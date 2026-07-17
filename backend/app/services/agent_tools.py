"""Agent tools — unified file-based tools that give digital employees
access to their own structured workspace.

Design principle:  ONE set of file tools covers EVERYTHING.
The agent's workspace uses well-known paths:
  - soul.md             → personality definition
  - memory/memory.md    → long-term memory / notes
  - skills/             → skill definitions (markdown files)
  - workspace/          → general working files, reports, etc.

The agent reads/writes these files directly. No per-concept tools needed.
"""

import asyncio
from collections.abc import Mapping
from copy import deepcopy
from dataclasses import dataclass, replace
import fnmatch
import hashlib
import json
import math
import multiprocessing as mp
import os
import queue
import tempfile
import uuid
import unicodedata
from contextvars import ContextVar
from datetime import date, datetime, timedelta, timezone
from pathlib import Path
from typing import Optional, Any, cast
import re

from loguru import logger
from sqlalchemy import select, or_

from app.core.permissions import (
    evaluate_roster_agent_visibility,
    evaluate_roster_human_visibility,
)
from app.database import async_session
from app.models.agent import Agent as AgentModel
from app.models.audit import ChatMessage
from app.models.chat_session import ChatSession
from app.models.channel_config import ChannelConfig
from app.models.identity import IdentityProvider
from app.models.org import (
    OrgDepartment,
    OrgMember,
)
from app.models.task import Task
from app.models.user import User as UserModel
from app.services.channel_session import find_or_create_channel_session
from app.services.channel_user_service import get_platform_user_by_org_member
from app.services.document_conversion import (
    convert_html_to_pdf as convert_html_file_to_pdf,
    convert_html_to_pptx as convert_html_file_to_pptx,
)
from app.services.focus_service import (
    complete_focus_item,
    ensure_focus_item,
    is_focus_file_path,
    list_focus_items,
    upsert_focus_item,
)
from app.services import agent_directory
from app.services.workspace_collaboration import (
    delete_workspace_file,
    move_workspace_path,
    normalize_workspace_path,
    write_workspace_file,
)
from app.services.storage import get_storage_backend, normalize_storage_key
from app.services.storage_runtime.base import WriteCondition, content_hash_bytes
from app.services.workspace_locking import workspace_locks
from app.config import get_settings
from app.services.llm.finish import (
    FINISH_TOOL_NAME,
)
from app.services.builtin_tool_definitions import (
    BUILTIN_TOOL_DEFINITIONS,
    BUILTIN_TOOL_NAMES,
    builtin_model_definition,
    builtin_model_definitions,
    builtin_readiness,
    builtin_sensitive_paths,
    is_reserved_custom_tool_name,
)
from app.services.agent_runtime.tool_execution import (
    ToolExecutionOutcome,
    sanitize_tool_arguments,
)


_settings = get_settings()
WORKSPACE_ROOT = Path(_settings.STORAGE_LOCAL_ROOT or _settings.AGENT_DATA_DIR)
TOOL_MATERIALIZE_MAX_FILE_BYTES = 10 * 1024 * 1024
TOOL_MATERIALIZE_MAX_TOTAL_BYTES = 100 * 1024 * 1024
TEMP_WORKSPACE_DEFAULT_PATHS = ["workspace", "memory", "skills", "focus.md", "soul.md", "HEARTBEAT.md"]
MAX_EXEC_STDOUT_CAPTURE_BYTES = 1_000_000
MAX_EXEC_STDERR_CAPTURE_BYTES = 500_000
_READ_FILE_BINARY_EXTENSIONS = frozenset(
    {
        ".7z",
        ".avi",
        ".bin",
        ".bmp",
        ".doc",
        ".docx",
        ".exe",
        ".gif",
        ".gz",
        ".ico",
        ".jpeg",
        ".jpg",
        ".mov",
        ".mp3",
        ".mp4",
        ".pdf",
        ".png",
        ".ppt",
        ".pptx",
        ".rar",
        ".tar",
        ".wav",
        ".webp",
        ".xls",
        ".xlsb",
        ".xlsm",
        ".xlsx",
        ".zip",
    }
)


def _read_file_binary_extension(path: str) -> str | None:
    suffix = Path(path.strip()).suffix.lower()
    return suffix if suffix in _READ_FILE_BINARY_EXTENSIONS else None


def _read_file_binary_error(path: str) -> str | None:
    suffix = _read_file_binary_extension(path)
    if suffix is None:
        return None
    return (
        f"read_file supports text files only; binary file type '{suffix}' "
        "must be opened with read_document instead."
    )


def _observability_arguments(tool_name: str, arguments: dict) -> dict:
    """Return a fail-closed, canonical-path-aware copy for logs/UI errors."""
    try:
        return sanitize_tool_arguments(
            arguments,
            sensitive_paths=builtin_sensitive_paths(tool_name),
        )
    except Exception:
        return {"_redacted": "tool arguments could not be safely serialized"}


def _observability_text(value: object) -> str:
    try:
        sanitized = sanitize_tool_arguments({"value": str(value)})
        return str(sanitized["value"])
    except Exception:
        return "[REDACTED: result could not be safely serialized]"

# ─── Tool Config Cache ──────────────────────────────────────────
# Cache tool configurations to avoid frequent DB queries
# Key: (agent_id, tool_name), Value: (config, expiry_time)
_tool_config_cache: dict[tuple, tuple[dict, datetime]] = {}
_TOOL_CONFIG_CACHE_TTL_SECONDS = 60

# Sensitive field keys that should be encrypted/decrypted
SENSITIVE_FIELD_KEYS = {"api_key", "private_key", "auth_code", "password", "secret", "atlassian_api_key"}

def _decrypt_sensitive_fields(config: dict, config_schema: dict | None = None) -> dict:
    """Decrypt sensitive fields in config dict.

    When config_schema is provided, also decrypts fields with type='password'
    (e.g. smithery_api_key) that are not in the hardcoded SENSITIVE_FIELD_KEYS.
    """
    if not config:
        return config

    from app.core.security import decrypt_data
    from app.config import get_settings

    settings = get_settings()
    result = dict(config)

    # Build the set of sensitive keys: hardcoded + schema-derived
    sensitive_keys = set(SENSITIVE_FIELD_KEYS)
    if config_schema:
        for field in config_schema.get("fields", []):
            if field.get("type") == "password":
                key = field.get("key", "")
                if key:
                    sensitive_keys.add(key)

    for key in sensitive_keys:
        if key in result and result[key]:
            value = result[key]
            if isinstance(value, str) and value:
                try:
                    result[key] = decrypt_data(value, settings.SECRET_KEY)
                except Exception:
                    # If decryption fails, assume it's plaintext
                    pass

    return result


def _get_cached_tool_config(agent_id: Optional[uuid.UUID], tool_name: str) -> Optional[dict]:
    """获取缓存的工具配置，过期返回 None。"""
    cache_key = (str(agent_id) if agent_id else None, tool_name)
    if cache_key in _tool_config_cache:
        config, expiry = _tool_config_cache[cache_key]
        if datetime.now() < expiry:
            return config
        # 过期，删除
        del _tool_config_cache[cache_key]
    return None


def _set_cached_tool_config(agent_id: Optional[uuid.UUID], tool_name: str, config: dict):
    """设置工具配置缓存。"""
    cache_key = (str(agent_id) if agent_id else None, tool_name)
    expiry = datetime.now() + timedelta(seconds=_TOOL_CONFIG_CACHE_TTL_SECONDS)
    _tool_config_cache[cache_key] = (config, expiry)


async def _get_tool_config(agent_id: Optional[uuid.UUID], tool_name: str) -> Optional[dict]:
    """Get merged tool config (with caching).

    Priority:
    1. agent_tools.config (per-agent override)
    2. tenant_settings tool_config:<tool_name> for builtin company config
    3. tools.config (tenant-specific/admin tool config or non-secret defaults)

    Both configs are decrypted using the tool's config_schema for
    schema-aware field detection (e.g. smithery_api_key with type=password).
    """
    # Check cache first
    cached = _get_cached_tool_config(agent_id, tool_name)
    if cached is not None:
        logger.debug(f"[ToolConfig] Cache hit for {tool_name}, agent_id={agent_id}")
        return cached

    from app.models.tool import Tool, AgentTool
    from app.models.agent import Agent as AgentModel
    from app.services.tool_config import get_tenant_tool_config

    async with async_session() as db:
        agent_tenant_id = None
        if agent_id:
            tenant_r = await db.execute(select(AgentModel.tenant_id).where(AgentModel.id == agent_id))
            agent_tenant_id = tenant_r.scalar_one_or_none()

        # 1. Try per-agent + global config together
        if agent_id:
            result = await db.execute(
                select(AgentTool.config, Tool.config, Tool.config_schema, Tool.source, Tool.name)
                .join(Tool, AgentTool.tool_id == Tool.id)
                .where(AgentTool.agent_id == agent_id, Tool.name == tool_name)
            )
            row = result.first()
            if row:
                agent_config, global_config, config_schema, tool_source, db_tool_name = row
                base_config = global_config or {}
                tenant_config = {}
                if tool_source == "builtin":
                    tenant_config = await get_tenant_tool_config(db, agent_tenant_id, db_tool_name, config_schema)
                # Merge: agent overrides global
                merged = {**base_config, **tenant_config, **(agent_config or {})}
                if merged:
                    # Decrypt with schema awareness
                    merged = _decrypt_sensitive_fields(merged, config_schema)
                    logger.info(f"[ToolConfig] DB merged config for {tool_name}, agent_id={agent_id}")
                    _set_cached_tool_config(agent_id, tool_name, merged)
                    return merged

        # 2. Fallback to global config only
        result = await db.execute(select(Tool).where(Tool.name == tool_name))
        tool = result.scalar_one_or_none()
        if tool:
            tenant_config = {}
            if tool.source == "builtin":
                tenant_config = await get_tenant_tool_config(db, agent_tenant_id, tool.name, tool.config_schema)
            base_config = tool.config or {}
            merged = {**base_config, **tenant_config}
        else:
            merged = {}
        if tool and merged:
            # Decrypt with schema awareness
            decrypted = _decrypt_sensitive_fields(merged, tool.config_schema)
            logger.info(f"[ToolConfig] DB global config for {tool_name}")
            _set_cached_tool_config(agent_id, tool_name, decrypted)
            return decrypted

    # Optional tools are resolved through this same path during every Runtime
    # workset build. An absent row/config is therefore an expected readiness
    # result, not a configuration failure. Database/query failures still
    # propagate to the caller and are logged as warnings by readiness gates.
    logger.debug(f"[ToolConfig] No DB config found for {tool_name}, agent_id={agent_id}")
    return None

# ContextVar set by each channel handler so send_channel_file knows where to send
# Value: async callable(file_path: Path) -> None  |  None for web chat (returns URL)
channel_file_sender: ContextVar = ContextVar('channel_file_sender', default=None)
# For web chat: agent_id needed to build download URL
channel_web_agent_id: ContextVar = ContextVar('channel_web_agent_id', default=None)
# Set by Feishu channel handler — open_id of the message sender so calendar tool
# can auto-invite them as attendee when no explicit attendee list is given
channel_feishu_sender_open_id: ContextVar = ContextVar('channel_feishu_sender_open_id', default=None)
# AgentBay execution identity is runtime context, not model-provided tool input.
agentbay_session_scope_id: ContextVar[str] = ContextVar(
    "agentbay_session_scope_id",
    default="",
)
agentbay_run_scope_id: ContextVar[str] = ContextVar(
    "agentbay_run_scope_id",
    default="",
)


def _agentbay_scope_ids(arguments: Mapping[str, Any]) -> tuple[str, str]:
    """Resolve exact scope without mutating durable/model arguments."""
    context_session_id = agentbay_session_scope_id.get().strip()
    legacy_session_id = arguments.get("_session_id", "")
    session_id = context_session_id or (
        legacy_session_id.strip()
        if isinstance(legacy_session_id, str)
        else ""
    )
    return session_id, agentbay_run_scope_id.get().strip()


async def _get_scoped_agentbay_client(
    agent_id: uuid.UUID,
    image_type: str,
    arguments: Mapping[str, Any],
):
    from app.services.agentbay_client import get_agentbay_client_for_agent

    session_id, run_id = _agentbay_scope_ids(arguments)
    return await get_agentbay_client_for_agent(
        agent_id,
        image_type,
        session_id=session_id,
        run_id=run_id,
    )

# ─── Tool Definitions (OpenAI function-calling format) ──────────

_HIDDEN_FROM_LLM_TOOL_NAMES = {
    "query_roster",
    "send_feishu_message",
}

# Compatibility export for call sites that still expect an OpenAI tools list.
# The description and JSON Schema are derived from the canonical builtin data;
# legacy aliases that are never model-visible remain Seeder-only definitions.
AGENT_TOOLS = [
    tool
    for tool in builtin_model_definitions()
    if tool["function"]["name"] not in _HIDDEN_FROM_LLM_TOOL_NAMES
]

_OKR_AGENT_ONLY_TOOL_NAMES = frozenset(
    str(definition["name"])
    for definition in BUILTIN_TOOL_DEFINITIONS
    if (definition.get("config") or {}).get("okr_agent_only") is True
)

_OKR_TRANSACTION_TOOL_NAMES = frozenset(
    {
        "get_okr",
        "get_my_okr",
        "get_okr_settings",
        "update_kr_progress",
        "update_kr_content",
        "create_objective",
        "create_key_result",
        "update_objective",
        "update_any_kr_progress",
        "upsert_member_daily_report",
    }
)

_OKR_JOB_TOOL_NAMES = frozenset(
    {
        "collect_okr_progress",
        "generate_okr_report",
        "generate_monthly_okr_report",
    }
)

_VERCEL_READ_TOOL_NAMES = frozenset(
    {
        "vercel_get_deploy_logs",
        "vercel_list_deployments",
    }
)

_DEPLOY_SIMPLE_WRITE_TOOL_NAMES = frozenset(
    {
        "vercel_set_env",
        "vercel_manage_domain",
        "neon_create_database",
    }
)

_AGENTBAY_A1_READ_TOOL_NAMES = frozenset(
    {
        "agentbay_browser_screenshot",
        "agentbay_browser_extract",
        "agentbay_browser_observe",
        "agentbay_code_read_file",
        "agentbay_computer_screenshot",
        "agentbay_computer_precision_screenshot",
        "agentbay_computer_get_screen_size",
        "agentbay_computer_get_installed_apps",
        "agentbay_computer_get_cursor_position",
        "agentbay_computer_get_active_window",
        "agentbay_computer_list_windows",
        "agentbay_computer_list_visible_apps",
    }
)

_IMAGE_GENERATION_TOOL_NAMES = frozenset(
    {
        "generate_image_siliconflow",
        "generate_image_openai",
        "generate_image_google",
        "generate_image_custom",
    }
)

_IMAGE_GENERATION_PROVIDER_BY_TOOL = {
    "generate_image_siliconflow": "siliconflow",
    "generate_image_openai": "openai",
    "generate_image_google": "google",
    "generate_image_custom": "custom",
}

_IMAGE_GENERATION_SIZES = frozenset(
    {
        "1024x1024",
        "1024x768",
        "768x1024",
        "1366x768",
        "768x1366",
        "1536x1024",
        "1024x1536",
    }
)
_MAX_GENERATED_IMAGE_BYTES = 25 * 1024 * 1024

# Application tools that have a native typed execution fact in Durable Runtime.
# `send_message_to_agent` is settled by RuntimeA2AService before the generic
# executor. Tools absent from this set remain available to legacy callers but
# are deterministically hidden from Durable Runtime until their real business
# boundary has a typed adapter.
RUNTIME_TYPED_APPLICATION_TOOL_NAMES = frozenset(
    {
        "list_files",
        "read_file",
        "search_files",
        "find_files",
        "list_focus_items",
        "upsert_focus_item",
        "complete_focus_item",
        "write_file",
        "move_file",
        "delete_file",
        "edit_file",
        "update_trigger",
        "cancel_trigger",
        "list_triggers",
        "query_directory",
        "send_channel_message",
        "send_platform_message",
        "send_message_to_agent",
        "execute_code",
        "execute_code_e2b",
        "convert_csv_to_xlsx",
        "convert_html_to_pdf",
        "convert_html_to_pptx",
        "convert_markdown_to_docx",
        "convert_markdown_to_pdf",
        "read_document",
        "read_webpage",
        "upload_image",
        *_IMAGE_GENERATION_TOOL_NAMES,
        "publish_page",
        "list_published_pages",
        "set_trigger",
        "send_channel_file",
        "send_file_to_agent",
        "duckduckgo_search",
        "web_search",
        "jina_search",
        "jina_read",
        "exa_search",
        "tavily_search",
        "google_search",
        "bing_search",
        "search_experience",
        "read_experience",
        "propose_experience_draft",
        "discover_resources",
        "import_mcp_server",
        "get_okr",
        "get_my_okr",
        "get_okr_settings",
        "update_kr_progress",
        "update_kr_content",
        "collect_okr_progress",
        "generate_okr_report",
        "generate_monthly_okr_report",
        "create_objective",
        "create_key_result",
        "update_objective",
        "update_any_kr_progress",
        "upsert_member_daily_report",
        "search_clawhub",
        "install_skill",
        "feishu_calendar_list",
        "feishu_calendar_create",
        "feishu_calendar_update",
        "feishu_calendar_delete",
        "feishu_wiki_list",
        "feishu_doc_search",
        "feishu_doc_read",
        "feishu_doc_create",
        "feishu_doc_append",
        "feishu_drive_share",
        "feishu_drive_delete",
        "feishu_user_search",
        "feishu_approval_query",
        "feishu_approval_get",
        "read_emails",
        "send_email",
        "reply_email",
        "bitable_create_app",
        "bitable_list_tables",
        "bitable_list_fields",
        "bitable_query_records",
        "bitable_create_record",
        "bitable_update_record",
        "bitable_delete_record",
        "vercel_list_deployments",
        "vercel_get_deploy_logs",
        "vercel_deploy",
        "vercel_set_env",
        "vercel_manage_domain",
        "neon_create_database",
        *_AGENTBAY_A1_READ_TOOL_NAMES,
    }
)


# Core tools that should always be available to agents regardless of
# DB configuration.
# Note: send_channel_message is intentionally NOT here — it lives in
# _CHANNEL_MESSAGE_TOOL_NAMES and is only added when a channel is configured,
# to avoid sending duplicate tool definitions to the LLM.
_ALWAYS_INCLUDE_CORE = {
    "complete_focus_item",
    FINISH_TOOL_NAME,
    "list_focus_items",
    "query_directory",
    "send_channel_file",
    "send_file_to_agent",
    "upsert_focus_item",
    "write_file",
}
# Channel message tool - available when any channel (Feishu/DingTalk/WeCom) is configured
_CHANNEL_MESSAGE_TOOL_NAMES = {
    "send_channel_message",
}
_always_core_tools = [t for t in AGENT_TOOLS if t["function"]["name"] in _ALWAYS_INCLUDE_CORE]
_channel_tools = [t for t in AGENT_TOOLS if t["function"]["name"] in _CHANNEL_MESSAGE_TOOL_NAMES]


async def _get_computer_os_type(agent_id: uuid.UUID) -> str:
    """Return the configured OS type for the agent's computer tool.

    Reads from agentbay_browser_navigate tool config (which stores all AgentBay
    settings including os_type). Defaults to 'windows' to match AgentBay's default.
    """
    try:
        config = await _get_tool_config(agent_id, "agentbay_browser_navigate")
        return (config or {}).get("os_type", "windows")
    except Exception:
        return "windows"


def _patch_computer_tool_descriptions(tools: list[dict], os_type: str) -> list[dict]:
    """Rewrite path examples in agentbay_file_transfer to match the agent's OS.

    This ensures the Agent always sees the correct desktop and home-directory
    paths for its specific computer environment without having to guess.
    """
    import copy

    if os_type == "windows":
        # Windows paths used by AgentBay's windows_latest image
        desktop_path = r"C:\Users\Administrator\Desktop"
        home_path    = r"C:\Users\Administrator"
        computer_os_label = "Windows"
    else:
        # Linux paths used by AgentBay's linux_latest image
        desktop_path = "/home/wuying/Desktop"
        home_path    = "/home/wuying"
        computer_os_label = "Linux"

    # Build the OS-aware description for agentbay_file_transfer
    new_file_transfer_desc = (
        (
        "Transfer a file between any two endpoints: the agent workspace, "
        "the AgentBay browser environment, the cloud desktop (computer), or the code sandbox.\n\n"
        f"COMPUTER ENVIRONMENT OS: {computer_os_label}\n"
        f"VERIFIED PATH CONVENTIONS for the computer environment ({computer_os_label}):\n"
        f"- computer desktop: {desktop_path}\\<filename>  (e.g. {desktop_path}\\report.xlsx)\n"
        f"- computer home:    {home_path}\\<filename>\n\n"
        "Other environments (Linux-based, user 'wuying', HOME=/home/wuying/):\n"
        "- code env:     /home/wuying/<filename>  (e.g. /home/wuying/data.csv)\n"
        "- browser env:  /home/wuying/下载/<filename>  (download folder)\n"
        "- workspace:    relative path, e.g. 'workspace/data.csv'\n\n"
        "Transfer directions:\n"
        "- workspace -> env: upload a workspace file into a cloud environment\n"
        "- env -> workspace: download a file from a cloud environment into the workspace\n"
        "- env A -> env B:   transfer between environments (transparent backend temp)"
        )
        if os_type == "windows"
        else (
        "Transfer a file between any two endpoints: the agent workspace, "
        "the AgentBay browser environment, the cloud desktop (computer), or the code sandbox.\n\n"
        f"COMPUTER ENVIRONMENT OS: {computer_os_label}\n"
        f"VERIFIED PATH CONVENTIONS for the computer environment ({computer_os_label}):\n"
        f"- computer desktop: {desktop_path}/<filename>  (e.g. {desktop_path}/report.xlsx)\n"
        f"- computer home:    {home_path}/<filename>\n\n"
        "Other environments (also Linux, user 'wuying'):\n"
        "- code env:     /home/wuying/<filename>  (e.g. /home/wuying/data.csv)\n"
        "- browser env:  /home/wuying/下载/<filename>  (download folder)\n"
        "- workspace:    relative path, e.g. 'workspace/data.csv'\n\n"
        "Transfer directions:\n"
        "- workspace -> env: upload a workspace file into a cloud environment\n"
        "- env -> workspace: download a file from a cloud environment into the workspace\n"
        "- env A -> env B:   transfer between environments (transparent backend temp)"
        )
    )

    patched = []
    for tool in tools:
        fn = tool.get("function", {})
        name = fn.get("name", "")
        if name == "agentbay_file_transfer":
            # Deep copy to avoid mutating the shared AGENT_TOOLS constant
            tool = copy.deepcopy(tool)
            tool["function"]["description"] = new_file_transfer_desc
            # Also patch from_path and to_path parameter hints
            props = tool["function"].get("parameters", {}).get("properties", {})
            if "from_path" in props:
                if os_type == "windows":
                    props["from_path"]["description"] = (
                        r"Source path. Relative if workspace (e.g. 'workspace/data.csv'). "
                        r"Absolute if env: computer → C:\Users\Administrator\Desktop\file, "
                        r"code → /home/wuying/file, browser → /home/wuying/下载/file."
                    )
                else:
                    props["from_path"]["description"] = (
                        "Source path. Relative if workspace (e.g. 'workspace/data.csv'). "
                        "Absolute if env: computer → /home/wuying/Desktop/file, "
                        "code → /home/wuying/file, browser → /home/wuying/下载/file."
                    )
            if "to_path" in props:
                if os_type == "windows":
                    props["to_path"]["description"] = (
                        r"Destination path. Relative if workspace (e.g. 'workspace/output.csv'). "
                        r"Absolute if env: computer → C:\Users\Administrator\Desktop\file, "
                        r"code → /home/wuying/file, browser → /home/wuying/下载/file."
                    )
                else:
                    props["to_path"]["description"] = (
                        "Destination path. Relative if workspace (e.g. 'workspace/output.csv'). "
                        "Absolute if env: computer → /home/wuying/Desktop/file, "
                        "code → /home/wuying/file, browser → /home/wuying/下载/file."
                    )
        patched.append(tool)
    return patched


async def _agent_has_feishu(agent_id: uuid.UUID) -> bool:
    """Check deterministic local Feishu channel readiness."""
    try:
        from app.models.channel_config import ChannelConfig
        async with async_session() as db:
            r = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "feishu",
                    ChannelConfig.is_configured.is_(True),
                )
            )
            config = r.scalar_one_or_none()
            return bool(
                config
                and config.is_configured
                and isinstance(config.app_id, str)
                and bool(config.app_id.strip())
                and isinstance(config.app_secret, str)
                and bool(config.app_secret.strip())
            )
    except Exception:
        return False


async def _agent_has_any_channel(agent_id: uuid.UUID) -> bool:
    """Check if agent has any configured channel (Feishu/DingTalk/WeCom)."""
    try:
        from app.models.channel_config import ChannelConfig
        async with async_session() as db:
            r = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.is_configured == True,
                )
            )
            return r.scalar_one_or_none() is not None
    except Exception:
        return False


# ─── Dynamic Tool Loading from DB ──────────────────────────────


def _canonicalize_llm_tool(tool_def: dict, *, source: str = "builtin") -> dict:
    """Replace stale DB schemas with the current model-facing contract."""
    name = tool_def.get("function", {}).get("name")
    if source != "builtin" or name not in BUILTIN_TOOL_NAMES:
        return tool_def
    return builtin_model_definition(name)


async def get_agent_tools_for_llm(agent_id: uuid.UUID) -> list[dict]:
    """Load enabled tools for an agent from DB (OpenAI function-calling format).

    Falls back to hardcoded AGENT_TOOLS if DB not ready.
    Includes core system tools (send_channel_file, write_file) unless the user
    has explicitly disabled them via the Agent tool panel.
    Feishu tools are only included when the agent has a configured Feishu channel.
    send_channel_message is included when any channel (Feishu/DingTalk/WeCom) is configured.

    Also patches agentbay_file_transfer description with OS-specific paths based on
    the agent's computer tool configuration (os_type: 'windows' | 'linux').

    A2A always exposes notify, consult, and task_delegate; the durable Runtime
    owns their different wait/resume behavior.
    """
    has_feishu = await _agent_has_feishu(agent_id)
    has_any_channel = await _agent_has_any_channel(agent_id)
    # A configured channel satisfies a prerequisite; it does not assign every
    # tool in that provider family. Feishu application tools still require an
    # enabled AgentTool assignment or their explicit canonical default.
    _always_tools = _always_core_tools + (
        _channel_tools if has_any_channel else []
    )

    is_system_agent = False
    agent_tenant_id = None
    try:
        from app.models.agent import Agent as AgentModel
        async with async_session() as _flag_db:
            _ag_r = await _flag_db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            _agent = _ag_r.scalar_one_or_none()
            _tid = _agent.tenant_id if _agent else None
            agent_tenant_id = _tid
            is_system_agent = bool(_agent and _agent.is_system)
    except Exception:
        pass

    # Read os_type once; used to patch agentbay_file_transfer paths below
    computer_os_type = await _get_computer_os_type(agent_id)

    try:
        from app.models.tool import Tool, AgentTool

        async with async_session() as db:
            # Get agent-specific assignments
            agent_tools_r = await db.execute(select(AgentTool).where(AgentTool.agent_id == agent_id))
            assignments = {str(at.tool_id): at for at in agent_tools_r.scalars().all()}
            assigned_tool_ids = [uuid.UUID(tool_id) for tool_id in assignments]

            visible_clauses = [Tool.source == "builtin"]
            # Admin tools: visible if they are global (tenant_id is NULL) or belong to the agent's tenant
            admin_cond = (Tool.tenant_id == None)
            if agent_tenant_id:
                admin_cond = admin_cond | (Tool.tenant_id == agent_tenant_id)
            visible_clauses.append((Tool.source == "admin") & admin_cond)
            # Explicitly assigned tools: always visible regardless of source (builtin, admin, agent)
            if assigned_tool_ids:
                visible_clauses.append(Tool.id.in_(assigned_tool_ids))

            # Get all tools visible within this agent's tenant boundary.
            all_tools_r = await db.execute(
                select(Tool).where(Tool.enabled == True, or_(*visible_clauses))
            )
            all_tools = all_tools_r.scalars().all()

            result = []
            db_tool_names = set()
            # Track tool names that were explicitly disabled by the user
            # (have an AgentTool record with enabled=False). These must NOT
            # be re-added by the _always_tools fallback below.
            explicitly_disabled_names = set()
            # Track tools included via is_default fallback (no AgentTool record)
            default_included_names = []

            for t in all_tools:
                # ORM rows always carry `source`; lightweight compatibility
                # fixtures and pre-source legacy rows are builtin by default.
                source = getattr(t, "source", "builtin")
                if t.name in _HIDDEN_FROM_LLM_TOOL_NAMES:
                    continue
                if source == "builtin" and t.name not in BUILTIN_TOOL_NAMES:
                    logger.warning(
                        "[Tools] Ignoring builtin row without a canonical definition: {}",
                        t.name,
                    )
                    continue
                if source != "builtin" and (
                    t.name in BUILTIN_TOOL_NAMES
                    or is_reserved_custom_tool_name(t.name)
                ):
                    logger.warning(
                        "[Tools] Ignoring custom tool with canonical or "
                        "Runtime-reserved name: {}",
                        t.name,
                    )
                    continue

                tid = str(t.id)
                at = assignments.get(tid)

                # If no explicit assignment, fallback to t.is_default
                enabled = at.enabled if at is not None else t.is_default

                if at is None and t.is_default:
                    default_included_names.append(t.name)

                if not enabled:
                    if at and not at.enabled:
                        explicitly_disabled_names.add(t.name)
                    continue

                # Skip feishu tools if the agent has no Feishu channel configured
                if t.category == "feishu" and not has_feishu:
                    continue
                if t.name in _CHANNEL_MESSAGE_TOOL_NAMES and not has_any_channel:
                    continue
                # Match the Agent Tools UI: regular agents must not receive
                # OKR-system-only tools, even if the DB default says enabled.
                if (t.config or {}).get("okr_agent_only") and not is_system_agent:
                    continue
                # Build OpenAI function-calling format
                tool_def = {
                    "type": "function",
                    "function": {
                        "name": t.name,
                        "description": t.description,
                        "parameters": t.parameters_schema or {"type": "object", "properties": {}},
                    },
                }
                tool_def = _canonicalize_llm_tool(tool_def, source=source)
                # Defensive dedup: skip if this name was already added.
                # Normally the UNIQUE constraint on tool.name prevents duplicate
                # rows, but old DB dumps (pre-constraint) may have them. Without
                # this guard, the LLM would receive duplicate tool names and
                # return HTTP 400 "Tool names must be unique".
                if t.name in db_tool_names:
                    logger.warning(
                        f"[Tools] Duplicate tool name '{t.name}' found in DB "
                        f"(id={t.id}). Skipping to avoid LLM error. "
                        "Run: DELETE FROM tools WHERE id IN (SELECT id FROM "
                        "(SELECT id, ROW_NUMBER() OVER (PARTITION BY name "
                        "ORDER BY created_at DESC) AS rn FROM tools) t WHERE rn > 1);"
                    )
                    continue

                result.append(tool_def)
                db_tool_names.add(t.name)

            if default_included_names:
                logger.info(
                    f"[Tools] agent={agent_id} included via default fallback (no AgentTool record): "
                    f"{sorted(default_included_names)}"
                )

            if result:
                # Append always-available system tools that aren't already in
                # the DB list — but respect explicit user disabling.
                always_added = []
                for t in _always_tools:
                    fn_name = t["function"]["name"]
                    if fn_name not in db_tool_names and fn_name not in explicitly_disabled_names:
                        result.append(t)
                        always_added.append(fn_name)
                if always_added:
                    logger.debug(
                        f"[Tools] agent={agent_id} added from _always_tools: {always_added}"
                    )
                # Inject OS-aware paths into computer-related tool descriptions
                result = _patch_computer_tool_descriptions(result, computer_os_type)
                # Final diagnostic: log the complete tool list and assignment stats
                final_names = sorted(t["function"]["name"] for t in result)
                logger.info(
                    f"[Tools] agent={agent_id} FINAL {len(result)} tools "
                    f"(assignments={len(assignments)}, "
                    f"disabled={len(explicitly_disabled_names)}, "
                    f"default_fallback={len(default_included_names)}): "
                    f"{final_names}"
                )
                return result
            # If DB loading fails, do not expose the full hardcoded tool catalog: that
            # can leak disabled tools (for example search tools) into the LLM. Keep only
            # the minimal always-available core/channel tools.
            # (Note: we fall through to the except-clause fallback below if result is empty or exception is raised)
            raise ValueError("No tools found for agent in DB")
    except Exception as e:
        logger.error(f"[Tools] DB load failed, using fallback: {e}")

    # If DB loading fails, do not expose the full hardcoded tool catalog: that
    # can leak disabled tools (for example search tools) into the LLM. Keep only
    # the minimal always-available core/channel tools.
    fallback = _patch_computer_tool_descriptions(_always_tools, computer_os_type)
    return fallback


def _runtime_typed_tools(
    tools: list[dict],
    *,
    dynamic_mcp_names: set[str] | frozenset[str] = frozenset(),
) -> list[dict]:
    """Keep only tools with a native Runtime execution fact.

    Canonical builtin names always use the explicit typed-name gate.  A
    dynamic MCP row may enter only through the separately resolved exact-name
    workset and may not replace Runtime control, Group, or builtin contracts.
    """
    return [
        tool
        for tool in tools
        if (
            (name := str(tool.get("function", {}).get("name") or ""))
            in RUNTIME_TYPED_APPLICATION_TOOL_NAMES
            or (
                name in dynamic_mcp_names
                and name not in BUILTIN_TOOL_NAMES
                and not is_reserved_custom_tool_name(name)
            )
        )
    ]


async def _agent_is_designated_okr_agent(agent_id: uuid.UUID) -> bool:
    """Fail closed unless tenant OKR settings designate this exact Agent."""
    from app.models.okr import OKRSettings

    try:
        async with async_session() as db:
            result = await db.execute(
                select(OKRSettings.okr_agent_id).where(
                    OKRSettings.okr_agent_id == agent_id
                )
            )
            return result.scalar_one_or_none() == agent_id
    except Exception as exc:
        logger.warning(
            "[Tools] Designated OKR Agent lookup failed: {}",
            type(exc).__name__,
        )
        return False


async def _get_runtime_dynamic_mcp_tool_names(
    agent_id: uuid.UUID,
) -> set[str]:
    """Resolve locally ready dynamic MCP names without provider I/O."""
    from urllib.parse import urlparse

    from app.models.tool import AgentTool, Tool

    try:
        async with async_session() as db:
            result = await db.execute(
                select(Tool)
                .join(AgentTool, AgentTool.tool_id == Tool.id)
                .where(
                    AgentTool.agent_id == agent_id,
                    AgentTool.enabled.is_(True),
                    Tool.enabled.is_(True),
                    Tool.type == "mcp",
                )
            )
            tools = result.scalars().all()
    except Exception as exc:
        logger.warning(
            "[Tools] Dynamic MCP readiness lookup failed: {}",
            type(exc).__name__,
        )
        return set()

    ready: set[str] = set()
    for tool in tools:
        name = str(tool.name or "")
        server_url = str(tool.mcp_server_url or "").strip()
        parsed = urlparse(server_url)
        if (
            not name
            or name in BUILTIN_TOOL_NAMES
            or is_reserved_custom_tool_name(name)
            or not str(tool.mcp_tool_name or "").strip()
            or parsed.scheme not in {"http", "https"}
            or not parsed.netloc
        ):
            logger.info(
                "[Tools] Durable Runtime hid locally unready MCP tool {}",
                name or "<unnamed>",
            )
            continue
        ready.add(name)
    return ready


async def get_runtime_agent_tools_for_llm(agent_id: uuid.UUID) -> list[dict]:
    """Resolve the current Durable Runtime workset with typed-outcome gating."""
    tools = await get_agent_tools_for_llm(agent_id)
    dynamic_mcp_names = await _get_runtime_dynamic_mcp_tool_names(agent_id)
    resolved = _runtime_typed_tools(
        tools,
        dynamic_mcp_names=dynamic_mcp_names,
    )
    ready: list[dict] = []
    is_designated_okr_agent: bool | None = None
    for tool in resolved:
        name = str(tool.get("function", {}).get("name") or "")
        if name in _OKR_AGENT_ONLY_TOOL_NAMES:
            if is_designated_okr_agent is None:
                is_designated_okr_agent = (
                    await _agent_is_designated_okr_agent(agent_id)
                )
            if not is_designated_okr_agent:
                logger.info(
                    "[Tools] Durable Runtime hid {} because this Agent is "
                    "not the tenant's designated OKR Agent",
                    name,
                )
                continue
        readiness = builtin_readiness(name)
        if name == "web_search":
            try:
                search_config = await _get_tool_config(agent_id, name) or {}
            except Exception as exc:
                logger.warning(
                    "[Tools] Web search readiness lookup failed: {}",
                    type(exc).__name__,
                )
                continue
            engine = str(
                search_config.get("search_engine") or "duckduckgo"
            ).strip().lower()
            if engine == "duckduckgo" or (
                engine in {"tavily", "google", "bing", "exa"}
                and bool(search_config.get("api_key"))
            ):
                ready.append(tool)
            else:
                logger.info(
                    "[Tools] Durable Runtime hid web_search because its local "
                    "engine configuration is not ready"
                )
            continue
        if readiness == "e2b_configuration":
            try:
                e2b_config = await _get_tool_config(agent_id, name) or {}
            except Exception as exc:
                logger.warning(
                    "[Tools] E2B readiness lookup failed: {}",
                    type(exc).__name__,
                )
                continue
            if (
                e2b_config.get("sandbox_type") == "e2b"
                and isinstance(e2b_config.get("api_key"), str)
                and bool(e2b_config["api_key"].strip())
            ):
                ready.append(tool)
            else:
                logger.info(
                    "[Tools] Durable Runtime hid execute_code_e2b because "
                    "its local E2B configuration is not ready"
                )
            continue
        if readiness == "configured_channel":
            try:
                if await _agent_has_any_channel(agent_id):
                    ready.append(tool)
                else:
                    logger.info(
                        "[Tools] Durable Runtime hid {} because no channel is configured",
                        name,
                    )
            except Exception as exc:
                logger.warning(
                    "[Tools] Channel readiness lookup failed for {}: {}",
                    name,
                    type(exc).__name__,
                )
            continue
        if readiness == "feishu_channel":
            try:
                if await _agent_has_feishu(agent_id):
                    ready.append(tool)
                else:
                    logger.info(
                        "[Tools] Durable Runtime hid {} because the local "
                        "Feishu channel credentials are incomplete",
                        name,
                    )
            except Exception as exc:
                logger.warning(
                    "[Tools] Feishu readiness lookup failed for {}: {}",
                    name,
                    type(exc).__name__,
                )
            continue
        if readiness == "email_configuration":
            try:
                email_config = await _get_email_config(agent_id)
            except Exception as exc:
                logger.warning(
                    "[Tools] Email readiness lookup failed for {}: {}",
                    name,
                    type(exc).__name__,
                )
                continue
            _, ready_protocols = _resolve_local_email_configuration(
                email_config
            )
            required_protocols = {
                "send_email": frozenset({"smtp"}),
                "read_emails": frozenset({"imap"}),
                "reply_email": frozenset({"imap", "smtp"}),
            }.get(name, frozenset())
            if required_protocols and required_protocols <= ready_protocols:
                ready.append(tool)
            else:
                logger.info(
                    "[Tools] Durable Runtime hid {} because its local Email "
                    "protocol configuration is incomplete",
                    name,
                )
            continue
        if readiness == "agentbay_configuration":
            try:
                # AgentBay stores the family configuration on one canonical
                # representative. Readiness is local-only and never constructs
                # the SDK or calls the Provider.
                agentbay_config = await _get_tool_config(
                    agent_id,
                    "agentbay_browser_navigate",
                ) or {}
            except Exception as exc:
                logger.warning(
                    "[Tools] AgentBay readiness lookup failed for {}: {}",
                    name,
                    type(exc).__name__,
                )
                continue
            from app.services.agentbay_client import (
                _is_plausible_agentbay_api_key,
            )

            if (
                _is_plausible_agentbay_api_key(agentbay_config.get("api_key"))
                and str(agentbay_config.get("os_type") or "").strip()
                in {"linux", "windows"}
            ):
                ready.append(tool)
            else:
                logger.info(
                    "[Tools] Durable Runtime hid {} because its local "
                    "AgentBay configuration is incomplete",
                    name,
                )
            continue
        if readiness != "configured_credentials":
            ready.append(tool)
            continue
        # D-020 requires deterministic local prerequisites to be checked at
        # model-step resolution without pinging the provider. Vercel siblings
        # deliberately share the credential stored by vercel_deploy; image
        # generators remain isolated to their own configuration.
        config_tool_name = (
            "vercel_deploy" if name.startswith("vercel_") else name
        )
        try:
            config = await _get_tool_config(agent_id, config_tool_name) or {}
        except Exception as exc:
            logger.warning(
                "[Tools] Readiness config lookup failed for {}: {}",
                name,
                type(exc).__name__,
            )
            config = {}
        if name.startswith("vercel_") and str(
            config.get("vercel_token") or ""
        ).strip():
            ready.append(tool)
        elif name == "neon_create_database" and str(
            config.get("neon_api_key") or ""
        ).strip():
            ready.append(tool)
        elif name == "upload_image" and str(
            config.get("private_key") or ""
        ).strip():
            ready.append(tool)
        elif name in {
            "generate_image_siliconflow",
            "generate_image_openai",
            "generate_image_google",
        } and str(config.get("api_key") or "").strip():
            ready.append(tool)
        elif name == "generate_image_custom" and all(
            str(config.get(field) or "").strip()
            for field in (
                "api_key",
                "base_url",
                "model",
                "response_image_path",
            )
        ):
            ready.append(tool)
        elif name == "discover_resources" and (
            config.get("smithery_api_key")
            or config.get("modelscope_api_token")
        ):
            ready.append(tool)
        elif name == "import_mcp_server" and config.get("smithery_api_key"):
            ready.append(tool)
        elif name in {"tavily_search", "google_search", "bing_search"} and (
            config.get("api_key")
        ):
            ready.append(tool)
        elif name == "exa_search" and (
            config.get("api_key") or get_settings().EXA_API_KEY
        ):
            ready.append(tool)
        else:
            logger.info(
                "[Tools] Durable Runtime hid {} because credentials are not configured",
                name,
            )
    hidden = sorted(
        {
            str(tool.get("function", {}).get("name") or "")
            for tool in tools
        }
        - RUNTIME_TYPED_APPLICATION_TOOL_NAMES
        - dynamic_mcp_names
        - {""}
    )
    if hidden:
        logger.info(
            "[Tools] Durable Runtime hid tools without typed outcomes: {}",
            hidden,
        )
    return ready


# ─── Workspace initialization ──────────────────────────────────


async def initialize_agent_workspace(agent_id: uuid.UUID) -> None:
    """Seed default workspace files into shared storage once at agent creation time."""
    storage = get_storage_backend()
    mem_key = normalize_storage_key(f"{agent_id}/memory/memory.md")
    if not await storage.is_file(mem_key):
        await storage.write_text(
            mem_key,
            "# Memory\n\n_Record important information and knowledge here._\n",
            encoding="utf-8",
        )

    soul_key = normalize_storage_key(f"{agent_id}/soul.md")
    if not await storage.is_file(soul_key):
        # Soul is an independently editable personality artifact. The Agent
        # role enters the prompt through Identity and must not be duplicated
        # into Soul as a fallback.
        await storage.write_text(
            soul_key,
            "# Personality\n\n_Describe personality, values, and working style here._\n",
            encoding="utf-8",
        )


@dataclass
class TempWorkspaceManifestEntry:
    rel_path: str
    storage_key: str
    base_version_token: str
    base_hash: str
    size: int


@dataclass
class TempWorkspace:
    temp_dir: tempfile.TemporaryDirectory
    root: Path
    agent_id: uuid.UUID
    tenant_id: str | None
    selected_paths: list[str]
    manifest: dict[str, TempWorkspaceManifestEntry]

    def cleanup(self) -> None:
        self.temp_dir.cleanup()


async def _materialize_storage_workspace(storage, storage_key: str, local_root: Path) -> None:
    if not await storage.is_dir(storage_key):
        return
    for entry in await storage.list_dir(storage_key):
        await _materialize_storage_entry(storage, entry.key, storage_key, local_root)


async def _materialize_storage_entry(storage, entry_key: str, root_key: str, local_root: Path) -> None:
    rel = entry_key.removeprefix(root_key.rstrip("/") + "/")
    target = (local_root / rel).resolve()
    if not target.is_relative_to(local_root.resolve()):
        return
    if await storage.is_dir(entry_key):
        target.mkdir(parents=True, exist_ok=True)
        for child in await storage.list_dir(entry_key):
            await _materialize_storage_entry(storage, child.key, root_key, local_root)
        return
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_bytes(await storage.read_bytes(entry_key))


async def _prepare_temp_workspace(
    agent_id: uuid.UUID,
    tenant_id: str | None = None,
    paths: list[str] | None = None,
) -> TempWorkspace:
    tmp = tempfile.TemporaryDirectory(prefix=f"clawith-agent-{str(agent_id)[:8]}-")
    temp_ws = Path(tmp.name)
    for folder in ("workspace", "memory", "skills"):
        (temp_ws / folder).mkdir(parents=True, exist_ok=True)

    storage = get_storage_backend()
    budget = {"total": 0}
    selected = TEMP_WORKSPACE_DEFAULT_PATHS if paths is None else [path for path in paths if path]
    manifest: dict[str, TempWorkspaceManifestEntry] = {}
    for rel_path in selected:
        storage_key, normalized, is_enterprise = _tool_storage_key(agent_id, rel_path, tenant_id)
        if is_enterprise:
            continue
        await _materialize_storage_path_with_budget(storage, storage_key, normalized, temp_ws, budget, manifest)
    return TempWorkspace(
        temp_dir=tmp,
        root=temp_ws,
        agent_id=agent_id,
        tenant_id=tenant_id,
        selected_paths=list(selected),
        manifest=manifest,
    )


async def _materialize_storage_path_with_budget(
    storage,
    storage_key: str,
    rel_path: str,
    local_root: Path,
    budget: dict,
    manifest: dict[str, TempWorkspaceManifestEntry],
) -> None:
    if await storage.is_file(storage_key):
        version = await storage.get_version(storage_key)
        if version.size > TOOL_MATERIALIZE_MAX_FILE_BYTES:
            return
        if budget["total"] + version.size > TOOL_MATERIALIZE_MAX_TOTAL_BYTES:
            return
        target = (local_root / rel_path).resolve()
        if not target.is_relative_to(local_root.resolve()):
            return
        target.parent.mkdir(parents=True, exist_ok=True)
        data = await storage.read_bytes(storage_key)
        target.write_bytes(data)
        normalized_rel = normalize_workspace_path(rel_path)
        manifest[normalized_rel] = TempWorkspaceManifestEntry(
            rel_path=normalized_rel,
            storage_key=storage_key,
            base_version_token=version.token,
            base_hash=content_hash_bytes(data),
            size=version.size,
        )
        budget["total"] += version.size
        return
    if await storage.is_dir(storage_key):
        (local_root / rel_path).mkdir(parents=True, exist_ok=True)
        for entry in await storage.list_dir(storage_key):
            child_rel = f"{rel_path.rstrip('/')}/{entry.name}" if rel_path else entry.name
            await _materialize_storage_path_with_budget(storage, entry.key, child_rel, local_root, budget, manifest)


async def _sync_tasks_to_file(agent_id: uuid.UUID, ws: Path):
    """Sync tasks from DB to legacy tasks.json, if the file already exists."""
    tasks_path = ws / "tasks.json"
    if not tasks_path.exists():
        return

    try:
        async with async_session() as db:
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id).order_by(Task.created_at.desc())
            )
            tasks = result.scalars().all()

        task_list = []
        for t in tasks:
            task_list.append({
                "title": t.title,
                "status": t.status,
                "priority": t.priority,
                "description": t.description or "",
                "created_at": t.created_at.isoformat() if t.created_at else "",
                "completed_at": t.completed_at.isoformat() if t.completed_at else "",
            })

        tasks_path.write_text(
            json.dumps(task_list, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    except Exception as e:
        logger.error(f"[AgentTools] Failed to sync tasks: {e}")


async def flush_temp_workspace(temp_workspace: TempWorkspace, conflict_mode: str = "fail") -> dict[str, list[str]]:
    """Flush local changes back to storage using manifest-based conflict checks."""
    storage = get_storage_backend()
    selected_paths = [normalize_workspace_path(path) for path in temp_workspace.selected_paths]
    manifest = temp_workspace.manifest
    local_files = _collect_temp_workspace_files(temp_workspace.root, selected_paths)

    updated: list[str] = []
    conflicted: list[str] = []
    deleted: list[str] = []
    skipped: list[str] = []

    async with workspace_locks(temp_workspace.agent_id, selected_paths):
        for rel_path, local_path in local_files.items():
            if local_path.name.startswith("_exec_tmp") or "__pycache__" in local_path.parts:
                continue
            data = local_path.read_bytes()
            current_hash = content_hash_bytes(data)
            entry = manifest.get(rel_path)
            if entry and entry.base_hash == current_hash:
                skipped.append(rel_path)
                continue
            condition = (
                WriteCondition(version_token=entry.base_version_token)
                if entry
                else WriteCondition(require_absent=True)
            )
            storage_key = entry.storage_key if entry else normalize_storage_key(f"{temp_workspace.agent_id}/{rel_path}")
            result = await storage.write_bytes_if_match(
                storage_key,
                data,
                condition=condition,
            )
            if not result.ok:
                conflicted.append(rel_path)
                if conflict_mode == "fail":
                    return {"updated": updated, "deleted": deleted, "conflicted": conflicted, "skipped": skipped}
                continue
            updated.append(rel_path)

        for rel_path, entry in manifest.items():
            if rel_path in local_files:
                continue
            result = await storage.delete_if_match(
                entry.storage_key,
                condition=WriteCondition(version_token=entry.base_version_token),
            )
            if not result.ok:
                conflicted.append(rel_path)
                if conflict_mode == "fail":
                    return {"updated": updated, "deleted": deleted, "conflicted": conflicted, "skipped": skipped}
                continue
            deleted.append(rel_path)

    return {"updated": updated, "deleted": deleted, "conflicted": conflicted, "skipped": skipped}


def _collect_temp_workspace_files(root: Path, selected_paths: list[str]) -> dict[str, Path]:
    files: dict[str, Path] = {}
    root_resolved = root.resolve()
    for selected in selected_paths:
        if not selected:
            continue
        target = (root_resolved / selected).resolve()
        if not target.is_relative_to(root_resolved):
            continue
        if target.is_file():
            files[normalize_workspace_path(selected)] = target
            continue
        if not target.exists() or not target.is_dir():
            continue
        for path in target.rglob("*"):
            if not path.is_file():
                continue
            rel = path.resolve().relative_to(root_resolved).as_posix()
            files[normalize_workspace_path(rel)] = path
    return files


# ─── Tool Executors ─────────────────────────────────────────────

# Mapping from tool_name to autonomy action_type used for policy lookup and notifications.
# Each tool name maps to the action_type key in the agent's autonomy_policy dict.
# Using the tool's own name avoids misleading notification titles (e.g. showing
# "send_feishu_message" when the agent actually called send_message_to_agent).
_TOOL_AUTONOMY_MAP = {
    "write_file": "write_workspace_files",
    "move_file": "write_workspace_files",
    "delete_file": "delete_files",
    "send_feishu_message": "send_feishu_message",
    "send_message_to_agent": "send_message_to_agent",  # A2A messaging — distinct from feishu
    "send_file_to_agent": "send_file_to_agent",          # A2A file transfer
    "web_search": "web_search",
    "execute_code": "execute_code",
    "execute_code_e2b": "execute_code",
}


def _is_enterprise_info_path(path: str | None) -> bool:
    normalized = str(path or "").replace("\\", "/").strip().strip("/")
    return normalized == "enterprise_info" or normalized.startswith("enterprise_info/")


async def _get_agent_tenant_id(agent_id: uuid.UUID) -> str | None:
    """Get the agent tenant ID for tenant-scoped shared paths."""
    try:
        async with async_session() as db:

            r = await db.execute(select(AgentModel.tenant_id).where(AgentModel.id == agent_id))

            tenant_id = r.scalar_one_or_none()
            if tenant_id:
                return str(tenant_id)
    except Exception:
        pass
    return None


def _agent_workspace_root(agent_id: uuid.UUID) -> Path:
    """Return the per-agent local path without creating or hydrating it."""
    return WORKSPACE_ROOT / str(agent_id)


def _non_empty_paths(*paths: str | None) -> list[str] | None:
    selected = [path for path in paths if path]
    return selected or None


async def _run_with_temp_workspace(
    agent_id: uuid.UUID,
    tenant_id: str | None,
    runner,
    *,
    paths: list[str] | None = None,
    sync_back: bool = False,
) -> str:
    """Materialize a temporary workspace for tools that require local files."""
    temp_workspace = await _prepare_temp_workspace(agent_id, tenant_id=tenant_id, paths=paths)
    try:
        result = await runner(temp_workspace.root)
        if sync_back:
            flush_result = await flush_temp_workspace(temp_workspace, conflict_mode="fail")
            if flush_result["conflicted"]:
                conflict_list = ", ".join(flush_result["conflicted"][:5])
                return f"❌ Workspace sync conflict for: {conflict_list}"
        return result
    finally:
        temp_workspace.cleanup()


def _workspace_artifact_ref(agent_id: uuid.UUID, path: str) -> str:
    return f"workspace://{agent_id}/{normalize_workspace_path(path)}"


async def _run_with_temp_workspace_outcome(
    agent_id: uuid.UUID,
    tenant_id: str | None,
    runner,
    *,
    paths: list[str] | None = None,
    sync_back: bool = False,
    sync_back_on_non_success: bool = False,
) -> ToolExecutionOutcome:
    """Run a typed local-content tool and preserve explicit sync facts."""
    try:
        temp_workspace = await _prepare_temp_workspace(
            agent_id,
            tenant_id=tenant_id,
            paths=paths,
        )
    except Exception as exc:
        return _typed_failure(
            f"Local content could not be materialized: {type(exc).__name__}.",
            "local_content_materialize_failed",
        )
    try:
        outcome = await runner(temp_workspace.root)
        if not isinstance(outcome, ToolExecutionOutcome):
            return _typed_failure(
                "Local content adapter returned an invalid outcome.",
                "invalid_local_content_outcome",
            )
        if not sync_back or (
            outcome.status != "succeeded" and not sync_back_on_non_success
        ):
            return outcome
        try:
            flush_result = await flush_temp_workspace(
                temp_workspace,
                conflict_mode="fail",
            )
        except Exception as exc:
            return _typed_unknown(
                f"Local execution completed but workspace sync is unknown: {type(exc).__name__}.",
                "workspace_sync_outcome_unknown",
            )
        if flush_result["conflicted"]:
            conflict_list = ", ".join(flush_result["conflicted"][:5])
            return _typed_unknown(
                f"Local execution completed but workspace sync conflicted for: {conflict_list}",
                "workspace_sync_conflict",
            )
        changed_refs = tuple(
            _workspace_artifact_ref(agent_id, path)
            for path in flush_result["updated"]
        )
        return replace(
            outcome,
            artifact_refs=tuple(
                dict.fromkeys((*outcome.artifact_refs, *changed_refs))
            ),
        )
    finally:
        temp_workspace.cleanup()


async def _execute_workspace_mutation(
    tool_name: str,
    arguments: dict,
    *,
    agent_id: uuid.UUID,
    base_dir: Path,
    session_id: str | None,
) -> str:
    """Handle shared workspace mutations for both direct and normal tool execution."""
    if tool_name == "write_file":
        path = arguments.get("path")
        content = arguments.get("content")
        if not path:
            return "❌ Missing required argument 'path' for write_file. Please provide a file path like 'skills/my-skill/SKILL.md'"
        if content is None:
            return "❌ Missing required argument 'content' for write_file"
        if is_focus_file_path(path):
            return "❌ Focus is no longer stored in focus.md. Use upsert_focus_item or complete_focus_item."
        if _is_enterprise_info_path(path):
            return "❌ enterprise_info is shared company context and is read-only for agents. Ask an admin to update it."
        async with async_session() as _wdb:
            write_result = await write_workspace_file(
                _wdb,
                agent_id=agent_id,
                base_dir=base_dir,
                path=path,
                content=content,
                actor_type="agent",
                actor_id=agent_id,
                operation="write",
                session_id=session_id,
                enforce_human_lock=True,
            )
            await _wdb.commit()
        return (
            f"✅ Written to {write_result.path} ({len(content)} chars)"
            if write_result.ok
            else f"❌ {write_result.message}"
        )

    if tool_name == "move_file":
        source_path = arguments.get("source_path")
        destination_path = arguments.get("destination_path")
        if not source_path:
            return "❌ Missing required argument 'source_path' for move_file"
        if not destination_path:
            return "❌ Missing required argument 'destination_path' for move_file"
        if is_focus_file_path(source_path) or is_focus_file_path(destination_path):
            return "❌ Focus is no longer stored in focus.md. Use Focus tools instead."
        if str(source_path).strip("/") in {"tasks.json", "soul.md"}:
            return f"❌ {source_path} cannot be moved (protected)"
        if _is_enterprise_info_path(source_path) or _is_enterprise_info_path(destination_path):
            return "❌ enterprise_info is shared company context and is read-only for agents. Ask an admin to update it."
        async with async_session() as _wdb:
            move_result = await move_workspace_path(
                _wdb,
                agent_id=agent_id,
                base_dir=base_dir,
                source_path=source_path,
                destination_path=destination_path,
                actor_type="agent",
                actor_id=agent_id,
                session_id=session_id,
                enforce_human_lock=True,
                overwrite=bool(arguments.get("overwrite", False)),
            )
            await _wdb.commit()
        return f"✅ {move_result.message}" if move_result.ok else f"❌ {move_result.message}"

    if tool_name == "delete_file":
        path = arguments.get("path", "")
        if is_focus_file_path(path):
            return "❌ Focus is no longer stored in focus.md. Use Focus tools instead."
        if _is_enterprise_info_path(path):
            return "❌ enterprise_info is shared company context and is read-only for agents. Ask an admin to update it."
        async with async_session() as _wdb:
            delete_result = await delete_workspace_file(
                _wdb,
                agent_id=agent_id,
                base_dir=base_dir,
                path=path,
                actor_type="agent",
                actor_id=agent_id,
                session_id=session_id,
                enforce_human_lock=True,
            )
            await _wdb.commit()
        return f"✅ Deleted {delete_result.path}" if delete_result.ok else f"❌ {delete_result.message}"

    if tool_name == "edit_file":
        path = arguments.get("path")
        old_string = arguments.get("old_string")
        new_string = arguments.get("new_string")
        if not path:
            return "❌ Missing required argument 'path' for edit_file"
        if old_string is None:
            return "❌ Missing required argument 'old_string' for edit_file"
        if new_string is None:
            return "❌ Missing required argument 'new_string' for edit_file"
        if is_focus_file_path(path):
            return "❌ Focus is no longer stored in focus.md. Use upsert_focus_item or complete_focus_item."
        if _is_enterprise_info_path(path):
            return "❌ enterprise_info is shared company context and is read-only for agents. Ask an admin to update it."

        replace_all = arguments.get("replace_all", False)
        storage = get_storage_backend()
        storage_key, normalized_path, _ = _tool_storage_key(agent_id, path, None)
        if not await storage.is_file(storage_key):
            return f"File not found: {path}"

        content = await storage.read_text(storage_key, encoding="utf-8", errors="replace")
        if old_string not in content:
            return f"❌ 'old_string' not found in {path}. Please check the exact text including whitespace and newlines."
        count = content.count(old_string)
        if count > 1 and not replace_all:
            return f"❌ 'old_string' appears {count} times in {path}. Use replace_all=true or provide more context to make the match unique."

        new_content = content.replace(old_string, new_string) if replace_all else content.replace(old_string, new_string, 1)
        async with async_session() as _wdb:
            write_result = await write_workspace_file(
                _wdb,
                agent_id=agent_id,
                base_dir=base_dir,
                path=normalized_path,
                content=new_content,
                actor_type="agent",
                actor_id=agent_id,
                operation="edit",
                session_id=session_id,
                enforce_human_lock=True,
            )
            await _wdb.commit()
        replaced = count if replace_all else 1
        return (
            f"✅ Replaced {replaced} occurrence(s) in {write_result.path}"
            if write_result.ok
            else f"❌ {write_result.message}"
        )

    return f"Tool {tool_name} does not support workspace mutation execution"


def _typed_failure(
    summary: str,
    error_code: str,
    *,
    retryable: bool = False,
    result_ref: str | None = None,
    metadata: dict | None = None,
) -> ToolExecutionOutcome:
    return ToolExecutionOutcome(
        status="failed",
        result_summary=summary,
        result_ref=result_ref,
        error_code=error_code,
        retryable=retryable,
        metadata=metadata or {},
    )


def _typed_success(
    summary: str,
    *,
    result_ref: str | None = None,
    artifact_refs: tuple[str, ...] = (),
    evidence_refs: tuple[str, ...] = (),
    metadata: dict | None = None,
    private_binary: bytes | None = None,
) -> ToolExecutionOutcome:
    return ToolExecutionOutcome(
        status="succeeded",
        result_summary=summary,
        result_ref=result_ref,
        artifact_refs=artifact_refs,
        evidence_refs=evidence_refs,
        metadata=metadata or {},
        private_binary=private_binary,
    )


def _typed_unknown(
    summary: str,
    error_code: str,
    *,
    result_ref: str | None = None,
    metadata: dict | None = None,
) -> ToolExecutionOutcome:
    return ToolExecutionOutcome(
        status="unknown",
        result_summary=summary,
        result_ref=result_ref,
        error_code=error_code,
        metadata=metadata or {},
    )


def _typed_pending(summary: str, *, metadata: dict) -> ToolExecutionOutcome:
    return ToolExecutionOutcome(
        status="pending",
        result_summary=summary,
        result_ref=None,
        metadata=metadata,
    )


def _legacy_tool_outcome_text(
    outcome: ToolExecutionOutcome,
    *,
    fallback: str,
) -> str:
    """Serialize a typed outcome only at a legacy text-consumer boundary."""
    prefix = {
        "succeeded": "✅",
        "failed": "❌",
        "pending": "⏳",
        "unknown": "⚠️",
    }[outcome.status]
    return f"{prefix} {outcome.result_summary or fallback}"


def _propose_experience_draft_outcome(
    arguments: dict,
) -> ToolExecutionOutcome:
    """Validate the human-gated draft without claiming a storage write."""
    for field in ("title", "body", "applicability"):
        value = arguments.get(field)
        if not isinstance(value, str) or not value.strip():
            return _typed_failure(
                "propose_experience_draft requires non-empty title, body, and applicability.",
                "invalid_tool_arguments",
            )
    tags = arguments.get("tags")
    if tags is not None and (
        not isinstance(tags, list)
        or any(not isinstance(tag, str) or not tag.strip() for tag in tags)
    ):
        return _typed_failure(
            "propose_experience_draft tags must be an array of non-empty strings.",
            "invalid_tool_arguments",
        )
    return _typed_success(
        "The structured experience draft is ready for human review. "
        "Nothing was written to the experience library; the user must confirm it."
    )


async def _list_focus_items_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Return a typed Focus read without interpreting display strings."""
    try:
        items = await list_focus_items(
            agent_id,
            include_completed=bool(arguments.get("include_completed", False)),
        )
    except Exception as exc:
        return _typed_failure(
            f"Focus items could not be read: {type(exc).__name__}",
            "focus_read_failed",
            retryable=True,
        )
    if not items:
        return _typed_success("No Focus items.")
    lines = ["Focus items:"]
    for item in items:
        label = "completed" if item["status"] == "completed" else "in_progress"
        kind = f", {item['kind']}" if item.get("kind") == "system" else ""
        title = item.get("title")
        if title:
            lines.append(
                f"- {title} ({item['key']}) [{label}{kind}]: {item['description']}"
            )
        else:
            lines.append(
                f"- {item['key']} [{label}{kind}]: {item['description']}"
            )
    return _typed_success("\n".join(lines))


async def _upsert_focus_item_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    description = (arguments.get("description") or "").strip()
    if not description:
        return _typed_failure(
            "Missing required argument 'description' for upsert_focus_item.",
            "invalid_tool_arguments",
        )
    try:
        item = await upsert_focus_item(
            agent_id,
            key=arguments.get("key"),
            title=arguments.get("title"),
            description=description,
            status="in_progress",
            kind=arguments.get("kind") or "normal",
            source=arguments.get("source") or "user",
            metadata={"tool": "upsert_focus_item"},
        )
    except Exception as exc:
        return _typed_failure(
            f"Focus item could not be saved: {type(exc).__name__}",
            "focus_write_failed",
        )
    title = f" (title: {item['title']})" if item.get("title") else ""
    return _typed_success(
        f"Focus item saved: {item['key']}{title} — {item['description']}"
    )


async def _complete_focus_item_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    key = (arguments.get("key") or "").strip()
    if not key:
        return _typed_failure(
            "Missing required argument 'key' for complete_focus_item.",
            "invalid_tool_arguments",
        )
    try:
        item = await complete_focus_item(agent_id, key=key)
    except Exception as exc:
        return _typed_failure(
            f"Focus item could not be completed: {type(exc).__name__}",
            "focus_write_failed",
        )
    if item is None:
        return _typed_failure(
            f"Focus item not found: {key}",
            "focus_item_not_found",
        )
    return _typed_success(f"Focus item completed: {key}")


async def _read_file_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    tenant_id: str | None,
) -> ToolExecutionOutcome:
    """Read one text file from StorageBackend with explicit typed branches."""
    path = arguments.get("path")
    if not isinstance(path, str) or not path.strip():
        return _typed_failure(
            "Missing required argument 'path' for read_file.",
            "invalid_tool_arguments",
        )
    if is_focus_file_path(path):
        return _typed_failure(
            "Focus is structured data; use list_focus_items.",
            "focus_file_path_removed",
        )
    binary_error = _read_file_binary_error(path)
    if binary_error is not None:
        return _typed_failure(
            binary_error,
            "workspace_binary_file_unsupported",
        )
    try:
        offset = int(arguments.get("offset", 0))
        limit = int(arguments.get("limit", 2000))
    except (TypeError, ValueError):
        return _typed_failure(
            "read_file offset and limit must be integers.",
            "invalid_tool_arguments",
        )
    if offset < 0 or limit <= 0:
        return _typed_failure(
            "read_file offset must be non-negative and limit must be positive.",
            "invalid_tool_arguments",
        )
    storage = get_storage_backend()
    try:
        storage_key, normalized, _ = _tool_storage_key(agent_id, path, tenant_id)
        if not normalized or not await storage.is_file(storage_key):
            return _typed_failure(
                f"File not found: {path}",
                "workspace_file_not_found",
            )
        content = await storage.read_text(
            storage_key,
            encoding="utf-8",
            errors="replace",
        )
    except Exception as exc:
        return _typed_failure(
            f"File read failed: {type(exc).__name__}",
            "workspace_read_failed",
            retryable=True,
        )
    lines = content.splitlines()
    end = min(len(lines), offset + limit)
    if offset >= len(lines) and lines:
        return _typed_failure(
            f"Offset {offset} exceeds file length ({len(lines)} lines total).",
            "workspace_read_offset_invalid",
        )
    selected = "\n".join(
        f"{index + 1:6}\t{line}"
        for index, line in enumerate(lines[offset:end], start=offset)
    )
    if len(lines) > end:
        selected += (
            f"\n\n... [{len(lines) - end} more lines not shown, "
            f"lines {end + 1}-{len(lines)}]"
        )
    return _typed_success(
        f"📄 {path} (lines {offset + 1 if lines else 0}-{end} of {len(lines)})\n"
        f"{selected}"
    )


async def _write_file_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    base_dir: Path,
    session_id: str | None,
) -> ToolExecutionOutcome:
    """Write one workspace file using the structured collaboration result."""
    path = arguments.get("path")
    content = arguments.get("content")
    if not isinstance(path, str) or not path.strip() or content is None:
        return _typed_failure(
            "write_file requires non-empty path and content.",
            "invalid_tool_arguments",
        )
    if not isinstance(content, str):
        return _typed_failure(
            "write_file content must be a string.",
            "invalid_tool_arguments",
        )
    if is_focus_file_path(path):
        return _typed_failure(
            "Focus is structured data; use upsert_focus_item.",
            "focus_file_path_removed",
        )
    if _is_enterprise_info_path(path):
        return _typed_failure(
            "enterprise_info is read-only for Agents.",
            "workspace_path_read_only",
        )
    write_started = False
    try:
        async with async_session() as db:
            write_started = True
            write_result = await write_workspace_file(
                db,
                agent_id=agent_id,
                base_dir=base_dir,
                path=path,
                content=content,
                actor_type="agent",
                actor_id=agent_id,
                operation="write",
                session_id=session_id,
                enforce_human_lock=True,
            )
            if not write_result.ok:
                return _typed_failure(
                    write_result.message,
                    "workspace_write_rejected",
                )
            await db.commit()
    except Exception as exc:
        if write_started:
            return _typed_unknown(
                "Workspace write outcome is unknown; reconcile before retrying.",
                "workspace_write_outcome_unknown",
            )
        return _typed_failure(
            f"Workspace write failed: {type(exc).__name__}",
            "workspace_write_failed",
        )
    return _typed_success(
        f"Written to {write_result.path} ({len(content)} chars)."
    )


async def _list_files_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    tenant_id: str | None,
) -> ToolExecutionOutcome:
    path = arguments.get("path", "")
    if not isinstance(path, str):
        return _typed_failure(
            "list_files path must be a string.",
            "invalid_tool_arguments",
        )
    try:
        storage = get_storage_backend()
        storage_key, normalized, _ = _tool_storage_key(agent_id, path, tenant_id)
        exists = await storage.exists(storage_key)
        is_dir = await storage.is_dir(storage_key)
        if exists and not is_dir:
            return _typed_failure(
                f"Path is not a directory: {path}",
                "workspace_path_not_directory",
            )
        if not exists and not is_dir and normalized:
            return _typed_failure(
                f"Directory not found: {path or '/'}",
                "workspace_directory_not_found",
            )
        summary = await _storage_list_dir(agent_id, path, tenant_id=tenant_id)
    except Exception as exc:
        return _typed_failure(
            f"Directory could not be listed: {type(exc).__name__}.",
            "workspace_list_failed",
            retryable=True,
        )
    return _typed_success(summary)


async def _search_files_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    tenant_id: str | None,
) -> ToolExecutionOutcome:
    pattern = arguments.get("pattern")
    if not isinstance(pattern, str) or not pattern:
        return _typed_failure(
            "search_files requires a non-empty pattern.",
            "invalid_tool_arguments",
        )
    try:
        re.compile(pattern, re.IGNORECASE if arguments.get("ignore_case", False) else 0)
    except re.error as exc:
        return _typed_failure(
            f"Invalid regex pattern: {exc}",
            "invalid_tool_arguments",
        )
    path = arguments.get("path", ".")
    file_pattern = arguments.get("file_pattern", "*")
    if not isinstance(path, str) or not isinstance(file_pattern, str):
        return _typed_failure(
            "search_files path and file_pattern must be strings.",
            "invalid_tool_arguments",
        )
    try:
        storage = get_storage_backend()
        rel_path = "" if path in ("", ".") else path
        base_key, normalized, _ = _tool_storage_key(agent_id, rel_path, tenant_id)
        if normalized and not await storage.is_dir(base_key):
            return _typed_failure(
                f"Directory not found: {path}",
                "workspace_directory_not_found",
            )
        summary = await _storage_search_files(
            agent_id,
            pattern,
            path=path,
            file_pattern=file_pattern,
            ignore_case=bool(arguments.get("ignore_case", False)),
            tenant_id=tenant_id,
        )
    except Exception as exc:
        return _typed_failure(
            f"Workspace search failed: {type(exc).__name__}.",
            "workspace_search_failed",
            retryable=True,
        )
    return _typed_success(summary)


async def _find_files_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    tenant_id: str | None,
) -> ToolExecutionOutcome:
    pattern = arguments.get("pattern")
    path = arguments.get("path", ".")
    if not isinstance(pattern, str) or not pattern or not isinstance(path, str):
        return _typed_failure(
            "find_files requires a non-empty pattern and string path.",
            "invalid_tool_arguments",
        )
    try:
        storage = get_storage_backend()
        rel_path = "" if path in ("", ".") else path
        base_key, normalized, _ = _tool_storage_key(agent_id, rel_path, tenant_id)
        if normalized and not await storage.is_dir(base_key):
            return _typed_failure(
                f"Directory not found: {path}",
                "workspace_directory_not_found",
            )
        summary = await _storage_find_files(
            agent_id,
            pattern,
            path=path,
            tenant_id=tenant_id,
        )
    except Exception as exc:
        return _typed_failure(
            f"Workspace file lookup failed: {type(exc).__name__}.",
            "workspace_find_failed",
            retryable=True,
        )
    return _typed_success(summary)


async def _move_file_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    base_dir: Path,
    session_id: str | None,
) -> ToolExecutionOutcome:
    source_path = arguments.get("source_path")
    destination_path = arguments.get("destination_path")
    if not isinstance(source_path, str) or not source_path or not isinstance(destination_path, str) or not destination_path:
        return _typed_failure(
            "move_file requires source_path and destination_path.",
            "invalid_tool_arguments",
        )
    if is_focus_file_path(source_path) or is_focus_file_path(destination_path):
        return _typed_failure(
            "Focus is structured data and cannot be moved as a file.",
            "focus_file_path_removed",
        )
    if _is_enterprise_info_path(source_path) or _is_enterprise_info_path(destination_path):
        return _typed_failure(
            "enterprise_info is read-only for Agents.",
            "workspace_path_read_only",
        )
    mutation_started = False
    try:
        async with async_session() as db:
            mutation_started = True
            result = await move_workspace_path(
                db,
                agent_id=agent_id,
                base_dir=base_dir,
                source_path=source_path,
                destination_path=destination_path,
                actor_type="agent",
                actor_id=agent_id,
                session_id=session_id,
                enforce_human_lock=True,
                overwrite=bool(arguments.get("overwrite", False)),
            )
            if not result.ok:
                return _typed_failure(result.message, "workspace_move_rejected")
            await db.commit()
    except Exception as exc:
        if mutation_started:
            return _typed_unknown(
                "Workspace move outcome is unknown; reconcile before retrying.",
                "workspace_move_outcome_unknown",
            )
        return _typed_failure(
            f"Workspace move failed: {type(exc).__name__}.",
            "workspace_move_failed",
        )
    return _typed_success(result.message)


async def _delete_file_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    base_dir: Path,
    session_id: str | None,
) -> ToolExecutionOutcome:
    path = arguments.get("path")
    if not isinstance(path, str) or not path:
        return _typed_failure(
            "delete_file requires a non-empty path.",
            "invalid_tool_arguments",
        )
    if is_focus_file_path(path):
        return _typed_failure(
            "Focus is structured data and cannot be deleted as a file.",
            "focus_file_path_removed",
        )
    if _is_enterprise_info_path(path):
        return _typed_failure(
            "enterprise_info is read-only for Agents.",
            "workspace_path_read_only",
        )
    mutation_started = False
    try:
        async with async_session() as db:
            mutation_started = True
            result = await delete_workspace_file(
                db,
                agent_id=agent_id,
                base_dir=base_dir,
                path=path,
                actor_type="agent",
                actor_id=agent_id,
                session_id=session_id,
                enforce_human_lock=True,
            )
            if not result.ok:
                return _typed_failure(result.message, "workspace_delete_rejected")
            await db.commit()
    except Exception as exc:
        if mutation_started:
            return _typed_unknown(
                "Workspace delete outcome is unknown; reconcile before retrying.",
                "workspace_delete_outcome_unknown",
            )
        return _typed_failure(
            f"Workspace delete failed: {type(exc).__name__}.",
            "workspace_delete_failed",
        )
    return _typed_success(result.message)


async def _edit_file_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    base_dir: Path,
    session_id: str | None,
) -> ToolExecutionOutcome:
    path = arguments.get("path")
    old_string = arguments.get("old_string")
    new_string = arguments.get("new_string")
    if not isinstance(path, str) or not path or not isinstance(old_string, str) or not isinstance(new_string, str):
        return _typed_failure(
            "edit_file requires string path, old_string, and new_string.",
            "invalid_tool_arguments",
        )
    if is_focus_file_path(path):
        return _typed_failure(
            "Focus is structured data and cannot be edited as a file.",
            "focus_file_path_removed",
        )
    if _is_enterprise_info_path(path):
        return _typed_failure(
            "enterprise_info is read-only for Agents.",
            "workspace_path_read_only",
        )
    try:
        storage = get_storage_backend()
        storage_key, normalized_path, _ = _tool_storage_key(agent_id, path, None)
        if not await storage.is_file(storage_key):
            return _typed_failure(
                f"File not found: {path}",
                "workspace_file_not_found",
            )
        content = await storage.read_text(storage_key, encoding="utf-8", errors="replace")
    except Exception as exc:
        return _typed_failure(
            f"Workspace file could not be read for editing: {type(exc).__name__}.",
            "workspace_read_failed",
            retryable=True,
        )
    count = content.count(old_string)
    replace_all = bool(arguments.get("replace_all", False))
    if count == 0:
        return _typed_failure(
            f"old_string was not found in {path}.",
            "workspace_edit_text_not_found",
        )
    if count > 1 and not replace_all:
        return _typed_failure(
            f"old_string appears {count} times in {path}; provide a unique match or set replace_all.",
            "workspace_edit_text_ambiguous",
        )
    new_content = (
        content.replace(old_string, new_string)
        if replace_all
        else content.replace(old_string, new_string, 1)
    )
    mutation_started = False
    try:
        async with async_session() as db:
            mutation_started = True
            result = await write_workspace_file(
                db,
                agent_id=agent_id,
                base_dir=base_dir,
                path=normalized_path,
                content=new_content,
                actor_type="agent",
                actor_id=agent_id,
                operation="edit",
                session_id=session_id,
                enforce_human_lock=True,
            )
            if not result.ok:
                return _typed_failure(result.message, "workspace_edit_rejected")
            await db.commit()
    except Exception as exc:
        if mutation_started:
            return _typed_unknown(
                "Workspace edit outcome is unknown; reconcile before retrying.",
                "workspace_edit_outcome_unknown",
            )
        return _typed_failure(
            f"Workspace edit failed: {type(exc).__name__}.",
            "workspace_edit_failed",
        )
    replaced = count if replace_all else 1
    return _typed_success(
        f"Replaced {replaced} occurrence(s) in {result.path}."
    )


async def execute_builtin_tool_outcome(
    tool_name: str,
    arguments: dict,
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    session_id: str = "",
    on_output=None,
) -> ToolExecutionOutcome | str:
    """Execute only explicitly migrated builtin branches as typed outcomes.

    Unmigrated builtin and dynamic handlers intentionally remain strings.  The
    Durable Runtime rejects those as ``untyped_tool_outcome``; this function
    never infers success from display text or from a non-raising handler.
    """
    tenant_id: str | None = None
    if tool_name in {
        "list_files",
        "read_file",
        "search_files",
        "find_files",
        "read_document",
        "execute_code",
        "execute_code_e2b",
        "convert_csv_to_xlsx",
        "convert_html_to_pdf",
        "convert_html_to_pptx",
        "convert_markdown_to_docx",
        "convert_markdown_to_pdf",
        "upload_image",
        *_IMAGE_GENERATION_TOOL_NAMES,
    }:
        tenant_id = await _get_agent_tenant_id(agent_id)
    if tool_name == "list_files":
        return await _list_files_outcome(
            agent_id,
            arguments,
            tenant_id=tenant_id,
        )
    if tool_name == "list_focus_items":
        return await _list_focus_items_outcome(agent_id, arguments)
    if tool_name == "upsert_focus_item":
        return await _upsert_focus_item_outcome(agent_id, arguments)
    if tool_name == "complete_focus_item":
        return await _complete_focus_item_outcome(agent_id, arguments)
    if tool_name == "read_file":
        return await _read_file_outcome(
            agent_id,
            arguments,
            tenant_id=tenant_id,
        )
    if tool_name == "search_files":
        return await _search_files_outcome(
            agent_id,
            arguments,
            tenant_id=tenant_id,
        )
    if tool_name == "find_files":
        return await _find_files_outcome(
            agent_id,
            arguments,
            tenant_id=tenant_id,
        )
    if tool_name == "read_document":
        return await _read_document_outcome(
            agent_id,
            arguments,
            tenant_id=tenant_id,
        )
    if tool_name == "write_file":
        return await _write_file_outcome(
            agent_id,
            arguments,
            base_dir=_agent_workspace_root(agent_id),
            session_id=session_id or None,
        )
    if tool_name == "move_file":
        return await _move_file_outcome(
            agent_id,
            arguments,
            base_dir=_agent_workspace_root(agent_id),
            session_id=session_id or None,
        )
    if tool_name == "delete_file":
        return await _delete_file_outcome(
            agent_id,
            arguments,
            base_dir=_agent_workspace_root(agent_id),
            session_id=session_id or None,
        )
    if tool_name == "edit_file":
        return await _edit_file_outcome(
            agent_id,
            arguments,
            base_dir=_agent_workspace_root(agent_id),
            session_id=session_id or None,
        )
    if tool_name in {
        "convert_csv_to_xlsx",
        "convert_html_to_pdf",
        "convert_html_to_pptx",
        "convert_markdown_to_docx",
        "convert_markdown_to_pdf",
    }:
        return await _run_with_temp_workspace_outcome(
            agent_id,
            tenant_id,
            lambda temp_ws: _convert_file_outcome(
                agent_id,
                temp_ws,
                arguments,
                tool_name=tool_name,
            ),
            paths=_non_empty_paths(
                arguments.get("source_path"),
                arguments.get("target_path"),
            ),
            sync_back=True,
        )
    if tool_name in {"execute_code", "execute_code_e2b"}:
        return await _run_with_temp_workspace_outcome(
            agent_id,
            tenant_id,
            lambda temp_ws: _execute_code_outcome(
                agent_id,
                temp_ws,
                arguments,
                tool_name=tool_name,
                on_output=on_output,
            ),
            sync_back=True,
            sync_back_on_non_success=True,
        )
    if tool_name == "read_webpage":
        return await _read_webpage_outcome(arguments)
    if tool_name == "upload_image":
        file_path = arguments.get("file_path")
        return await _run_with_temp_workspace_outcome(
            agent_id,
            tenant_id,
            lambda temp_ws: _upload_image_outcome(
                agent_id,
                temp_ws,
                arguments,
            ),
            paths=_non_empty_paths(file_path),
        )
    if tool_name in _IMAGE_GENERATION_TOOL_NAMES:
        return await _run_with_temp_workspace_outcome(
            agent_id,
            tenant_id,
            lambda temp_ws: _generate_image_outcome(
                agent_id,
                temp_ws,
                arguments,
                _IMAGE_GENERATION_PROVIDER_BY_TOOL[tool_name],
            ),
            sync_back=True,
        )
    if tool_name == "publish_page":
        return await _publish_page_outcome(
            agent_id,
            user_id,
            _agent_workspace_root(agent_id),
            arguments,
        )
    if tool_name == "list_published_pages":
        return await _list_published_pages_outcome(agent_id)
    if tool_name == "set_trigger":
        return await _handle_set_trigger_outcome(
            agent_id,
            arguments,
            session_id=session_id,
            user_id=user_id,
        )
    if tool_name == "send_channel_file":
        file_path = arguments.get("file_path")
        if not isinstance(file_path, str) or not file_path.strip():
            return _typed_failure(
                "send_channel_file requires file_path.",
                "invalid_tool_arguments",
            )
        tenant_id = await _get_agent_tenant_id(agent_id)
        return await _run_with_temp_workspace_outcome(
            agent_id,
            tenant_id,
            lambda temp_ws: _send_channel_file_outcome(
                agent_id,
                temp_ws,
                arguments,
            ),
            paths=[file_path],
        )
    if tool_name == "send_file_to_agent":
        return await _send_file_to_agent_outcome(agent_id, arguments)
    if tool_name == "duckduckgo_search":
        return await _duckduckgo_search_outcome(arguments)
    if tool_name == "web_search":
        return await _web_search_outcome(arguments, agent_id)
    if tool_name == "jina_search":
        return await _jina_search_outcome(arguments, agent_id)
    if tool_name == "jina_read":
        return await _jina_read_outcome(arguments, agent_id)
    if tool_name == "exa_search":
        return await _exa_search_outcome(arguments, agent_id)
    if tool_name == "tavily_search":
        return await _tavily_search_outcome(arguments, agent_id)
    if tool_name == "google_search":
        return await _google_search_outcome(arguments, agent_id)
    if tool_name == "bing_search":
        return await _bing_search_outcome(arguments, agent_id)
    if tool_name == "search_experience":
        from app.services.experience_retrieval import search_experience_outcome

        return await search_experience_outcome(agent_id, arguments)
    if tool_name == "read_experience":
        from app.services.experience_retrieval import read_experience_outcome

        return await read_experience_outcome(agent_id, arguments)
    if tool_name == "propose_experience_draft":
        return _propose_experience_draft_outcome(arguments)
    if tool_name == "discover_resources":
        return await _discover_resources_outcome(agent_id, arguments)
    if tool_name == "import_mcp_server":
        return await _import_mcp_server_outcome(agent_id, arguments)
    if tool_name in _VERCEL_READ_TOOL_NAMES:
        return await _vercel_read_outcome(tool_name, agent_id, arguments)
    if tool_name == "vercel_deploy":
        return await _vercel_deploy_outcome(
            agent_id,
            _agent_workspace_root(agent_id),
            arguments,
        )
    if tool_name in _DEPLOY_SIMPLE_WRITE_TOOL_NAMES:
        return await _deploy_simple_write_outcome(
            tool_name,
            agent_id,
            arguments,
        )
    if tool_name in _AGENTBAY_A1_READ_TOOL_NAMES:
        return await _agentbay_read_outcome(
            tool_name,
            agent_id,
            arguments,
            session_id=session_id,
        )
    if tool_name in _OKR_TRANSACTION_TOOL_NAMES:
        return await _okr_transaction_outcome(
            tool_name,
            agent_id,
            user_id,
            arguments,
        )
    if tool_name in _OKR_JOB_TOOL_NAMES:
        return await _okr_job_outcome(
            tool_name,
            agent_id,
            arguments,
        )
    if tool_name == "search_clawhub":
        return await _search_clawhub_outcome(agent_id, arguments)
    if tool_name == "install_skill":
        source = arguments.get("source")
        if not isinstance(source, str) or not source.strip():
            return _typed_failure(
                "install_skill requires source.",
                "invalid_tool_arguments",
            )
        tenant_id = await _get_agent_tenant_id(agent_id)
        return await _run_with_temp_workspace_outcome(
            agent_id,
            tenant_id,
            lambda temp_ws: _install_skill_outcome(
                agent_id,
                temp_ws,
                arguments,
            ),
            paths=["skills"],
            sync_back=True,
        )
    if tool_name == "send_channel_message":
        return await _send_channel_message_outcome(agent_id, arguments)
    if tool_name == "send_platform_message":
        return await _send_platform_message_outcome(agent_id, arguments)
    if tool_name == "query_directory":
        return await _query_directory_outcome(agent_id, arguments)
    if tool_name == "update_trigger":
        return await _handle_update_trigger_outcome(agent_id, arguments)
    if tool_name == "cancel_trigger":
        return await _handle_cancel_trigger_outcome(agent_id, arguments)
    if tool_name == "list_triggers":
        return await _handle_list_triggers_outcome(agent_id)
    if tool_name == "read_emails":
        return await _read_emails_outcome(agent_id, arguments)
    if tool_name in {"send_email", "reply_email"}:
        return await _email_write_outcome(tool_name, agent_id, arguments)
    if tool_name == "feishu_calendar_list":
        return await _feishu_calendar_list_outcome(agent_id, arguments)
    if tool_name == "feishu_calendar_create":
        return await _feishu_calendar_create_outcome(agent_id, arguments)
    if tool_name in {"feishu_calendar_update", "feishu_calendar_delete"}:
        return await _feishu_calendar_mutation_outcome(
            tool_name,
            agent_id,
            arguments,
        )
    if tool_name == "feishu_wiki_list":
        return await _feishu_wiki_list_outcome(agent_id, arguments)
    if tool_name == "feishu_doc_search":
        return await _feishu_doc_search_outcome(agent_id, arguments)
    if tool_name == "feishu_doc_read":
        return await _feishu_doc_read_outcome(agent_id, arguments)
    if tool_name == "feishu_doc_create":
        return await _feishu_doc_create_outcome(agent_id, arguments)
    if tool_name == "feishu_doc_append":
        return await _feishu_doc_append_outcome(agent_id, arguments)
    if tool_name == "feishu_drive_share":
        return await _feishu_drive_share_outcome(agent_id, arguments)
    if tool_name == "feishu_drive_delete":
        return await _feishu_drive_delete_outcome(agent_id, arguments)
    if tool_name == "feishu_user_search":
        return await _feishu_user_search_outcome(agent_id, arguments)
    if tool_name == "feishu_approval_query":
        return await _feishu_approval_query_outcome(agent_id, arguments)
    if tool_name == "feishu_approval_get":
        return await _feishu_approval_get_outcome(agent_id, arguments)
    if tool_name in {
        "bitable_list_tables",
        "bitable_list_fields",
        "bitable_query_records",
    }:
        return await _bitable_read_outcome(tool_name, agent_id, arguments)
    if tool_name in {
        "bitable_create_app",
        "bitable_create_record",
        "bitable_update_record",
        "bitable_delete_record",
    }:
        return await _bitable_write_outcome(tool_name, agent_id, arguments)

    # Dynamic MCP tools are not members of the canonical builtin registry.
    # Resolve an exact, enabled AgentTool assignment before selecting their
    # typed adapter.  A name with no MCP row remains on the legacy untyped path
    # so arbitrary custom handlers are never promoted from display text.
    if (
        agent_id is not None
        and tool_name not in BUILTIN_TOOL_NAMES
        and not is_reserved_custom_tool_name(tool_name)
    ):
        mcp_target = await _resolve_mcp_execution_target(tool_name, agent_id)
        if mcp_target is not None:
            return await _execute_resolved_mcp_target_outcome(
                mcp_target,
                arguments,
                agent_id=agent_id,
            )
    return await execute_tool(
        tool_name,
        arguments,
        agent_id,
        user_id,
        session_id,
        on_output,
    )


async def _execute_tool_direct(
    tool_name: str,
    arguments: dict,
    agent_id: uuid.UUID,
) -> str:
    """Execute a tool directly, bypassing autonomy checks.

    Used by the approval post-processing hook after an action
    has been approved and needs to actually run.
    """
    _agent_tenant_id = await _get_agent_tenant_id(agent_id)
    ws = _agent_workspace_root(agent_id)
    try:
        if tool_name in {"delete_file", "write_file", "move_file", "edit_file"}:
            return await _execute_workspace_mutation(
                tool_name,
                arguments,
                agent_id=agent_id,
                base_dir=ws,
                session_id=None,
            )
        elif tool_name in ("execute_code", "execute_code_e2b"):
            logger.info(
                "[DirectTool] Executing code ({}) with arguments: {}",
                tool_name,
                _observability_arguments(tool_name, arguments),
            )
            return await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _execute_code(agent_id, temp_ws, arguments, tool_name=tool_name),
                sync_back=True,
            )
        elif tool_name == "web_search":
            return await _web_search(arguments, agent_id)
        elif tool_name == "jina_search":
            return await _jina_search(arguments, agent_id)
        elif tool_name == "read_webpage":
            return await _read_webpage(arguments)
        elif tool_name == "exa_search":
            return await _exa_search(arguments, agent_id)
        elif tool_name == "duckduckgo_search":
            return await _duckduckgo_search_tool(arguments)
        elif tool_name == "tavily_search":
            return await _tavily_search_tool(arguments, agent_id)
        elif tool_name == "google_search":
            return await _google_search_tool(arguments, agent_id)
        elif tool_name == "bing_search":
            return await _bing_search_tool(arguments, agent_id)
        elif tool_name == "send_feishu_message":
            return await _send_feishu_message(agent_id, arguments)
        elif tool_name == "query_directory":
            return await _query_directory(agent_id, arguments)
        elif tool_name == "send_message_to_agent":
            return await _send_message_to_agent(
                agent_id,
                arguments,
                user_id=None,
                origin_session_id=None,
            )
        elif tool_name == "send_file_to_agent":
            return await _send_file_to_agent(agent_id, arguments)
        else:
            return f"Tool {tool_name} does not support post-approval execution"
    except Exception as e:
        logger.exception(f"[DirectTool] Error executing {tool_name}: {e}")
        return f"Error executing {tool_name}: {e}"


async def execute_tool(
    tool_name: str,
    arguments: dict,
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    session_id: str = "",
    on_output=None,
) -> str:
    """Execute a tool call and return the result as a string.

    Args:
        session_id: The ChatSession ID, used to isolate AgentBay instances
                    per conversation. Passed through to agentbay_* tools.
    """
    if not isinstance(tool_name, str):
        tool_name = str(tool_name or "")
    tool_name = (
        tool_name
        .replace("`", "")
        .replace("\u200b", "")
        .replace("\u200c", "")
        .replace("\u200d", "")
        .replace("\ufeff", "")
        .strip()
    )
    if tool_name == FINISH_TOOL_NAME:
        content = arguments.get("content", "")
        return content if isinstance(content, str) else str(content)

    _agent_tenant_id = await _get_agent_tenant_id(agent_id)

    ws = _agent_workspace_root(agent_id)

    # ── Autonomy boundary check ──
    action_type = _TOOL_AUTONOMY_MAP.get(tool_name)
    if action_type:
        try:
            from app.services.autonomy_service import autonomy_service
            from app.models.agent import Agent as AgentModel
            async with async_session() as _adb:
                _ar = await _adb.execute(select(AgentModel).where(AgentModel.id == agent_id))
                _agent = _ar.scalar_one_or_none()
                if _agent:
                    result_check = await autonomy_service.check_and_enforce(
                        _adb,
                        _agent,
                        action_type,
                        {
                            "tool": tool_name,
                            "args": str(
                                _observability_arguments(tool_name, arguments)
                            )[:200],
                            "requested_by": str(user_id),
                        },
                    )
                    await _adb.commit()
                    if not result_check.get("allowed"):
                        level = result_check.get("level", "L3")
                        logger.info(f"[Autonomy] Tool {tool_name} denied, level: {level}")
                        if level == "L3":
                            return f"⏳ This action requires approval. An approval request has been sent. Please wait for approval before retrying. (Approval ID: {result_check.get('approval_id', 'N/A')})"
                        return f"❌ Action denied: {result_check.get('message', 'unknown reason')}"
        except Exception as e:
            logger.exception(f"[Autonomy] Check failed: {e}")
            return f"⚠️ Autonomy check failed ({e}). Operation blocked for safety. Please retry or contact admin."

    agentbay_scope_token = None
    if tool_name.startswith("agentbay_"):
        # Take Control lock: block automatic tool execution while a human
        # is manually controlling the browser/desktop session. This prevents
        # input collisions between human clicks and agent-initiated actions.
        from app.api.agentbay_control import is_session_locked
        if is_session_locked(str(agent_id), session_id):
            return (
                "⏸️ A human operator is currently controlling this browser session "
                "(Take Control mode). Please wait for them to finish before retrying "
                "browser/computer operations."
            )
        # Keep execution identity out of durable/model arguments. A private
        # copy also prevents legacy handlers from mutating the caller's input.
        arguments = deepcopy(arguments)
        agentbay_scope_token = agentbay_session_scope_id.set(session_id)

    try:
        if tool_name == "list_files":
            result = await _storage_list_dir(agent_id, arguments.get("path", ""), tenant_id=_agent_tenant_id)
        elif tool_name == "list_focus_items":
            items = await list_focus_items(agent_id, include_completed=bool(arguments.get("include_completed", False)))
            if not items:
                result = "No Focus items."
            else:
                lines = ["Focus items:"]
                for item in items:
                    label = "completed" if item["status"] == "completed" else "in_progress"
                    kind = f", {item['kind']}" if item.get("kind") == "system" else ""
                    if item.get("title"):
                        lines.append(f"- {item['title']} ({item['key']}) [{label}{kind}]: {item['description']}")
                    else:
                        lines.append(f"- {item['key']} [{label}{kind}]: {item['description']}")
                result = "\n".join(lines)
        elif tool_name == "upsert_focus_item":
            description = (arguments.get("description") or "").strip()
            if not description:
                return "❌ Missing required argument 'description' for upsert_focus_item"
            item = await upsert_focus_item(
                agent_id,
                key=arguments.get("key"),
                title=arguments.get("title"),
                description=description,
                status="in_progress",
                kind=arguments.get("kind") or "normal",
                source=arguments.get("source") or "user",
                metadata={"tool": "upsert_focus_item"},
            )
            result = f"✅ Focus item saved: {item['key']} (title: {item['title']}) — {item['description']}" if item.get("title") else f"✅ Focus item saved: {item['key']} — {item['description']}"
        elif tool_name == "complete_focus_item":
            key = (arguments.get("key") or "").strip()
            if not key:
                return "❌ Missing required argument 'key' for complete_focus_item"
            item = await complete_focus_item(agent_id, key=key)
            result = f"✅ Focus item completed: {key}" if item else f"❌ Focus item not found: {key}"
        elif tool_name == "read_file":
            path = arguments.get("path")
            if not path:
                return "❌ Missing required argument 'path' for read_file"
            if is_focus_file_path(path):
                return "❌ Focus is no longer stored in focus.md. Use list_focus_items, upsert_focus_item, and complete_focus_item."
            offset = int(arguments.get("offset", 0))
            limit = int(arguments.get("limit", 2000))
            result = await _storage_read_file(agent_id, path, tenant_id=_agent_tenant_id, offset=offset, limit=limit)
        elif tool_name == "read_document":
            path = arguments.get("path")
            if not path:
                return "❌ Missing required argument 'path' for read_document"
            max_chars = min(int(arguments.get("max_chars", 8000)), 20000)
            result = await _read_document_from_storage(agent_id, path, max_chars=max_chars, tenant_id=_agent_tenant_id)
        elif tool_name in {"write_file", "move_file", "delete_file", "edit_file"}:
            result = await _execute_workspace_mutation(
                tool_name,
                arguments,
                agent_id=agent_id,
                base_dir=ws,
                session_id=session_id,
            )
        # --- Enhanced file management tools ---
        elif tool_name == "convert_csv_to_xlsx":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _convert_csv_to_xlsx(agent_id, temp_ws, arguments),
                paths=_non_empty_paths(arguments.get("source_path", ""), arguments.get("target_path", "")),
                sync_back=True,
            )
        elif tool_name == "convert_html_to_pdf":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _convert_html_to_pdf(agent_id, temp_ws, arguments),
                paths=_non_empty_paths(arguments.get("source_path", ""), arguments.get("target_path", "")),
                sync_back=True,
            )
        elif tool_name == "convert_html_to_pptx":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _convert_html_to_pptx(agent_id, temp_ws, arguments),
                paths=_non_empty_paths(arguments.get("source_path", ""), arguments.get("target_path", "")),
                sync_back=True,
            )
        elif tool_name == "convert_markdown_to_docx":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _convert_markdown_to_docx(agent_id, temp_ws, arguments),
                paths=_non_empty_paths(arguments.get("source_path", ""), arguments.get("target_path", "")),
                sync_back=True,
            )
        elif tool_name == "convert_markdown_to_pdf":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _convert_markdown_to_pdf(agent_id, temp_ws, arguments),
                paths=_non_empty_paths(arguments.get("source_path", ""), arguments.get("target_path", "")),
                sync_back=True,
            )
        elif tool_name == "search_files":
            pattern = arguments.get("pattern")
            if not pattern:
                return "❌ Missing required argument 'pattern' for search_files"
            result = await _storage_search_files(
                agent_id,
                pattern,
                path=arguments.get("path", "."),
                file_pattern=arguments.get("file_pattern", "*"),
                ignore_case=arguments.get("ignore_case", False),
                tenant_id=_agent_tenant_id
            )
        elif tool_name == "find_files":
            pattern = arguments.get("pattern")
            if not pattern:
                return "❌ Missing required argument 'pattern' for find_files"
            result = await _storage_find_files(
                agent_id,
                pattern,
                path=arguments.get("path", "."),
                tenant_id=_agent_tenant_id
            )
        elif tool_name == "manage_tasks":
            result = await _manage_tasks(agent_id, user_id, ws, arguments)
        elif tool_name == "set_trigger":
            result = await _handle_set_trigger(
                agent_id,
                arguments,
                session_id=session_id,
                user_id=user_id,
            )
        elif tool_name == "update_trigger":
            result = await _handle_update_trigger(agent_id, arguments)
        elif tool_name == "cancel_trigger":
            result = await _handle_cancel_trigger(agent_id, arguments)
        elif tool_name == "list_triggers":
            result = await _handle_list_triggers(agent_id)
        elif tool_name == "query_directory":
            result = await _query_directory(agent_id, arguments)
        elif tool_name == "send_feishu_message":
            result = await _send_feishu_message(agent_id, arguments)
        elif tool_name == "send_platform_message":
            result = await _send_platform_message(agent_id, arguments)
        elif tool_name == "send_channel_message":
            result = await _send_channel_message(agent_id, arguments)
        elif tool_name == "send_message_to_agent":
            result = await _send_message_to_agent(
                agent_id,
                arguments,
                user_id=user_id,
                origin_session_id=session_id,
            )
        elif tool_name == "send_file_to_agent":
            result = await _send_file_to_agent(agent_id, arguments)
        elif tool_name == "send_channel_file":
            file_path = (arguments.get("file_path") or "").strip()
            if not file_path:
                result = "Error: file_path is required"
            else:
                result = await _run_with_temp_workspace(
                    agent_id,
                    _agent_tenant_id,
                    lambda temp_ws: _send_channel_file(agent_id, temp_ws, arguments),
                    paths=[file_path],
                )
        elif tool_name == "web_search":
            result = await _web_search(arguments, agent_id)
        elif tool_name == "jina_search":
            result = await _jina_search(arguments, agent_id)
        elif tool_name == "exa_search":
            result = await _exa_search(arguments, agent_id)
        elif tool_name == "duckduckgo_search":
            result = await _duckduckgo_search_tool(arguments)
        elif tool_name == "tavily_search":
            result = await _tavily_search_tool(arguments, agent_id)
        elif tool_name == "google_search":
            result = await _google_search_tool(arguments, agent_id)
        elif tool_name == "bing_search":
            result = await _bing_search_tool(arguments, agent_id)
        elif tool_name == "jina_read":
            result = await _jina_read(arguments, agent_id)
        elif tool_name == "read_webpage":
            result = await _read_webpage(arguments)
        elif tool_name in ("plaza_get_new_posts", "plaza_create_post", "plaza_add_comment"):
            # Deprecated: Plaza social feed replaced by the human-curated experience library.
            result = "[DISABLED] Plaza is now a human-curated experience library. Agents no longer post; contribute via the human-led distillation flow instead."
        elif tool_name == "search_experience":
            from app.services.experience_retrieval import search_experience
            result = await search_experience(agent_id, arguments)
        elif tool_name == "read_experience":
            from app.services.experience_retrieval import read_experience
            result = await read_experience(agent_id, arguments)
        elif tool_name == "propose_experience_draft":
            # No-op by design: writes nothing. The structured args are rendered as a
            # human-gated review card in the UI; a row is created only if the human confirms.
            result = (
                "[已呈现草稿] 已把这条经验的结构化草稿展示给用户，等待其点击『沉淀为经验』人工确认后入库。"
                "本工具未写入任何存储；请如实告诉用户你无法直接入库、需要他确认。"
            )
        elif tool_name in ("execute_code", "execute_code_e2b"):
            logger.info(
                "[DirectTool] Executing code ({}) with arguments: {}",
                tool_name,
                _observability_arguments(tool_name, arguments),
            )
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _execute_code(agent_id, temp_ws, arguments, tool_name=tool_name, on_output=on_output),
                sync_back=True,
            )
        elif tool_name == "upload_image":
            file_path = (arguments.get("file_path") or "").strip()
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _upload_image(agent_id, temp_ws, arguments),
                paths=_non_empty_paths(file_path),
            )
        elif tool_name == "generate_image_siliconflow":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _generate_image(agent_id, temp_ws, arguments, "siliconflow"),
                sync_back=True,
            )
        elif tool_name == "generate_image_openai":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _generate_image(agent_id, temp_ws, arguments, "openai"),
                sync_back=True,
            )
        elif tool_name == "generate_image_google":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _generate_image(agent_id, temp_ws, arguments, "google"),
                sync_back=True,
            )
        elif tool_name == "generate_image_custom":
            result = await _run_with_temp_workspace(
                agent_id,
                _agent_tenant_id,
                lambda temp_ws: _generate_image(agent_id, temp_ws, arguments, "custom"),
                sync_back=True,
            )
        elif tool_name == "discover_resources":
            result = await _discover_resources(agent_id, arguments)
        elif tool_name == "import_mcp_server":
            result = await _import_mcp_server(agent_id, arguments)
        # ── Feishu Bitable Tools ──
        elif tool_name == "bitable_create_app":
            result = await _bitable_create_app(agent_id, arguments)
        elif tool_name == "bitable_list_tables":
            result = await _bitable_list_tables(agent_id, arguments)
        elif tool_name == "bitable_list_fields":
            result = await _bitable_list_fields(agent_id, arguments)
        elif tool_name == "bitable_query_records":
            result = await _bitable_query_records(agent_id, arguments)
        elif tool_name == "bitable_create_record":
            result = await _bitable_create_record(agent_id, arguments)
        elif tool_name == "bitable_update_record":
            result = await _bitable_update_record(agent_id, arguments)
        elif tool_name == "bitable_delete_record":
            result = await _bitable_delete_record(agent_id, arguments)
        # ── Feishu Document Tools ──
        elif tool_name == "feishu_doc_search":
            result = await _feishu_doc_search(agent_id, arguments)
        elif tool_name == "feishu_wiki_list":
            result = await _feishu_wiki_list(agent_id, arguments)
        elif tool_name == "feishu_doc_read":
            result = await _feishu_doc_read(agent_id, arguments)
        elif tool_name == "feishu_doc_create":
            result = await _feishu_doc_create(agent_id, arguments)
        elif tool_name == "feishu_doc_append":
            result = await _feishu_doc_append(agent_id, arguments)
        # ── Feishu Calendar Tools ──
        elif tool_name == "feishu_drive_share":
            result = await _feishu_drive_share(agent_id, arguments)
        elif tool_name == "feishu_drive_delete":
            result = await _feishu_drive_delete(agent_id, arguments)
        elif tool_name == "feishu_user_search":
            result = await _feishu_user_search(agent_id, arguments)
        elif tool_name == "feishu_calendar_list":
            result = await _feishu_calendar_list(agent_id, arguments)
        elif tool_name == "feishu_calendar_create":
            result = await _feishu_calendar_create(agent_id, arguments)
        elif tool_name == "feishu_calendar_update":
            result = await _feishu_calendar_update(agent_id, arguments)
        elif tool_name == "feishu_calendar_delete":
            result = await _feishu_calendar_delete(agent_id, arguments)
        elif tool_name == "feishu_approval_create":
            result = await _feishu_approval_create(agent_id, arguments)
        elif tool_name == "feishu_approval_query":
            result = await _feishu_approval_query(agent_id, arguments)
        elif tool_name == "feishu_approval_get":
            result = await _feishu_approval_get(agent_id, arguments)
        # ── Email Tools ──
        elif tool_name in ("send_email", "read_emails", "reply_email"):
            result = await _handle_email_tool(tool_name, agent_id, ws, arguments)
        # ── Pages: public HTML hosting ──
        elif tool_name == "publish_page":
            result = await _publish_page(agent_id, user_id, ws, arguments)
        elif tool_name == "list_published_pages":
            result = await _list_published_pages(agent_id)
        # ── AgentBay Tools ──
        elif tool_name == "agentbay_browser_navigate":
            result = await _agentbay_browser_navigate(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_screenshot":
            result = await _agentbay_browser_screenshot(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_save_screenshot":
            result = await _agentbay_browser_save_screenshot(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_click":
            result = await _agentbay_browser_click(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_type":
            result = await _agentbay_browser_type(agent_id, ws, arguments)
        elif tool_name == "agentbay_code_execute":
            result = await _agentbay_code_execute(agent_id, ws, arguments)
        elif tool_name == "agentbay_code_write_file":
            result = await _agentbay_code_write_file(agent_id, ws, arguments)
        elif tool_name == "agentbay_code_read_file":
            result = await _agentbay_code_read_file(agent_id, ws, arguments)
        elif tool_name == "agentbay_code_edit_file":
            result = await _agentbay_code_edit_file(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_extract":
            result = await _agentbay_browser_extract(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_observe":
            result = await _agentbay_browser_observe(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_login":
            result = await _agentbay_browser_login(agent_id, ws, arguments)
        elif tool_name == "agentbay_command_exec":
            result = await _agentbay_command_exec(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_screenshot":
            result = await _agentbay_computer_screenshot(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_save_screenshot":
            result = await _agentbay_computer_save_screenshot(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_precision_screenshot":
            result = await _agentbay_computer_precision_screenshot(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_click":
            result = await _agentbay_computer_click(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_input_text":
            result = await _agentbay_computer_input_text(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_press_keys":
            result = await _agentbay_computer_press_keys(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_scroll":
            result = await _agentbay_computer_scroll(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_move_mouse":
            result = await _agentbay_computer_move_mouse(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_drag_mouse":
            result = await _agentbay_computer_drag_mouse(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_get_screen_size":
            result = await _agentbay_computer_get_screen_size(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_start_app":
            result = await _agentbay_computer_start_app(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_get_installed_apps":
            result = await _agentbay_computer_get_installed_apps(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_get_cursor_position":
            result = await _agentbay_computer_get_cursor_position(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_get_active_window":
            result = await _agentbay_computer_get_active_window(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_list_windows":
            result = await _agentbay_computer_list_windows(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_activate_window":
            result = await _agentbay_computer_activate_window(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_close_window":
            result = await _agentbay_computer_close_window(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_dismiss_dialog":
            result = await _agentbay_computer_dismiss_dialog(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_list_visible_apps":
            result = await _agentbay_computer_list_visible_apps(agent_id, ws, arguments)
        elif tool_name == "agentbay_file_transfer":
            result = await _agentbay_file_transfer(agent_id, ws, arguments)
        # ── Skill Management ──
        elif tool_name == "search_clawhub":
            result = await _search_clawhub(agent_id, arguments)
        elif tool_name == "install_skill":
            result = await _install_skill(agent_id, ws, arguments)
        # ── OKR Tools ──
        elif tool_name == "get_okr":
            result = await _get_okr(agent_id, arguments)
        elif tool_name == "get_my_okr":
            result = await _get_my_okr(agent_id, arguments)
        elif tool_name == "update_kr_content":
            result = await _update_kr_content(agent_id, user_id, arguments)
        elif tool_name == "update_kr_progress":
            result = await _update_kr_progress(agent_id, user_id, arguments)
        # collect_okr_progress: legacy batch progress collection
        elif tool_name == "collect_okr_progress":
            result = await _collect_okr_progress(agent_id)
        # generate_okr_report: build daily/weekly structured report and store it
        elif tool_name == "generate_okr_report":
            result = await _generate_okr_report(agent_id, arguments)
        # get_okr_settings: read tenant OKR configuration for scheduling decisions
        elif tool_name == "get_okr_settings":
            result = await _get_okr_settings_tool(agent_id)
        # ── OKR Management Tools (OKR Agent exclusive) ──
        elif tool_name == "create_objective":
            result = await _create_objective(agent_id, user_id, arguments)
        elif tool_name == "create_key_result":
            result = await _create_key_result(agent_id, user_id, arguments)
        elif tool_name == "update_objective":
            result = await _update_objective(agent_id, user_id, arguments)
        elif tool_name == "update_any_kr_progress":
            result = await _update_any_kr_progress(agent_id, user_id, arguments)
        # generate_monthly_okr_report: produce the monthly summary report
        elif tool_name == "generate_monthly_okr_report":
            result = await _generate_monthly_okr_report(agent_id)
        elif tool_name == "upsert_member_daily_report":
            result = await _upsert_member_daily_report(agent_id, arguments)
        # ── Vercel & Neon Deploy Tools ──
        elif tool_name == "vercel_deploy":
            result = await _vercel_deploy(agent_id, ws, arguments)
        elif tool_name == "vercel_list_deployments":
            result = await _vercel_list_deployments(agent_id, arguments)
        elif tool_name == "vercel_get_deploy_logs":
            result = await _vercel_get_deploy_logs(agent_id, arguments)
        elif tool_name == "vercel_set_env":
            result = await _vercel_set_env(agent_id, arguments)
        elif tool_name == "vercel_manage_domain":
            result = await _vercel_manage_domain(agent_id, arguments)
        elif tool_name == "neon_create_database":
            result = await _neon_create_database(agent_id, arguments)
        else:

            # Try MCP tool execution
            result = await _execute_mcp_tool(tool_name, arguments, agent_id=agent_id)

        # Log tool call activity (skip noisy read operations)
        if tool_name not in ("list_files", "read_file", "read_document"):
            from app.services.activity_logger import log_activity
            safe_arguments = _observability_arguments(tool_name, arguments)
            safe_result = _observability_text(result)
            await log_activity(
                agent_id, "tool_call",
                f"Called tool {tool_name}: {safe_result[:80]}",
                detail={
                    "tool": tool_name,
                    "args": safe_arguments,
                    "result": safe_result[:300],
                },
            )
        # Save error message to current session if a messaging tool fails, so the user is notified
        if session_id and tool_name in ("send_channel_message", "send_feishu_message", "send_platform_message", "send_message_to_agent") and isinstance(result, str) and result.startswith("❌"):
            try:
                async with async_session() as _err_db:
                    from app.models.audit import ChatMessage as _CM
                    _err_db.add(_CM(
                        agent_id=agent_id,
                        user_id=user_id,
                        role="assistant",
                        content=(
                            "⚠️ [系统提示] 数字员工工具调用失败！\n"
                            f"工具名: `{tool_name}`\n"
                            "参数: `"
                            f"{json.dumps(_observability_arguments(tool_name, arguments), ensure_ascii=False)}"
                            "`\n"
                            f"错误信息: {_observability_text(result)}"
                        ),
                        conversation_id=session_id,
                    ))
                    await _err_db.commit()
            except Exception as _e:
                logger.warning(f"Failed to save tool error message to session: {_e}")

        if agentbay_scope_token is not None:
            agentbay_session_scope_id.reset(agentbay_scope_token)
        return result
    except Exception as e:
        if agentbay_scope_token is not None:
            agentbay_session_scope_id.reset(agentbay_scope_token)
        logger.exception(f"[Tool] Execution failed: {tool_name}")
        return f"Tool execution error ({tool_name}): {type(e).__name__}: {str(e)[:200]}"


def _read_http_status_retryable(status_code: int) -> bool:
    """Record retry eligibility for canonical read/safe HTTP tools."""
    return status_code in {408, 429} or status_code >= 500


async def _web_search_outcome(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> ToolExecutionOutcome:
    """Route the deprecated unified search tool to one native provider fact."""
    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "web_search requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()
    config = await _get_tool_config(agent_id, "web_search") or {}
    try:
        max_results = int(
            arguments.get("max_results", config.get("max_results", 5))
        )
    except (TypeError, ValueError):
        return _typed_failure(
            "web_search max_results must be an integer.",
            "invalid_tool_arguments",
        )
    if max_results < 1:
        return _typed_failure(
            "web_search max_results must be positive.",
            "invalid_tool_arguments",
        )
    max_results = min(max_results, 10)
    engine = str(config.get("search_engine") or "duckduckgo").strip().lower()
    api_key = config.get("api_key")
    if not isinstance(api_key, str):
        return _typed_failure(
            "web_search API key configuration is invalid.",
            "search_configuration_invalid",
        )
    api_key = api_key.strip()
    language = str(config.get("language") or "en")

    if engine == "duckduckgo":
        return await _duckduckgo_search_outcome(
            {"query": query, "max_results": max_results}
        )
    if engine not in {"tavily", "google", "bing", "exa"}:
        return _typed_failure(
            f"web_search engine '{engine}' is not supported.",
            "search_configuration_invalid",
        )
    if not api_key:
        return _typed_failure(
            f"web_search engine '{engine}' requires configured credentials.",
            "search_credentials_missing",
        )
    if engine == "tavily":
        return await _search_tavily_outcome(query, api_key, max_results)
    if engine == "google":
        return await _search_google_outcome(
            query,
            api_key,
            max_results,
            language,
        )
    if engine == "bing":
        return await _search_bing_outcome(
            query,
            api_key,
            max_results,
            language,
        )
    return await _exa_search_outcome(
        {"query": query, "max_results": max_results},
        agent_id,
        api_key_override=api_key,
    )


async def _web_search(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> str:
    """Legacy display adapter for typed unified web search."""
    outcome = await _web_search_outcome(arguments, agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Web search returned no summary.",
    )

async def _get_jina_api_key() -> str:
    """Read Jina API key from DB system_settings first, then fall back to env."""
    try:
        from app.database import async_session
        from app.models.system_settings import SystemSetting
        from sqlalchemy import select
        async with async_session() as db:
            result = await db.execute(select(SystemSetting).where(SystemSetting.key == "jina_api_key"))
            setting = result.scalar_one_or_none()
            if setting and setting.value.get("api_key"):
                return setting.value["api_key"]
    except Exception:
        pass
    from app.config import get_settings
    return get_settings().JINA_API_KEY


async def _jina_search_outcome(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> ToolExecutionOutcome:
    """Search Jina using HTTP status and decoded response facts."""
    import httpx

    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "jina_search requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()
    try:
        max_results = int(arguments.get("max_results", 5))
    except (TypeError, ValueError):
        return _typed_failure(
            "jina_search max_results must be an integer.",
            "invalid_tool_arguments",
        )
    if max_results < 1:
        return _typed_failure(
            "jina_search max_results must be positive.",
            "invalid_tool_arguments",
        )
    max_results = min(max_results, 10)
    config = await _get_tool_config(agent_id, "jina_search") or {}
    configured_key = config.get("api_key", "")
    if not isinstance(configured_key, str):
        return _typed_failure(
            "jina_search API key configuration is invalid.",
            "search_configuration_invalid",
        )
    api_key = configured_key.strip() or await _get_jina_api_key()

    headers: dict = {
        "Accept": "application/json",
        "X-Respond-With": "no-content",  # return snippets/descriptions, not full pages (faster)
        "X-Return-Format": "markdown",
    }
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
            resp = await client.get(
                f"https://s.jina.ai/{__import__('urllib.parse', fromlist=['quote']).quote(query)}",
                headers=headers,
            )
    except httpx.TimeoutException:
        return _typed_failure(
            "Jina Search timed out.",
            "jina_search_timeout",
            retryable=True,
        )
    except httpx.TransportError as exc:
        return _typed_failure(
            f"Jina Search transport failed: {type(exc).__name__}.",
            "jina_search_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"Jina Search failed: {type(exc).__name__}.",
            "jina_search_failed",
        )

    if resp.status_code != 200:
        return _typed_failure(
            f"Jina Search returned HTTP {resp.status_code}.",
            "jina_search_http_error",
            retryable=_read_http_status_retryable(resp.status_code),
        )
    try:
        data = resp.json()
    except Exception:
        return _typed_failure(
            "Jina Search returned invalid JSON.",
            "jina_search_response_invalid",
            retryable=True,
        )
    if not isinstance(data, Mapping) or not isinstance(data.get("data"), list):
        return _typed_failure(
            "Jina Search returned an invalid result collection.",
            "jina_search_response_invalid",
            retryable=True,
        )
    items = data["data"][:max_results]
    if any(not isinstance(item, Mapping) for item in items):
        return _typed_failure(
            "Jina Search returned an invalid result entry.",
            "jina_search_response_invalid",
            retryable=True,
        )
    if not items:
        return _typed_success(f'No Jina Search results found for "{query}".')

    parts = []
    for index, item in enumerate(items, 1):
        title = item.get("title", "Untitled")
        url = item.get("url", "")
        description = item.get("description", "") or str(
            item.get("content", "")
        )[:500]
        parts.append(f"**{index}. {title}**\n{url}\n{description}")
    return _typed_success(
        f'Jina Search results for "{query}" ({len(items)} items):\n\n'
        + "\n\n---\n\n".join(parts)
    )


async def _jina_search(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> str:
    """Legacy display adapter for typed Jina Search."""
    outcome = await _jina_search_outcome(arguments, agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Jina Search returned no summary.",
    )


async def _jina_read_outcome(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> ToolExecutionOutcome:
    """Read one page through Jina using HTTP and bounded-content facts."""
    import httpx
    from urllib.parse import urlparse

    url = arguments.get("url")
    if not isinstance(url, str) or not url.strip():
        return _typed_failure(
            "jina_read requires url.",
            "invalid_tool_arguments",
        )
    url = url.strip()
    if "://" not in url:
        url = "https://" + url
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        return _typed_failure(
            "jina_read url must be a valid HTTP(S) URL.",
            "invalid_tool_arguments",
        )
    try:
        max_chars = int(arguments.get("max_chars", 8000))
    except (TypeError, ValueError):
        return _typed_failure(
            "jina_read max_chars must be an integer.",
            "invalid_tool_arguments",
        )
    if max_chars < 1:
        return _typed_failure(
            "jina_read max_chars must be positive.",
            "invalid_tool_arguments",
        )
    max_chars = min(max_chars, 20000)
    config = await _get_tool_config(agent_id, "jina_read") or {}
    configured_key = config.get("api_key", "")
    if not isinstance(configured_key, str):
        return _typed_failure(
            "jina_read API key configuration is invalid.",
            "search_configuration_invalid",
        )
    api_key = configured_key.strip() or await _get_jina_api_key()

    headers: dict = {
        "Accept": "text/plain, text/markdown, */*",
        "X-Return-Format": "markdown",
        "X-Remove-Selector": "header, footer, nav, aside, .ads, .advertisement",
    }
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
            resp = await client.get(
                f"https://r.jina.ai/{url}",
                headers=headers,
            )
    except httpx.TimeoutException:
        return _typed_failure(
            "Jina Reader timed out.",
            "jina_read_timeout",
            retryable=True,
        )
    except httpx.TransportError as exc:
        return _typed_failure(
            f"Jina Reader transport failed: {type(exc).__name__}.",
            "jina_read_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"Jina Reader failed: {type(exc).__name__}.",
            "jina_read_failed",
        )

    if resp.status_code != 200:
        return _typed_failure(
            f"Jina Reader returned HTTP {resp.status_code}.",
            "jina_read_http_error",
            retryable=_read_http_status_retryable(resp.status_code),
        )
    text = resp.text.strip()
    if len(text) < 100:
        return _typed_failure(
            "Jina Reader returned no usable content.",
            "jina_read_content_empty",
            retryable=True,
        )
    if len(text) > max_chars:
        text = text[:max_chars] + f"\n\n[... truncated at {max_chars} chars]"
    return _typed_success(f"Content from: {url}\n\n{text}")


async def _jina_read(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> str:
    """Legacy display adapter for typed Jina Reader."""
    outcome = await _jina_read_outcome(arguments, agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Jina Reader returned no summary.",
    )


async def _validate_public_http_url(url: str) -> tuple[str | None, str | None]:
    """Normalize a URL and reject local/private network targets."""
    import ipaddress
    import socket
    from urllib.parse import urlparse

    url = (url or "").strip()
    if not url:
        return None, "❌ Please provide a URL"
    if "://" not in url:
        url = "https://" + url

    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"}:
        return None, "❌ Only HTTP and HTTPS URLs are supported"
    if not parsed.hostname:
        return None, "❌ URL must include a hostname"

    hostname = parsed.hostname
    try:
        ipaddress.ip_address(hostname)
        host_is_ip = True
    except ValueError:
        host_is_ip = False

    if hostname.lower() in {"localhost", "localhost.localdomain"}:
        return None, "❌ Localhost URLs are blocked for safety"

    try:
        if host_is_ip:
            addresses = [hostname]
        else:
            loop = asyncio.get_running_loop()
            infos = await loop.run_in_executor(
                None,
                lambda: socket.getaddrinfo(hostname, parsed.port or (443 if parsed.scheme == "https" else 80), type=socket.SOCK_STREAM),
            )
            addresses = [info[4][0] for info in infos]
    except Exception as exc:
        return None, f"❌ Could not resolve hostname {hostname}: {str(exc)[:160]}"

    for address in set(addresses):
        try:
            ip = ipaddress.ip_address(address)
        except ValueError:
            return None, f"❌ Could not validate resolved address: {address}"
        is_proxy_test_range = (not host_is_ip) and ip in ipaddress.ip_network("198.18.0.0/15")
        if (
            ip.is_loopback
            or ip.is_link_local
            or ip.is_multicast
            or ip.is_unspecified
            or ip.is_reserved
            or (ip.is_private and not is_proxy_test_range)
        ):
            return None, f"❌ Private, local, reserved, or internal network URLs are blocked ({address})"

    return url, None


def _fallback_extract_visible_text(html: str) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript", "template", "svg", "canvas", "header", "footer", "nav", "aside"]):
        tag.decompose()
    text = soup.get_text("\n")
    lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def _extract_page_links(html: str, base_url: str, limit: int = 30) -> list[str]:
    from bs4 import BeautifulSoup
    from urllib.parse import urljoin

    soup = BeautifulSoup(html, "html.parser")
    links: list[str] = []
    seen: set[str] = set()
    for anchor in soup.find_all("a", href=True):
        href = urljoin(base_url, anchor["href"].strip())
        if not href.startswith(("http://", "https://")) or href in seen:
            continue
        label = re.sub(r"\s+", " ", anchor.get_text(" ", strip=True))[:80] or href
        seen.add(href)
        links.append(f"- {label}: {href}")
        if len(links) >= limit:
            break
    return links


async def _read_webpage_outcome(arguments: dict) -> ToolExecutionOutcome:
    """Fetch and extract readable content from a public webpage without a third-party reader API."""
    import httpx
    import trafilatura
    from bs4 import BeautifulSoup

    url, validation_error = await _validate_public_http_url(arguments.get("url", ""))
    if validation_error:
        return _typed_failure(validation_error, "webpage_url_invalid")

    try:
        max_chars = min(max(int(arguments.get("max_chars", 12000)), 500), 50000)
    except (TypeError, ValueError):
        return _typed_failure(
            "read_webpage max_chars must be an integer.",
            "invalid_tool_arguments",
        )
    include_links = bool(arguments.get("include_links", False))
    max_bytes = 2_000_000
    headers = {
        "User-Agent": "ClawithBot/1.0 (+https://clawith.ai) Mozilla/5.0",
        "Accept": "text/html, text/plain, application/json, application/xml;q=0.9, text/*;q=0.8, */*;q=0.5",
    }

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=15) as client:
            async with client.stream("GET", url, headers=headers) as resp:
                content_length = resp.headers.get("content-length")
                if content_length and content_length.isdigit() and int(content_length) > max_bytes:
                    return _typed_failure(
                        f"Page is too large to read safely ({content_length} bytes, limit {max_bytes} bytes).",
                        "webpage_too_large",
                    )

                chunks: list[bytes] = []
                total = 0
                truncated_bytes = False
                async for chunk in resp.aiter_bytes():
                    total += len(chunk)
                    if total > max_bytes:
                        remaining = max_bytes - sum(len(part) for part in chunks)
                        if remaining > 0:
                            chunks.append(chunk[:remaining])
                        truncated_bytes = True
                        break
                    chunks.append(chunk)

                status_code = resp.status_code
                final_url = str(resp.url)
                content_type = (resp.headers.get("content-type") or "").split(";")[0].strip().lower()
                encoding = resp.encoding or "utf-8"

        if status_code >= 400:
            return _typed_failure(
                f"Webpage fetch failed HTTP {status_code}: {final_url}",
                "webpage_http_error",
                retryable=status_code >= 500,
            )
        validated_final_url, final_url_error = await _validate_public_http_url(
            final_url
        )
        if final_url_error or not validated_final_url:
            return _typed_failure(
                final_url_error or "Webpage redirect target is invalid.",
                "webpage_redirect_target_invalid",
            )
        final_url = validated_final_url

        raw = b"".join(chunks)
        text = raw.decode(encoding, errors="replace").strip()
        if not text:
            return _typed_failure(
                f"Empty response from {final_url}",
                "webpage_empty_response",
                retryable=True,
            )

        title = ""
        description = ""
        extracted = text
        links: list[str] = []

        if content_type in {"", "text/html", "application/xhtml+xml"} or "<html" in text[:500].lower():
            soup = BeautifulSoup(text, "html.parser")
            if soup.title and soup.title.string:
                title = soup.title.string.strip()
            meta_description = soup.find("meta", attrs={"name": "description"})
            if meta_description and meta_description.get("content"):
                description = meta_description["content"].strip()

            extracted = trafilatura.extract(
                text,
                url=final_url,
                output_format="markdown",
                include_links=include_links,
                include_comments=False,
                include_tables=True,
            ) or _fallback_extract_visible_text(text)
            if include_links:
                links = _extract_page_links(text, final_url)
        elif content_type.startswith("text/") or content_type in {"application/json", "application/xml", "text/xml"}:
            title = final_url
        else:
            return _typed_failure(
                f"Unsupported content type: {content_type or 'unknown'}",
                "webpage_content_type_unsupported",
            )

        extracted = extracted.strip()
        if not extracted:
            return _typed_failure(
                f"Could not extract readable content from {final_url}",
                "webpage_content_unreadable",
            )

        truncated_chars = len(extracted) > max_chars
        if truncated_chars:
            extracted = extracted[:max_chars].rstrip() + f"\n\n[... truncated at {max_chars} chars]"

        meta_lines = [
            f"URL: {final_url}",
            f"Status: HTTP {status_code}",
        ]
        if title:
            meta_lines.append(f"Title: {title}")
        if description:
            meta_lines.append(f"Description: {description}")
        if truncated_bytes:
            meta_lines.append(f"Note: response body truncated at {max_bytes} bytes before extraction")
        if truncated_chars:
            meta_lines.append(f"Note: extracted text truncated at {max_chars} characters")

        result = "🌐 **Webpage content**\n\n" + "\n".join(meta_lines) + "\n\n---\n\n" + extracted
        if links:
            result += "\n\n---\n\nLinks:\n" + "\n".join(links)
        return _typed_success(result, evidence_refs=(final_url,))

    except httpx.TimeoutException:
        return _typed_failure(
            f"Webpage fetch timed out: {url}",
            "webpage_timeout",
            retryable=True,
        )
    except Exception as e:
        return _typed_failure(
            f"Webpage read error: {type(e).__name__}.",
            "webpage_read_failed",
            retryable=True,
        )


async def _read_webpage(arguments: dict) -> str:
    outcome = await _read_webpage_outcome(arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Webpage read returned no summary.",
    )



async def _search_tavily_outcome(
    query: str,
    api_key: str,
    max_results: int,
) -> ToolExecutionOutcome:
    """Search Tavily using HTTP status and its results collection."""
    import httpx

    try:
        async with httpx.AsyncClient() as client:
            resp = await client.post(
                "https://api.tavily.com/search",
                json={
                    "query": query,
                    "max_results": max_results,
                    "search_depth": "basic",
                },
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
                timeout=15,
            )
    except httpx.TimeoutException:
        return _typed_failure(
            "Tavily search timed out.",
            "tavily_search_timeout",
            retryable=True,
        )
    except httpx.TransportError as exc:
        return _typed_failure(
            f"Tavily search transport failed: {type(exc).__name__}.",
            "tavily_search_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"Tavily search failed: {type(exc).__name__}.",
            "tavily_search_failed",
        )
    if resp.status_code != 200:
        return _typed_failure(
            f"Tavily search returned HTTP {resp.status_code}.",
            "tavily_search_http_error",
            retryable=_read_http_status_retryable(resp.status_code),
        )
    try:
        data = resp.json()
    except Exception:
        return _typed_failure(
            "Tavily search returned invalid JSON.",
            "tavily_search_response_invalid",
            retryable=True,
        )
    if (
        not isinstance(data, Mapping)
        or "error" in data
        or not isinstance(data.get("results"), list)
    ):
        return _typed_failure(
            "Tavily search returned an invalid result collection.",
            "tavily_search_response_invalid",
            retryable=True,
        )
    items = data["results"][:max_results]
    if any(not isinstance(item, Mapping) for item in items):
        return _typed_failure(
            "Tavily search returned an invalid result entry.",
            "tavily_search_response_invalid",
            retryable=True,
        )
    results = []
    for item in items:
        results.append(
            f"**{item.get('title', '')}**\n{item.get('url', '')}\n"
            f"{str(item.get('content', ''))[:200]}"
        )
    if not results:
        return _typed_success(f'No Tavily results found for "{query}".')
    return _typed_success(
        f'Tavily search for "{query}" ({len(results)} items):\n\n'
        + "\n\n---\n\n".join(results)
    )


async def _search_tavily(query: str, api_key: str, max_results: int) -> str:
    """Legacy display adapter for typed Tavily search."""
    outcome = await _search_tavily_outcome(query, api_key, max_results)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Tavily search returned no summary.",
    )


async def _search_google_outcome(
    query: str,
    api_key: str,
    max_results: int,
    language: str,
) -> ToolExecutionOutcome:
    """Search Google Custom Search using HTTP and decoded response facts."""
    import httpx

    parts = api_key.split(":", 1)
    if len(parts) != 2 or not all(part.strip() for part in parts):
        return _typed_failure(
            "Google search credentials must use API_KEY:SEARCH_ENGINE_ID format.",
            "search_configuration_invalid",
        )

    gapi_key, cx = parts
    try:
        async with httpx.AsyncClient() as client:
            resp = await client.get(
                "https://www.googleapis.com/customsearch/v1",
                params={
                    "key": gapi_key,
                    "cx": cx,
                    "q": query,
                    "num": max_results,
                    "lr": f"lang_{language[:2]}",
                },
                timeout=10,
            )
    except httpx.TimeoutException:
        return _typed_failure(
            "Google search timed out.",
            "google_search_timeout",
            retryable=True,
        )
    except httpx.TransportError as exc:
        return _typed_failure(
            f"Google search transport failed: {type(exc).__name__}.",
            "google_search_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"Google search failed: {type(exc).__name__}.",
            "google_search_failed",
        )
    if resp.status_code != 200:
        return _typed_failure(
            f"Google search returned HTTP {resp.status_code}.",
            "google_search_http_error",
            retryable=_read_http_status_retryable(resp.status_code),
        )
    try:
        data = resp.json()
    except Exception:
        return _typed_failure(
            "Google search returned invalid JSON.",
            "google_search_response_invalid",
            retryable=True,
        )
    if not isinstance(data, Mapping) or "error" in data:
        return _typed_failure(
            "Google search returned an invalid response.",
            "google_search_response_invalid",
        )
    raw_items = data.get("items")
    if raw_items is None:
        if not any(key in data for key in ("queries", "searchInformation")):
            return _typed_failure(
                "Google search response did not prove a completed search.",
                "google_search_response_invalid",
                retryable=True,
            )
        raw_items = []
    if not isinstance(raw_items, list) or any(
        not isinstance(item, Mapping) for item in raw_items
    ):
        return _typed_failure(
            "Google search returned an invalid result collection.",
            "google_search_response_invalid",
            retryable=True,
        )
    results = []
    for item in raw_items[:max_results]:
        results.append(
            f"**{item.get('title', '')}**\n{item.get('link', '')}\n"
            f"{item.get('snippet', '')}"
        )
    if not results:
        return _typed_success(f'No Google results found for "{query}".')
    return _typed_success(
        f'Google search for "{query}" ({len(results)} items):\n\n'
        + "\n\n---\n\n".join(results)
    )


async def _search_google(
    query: str,
    api_key: str,
    max_results: int,
    language: str,
) -> str:
    """Legacy display adapter for typed Google search."""
    outcome = await _search_google_outcome(
        query,
        api_key,
        max_results,
        language,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Google search returned no summary.",
    )


async def _search_bing_outcome(
    query: str,
    api_key: str,
    max_results: int,
    language: str,
) -> ToolExecutionOutcome:
    """Search Bing using HTTP and its webPages result collection."""
    import httpx

    try:
        async with httpx.AsyncClient() as client:
            resp = await client.get(
                "https://api.bing.microsoft.com/v7.0/search",
                params={"q": query, "count": max_results, "mkt": language},
                headers={"Ocp-Apim-Subscription-Key": api_key},
                timeout=10,
            )
    except httpx.TimeoutException:
        return _typed_failure(
            "Bing search timed out.",
            "bing_search_timeout",
            retryable=True,
        )
    except httpx.TransportError as exc:
        return _typed_failure(
            f"Bing search transport failed: {type(exc).__name__}.",
            "bing_search_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"Bing search failed: {type(exc).__name__}.",
            "bing_search_failed",
        )
    if resp.status_code != 200:
        return _typed_failure(
            f"Bing search returned HTTP {resp.status_code}.",
            "bing_search_http_error",
            retryable=_read_http_status_retryable(resp.status_code),
        )
    try:
        data = resp.json()
    except Exception:
        return _typed_failure(
            "Bing search returned invalid JSON.",
            "bing_search_response_invalid",
            retryable=True,
        )
    if not isinstance(data, Mapping) or "errors" in data:
        return _typed_failure(
            "Bing search returned an invalid response.",
            "bing_search_response_invalid",
        )
    web_pages = data.get("webPages")
    if web_pages is None:
        if not isinstance(data.get("queryContext"), Mapping):
            return _typed_failure(
                "Bing search response did not prove a completed search.",
                "bing_search_response_invalid",
                retryable=True,
            )
        raw_items = []
    elif isinstance(web_pages, Mapping):
        raw_items = web_pages.get("value", [])
    else:
        raw_items = None
    if not isinstance(raw_items, list) or any(
        not isinstance(item, Mapping) for item in raw_items
    ):
        return _typed_failure(
            "Bing search returned an invalid result collection.",
            "bing_search_response_invalid",
            retryable=True,
        )
    results = []
    for item in raw_items[:max_results]:
        results.append(
            f"**{item.get('name', '')}**\n{item.get('url', '')}\n"
            f"{item.get('snippet', '')}"
        )
    if not results:
        return _typed_success(f'No Bing results found for "{query}".')
    return _typed_success(
        f'Bing search for "{query}" ({len(results)} items):\n\n'
        + "\n\n---\n\n".join(results)
    )


async def _search_bing(
    query: str,
    api_key: str,
    max_results: int,
    language: str,
) -> str:
    """Legacy display adapter for typed Bing search."""
    outcome = await _search_bing_outcome(
        query,
        api_key,
        max_results,
        language,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Bing search returned no summary.",
    )


async def _exa_search_outcome(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
    *,
    api_key_override: str | None = None,
) -> ToolExecutionOutcome:
    """Search Exa using HTTP status and its decoded results collection."""
    import httpx

    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "exa_search requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()

    if api_key_override is None:
        config = await _get_tool_config(agent_id, "exa_search") or {}
    else:
        config = {}
    configured_key = config.get("api_key", "")
    if not isinstance(configured_key, str) or (
        api_key_override is not None and not isinstance(api_key_override, str)
    ):
        return _typed_failure(
            "Exa API key configuration is invalid.",
            "search_configuration_invalid",
        )
    api_key = (
        (api_key_override or "").strip()
        or configured_key.strip()
        or get_settings().EXA_API_KEY
    )
    if not api_key:
        return _typed_failure(
            "Exa search credentials are not configured.",
            "search_credentials_missing",
        )

    try:
        max_results = int(arguments.get("max_results", 5))
    except (TypeError, ValueError):
        return _typed_failure(
            "exa_search max_results must be an integer.",
            "invalid_tool_arguments",
        )
    if max_results < 1:
        return _typed_failure(
            "exa_search max_results must be positive.",
            "invalid_tool_arguments",
        )
    max_results = min(max_results, 10)
    search_type = arguments.get("search_type", "auto")
    content_mode = arguments.get("content_mode", "text")
    if search_type not in {"auto", "neural", "fast"}:
        return _typed_failure(
            "exa_search search_type is invalid.",
            "invalid_tool_arguments",
        )
    if content_mode not in {"text", "highlights", "summary"}:
        return _typed_failure(
            "exa_search content_mode is invalid.",
            "invalid_tool_arguments",
        )
    category = arguments.get("category") or None
    include_domains = arguments.get("include_domains")
    exclude_domains = arguments.get("exclude_domains")
    if category is not None and not isinstance(category, str):
        return _typed_failure(
            "exa_search category must be a string.",
            "invalid_tool_arguments",
        )
    if any(
        value is not None and not isinstance(value, str)
        for value in (include_domains, exclude_domains)
    ):
        return _typed_failure(
            "exa_search domain filters must be comma-separated strings.",
            "invalid_tool_arguments",
        )

    body: dict = {
        "query": query,
        "type": search_type,
        "numResults": max_results,
        "contents": {},
    }

    if category:
        body["category"] = category
    if include_domains:
        body["includeDomains"] = [d.strip() for d in include_domains.split(",") if d.strip()]
    if exclude_domains:
        body["excludeDomains"] = [d.strip() for d in exclude_domains.split(",") if d.strip()]

    if content_mode == "highlights":
        body["contents"]["highlights"] = {"numSentences": 3}
    elif content_mode == "summary":
        body["contents"]["summary"] = {}
    else:
        body["contents"]["text"] = {"maxCharacters": 1000}

    try:
        async with httpx.AsyncClient() as client:
            resp = await client.post(
                "https://api.exa.ai/search",
                json=body,
                headers={
                    "x-api-key": api_key,
                    "Content-Type": "application/json",
                    "x-exa-integration": "clawith",
                },
                timeout=15,
            )
    except httpx.TimeoutException:
        return _typed_failure(
            "Exa search timed out.",
            "exa_search_timeout",
            retryable=True,
        )
    except httpx.TransportError as exc:
        return _typed_failure(
            f"Exa search transport failed: {type(exc).__name__}.",
            "exa_search_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"Exa search failed: {type(exc).__name__}.",
            "exa_search_failed",
        )

    if resp.status_code != 200:
        return _typed_failure(
            f"Exa search returned HTTP {resp.status_code}.",
            "exa_search_http_error",
            retryable=_read_http_status_retryable(resp.status_code),
        )
    try:
        data = resp.json()
    except Exception:
        return _typed_failure(
            "Exa search returned invalid JSON.",
            "exa_search_response_invalid",
            retryable=True,
        )
    if (
        not isinstance(data, Mapping)
        or "error" in data
        or not isinstance(data.get("results"), list)
    ):
        return _typed_failure(
            "Exa search returned an invalid result collection.",
            "exa_search_response_invalid",
            retryable=True,
        )
    items = data["results"][:max_results]
    if any(not isinstance(item, Mapping) for item in items):
        return _typed_failure(
            "Exa search returned an invalid result entry.",
            "exa_search_response_invalid",
            retryable=True,
        )
    if not items:
        return _typed_success(f'No Exa results found for "{query}".')

    parts = []
    for index, item in enumerate(items, 1):
        title = item.get("title", "Untitled")
        url = item.get("url", "")
        content = ""
        if content_mode == "highlights" and item.get("highlights"):
            highlights = item["highlights"]
            if not isinstance(highlights, list) or any(
                not isinstance(value, str) for value in highlights
            ):
                return _typed_failure(
                    "Exa search returned invalid highlights.",
                    "exa_search_response_invalid",
                    retryable=True,
                )
            content = " ... ".join(highlights)
        elif content_mode == "summary" and item.get("summary"):
            content = str(item["summary"])
        elif item.get("text"):
            content = str(item["text"])[:500]
        parts.append(f"**{index}. {title}**\n{url}\n{content}")
    return _typed_success(
        f'Exa search for "{query}" ({len(items)} items):\n\n'
        + "\n\n---\n\n".join(parts)
    )


async def _exa_search(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> str:
    """Legacy display adapter for typed Exa search."""
    outcome = await _exa_search_outcome(arguments, agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Exa search returned no summary.",
    )



# ── Standalone search engine tool wrappers ───────────────────────────────────
# Each function reads its own tool config (agent > company > defaults) and
# delegates to the existing private search implementations above.


async def _duckduckgo_search_outcome(arguments: dict) -> ToolExecutionOutcome:
    """Search DuckDuckGo using HTTP and parsed-result facts."""
    import httpx

    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "duckduckgo_search requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()
    try:
        max_results = int(arguments.get("max_results", 5))
    except (TypeError, ValueError):
        return _typed_failure(
            "duckduckgo_search max_results must be an integer.",
            "invalid_tool_arguments",
        )
    if max_results < 1:
        return _typed_failure(
            "duckduckgo_search max_results must be positive.",
            "invalid_tool_arguments",
        )
    max_results = min(max_results, 10)

    try:
        async with httpx.AsyncClient(follow_redirects=True) as client:
            response = await client.get(
                "https://html.duckduckgo.com/html/",
                params={"q": query},
                headers={
                    "User-Agent": (
                        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"
                    )
                },
                timeout=10,
            )
    except httpx.TimeoutException:
        return _typed_failure(
            "DuckDuckGo search timed out.",
            "duckduckgo_timeout",
            retryable=True,
        )
    except httpx.HTTPError as exc:
        return _typed_failure(
            f"DuckDuckGo search transport failed: {type(exc).__name__}.",
            "duckduckgo_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"DuckDuckGo search failed: {type(exc).__name__}.",
            "duckduckgo_search_failed",
            retryable=True,
        )

    if response.status_code != 200:
        return _typed_failure(
            f"DuckDuckGo returned HTTP {response.status_code}.",
            "duckduckgo_http_error",
            retryable=response.status_code == 429 or response.status_code >= 500,
        )

    blocks = re.findall(
        r'<a[^>]*class="result__a"[^>]*href="([^"]*)"[^>]*>(.*?)</a>.*?'
        r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>',
        response.text,
        re.DOTALL,
    )
    results: list[str] = []
    for url, title, snippet in blocks[:max_results]:
        title = re.sub(r"<[^>]+>", "", title).strip()
        snippet = re.sub(r"<[^>]+>", "", snippet).strip()
        if "uddg=" in url:
            from urllib.parse import parse_qs, unquote, urlparse

            parsed = parse_qs(urlparse(url).query)
            url = unquote(parsed.get("uddg", [url])[0])
        results.append(f"**{title}**\n{url}\n{snippet}")

    if not results:
        return _typed_success(f'No DuckDuckGo results found for "{query}".')
    return _typed_success(
        f'DuckDuckGo results for "{query}" ({len(results)} items):\n\n'
        + "\n\n---\n\n".join(results)
    )


async def _duckduckgo_search_tool(arguments: dict) -> str:
    """Legacy display adapter for the typed DuckDuckGo result."""
    outcome = await _duckduckgo_search_outcome(arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="DuckDuckGo search returned no summary.",
    )


async def _tavily_search_outcome(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> ToolExecutionOutcome:
    """Validate standalone Tavily configuration before its HTTP boundary."""
    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "tavily_search requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()
    config = await _get_tool_config(agent_id, "tavily_search") or {}
    api_key = config.get("api_key", "")
    if not isinstance(api_key, str):
        return _typed_failure(
            "Tavily API key configuration is invalid.",
            "search_configuration_invalid",
        )
    api_key = api_key.strip()
    if not api_key:
        return _typed_failure(
            "Tavily search credentials are not configured.",
            "search_credentials_missing",
        )
    try:
        max_results = int(arguments.get("max_results", 5))
    except (TypeError, ValueError):
        return _typed_failure(
            "tavily_search max_results must be an integer.",
            "invalid_tool_arguments",
        )
    if max_results < 1:
        return _typed_failure(
            "tavily_search max_results must be positive.",
            "invalid_tool_arguments",
        )
    return await _search_tavily_outcome(query, api_key, min(max_results, 10))


async def _tavily_search_tool(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> str:
    """Legacy display adapter for typed standalone Tavily search."""
    outcome = await _tavily_search_outcome(arguments, agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Tavily search returned no summary.",
    )


async def _google_search_outcome(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> ToolExecutionOutcome:
    """Validate standalone Google configuration before its HTTP boundary."""
    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "google_search requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()
    config = await _get_tool_config(agent_id, "google_search") or {}
    api_key = config.get("api_key", "")
    if not isinstance(api_key, str):
        return _typed_failure(
            "Google API key configuration is invalid.",
            "search_configuration_invalid",
        )
    api_key = api_key.strip()
    if not api_key:
        return _typed_failure(
            "Google search credentials are not configured.",
            "search_credentials_missing",
        )
    language = arguments.get("language") or config.get("language", "en")
    if not isinstance(language, str) or not language.strip():
        return _typed_failure(
            "google_search language must be a string.",
            "invalid_tool_arguments",
        )
    try:
        max_results = int(arguments.get("max_results", 5))
    except (TypeError, ValueError):
        return _typed_failure(
            "google_search max_results must be an integer.",
            "invalid_tool_arguments",
        )
    if max_results < 1:
        return _typed_failure(
            "google_search max_results must be positive.",
            "invalid_tool_arguments",
        )
    return await _search_google_outcome(
        query,
        api_key,
        min(max_results, 10),
        language.strip(),
    )


async def _google_search_tool(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> str:
    """Legacy display adapter for typed standalone Google search."""
    outcome = await _google_search_outcome(arguments, agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Google search returned no summary.",
    )


async def _bing_search_outcome(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> ToolExecutionOutcome:
    """Validate standalone Bing configuration before its HTTP boundary."""
    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "bing_search requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()
    config = await _get_tool_config(agent_id, "bing_search") or {}
    api_key = config.get("api_key", "")
    if not isinstance(api_key, str):
        return _typed_failure(
            "Bing API key configuration is invalid.",
            "search_configuration_invalid",
        )
    api_key = api_key.strip()
    if not api_key:
        return _typed_failure(
            "Bing search credentials are not configured.",
            "search_credentials_missing",
        )
    language = arguments.get("language") or config.get("language", "en-US")
    if not isinstance(language, str) or not language.strip():
        return _typed_failure(
            "bing_search language must be a string.",
            "invalid_tool_arguments",
        )
    try:
        max_results = int(arguments.get("max_results", 5))
    except (TypeError, ValueError):
        return _typed_failure(
            "bing_search max_results must be an integer.",
            "invalid_tool_arguments",
        )
    if max_results < 1:
        return _typed_failure(
            "bing_search max_results must be positive.",
            "invalid_tool_arguments",
        )
    return await _search_bing_outcome(
        query,
        api_key,
        min(max_results, 10),
        language.strip(),
    )


async def _bing_search_tool(
    arguments: dict,
    agent_id: uuid.UUID | None = None,
) -> str:
    """Legacy display adapter for typed standalone Bing search."""
    outcome = await _bing_search_outcome(arguments, agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Bing search returned no summary.",
    )


async def _send_channel_file_outcome(
    agent_id: uuid.UUID,
    ws: Path,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Deliver one materialized file using provider or local artifact facts."""
    rel_path = arguments.get("file_path")
    if not isinstance(rel_path, str) or not rel_path.strip():
        return _typed_failure(
            "send_channel_file requires file_path.",
            "invalid_tool_arguments",
        )
    rel_path = rel_path.strip()
    message = arguments.get("message", "")
    if not isinstance(message, str):
        return _typed_failure(
            "send_channel_file message must be a string.",
            "invalid_tool_arguments",
        )
    target_member_id = arguments.get("target_member_id", "")
    if not isinstance(target_member_id, str):
        return _typed_failure(
            "send_channel_file target_member_id must be a string.",
            "invalid_tool_arguments",
        )
    target_member_id = target_member_id.strip()
    if arguments.get("member_name") and not target_member_id:
        return _typed_failure(
            "send_channel_file accepts stable target_member_id, not member_name.",
            "invalid_tool_arguments",
        )
    target_channel = _normalize_roster_provider_type(arguments.get("channel"))

    root = ws.resolve()
    file_path = (root / rel_path).resolve()
    if not file_path.is_relative_to(root) or not file_path.is_file():
        return _typed_failure(
            f"File not found: {rel_path}",
            "workspace_file_not_found",
        )

    if target_member_id:
        return await _send_file_to_human_target_outcome(
            agent_id,
            file_path,
            target_member_id,
            target_channel,
            message,
        )

    sender = channel_file_sender.get()
    if sender is not None:
        try:
            await sender(file_path, message)
        except Exception:
            return _typed_unknown(
                "Channel file delivery outcome is unknown; reconcile before retrying.",
                "channel_file_outcome_unknown",
            )
        return _typed_success(
            f"File '{file_path.name}' was accepted by the current channel sender."
        )

    aid = channel_web_agent_id.get() or str(agent_id)
    from app.config import get_settings as _gs

    base_url = (getattr(_gs(), "BASE_URL", "") or "").rstrip("/")
    download_url = (
        f"{base_url}/api/agents/{aid}/files/download?path={rel_path}"
    )
    summary = f"File ready: [{file_path.name}]({download_url})"
    if message:
        summary = f"{message}\n\n{summary}"
    return _typed_success(
        summary,
        artifact_refs=(_workspace_artifact_ref(agent_id, rel_path),),
    )


async def _send_channel_file(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Send a file to a person or back to the current channel.

    Priority:
    1. If target_member_id is provided, deliver via that Directory member's channel.
    2. If channel_file_sender ContextVar is set (channel-initiated), use it directly.
    3. Fall back to web chat download URL when no explicit recipient is requested.
    """
    rel_path = arguments.get("file_path", "").strip()
    accompany_msg = arguments.get("message", "")
    member_name = (arguments.get("member_name") or "").strip()
    target_member_id = (arguments.get("target_member_id") or "").strip()
    target_channel = _normalize_roster_provider_type(arguments.get("channel"))
    if not rel_path:
        return "Error: file_path is required"
    if member_name and not target_member_id:
        return (
            "❌ member_name is no longer supported for send_channel_file. "
            "Call query_directory(member_type=\"human\", query=\"...\") first, then retry with target_member_id."
        )

    # Resolve file path within agent workspace
    file_path = (ws / rel_path).resolve()
    ws_resolved = ws.resolve()
    if not str(file_path).startswith(str(ws_resolved)):
        file_path = (WORKSPACE_ROOT / str(agent_id) / rel_path).resolve()
        if not file_path.exists():
            return f"Error: File not found: {rel_path}"
    if not file_path.exists():
        return f"Error: File not found: {rel_path}"

    # Priority 1: explicit recipient from roster
    if target_member_id:
        return await _send_file_to_human_target(
            agent_id,
            file_path,
            target_member_id,
            target_channel,
            accompany_msg,
        )

    # Priority 2: channel-initiated (ContextVar set by channel webhook handler)
    sender = channel_file_sender.get()
    if sender is not None:
        try:
            await sender(file_path, accompany_msg)
            return f"File '{file_path.name}' sent to user via channel."
        except Exception as e:
            return f"Failed to send file: {e}"

    # Priority 3: Web chat fallback — return download URL
    aid = channel_web_agent_id.get() or str(agent_id)
    base_abs = (WORKSPACE_ROOT / str(agent_id)).resolve()
    try:
        file_rel = str(file_path.resolve().relative_to(base_abs))
    except ValueError:
        file_rel = rel_path
    from app.config import get_settings as _gs
    _s = _gs()
    base_url = getattr(_s, 'BASE_URL', '').rstrip('/') or ''
    download_url = f"{base_url}/api/agents/{aid}/files/download?path={file_rel}"
    msg = f"File ready: [{file_path.name}]({download_url})"
    if accompany_msg:
        msg = accompany_msg + "\n\n" + msg
    return msg


async def _send_file_to_human_target(
    agent_id: uuid.UUID,
    file_path: Path,
    target_member_id: str,
    target_channel: str | None,
    message: str = "",
) -> str:
    """Send a file to an already selected human roster target."""
    from app.models.channel_config import ChannelConfig

    async with async_session() as db:
        target, error = await _resolve_roster_human_target(
            db,
            agent_id,
            target_member_id=target_member_id,
            provider_type=target_channel,
        )
        if error:
            return error

        result = await db.execute(
            select(ChannelConfig).where(ChannelConfig.agent_id == agent_id)
        )
        configs = {c.channel_type: c for c in result.scalars().all()}

    target_member = target.member
    display_name = target_member.name or target_member_id
    provider_type = target.provider_type
    if not provider_type and (target_member.external_id or target_member.open_id):
        provider_type = "feishu"

    if provider_type == "feishu":
        config = configs.get("feishu")
        if not config:
            return "❌ This agent has no Feishu channel configured"
        if target_member.external_id:
            return await _send_file_via_feishu_resolved(
                agent_id, config, file_path, display_name, target_member.external_id, "user_id", message
            )
        if target_member.open_id:
            return await _send_file_via_feishu_resolved(
                agent_id, config, file_path, display_name, target_member.open_id, "open_id", message
            )
        return f"❌ {display_name} has no Feishu user_id/open_id."

    if provider_type == "slack":
        config = configs.get("slack")
        if not config:
            return "❌ This agent has no Slack channel configured"
        slack_user_id = target_member.external_id or target_member.open_id or target_member.unionid
        if not slack_user_id:
            return f"❌ {display_name} has no Slack user id."
        return await _send_file_via_slack_user_id(agent_id, config, file_path, display_name, slack_user_id, message)

    return (
        f"❌ File delivery via {provider_type or 'this channel'} is not supported yet. "
        "Use send_channel_message to send a download link, or omit target_member_id to return a link here."
    )


async def _send_file_to_human_target_outcome(
    agent_id: uuid.UUID,
    file_path: Path,
    target_member_id: str,
    target_channel: str | None,
    message: str = "",
) -> ToolExecutionOutcome:
    """Resolve a human recipient before the provider dispatch boundary."""
    try:
        async with async_session() as db:
            target, error = await _resolve_roster_human_target(
                db,
                agent_id,
                target_member_id=target_member_id,
                provider_type=target_channel,
            )
            if error:
                return _typed_failure(error, "channel_file_recipient_invalid")
            result = await db.execute(
                select(ChannelConfig).where(ChannelConfig.agent_id == agent_id)
            )
            configs = {config.channel_type: config for config in result.scalars().all()}
    except Exception as exc:
        return _typed_failure(
            f"File recipient could not be resolved: {type(exc).__name__}.",
            "channel_file_recipient_resolution_failed",
        )

    target_member = target.member
    display_name = target_member.name or target_member_id
    provider_type = target.provider_type
    if not provider_type and (target_member.external_id or target_member.open_id):
        provider_type = "feishu"

    if provider_type == "feishu":
        config = configs.get("feishu")
        if not config:
            return _typed_failure(
                "This Agent has no Feishu channel configured.",
                "feishu_channel_not_configured",
            )
        receive_id = target_member.external_id or target_member.open_id
        receive_id_type = "user_id" if target_member.external_id else "open_id"
        if not receive_id:
            return _typed_failure(
                f"{display_name} has no Feishu recipient id.",
                "feishu_recipient_not_linked",
            )
        return await _send_file_via_feishu_resolved_outcome(
            config,
            file_path,
            display_name,
            receive_id,
            receive_id_type,
            message,
        )

    if provider_type == "slack":
        config = configs.get("slack")
        if not config:
            return _typed_failure(
                "This Agent has no Slack channel configured.",
                "slack_channel_not_configured",
            )
        slack_user_id = (
            target_member.external_id
            or target_member.open_id
            or target_member.unionid
        )
        if not slack_user_id:
            return _typed_failure(
                f"{display_name} has no Slack user id.",
                "slack_recipient_not_linked",
            )
        return await _send_file_via_slack_user_id_outcome(
            config,
            file_path,
            display_name,
            slack_user_id,
            message,
        )

    return _typed_failure(
        f"File delivery via {provider_type or 'this channel'} is not supported.",
        "channel_file_provider_unsupported",
    )


async def _send_file_via_feishu_resolved_outcome(
    config,
    file_path: Path,
    display_name: str,
    receive_id: str,
    receive_id_type: str,
    message: str,
) -> ToolExecutionOutcome:
    from app.services.feishu_service import feishu_service

    try:
        response = await feishu_service.upload_and_send_file(
            config.app_id,
            config.app_secret,
            receive_id,
            file_path,
            receive_id_type=receive_id_type,
            accompany_msg=message,
        )
    except Exception:
        return _typed_unknown(
            "Feishu file delivery outcome is unknown; reconcile before retrying.",
            "feishu_file_outcome_unknown",
        )
    if not isinstance(response, Mapping):
        return _typed_unknown(
            "Feishu returned an unreadable file response; reconcile before retrying.",
            "feishu_file_response_invalid",
        )
    if response.get("code") != 0:
        if "code" not in response:
            return _typed_unknown(
                "Feishu returned an incomplete file response; reconcile before retrying.",
                "feishu_file_response_invalid",
            )
        if message:
            return _typed_unknown(
                "Feishu rejected the file after an accompanying message may have been sent; "
                "reconcile before retrying.",
                "feishu_file_partial_outcome_unknown",
            )
        return _typed_failure(
            f"Feishu rejected file delivery: {response.get('msg') or 'unknown error'}.",
            "feishu_file_rejected",
        )
    return _typed_success(
        f"File '{file_path.name}' sent to {display_name} via Feishu."
    )


async def _send_file_via_slack_user_id_outcome(
    config,
    file_path: Path,
    display_name: str,
    slack_user_id: str,
    message: str,
) -> ToolExecutionOutcome:
    import httpx

    bot_token = config.app_secret or ""
    if not bot_token:
        return _typed_failure(
            "This Agent has no Slack bot token configured.",
            "slack_channel_not_configured",
        )
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            dm_response = await client.post(
                "https://slack.com/api/conversations.open",
                headers={
                    "Authorization": f"Bearer {bot_token}",
                    "Content-Type": "application/json",
                },
                json={"users": slack_user_id},
            )
            dm_data = dm_response.json()
            if dm_response.status_code >= 400 or not dm_data.get("ok"):
                return _typed_failure(
                    f"Slack rejected DM setup: {dm_data.get('error') or 'unknown error'}.",
                    "slack_file_rejected",
                )
            channel_id = str((dm_data.get("channel") or {}).get("id") or "")
            if not channel_id:
                return _typed_failure(
                    "Slack did not return a DM channel id.",
                    "slack_file_rejected",
                )

            upload_response = await client.post(
                "https://slack.com/api/files.getUploadURLExternal",
                headers={"Authorization": f"Bearer {bot_token}"},
                data={
                    "filename": file_path.name,
                    "length": str(file_path.stat().st_size),
                },
            )
            upload_data = upload_response.json()
            if upload_response.status_code >= 400 or not upload_data.get("ok"):
                return _typed_failure(
                    f"Slack rejected file upload setup: {upload_data.get('error') or 'unknown error'}.",
                    "slack_file_rejected",
                )
            upload_url = upload_data.get("upload_url")
            file_id = upload_data.get("file_id")
            if not upload_url or not file_id:
                return _typed_unknown(
                    "Slack returned an incomplete upload response; reconcile before retrying.",
                    "slack_file_response_invalid",
                )
            binary_response = await client.post(
                upload_url,
                content=file_path.read_bytes(),
                headers={"Content-Type": "application/octet-stream"},
            )
            if binary_response.status_code >= 400:
                return _typed_failure(
                    f"Slack rejected the file bytes with HTTP {binary_response.status_code}.",
                    "slack_file_rejected",
                )
            complete_response = await client.post(
                "https://slack.com/api/files.completeUploadExternal",
                headers={"Authorization": f"Bearer {bot_token}"},
                json={
                    "files": [{"id": file_id}],
                    "channel_id": channel_id,
                    "initial_comment": message,
                },
            )
            complete_data = complete_response.json()
            if complete_response.status_code >= 400 or not complete_data.get("ok"):
                return _typed_failure(
                    f"Slack rejected file completion: {complete_data.get('error') or 'unknown error'}.",
                    "slack_file_rejected",
                )
    except Exception:
        return _typed_unknown(
            "Slack file delivery outcome is unknown; reconcile before retrying.",
            "slack_file_outcome_unknown",
        )
    return _typed_success(
        f"File '{file_path.name}' sent to {display_name} via Slack."
    )


async def _send_file_via_feishu_resolved(
    agent_id,
    config,
    file_path: Path,
    display_name: str,
    receive_id: str,
    id_type: str,
    message: str,
) -> str:
    """Send file to a resolved Feishu recipient."""
    from app.services.feishu_service import feishu_service
    try:
        await feishu_service.upload_and_send_file(
            config.app_id, config.app_secret,
            receive_id, file_path,
            receive_id_type=id_type,
            accompany_msg=message,
        )
        return f"File '{file_path.name}' sent to {display_name} via Feishu."
    except Exception as e:
        # If upload fails, try sending a download link as fallback
        import json as _j
        from app.config import get_settings as _gs
        _s = _gs()
        base_url = getattr(_s, 'BASE_URL', '').rstrip('/') or ''
        base_abs = (WORKSPACE_ROOT / str(agent_id)).resolve()
        try:
            _rel = str(file_path.resolve().relative_to(base_abs))
        except ValueError:
            _rel = file_path.name
        parts = []
        if message:
            parts.append(message)
        if base_url:
            dl_url = f"{base_url}/api/agents/{agent_id}/files/download?path={_rel}"
            parts.append(f"{file_path.name}\n{dl_url}")
        parts.append(f"File upload failed ({e}). If you need direct file sending, enable im:resource permission in Feishu.")
        try:
            await feishu_service.send_message(
                config.app_id, config.app_secret,
                receive_id, "text",
                _j.dumps({"text": "\n\n".join(parts)}, ensure_ascii=False),
                receive_id_type=id_type,
            )
            return f"File upload to Feishu failed, sent download link to {display_name} instead."
        except Exception:
            return f"Failed to send file to {display_name} via Feishu: {e}"


async def _send_file_via_slack_user_id(
    agent_id,
    config,
    file_path: Path,
    display_name: str,
    slack_user_id: str,
    message: str,
) -> str:
    """Send file to a resolved Slack user id."""
    import httpx
    bot_token = config.app_secret or ""
    if not bot_token:
        return "❌ This agent has no Slack bot token configured"

    try:
        async with httpx.AsyncClient(timeout=10) as client:
            # Open a DM channel
            dm_resp = await client.post(
                "https://slack.com/api/conversations.open",
                headers={"Authorization": f"Bearer {bot_token}", "Content-Type": "application/json"},
                json={"users": slack_user_id},
            )
            dm_data = dm_resp.json()
            if not dm_data.get("ok"):
                return f"Slack DM open failed: {dm_data.get('error')}"
            channel_id = dm_data["channel"]["id"]

            # Upload file
            upload_url_resp = await client.post(
                "https://slack.com/api/files.getUploadURLExternal",
                headers={"Authorization": f"Bearer {bot_token}"},
                data={"filename": file_path.name, "length": str(file_path.stat().st_size)},
            )
            ud = upload_url_resp.json()
            if not ud.get("ok"):
                return f"Slack file upload failed: {ud.get('error')}"
            await client.post(ud["upload_url"], content=file_path.read_bytes(),
                            headers={"Content-Type": "application/octet-stream"})
            complete = await client.post(
                "https://slack.com/api/files.completeUploadExternal",
                headers={"Authorization": f"Bearer {bot_token}"},
                json={"files": [{"id": ud["file_id"]}], "channel_id": channel_id,
                      "initial_comment": message or ""},
            )
            if not complete.json().get("ok"):
                return f"Slack file upload complete failed: {complete.json().get('error')}"
            return f"File '{file_path.name}' sent to {display_name} via Slack."
    except Exception as e:
        return f"Failed to send file via Slack: {e}"


def _bounded_mcp_text(value: object, *, max_chars: int = 4000) -> str:
    """Create a bounded, secret-sanitized provider summary."""
    sanitized = _observability_text(value)
    if len(sanitized) <= max_chars:
        return sanitized
    return sanitized[: max_chars - 20] + "...[truncated]"


def _safe_mcp_json(value: object) -> object:
    """Sanitize provider JSON before it reaches an outcome or log."""
    try:
        return sanitize_tool_arguments({"value": value})["value"]
    except Exception:
        return "[MCP payload could not be safely serialized]"


def _mcp_result_summary(result: dict) -> tuple[str, dict]:
    content = result.get("content") if "content" in result else None
    structured = (
        result.get("structuredContent")
        if "structuredContent" in result
        else None
    )
    if content is not None and not isinstance(content, list):
        raise ValueError("MCP result.content must be a list")
    if structured is not None and not isinstance(structured, dict):
        raise ValueError("MCP result.structuredContent must be an object")
    if content is None and structured is None:
        raise ValueError("MCP result has neither content nor structuredContent")

    parts: list[str] = []
    for block in content or []:
        if isinstance(block, str):
            parts.append(_bounded_mcp_text(block))
            continue
        if not isinstance(block, dict):
            parts.append(_bounded_mcp_text(block))
            continue
        block_type = str(block.get("type") or "content")
        if block_type == "text":
            parts.append(_bounded_mcp_text(block.get("text", "")))
        elif block_type in {"image", "audio"}:
            mime_type = _bounded_mcp_text(
                block.get("mimeType") or block_type,
                max_chars=120,
            )
            parts.append(f"[{block_type.title()}: {mime_type}]")
        else:
            parts.append(
                _bounded_mcp_text(
                    json.dumps(
                        _safe_mcp_json(block),
                        ensure_ascii=False,
                        sort_keys=True,
                    )
                )
            )

    metadata: dict = {
        "content_block_count": len(content or []),
        "has_structured_content": structured is not None,
    }
    if structured is not None:
        safe_structured = _safe_mcp_json(structured)
        serialized = json.dumps(
            safe_structured,
            ensure_ascii=False,
            sort_keys=True,
        )
        parts.append(f"Structured content: {_bounded_mcp_text(serialized)}")
        if len(serialized.encode("utf-8")) <= 4096:
            metadata["structured_content"] = safe_structured
        else:
            metadata["structured_content_truncated"] = True

    summary = "\n".join(part for part in parts if part).strip()
    if not summary:
        summary = "MCP tool completed without inline content."
    return _bounded_mcp_text(summary), metadata


class _MCPAsyncContractError(ValueError):
    """A trusted async declaration or its provider result is malformed."""


def _json_pointer_parts(pointer: object) -> tuple[str, ...]:
    if not isinstance(pointer, str) or not pointer.startswith("/"):
        raise _MCPAsyncContractError("JSON pointer must start with '/'")
    return tuple(
        part.replace("~1", "/").replace("~0", "~")
        for part in pointer[1:].split("/")
    )


def _json_pointer_get(document: object, pointer: object) -> object:
    current = document
    for part in _json_pointer_parts(pointer):
        if isinstance(current, Mapping):
            if part not in current:
                raise _MCPAsyncContractError("JSON pointer does not exist")
            current = current[part]
            continue
        if isinstance(current, list):
            try:
                index = int(part)
            except (TypeError, ValueError) as exc:
                raise _MCPAsyncContractError("JSON pointer index is invalid") from exc
            if index < 0 or index >= len(current):
                raise _MCPAsyncContractError("JSON pointer index is out of range")
            current = current[index]
            continue
        raise _MCPAsyncContractError("JSON pointer traverses a scalar")
    return current


def _json_pointer_set(document: dict, pointer: object, value: object) -> None:
    parts = _json_pointer_parts(pointer)
    if not parts:
        raise _MCPAsyncContractError("root replacement is not supported")
    current = document
    for part in parts[:-1]:
        child = current.get(part)
        if child is None:
            child = {}
            current[part] = child
        if not isinstance(child, dict):
            raise _MCPAsyncContractError("poll pointer traverses a scalar")
        current = child
    current[parts[-1]] = deepcopy(value)


def _mcp_async_operation_outcome(
    *,
    result: dict,
    summary: str,
    metadata: dict,
    full_tool_name: str,
    arguments: Mapping[str, object],
    contract: object,
) -> ToolExecutionOutcome:
    """Apply only an admin-owned structured async completion contract."""
    try:
        if not isinstance(contract, Mapping) or contract.get("version") != 1:
            raise _MCPAsyncContractError("unsupported async contract version")
        result_spec = contract.get("result")
        if (
            not isinstance(result_spec, Mapping)
            or result_spec.get("source") != "content_text_json"
        ):
            raise _MCPAsyncContractError("unsupported async result source")
        content_index = result_spec.get("content_index", 0)
        if (
            isinstance(content_index, bool)
            or not isinstance(content_index, int)
            or content_index < 0
        ):
            raise _MCPAsyncContractError("invalid async content index")
        content = result.get("content")
        if not isinstance(content, list) or content_index >= len(content):
            raise _MCPAsyncContractError("async result content is missing")
        block = content[content_index]
        if (
            not isinstance(block, Mapping)
            or block.get("type") != "text"
            or not isinstance(block.get("text"), str)
        ):
            raise _MCPAsyncContractError("async result must be a text block")
        try:
            payload = json.loads(cast(str, block["text"]))
        except (TypeError, ValueError, json.JSONDecodeError) as exc:
            raise _MCPAsyncContractError("async result text is not JSON") from exc
        if not isinstance(payload, Mapping):
            raise _MCPAsyncContractError("async result JSON must be an object")
        provider_state = _json_pointer_get(
            payload,
            result_spec.get("status_pointer"),
        )
        if not isinstance(provider_state, str) or not provider_state.strip():
            raise _MCPAsyncContractError("async status must be a non-empty string")
        provider_state = provider_state.strip()

        operation_spec = contract.get("operation_id")
        if not isinstance(operation_spec, Mapping):
            raise _MCPAsyncContractError("async operation ID declaration is missing")
        operation_source = operation_spec.get("source")
        operation_document = (
            arguments
            if operation_source == "argument"
            else payload
            if operation_source == "result"
            else None
        )
        if operation_document is None:
            raise _MCPAsyncContractError("unsupported async operation ID source")
        raw_operation_id = _json_pointer_get(
            operation_document,
            operation_spec.get("pointer"),
        )
        if isinstance(raw_operation_id, bool) or not isinstance(
            raw_operation_id,
            (str, int),
        ):
            raise _MCPAsyncContractError("async operation ID must be a scalar")
        operation_id = str(raw_operation_id).strip()
        if not operation_id:
            raise _MCPAsyncContractError("async operation ID is empty")

        states = contract.get("states")
        if not isinstance(states, Mapping):
            raise _MCPAsyncContractError("async states declaration is missing")
        classified: dict[str, str] = {}
        for classification in ("pending", "succeeded", "failed", "unknown"):
            values = states.get(classification, [])
            if not isinstance(values, list) or any(
                not isinstance(value, str) or not value.strip() for value in values
            ):
                raise _MCPAsyncContractError("async state lists are invalid")
            for value in values:
                normalized = value.strip()
                if normalized in classified:
                    raise _MCPAsyncContractError("async states overlap")
                classified[normalized] = classification
        if not all(states.get(name) for name in ("pending", "succeeded", "failed")):
            raise _MCPAsyncContractError("async terminal state lists are incomplete")

        poll_spec = contract.get("poll")
        if not isinstance(poll_spec, Mapping) or poll_spec.get("tool") != "$self":
            raise _MCPAsyncContractError("async polling must target the same tool")
        copy_arguments = poll_spec.get("copy_arguments", [])
        set_arguments = poll_spec.get("set_arguments", {})
        interval_ms = poll_spec.get("interval_ms", 1000)
        if (
            not isinstance(copy_arguments, list)
            or not isinstance(set_arguments, Mapping)
            or isinstance(interval_ms, bool)
            or not isinstance(interval_ms, int)
            or interval_ms < 0
            or interval_ms > 600_000
        ):
            raise _MCPAsyncContractError("async poll declaration is invalid")
        poll_arguments: dict = {}
        for pointer in copy_arguments:
            _json_pointer_set(
                poll_arguments,
                pointer,
                _json_pointer_get(arguments, pointer),
            )
        for pointer, value in set_arguments.items():
            _json_pointer_set(poll_arguments, pointer, value)

        classification = classified.get(provider_state)
        operation_key = hashlib.sha256(
            json.dumps(
                {"tool": full_tool_name, "operation_id": operation_id},
                ensure_ascii=False,
                sort_keys=True,
                separators=(",", ":"),
            ).encode("utf-8")
        ).hexdigest()
        operation = {
            "version": 1,
            "operation_key": operation_key,
            "operation_id": operation_id,
            "state": provider_state,
            "poll": {
                "tool": full_tool_name,
                "arguments": poll_arguments,
                "interval_ms": interval_ms,
            },
        }
        async_metadata = {
            **metadata,
            "runtime_async_pending": classification == "pending",
            "async_operation": operation,
        }
        if classification == "pending":
            poll_json = json.dumps(poll_arguments, ensure_ascii=False, sort_keys=True)
            return _typed_pending(
                f"{summary}\n\nAsync operation is still {provider_state}. "
                f"Do not finish yet. Poll {full_tool_name} with arguments "
                f"{poll_json} until it reaches a terminal state.",
                metadata=async_metadata,
            )
        if classification == "succeeded":
            return _typed_success(summary, metadata=async_metadata)
        if classification == "failed":
            return _typed_failure(
                summary,
                "mcp_async_operation_failed",
                metadata=async_metadata,
            )
        return _typed_unknown(
            (
                f"Async operation reported unclassified state {provider_state!r}; "
                "reconcile it before retrying or finishing."
            ),
            "mcp_async_operation_unknown",
            metadata=async_metadata,
        )
    except (TypeError, ValueError):
        return _typed_unknown(
            "MCP async completion facts did not match the configured contract; "
            "reconcile before retrying or finishing.",
            "mcp_async_contract_invalid",
        )


def _mcp_call_response_outcome(
    data: object,
    *,
    full_tool_name: str,
    arguments: Mapping[str, object] | None = None,
    async_completion: object | None = None,
) -> ToolExecutionOutcome:
    """Map protocol facts to a typed outcome without text-prefix inference."""
    if not isinstance(data, dict):
        return _typed_unknown(
            "MCP returned a malformed response; reconcile before retrying.",
            "mcp_malformed_response",
        )
    if "error" in data:
        error = data.get("error")
        message = error.get("message") if isinstance(error, dict) else error
        safe_message = _bounded_mcp_text(message or "provider rejected the call")
        return _typed_failure(
            f"MCP provider rejected the call: {safe_message}",
            "mcp_provider_rejected",
        )
    result = data.get("result")
    if not isinstance(result, dict):
        return _typed_unknown(
            "MCP returned a malformed response; reconcile before retrying.",
            "mcp_malformed_response",
        )
    if "isError" in result and not isinstance(result.get("isError"), bool):
        return _typed_unknown(
            "MCP returned a malformed isError fact; reconcile before retrying.",
            "mcp_malformed_response",
        )
    is_error = result.get("isError") is True
    try:
        summary, metadata = _mcp_result_summary(result)
    except (TypeError, ValueError):
        if is_error:
            return _typed_failure(
                "MCP tool reported failure without valid error details.",
                "mcp_tool_error",
            )
        return _typed_unknown(
            "MCP returned a malformed response; reconcile before retrying.",
            "mcp_malformed_response",
        )
    metadata["mcp_full_tool_name"] = full_tool_name
    if async_completion is not None:
        async_outcome = _mcp_async_operation_outcome(
            result=result,
            summary=summary,
            metadata=metadata,
            full_tool_name=full_tool_name,
            arguments=arguments or {},
            contract=async_completion,
        )
        if is_error and async_outcome.status in {"pending", "succeeded"}:
            return replace(
                async_outcome,
                status="unknown",
                result_summary=(
                    "MCP reported isError=true while the declared async status "
                    "reported a non-failure state; reconcile before continuing."
                ),
                error_code="mcp_async_protocol_conflict",
                retryable=False,
                metadata={
                    **async_outcome.metadata,
                    "runtime_async_pending": False,
                },
            )
        return async_outcome
    if is_error:
        return _typed_failure(summary, "mcp_tool_error")
    return _typed_success(summary, metadata=metadata)


async def _resolve_mcp_execution_target(
    tool_name: str,
    agent_id,
    *,
    allow_legacy_bare_name: bool = False,
) -> dict | None:
    """Resolve one assigned MCP target by its exact durable identity.

    Bare ``mcp_tool_name`` lookup is retained only for the legacy text path.
    It is never used by ``execute_builtin_tool_outcome`` and refuses ambiguous
    raw names shared by multiple servers.
    """
    from urllib.parse import urlparse

    from app.models.tool import AgentTool, Tool

    async with async_session() as db:
        result = await db.execute(
            select(Tool).where(Tool.name == tool_name, Tool.type == "mcp")
        )
        tool = result.scalar_one_or_none()

        if tool is None and allow_legacy_bare_name:
            legacy_result = await db.execute(
                select(Tool).where(
                    Tool.mcp_tool_name == tool_name,
                    Tool.type == "mcp",
                )
            )
            matches = legacy_result.scalars().all()
            if len(matches) > 1:
                logger.warning(
                    "[MCP] Refusing ambiguous legacy bare tool name: {}",
                    tool_name,
                )
                return {
                    "full_name": tool_name,
                    "unavailable_error_code": "mcp_tool_name_ambiguous",
                }
            tool = matches[0] if matches else None

        if tool is None:
            return None
        if (
            not tool.enabled
            or tool.name in BUILTIN_TOOL_NAMES
            or is_reserved_custom_tool_name(str(tool.name or ""))
        ):
            return {
                "full_name": str(tool.name or tool_name),
                "unavailable_error_code": "mcp_tool_not_available",
            }

        assignment = None
        if agent_id is not None:
            assignment_result = await db.execute(
                select(AgentTool).where(
                    AgentTool.agent_id == agent_id,
                    AgentTool.tool_id == tool.id,
                )
            )
            assignment = assignment_result.scalar_one_or_none()
        if assignment is None or not assignment.enabled:
            return {
                "full_name": str(tool.name or tool_name),
                "unavailable_error_code": "mcp_tool_not_available",
            }

        server_url = str(tool.mcp_server_url or "").strip()
        parsed = urlparse(server_url)
        if parsed.scheme not in {"http", "https"} or not parsed.netloc:
            return {
                "full_name": str(tool.name or tool_name),
                "unavailable_error_code": "mcp_configuration_missing",
            }
        raw_name = str(tool.mcp_tool_name or "").strip()
        if not raw_name and not allow_legacy_bare_name:
            return {
                "full_name": str(tool.name or tool_name),
                "unavailable_error_code": "mcp_configuration_missing",
            }
        trusted_async_completion = deepcopy(
            (tool.config or {}).get("async_completion")
        )
        merged_config = {
            **(tool.config or {}),
            **(assignment.config or {}),
        }
        merged_config = _decrypt_sensitive_fields(
            merged_config,
            tool.config_schema,
        )
        return {
            "full_name": str(tool.name),
            "raw_name": raw_name or str(tool.name),
            "server_url": server_url,
            "server_name": str(tool.mcp_server_name or ""),
            "config": merged_config,
            # Completion semantics are admin-owned Tool metadata. Per-Agent
            # config may supply credentials but cannot redefine completion.
            "async_completion": trusted_async_completion,
        }


async def _execute_resolved_mcp_target_outcome(
    target: dict,
    arguments: dict,
    *,
    agent_id,
) -> ToolExecutionOutcome:
    unavailable_error = target.get("unavailable_error_code")
    if unavailable_error:
        return _typed_failure(
            "MCP tool is not enabled, assigned, or locally configured.",
            str(unavailable_error),
        )

    from urllib.parse import urlparse

    import httpx

    from app.services.mcp_client import (
        MCPClient,
        MCPTransportDetectionError,
    )

    full_name = str(target["full_name"])
    raw_name = str(target["raw_name"])
    server_url = str(target["server_url"])
    server_name = str(target.get("server_name") or "")
    config = dict(target.get("config") or {})
    async_completion = target.get("async_completion")

    hostname = (urlparse(server_url).hostname or "").lower()
    if hostname.endswith(".run.tools"):
        return await _execute_via_smithery_connect_outcome(
            server_url,
            raw_name,
            arguments,
            config,
            agent_id=agent_id,
            full_tool_name=full_name,
            async_completion=async_completion,
        )

    direct_api_key = config.get("api_key") or config.get(
        "atlassian_api_key"
    )
    if not direct_api_key and server_name == "Atlassian Rovo":
        try:
            from app.api.atlassian import get_atlassian_api_key_for_agent

            direct_api_key = await get_atlassian_api_key_for_agent(agent_id)
        except Exception:
            direct_api_key = None

    client = MCPClient(server_url, api_key=direct_api_key)
    try:
        data = await client.call_tool_result(raw_name, arguments)
    except MCPTransportDetectionError:
        return _typed_failure(
            "MCP transport is not locally reachable; the tool was not dispatched.",
            "mcp_transport_unavailable",
        )
    except httpx.HTTPStatusError:
        return _typed_failure(
            "MCP provider explicitly rejected the call.",
            "mcp_provider_rejected",
        )
    except Exception:
        return _typed_unknown(
            "MCP call outcome is unknown after dispatch; reconcile before retrying.",
            "mcp_call_outcome_unknown",
        )
    return _mcp_call_response_outcome(
        data,
        full_tool_name=full_name,
        arguments=arguments,
        async_completion=async_completion,
    )


async def _execute_mcp_tool_outcome(
    tool_name: str,
    arguments: dict,
    agent_id=None,
) -> ToolExecutionOutcome:
    """Durable exact-name MCP execution adapter."""
    try:
        target = await _resolve_mcp_execution_target(
            tool_name,
            agent_id,
            allow_legacy_bare_name=False,
        )
    except Exception:
        logger.exception("[MCP] Exact tool resolution failed: {}", tool_name)
        return _typed_failure(
            "MCP tool assignment could not be resolved.",
            "mcp_tool_resolution_failed",
        )
    if target is None:
        return _typed_failure(
            "MCP tool is not enabled and assigned under that exact name.",
            "mcp_tool_not_available",
        )
    return await _execute_resolved_mcp_target_outcome(
        target,
        arguments,
        agent_id=agent_id,
    )


async def _execute_mcp_tool(tool_name: str, arguments: dict, agent_id=None) -> str:
    """Legacy text wrapper; bare-name compatibility is isolated here."""
    try:
        target = await _resolve_mcp_execution_target(
            tool_name,
            agent_id,
            allow_legacy_bare_name=True,
        )
        if target is None:
            return f"Unknown tool: {tool_name}"
        outcome = await _execute_resolved_mcp_target_outcome(
            target,
            arguments,
            agent_id=agent_id,
        )
        return _legacy_tool_outcome_text(
            outcome,
            fallback="MCP tool call did not return a summary.",
        )
    except Exception:
        logger.exception("[MCP] Legacy tool execution error: {}", tool_name)
        return "❌ MCP tool execution failed."


def _parse_mcp_json_or_sse(raw: str) -> dict:
    data = None
    for line in raw.splitlines():
        stripped = line.strip()
        if not stripped.startswith("data:"):
            continue
        candidate = stripped[5:].strip()
        if not candidate or candidate == "[DONE]":
            continue
        try:
            parsed = json.loads(candidate)
        except json.JSONDecodeError:
            continue
        if isinstance(parsed, dict):
            data = parsed
            if parsed.get("id") == 1:
                break
    if data is None:
        parsed = json.loads(raw)
        if not isinstance(parsed, dict):
            raise ValueError("MCP response must be an object")
        data = parsed
    return data


async def _execute_via_smithery_connect_outcome(
    mcp_url: str,
    tool_name: str,
    arguments: dict,
    config: dict,
    agent_id=None,
    *,
    full_tool_name: str,
    async_completion: object | None = None,
) -> ToolExecutionOutcome:
    """Execute one Smithery business call and preserve protocol status."""
    import httpx

    from app.services.resource_discovery import _get_smithery_api_key

    api_key = await _get_smithery_api_key(agent_id)
    if not api_key:
        return _typed_failure(
            "Smithery credentials are not configured for this agent.",
            "mcp_auth_required",
        )

    local_config = dict(config)
    namespace = local_config.get("smithery_namespace")
    connection_id = local_config.get("smithery_connection_id")
    if not namespace or not connection_id:
        try:
            from app.models.tool import Tool

            async with async_session() as db:
                result = await db.execute(
                    select(Tool).where(Tool.name == "discover_resources")
                )
                discovery_tool = result.scalar_one_or_none()
                discovery_config = (
                    discovery_tool.config
                    if discovery_tool and discovery_tool.config
                    else {}
                )
                namespace = namespace or discovery_config.get(
                    "smithery_namespace"
                )
                connection_id = connection_id or discovery_config.get(
                    "smithery_connection_id"
                )
        except Exception:
            pass
    if not namespace or not connection_id:
        return _typed_failure(
            "Smithery connection is not locally configured for this agent.",
            "mcp_configuration_missing",
        )

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
        "Accept": "application/json, text/event-stream",
    }
    try:
        async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
            response = await client.post(
                f"https://api.smithery.ai/connect/{namespace}/{connection_id}/mcp",
                json={
                    "jsonrpc": "2.0",
                    "id": 1,
                    "method": "tools/call",
                    "params": {"name": tool_name, "arguments": arguments},
                },
                headers=headers,
            )
    except Exception:
        return _typed_unknown(
            "Smithery MCP call outcome is unknown after dispatch; reconcile before retrying.",
            "mcp_call_outcome_unknown",
        )

    if response.status_code in {401, 403, 404}:
        try:
            await _smithery_auto_recover(
                api_key,
                mcp_url,
                str(namespace),
                str(connection_id),
                agent_id,
            )
        except Exception:
            pass
        return _typed_failure(
            "Smithery authorization is required before this MCP tool can run.",
            "mcp_auth_required",
        )
    if response.status_code >= 400:
        return _typed_failure(
            "Smithery explicitly rejected the MCP tool call.",
            "mcp_provider_rejected",
        )

    try:
        data = _parse_mcp_json_or_sse(response.text)
    except (TypeError, ValueError, json.JSONDecodeError):
        return _typed_unknown(
            "Smithery returned a malformed response; reconcile before retrying.",
            "mcp_malformed_response",
        )

    error = data.get("error") if isinstance(data, dict) else None
    error_message = (
        error.get("message") if isinstance(error, dict) else str(error or "")
    )
    auth_keywords = {
        "auth",
        "unauthorized",
        "forbidden",
        "expired",
        "not found",
        "connection",
    }
    normalized_error_message = error_message.lower()
    if error and (
        "http://" in normalized_error_message
        or "https://" in normalized_error_message
        or any(
            keyword in normalized_error_message for keyword in auth_keywords
        )
    ):
        try:
            await _smithery_auto_recover(
                api_key,
                mcp_url,
                str(namespace),
                str(connection_id),
                agent_id,
            )
        except Exception:
            pass
        return _typed_failure(
            "Smithery authorization is required before this MCP tool can run.",
            "mcp_auth_required",
        )
    return _mcp_call_response_outcome(
        data,
        full_tool_name=full_tool_name,
        arguments=arguments,
        async_completion=async_completion,
    )


async def _execute_via_smithery_connect(
    mcp_url: str,
    tool_name: str,
    arguments: dict,
    config: dict,
    agent_id=None,
) -> str:
    """Legacy text adapter for Smithery MCP execution."""
    outcome = await _execute_via_smithery_connect_outcome(
        mcp_url,
        tool_name,
        arguments,
        config,
        agent_id=agent_id,
        full_tool_name=tool_name,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Smithery MCP call did not return a summary.",
    )


async def _smithery_auto_recover(api_key: str, mcp_url: str, namespace: str, connection_id: str, agent_id=None) -> str | None:
    """Attempt to auto-recover a failed Smithery connection.

    Re-creates the Smithery Connect connection. If OAuth is needed,
    returns the auth URL for the user. Returns None if recovery fails silently.
    """
    try:
        from app.services.resource_discovery import _ensure_smithery_connection
        display_name = connection_id.replace("-", " ").title() if connection_id else "MCP Server"

        conn_result = await _ensure_smithery_connection(api_key, mcp_url, display_name)
        if "error" in conn_result:
            return (
                f"❌ MCP tool connection expired and auto-recovery failed: {conn_result['error']}\n\n"
                f"💡 Please re-authorize by telling me: `import_mcp_server(server_id=\"...\", reauthorize=true)`"
            )

        if conn_result.get("auth_url"):
            # A newly-created Smithery connection is not usable until the user
            # completes OAuth. Keep the existing stored connection in place so
            # a still-valid old connection is not overwritten by an unauthenticated
            # replacement. The user-facing auth URL is enough for recovery.
            return (
                f"🔐 MCP tool connection expired. Re-authorization needed.\n\n"
                f"Please visit the following URL to re-authorize:\n"
                f"{conn_result['auth_url']}\n\n"
                f"After completing authorization, the tools will work again automatically."
            )

        # Update stored config with new connection info
        new_config = {
            "smithery_namespace": conn_result["namespace"],
            "smithery_connection_id": conn_result["connection_id"],
        }
        if agent_id:
            try:
                from app.models.tool import Tool, AgentTool
                async with async_session() as db:
                    # Update all MCP tools for this server URL
                    r = await db.execute(
                        select(Tool).where(Tool.mcp_server_url == mcp_url, Tool.type == "mcp")
                    )
                    for tool in r.scalars().all():
                        at_r = await db.execute(
                            select(AgentTool).where(
                                AgentTool.agent_id == agent_id,
                                AgentTool.tool_id == tool.id,
                            )
                        )
                        at = at_r.scalar_one_or_none()
                        if at:
                            at.config = {**(at.config or {}), **new_config}
                    await db.commit()
            except Exception:
                pass  # Non-critical — connection may still work

        # Connection re-created without OAuth — should work now
        return None  # Signal caller to retry (but we don't retry here to avoid loops)

    except Exception as e:
        return f"❌ Auto-recovery failed: {str(e)[:200]}"


def _normalize_tool_rel_path(rel_path: str) -> str:
    normalized = unicodedata.normalize("NFC", (rel_path or "").strip()).replace("\\", "/")
    normalized = re.sub(r"/+", "/", normalized).lstrip("./")
    return normalized


def _collapse_filename_for_match(name: str) -> str:
    return re.sub(r"\s+", "", unicodedata.normalize("NFC", name or "")).casefold()


def _allowed_root_for_tool_path(ws: Path, rel_path: str, tenant_id: str | None = None) -> tuple[Path, str]:
    normalized = _normalize_tool_rel_path(rel_path)
    if normalized.startswith("enterprise_info"):
        enterprise_root = (
            (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
            if tenant_id
            else (WORKSPACE_ROOT / "enterprise_info").resolve()
        )
        sub = normalized[len("enterprise_info"):].lstrip("/")
        return enterprise_root, sub
    return ws.resolve(), normalized


def _resolve_tool_source_path(ws: Path, rel_path: str, tenant_id: str | None = None) -> Path:
    root, normalized = _allowed_root_for_tool_path(ws, rel_path, tenant_id=tenant_id)
    candidate = (root / normalized).resolve() if normalized else root
    if not candidate.is_relative_to(root):
        raise ValueError("Access denied for this path")
    if candidate.exists():
        return candidate

    parent = candidate.parent
    if parent.exists():
        wanted = _collapse_filename_for_match(candidate.name)
        for sibling in parent.iterdir():
            if _collapse_filename_for_match(sibling.name) == wanted:
                return sibling
    return candidate


def _resolve_tool_target_path(ws: Path, rel_path: str, tenant_id: str | None = None) -> Path:
    root, normalized = _allowed_root_for_tool_path(ws, rel_path, tenant_id=tenant_id)
    candidate = (root / normalized).resolve() if normalized else root
    if not candidate.is_relative_to(root):
        raise ValueError("❌ Access denied.")
    return candidate


def _tool_storage_key(agent_id: uuid.UUID, rel_path: str, tenant_id: str | None = None) -> tuple[str, str, bool]:
    normalized = normalize_workspace_path(_normalize_tool_rel_path(rel_path))
    if _is_enterprise_info_path(normalized):
        if not tenant_id:
            return normalize_storage_key("enterprise_info/" + normalized.removeprefix("enterprise_info").lstrip("/")), normalized, True
        sub = normalized[len("enterprise_info"):].lstrip("/")
        key = f"enterprise_info_{tenant_id}/{sub}" if sub else f"enterprise_info_{tenant_id}"
        return normalize_storage_key(key), normalized, True
    key = f"{agent_id}/{normalized}" if normalized else str(agent_id)
    return normalize_storage_key(key), normalized, False


def _display_size(size_bytes: int) -> str:
    return f"{size_bytes}B" if size_bytes < 1024 else f"{size_bytes / 1024:.1f}KB"


async def _storage_list_dir(agent_id: uuid.UUID, rel_path: str, tenant_id: str | None = None) -> str:
    storage = get_storage_backend()
    storage_key, normalized, is_enterprise = _tool_storage_key(agent_id, rel_path, tenant_id)

    exists = await storage.exists(storage_key)
    is_dir = await storage.is_dir(storage_key)
    if exists and not is_dir:
        return f"Path is not a directory: {rel_path}"
    if not exists and not is_dir and normalized:
        return f"Directory not found: {rel_path or '/'}"

    items: list[str] = []
    dir_count = 0
    file_count = 0
    if not normalized and tenant_id:
        items.append("  📁 enterprise_info/ (shared company info)")
        dir_count += 1

    entries = await storage.list_dir(storage_key) if exists or is_dir else []
    for entry in entries:
        if entry.name.startswith("."):
            continue
        if entry.is_dir:
            dir_count += 1
            try:
                child_count = len([c for c in await storage.list_dir(entry.key) if not c.name.startswith(".")])
            except Exception:
                child_count = 0
            items.append(f"  📁 {entry.name}/ ({child_count} items)")
        else:
            file_count += 1
            items.append(f"  📄 {entry.name} ({_display_size(entry.size)})")

    if not items:
        return f"📂 {rel_path or 'root'}: Empty directory (0 files, 0 folders)"
    header = f"📂 {rel_path or 'root'}: {dir_count} folder(s), {file_count} file(s)\n"
    return header + "\n".join(items)


async def _storage_read_file(
    agent_id: uuid.UUID,
    rel_path: str,
    tenant_id: str | None = None,
    offset: int = 0,
    limit: int = 2000,
) -> str:
    binary_error = _read_file_binary_error(rel_path)
    if binary_error is not None:
        return binary_error
    storage = get_storage_backend()
    storage_key, normalized, _ = _tool_storage_key(agent_id, rel_path, tenant_id)
    if not normalized:
        return "File not found: root"
    if not await storage.is_file(storage_key):
        return f"File not found: {rel_path}"
    try:
        content = await storage.read_text(storage_key, encoding="utf-8", errors="replace")
        lines = content.splitlines()
        total_lines = len(lines)
        start = max(0, offset)
        end = min(total_lines, start + limit)
        if start >= total_lines and total_lines > 0:
            return f"Offset {offset} exceeds file length ({total_lines} lines total)"
        selected_lines = lines[start:end]
        output = "\n".join(f"{i + 1:6}\t{line}" for i, line in enumerate(selected_lines, start=start))
        if total_lines > end:
            output += f"\n\n... [{total_lines - end} more lines not shown, lines {end + 1}-{total_lines}]"
        header = f"📄 {rel_path} (lines {start + 1 if total_lines else 0}-{end} of {total_lines})\n"
        return header + output
    except Exception as e:
        return f"Read failed: {e}"


async def _storage_walk_files(storage, root_key: str) -> list:
    out = []
    for entry in await storage.list_dir(root_key):
        if entry.name.startswith("."):
            continue
        out.append(entry)
        if entry.is_dir:
            out.extend(await _storage_walk_files(storage, entry.key))
    return out


def _relative_storage_display(entry_key: str, base_key: str, display_base: str) -> str:
    rel = entry_key.removeprefix(base_key.rstrip("/") + "/")
    return f"{display_base.rstrip('/')}/{rel}".strip("/") if display_base else rel


async def _storage_search_files(
    agent_id: uuid.UUID,
    pattern: str,
    path: str = ".",
    file_pattern: str = "*",
    ignore_case: bool = False,
    tenant_id: str | None = None,
) -> str:
    storage = get_storage_backend()
    rel_path = "" if path in ("", ".") else path
    base_key, normalized, _ = _tool_storage_key(agent_id, rel_path, tenant_id)
    if not await storage.is_dir(base_key) and normalized:
        return f"Directory not found: {path}"
    flags = re.IGNORECASE if ignore_case else 0
    try:
        regex = re.compile(pattern, flags)
    except re.error as e:
        return f"Invalid regex pattern: {e}"

    results: list[str] = []
    total_matches = 0
    files_searched = 0
    entries = await _storage_walk_files(storage, base_key) if await storage.is_dir(base_key) else []
    for entry in entries:
        if entry.is_dir:
            continue
        rel_display = _relative_storage_display(entry.key, base_key, normalized)
        if not fnmatch.fnmatch(Path(rel_display).name, file_pattern) and not fnmatch.fnmatch(rel_display, file_pattern):
            continue
        if Path(rel_display).suffix.lower() in {".pyc", ".pyo", ".so", ".dll", ".exe", ".bin", ".png", ".jpg", ".jpeg", ".gif", ".zip", ".tar", ".gz"}:
            continue
        files_searched += 1
        try:
            content = await storage.read_text(entry.key, encoding="utf-8", errors="ignore")
        except Exception:
            continue
        for i, line in enumerate(content.splitlines(), 1):
            if regex.search(line):
                results.append(f"{rel_display}:{i}: {line.strip()[:100]}")
                total_matches += 1
                if len(results) >= 50:
                    break
        if len(results) >= 50:
            break
    if not results:
        return f"No matches found for pattern '{pattern}' in {files_searched} file(s)"
    truncated = total_matches > len(results)
    truncation_note = f" (showing first {len(results)} of {total_matches}+ — refine pattern or path for more)" if truncated else ""
    return f"🔍 Found {total_matches}+ match(es) in {files_searched} file(s) for pattern '{pattern}'{truncation_note}:\n" + "\n".join(results)


async def _storage_find_files(
    agent_id: uuid.UUID,
    pattern: str,
    path: str = ".",
    tenant_id: str | None = None,
) -> str:
    storage = get_storage_backend()
    rel_path = "" if path in ("", ".") else path
    base_key, normalized, _ = _tool_storage_key(agent_id, rel_path, tenant_id)
    if not await storage.is_dir(base_key) and normalized:
        return f"Directory not found: {path}"
    entries = await _storage_walk_files(storage, base_key) if await storage.is_dir(base_key) else []
    matches = []
    for entry in entries:
        rel_display = _relative_storage_display(entry.key, base_key, normalized)
        if fnmatch.fnmatch(rel_display, pattern) or fnmatch.fnmatch(Path(rel_display).name, pattern):
            matches.append((entry, rel_display))
    if not matches:
        return f"No files matching pattern: {pattern}"
    results = []
    dir_count = 0
    file_count = 0
    for entry, rel_display in matches[:100]:
        if entry.is_dir:
            dir_count += 1
            results.append(f"📁 {rel_display}/")
        else:
            file_count += 1
            results.append(f"📄 {rel_display} ({_display_size(entry.size)})")
    return f"📂 Found {len(matches)} item(s) ({dir_count} dirs, {file_count} files) matching '{pattern}':\n" + "\n".join(results)


def _list_files(ws: Path, rel_path: str, tenant_id: str | None = None) -> str:
    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        # Remap: enterprise_info/... → enterprise_info_{tenant_id}/...
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        target = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(target).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        target = (ws / rel_path) if rel_path else ws
        target = target.resolve()
        if not str(target).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not target.exists():
        return f"Directory not found: {rel_path or '/'}"

    items = []
    # If listing root, also show enterprise_info entry
    if not rel_path:
        if tenant_id:
            enterprise_dir = WORKSPACE_ROOT / f"enterprise_info_{tenant_id}"
        else:
            enterprise_dir = WORKSPACE_ROOT / "enterprise_info"
        if enterprise_dir.exists():
            items.append("  📁 enterprise_info/ (shared company info)")

    dir_count = 0
    file_count = 0
    for p in sorted(target.iterdir()):
        if p.name.startswith("."):
            continue
        if p.is_dir():
            dir_count += 1
            child_count = len([c for c in p.iterdir() if not c.name.startswith(".")])
            items.append(f"  📁 {p.name}/ ({child_count} items)")
        elif p.is_file():
            file_count += 1
            size_bytes = p.stat().st_size
            if size_bytes < 1024:
                size_str = f"{size_bytes}B"
            else:
                size_str = f"{size_bytes/1024:.1f}KB"
            items.append(f"  📄 {p.name} ({size_str})")

    if not items:
        return f"📂 {rel_path or 'root'}: Empty directory (0 files, 0 folders)"

    header = f"📂 {rel_path or 'root'}: {dir_count} folder(s), {file_count} file(s)\n"
    return header + "\n".join(items)


def _read_file(ws: Path, rel_path: str, tenant_id: str | None = None, offset: int = 0, limit: int = 2000) -> str:
    """Read file contents with optional line range support.

    Args:
        ws: Workspace root path
        rel_path: Relative file path
        tenant_id: Optional tenant ID for enterprise_info
        offset: Starting line number (0-indexed)
        limit: Maximum number of lines to read

    Returns:
        File content with line numbers, or error message
    """
    binary_error = _read_file_binary_error(rel_path)
    if binary_error is not None:
        return binary_error

    try:
        file_path = _resolve_tool_source_path(ws, rel_path, tenant_id=tenant_id)
    except ValueError as exc:
        return str(exc)

    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        content = file_path.read_text(encoding="utf-8", errors="replace")
        lines = content.splitlines()
        total_lines = len(lines)

        # Apply offset and limit
        start = max(0, offset)
        end = min(total_lines, start + limit)

        if start >= total_lines:
            return f"Offset {offset} exceeds file length ({total_lines} lines total)"

        selected_lines = lines[start:end]

        # Format with line numbers (like cat -n)
        result = []
        for i, line in enumerate(selected_lines, start=start):
            result.append(f"{i+1:6}\t{line}")

        output = "\n".join(result)

        # Add pagination info if file is larger than what we show
        if total_lines > end:
            output += f"\n\n... [{total_lines - end} more lines not shown, lines {end+1}-{total_lines}]"

        # Add header with file info
        header = f"📄 {rel_path} (lines {start+1}-{end} of {total_lines})\n"
        return header + output

    except Exception as e:
        return f"Read failed: {e}"


_READ_DOCUMENT_MAX_FILE_BYTES = 50 * 1024 * 1024
_READ_DOCUMENT_TIMEOUT_SECONDS = 25
_READ_DOCUMENT_FALLBACK_TIMEOUT_SECONDS = 10
_READ_DOCUMENT_MAX_CELL_CHARS = 500
_READ_DOCUMENT_MAX_COLUMNS = 80
_READ_DOCUMENT_MAX_XLSX_CELLS = 20000


@dataclass(frozen=True, slots=True)
class DocumentReadResult:
    ok: bool
    content: str
    error_code: str | None = None
    retryable: bool = False


def _safe_document_cell_text(value: Any) -> str:
    """Convert spreadsheet/table values without letting pathological cells dominate CPU."""
    if value is None:
        return ""
    if isinstance(value, int) and value.bit_length() > 4096:
        return "[large integer omitted]"
    text = str(value)
    if len(text) > _READ_DOCUMENT_MAX_CELL_CHARS:
        return text[:_READ_DOCUMENT_MAX_CELL_CHARS] + "...[cell truncated]"
    return text


def _read_document_sync(
    ws: Path,
    rel_path: str,
    max_chars: int = 8000,
    tenant_id: str | None = None,
) -> DocumentReadResult:
    """Synchronous document extraction. Must run outside the uvicorn event loop."""
    max_chars = min(max(int(max_chars), 1), 20000)
    try:
        file_path = _resolve_tool_source_path(ws, rel_path, tenant_id=tenant_id)
    except ValueError as exc:
        return DocumentReadResult(False, str(exc), "workspace_path_invalid")

    if not file_path.exists():
        return DocumentReadResult(
            False,
            f"File not found: {rel_path}",
            "document_not_found",
        )
    if file_path.is_dir():
        return DocumentReadResult(
            False,
            f"Path is a directory, not a document: {rel_path}",
            "document_path_is_directory",
        )
    try:
        file_size = file_path.stat().st_size
    except OSError:
        file_size = 0
    if file_size > _READ_DOCUMENT_MAX_FILE_BYTES:
        return DocumentReadResult(
            False,
            (
                f"Document is too large to read safely ({file_size / 1024 / 1024:.1f} MB). "
                "Please split or convert it to a smaller text/Markdown excerpt first."
            ),
            "document_too_large",
        )

    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            import pdfplumber
            text_parts = []
            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages[:50]):  # Limit to 50 pages
                    page_text = page.extract_text() or ""
                    if page_text:
                        text_parts.append(f"--- Page {i+1} ---\n{page_text}")
                    if sum(len(part) for part in text_parts) >= max_chars:
                        break
            content = "\n\n".join(text_parts) if text_parts else "(PDF is empty or text extraction failed)"

        elif ext == ".docx":
            from docx import Document
            from docx.oxml.ns import qn
            doc = Document(str(file_path))
            lines: list[str] = []

            def _extract_para_text(para) -> str:
                return para.text.strip()

            def _extract_table(table) -> str:
                """Flatten a table into readable text."""
                rows = []
                for row in table.rows:
                    cells = [_safe_document_cell_text(cell.text).strip() for cell in row.cells[:_READ_DOCUMENT_MAX_COLUMNS]]
                    if not cells:
                        continue
                    # Remove duplicate adjacent cells (merged cells repeat)
                    deduped = [cells[0]] + [c for i, c in enumerate(cells[1:]) if c != cells[i]]
                    row_str = " | ".join(c for c in deduped if c)
                    if row_str:
                        rows.append(row_str)
                return "\n".join(rows)

            # 1. Main paragraphs
            for para in doc.paragraphs:
                t = _extract_para_text(para)
                if t:
                    lines.append(t)

            # 2. Tables in main body
            for table in doc.tables:
                t = _extract_table(table)
                if t:
                    lines.append(t)

            # 3. Text boxes / drawing shapes (wmf/shapes in body XML)
            for shape in doc.element.body.iter(qn("w:txbxContent")):
                for child in shape.iter(qn("w:t")):
                    if child.text and child.text.strip():
                        lines.append(child.text.strip())

            # 4. Headers and footers
            for section in doc.sections:
                for hf in [section.header, section.footer]:
                    if hf and hf.is_linked_to_previous is False:
                        for para in hf.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)

            content = "\n".join(lines) if lines else "(Document is empty or uses unsupported formatting)"

        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            sheets = []
            cell_count = 0
            for ws_name in wb.sheetnames[:10]:  # Limit to 10 sheets
                sheet = wb[ws_name]
                rows = []
                for row in sheet.iter_rows(max_row=200, max_col=_READ_DOCUMENT_MAX_COLUMNS, values_only=True):
                    visible = row
                    cell_count += len(visible)
                    if cell_count > _READ_DOCUMENT_MAX_XLSX_CELLS:
                        rows.append("[cell limit reached; remaining cells omitted]")
                        break
                    row_str = "\t".join(_safe_document_cell_text(c) for c in visible)
                    if row_str.strip():
                        rows.append(row_str)
                if rows:
                    sheets.append(f"=== Sheet: {ws_name} ===\n" + "\n".join(rows))
                if cell_count > _READ_DOCUMENT_MAX_XLSX_CELLS or sum(len(part) for part in sheets) >= max_chars:
                    break
            wb.close()
            content = "\n\n".join(sheets) if sheets else "(Excel is empty)"

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides[:50]):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
                if texts:
                    slides.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
            content = "\n\n".join(slides) if slides else "(PPT is empty)"

        elif ext in (".txt", ".md", ".json", ".csv", ".log"):
            content = file_path.read_text(encoding="utf-8", errors="replace")

        else:
            return DocumentReadResult(
                False,
                f"Unsupported file format: {ext}. Supported: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV",
                "document_format_unsupported",
            )

        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n...[truncated, {len(content)} chars total]"
        return DocumentReadResult(True, content)

    except ImportError as e:
        return DocumentReadResult(
            False,
            f"Missing dependency: {e}. Install: pip install pdfplumber python-docx openpyxl python-pptx",
            "document_dependency_missing",
        )
    except Exception as e:
        return DocumentReadResult(
            False,
            f"Document read failed: {str(e)[:200]}",
            "document_read_failed",
            retryable=True,
        )


def _read_document_worker(
    out_queue: mp.Queue,
    ws_str: str,
    rel_path: str,
    max_chars: int,
    tenant_id: str | None,
) -> None:
    try:
        out_queue.put(
            _read_document_sync(
                Path(ws_str),
                rel_path,
                max_chars=max_chars,
                tenant_id=tenant_id,
            )
        )
    except BaseException as exc:
        out_queue.put(
            DocumentReadResult(
                False,
                f"Document read failed: {str(exc)[:200]}",
                "document_read_failed",
                retryable=True,
            )
        )


def _read_pdf_fast_sync(
    ws: Path,
    rel_path: str,
    max_chars: int = 8000,
    tenant_id: str | None = None,
) -> DocumentReadResult:
    """Fast PDF text extraction fallback for files that make pdfplumber/pdfminer hang."""
    max_chars = min(max(int(max_chars), 1), 20000)
    try:
        file_path = _resolve_tool_source_path(ws, rel_path, tenant_id=tenant_id)
    except ValueError as exc:
        return DocumentReadResult(False, str(exc), "workspace_path_invalid")

    if not file_path.exists():
        return DocumentReadResult(False, f"File not found: {rel_path}", "document_not_found")
    if file_path.is_dir():
        return DocumentReadResult(
            False,
            f"Path is a directory, not a document: {rel_path}",
            "document_path_is_directory",
        )

    try:
        import fitz

        text_parts = []
        with fitz.open(str(file_path)) as doc:
            for i, page in enumerate(doc[:50]):
                page_text = page.get_text("text") or ""
                if page_text:
                    text_parts.append(f"--- Page {i+1} ---\n{page_text}")
                if sum(len(part) for part in text_parts) >= max_chars:
                    break
        content = "\n\n".join(text_parts) if text_parts else "(PDF is empty or text extraction failed)"
        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n...[truncated, {len(content)} chars total]"
        return DocumentReadResult(True, content)
    except ImportError as exc:
        return DocumentReadResult(
            False,
            f"PDF fallback extractor unavailable: {exc}. Install: pip install PyMuPDF",
            "document_dependency_missing",
        )
    except Exception as exc:
        return DocumentReadResult(
            False,
            f"PDF fallback extraction failed: {str(exc)[:200]}",
            "document_read_failed",
            retryable=True,
        )


def _read_pdf_fast_worker(
    out_queue: mp.Queue,
    ws_str: str,
    rel_path: str,
    max_chars: int,
    tenant_id: str | None,
) -> None:
    try:
        out_queue.put(
            _read_pdf_fast_sync(
                Path(ws_str),
                rel_path,
                max_chars=max_chars,
                tenant_id=tenant_id,
            )
        )
    except BaseException as exc:
        out_queue.put(
            DocumentReadResult(
                False,
                f"PDF fallback extraction failed: {str(exc)[:200]}",
                "document_read_failed",
                retryable=True,
            )
        )


def _read_pdf_fast_with_timeout(
    ws: Path,
    rel_path: str,
    max_chars: int = 8000,
    tenant_id: str | None = None,
) -> DocumentReadResult:
    ctx = mp.get_context("spawn")
    out_queue: mp.Queue = ctx.Queue(maxsize=1)
    proc = ctx.Process(
        target=_read_pdf_fast_worker,
        args=(out_queue, str(ws), rel_path, max_chars, tenant_id),
        daemon=True,
    )
    proc.start()
    proc.join(_READ_DOCUMENT_FALLBACK_TIMEOUT_SECONDS)
    if proc.is_alive():
        proc.terminate()
        proc.join(2)
        if proc.is_alive():
            proc.kill()
            proc.join(1)
        return DocumentReadResult(
            False,
            (
                f"Document read timed out after {_READ_DOCUMENT_TIMEOUT_SECONDS}s, "
                f"and PDF fallback also timed out after {_READ_DOCUMENT_FALLBACK_TIMEOUT_SECONDS}s. "
                "The file may be too large or too complex to extract safely."
            ),
            "document_read_timeout",
            retryable=True,
        )
    try:
        result = out_queue.get_nowait()
    except queue.Empty:
        if proc.exitcode:
            return DocumentReadResult(
                False,
                f"PDF fallback extraction failed: extractor exited with code {proc.exitcode}",
                "document_extractor_failed",
                retryable=True,
            )
        return DocumentReadResult(
            False,
            "PDF fallback extraction failed: extractor returned no content",
            "document_extractor_failed",
            retryable=True,
        )
    if isinstance(result, DocumentReadResult):
        return result
    return DocumentReadResult(
        False,
        "PDF fallback extractor returned an invalid result",
        "document_extractor_invalid",
        retryable=True,
    )


def _read_document_with_timeout(
    ws: Path,
    rel_path: str,
    max_chars: int = 8000,
    tenant_id: str | None = None,
) -> DocumentReadResult:
    """Run document parsing in a killable child process so one bad file cannot freeze the site."""
    ctx = mp.get_context("spawn")
    out_queue: mp.Queue = ctx.Queue(maxsize=1)
    proc = ctx.Process(
        target=_read_document_worker,
        args=(out_queue, str(ws), rel_path, max_chars, tenant_id),
        daemon=True,
    )
    proc.start()
    proc.join(_READ_DOCUMENT_TIMEOUT_SECONDS)
    if proc.is_alive():
        proc.terminate()
        proc.join(2)
        if proc.is_alive():
            proc.kill()
            proc.join(1)
        if Path(rel_path).suffix.lower() == ".pdf":
            return _read_pdf_fast_with_timeout(ws, rel_path, max_chars=max_chars, tenant_id=tenant_id)
        return DocumentReadResult(
            False,
            (
                f"Document read timed out after {_READ_DOCUMENT_TIMEOUT_SECONDS}s. "
                "The file may be too large or too complex to extract safely. "
                "Please split it, convert it to text/Markdown, or read a smaller excerpt."
            ),
            "document_read_timeout",
            retryable=True,
        )
    try:
        result = out_queue.get_nowait()
    except queue.Empty:
        if proc.exitcode:
            return DocumentReadResult(
                False,
                f"Document read failed: extractor exited with code {proc.exitcode}",
                "document_extractor_failed",
                retryable=True,
            )
        return DocumentReadResult(
            False,
            "Document read failed: extractor returned no content",
            "document_extractor_failed",
            retryable=True,
        )
    if isinstance(result, DocumentReadResult):
        return result
    return DocumentReadResult(
        False,
        "Document extractor returned an invalid result",
        "document_extractor_invalid",
        retryable=True,
    )


async def _read_document_result(
    ws: Path,
    rel_path: str,
    max_chars: int = 8000,
    tenant_id: str | None = None,
) -> DocumentReadResult:
    return await asyncio.to_thread(_read_document_with_timeout, ws, rel_path, max_chars, tenant_id)


async def _read_document(
    ws: Path,
    rel_path: str,
    max_chars: int = 8000,
    tenant_id: str | None = None,
) -> str:
    """Legacy display adapter for office document extraction."""
    return (
        await _read_document_result(
            ws,
            rel_path,
            max_chars=max_chars,
            tenant_id=tenant_id,
        )
    ).content


async def _read_document_from_storage(
    agent_id: uuid.UUID,
    rel_path: str,
    max_chars: int = 8000,
    tenant_id: str | None = None,
) -> str:
    temp_workspace = await _prepare_temp_workspace(agent_id, tenant_id=tenant_id, paths=[rel_path])
    try:
        return await _read_document(temp_workspace.root, rel_path, max_chars=max_chars, tenant_id=None)
    finally:
        temp_workspace.cleanup()


async def _read_document_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    tenant_id: str | None,
) -> ToolExecutionOutcome:
    path = arguments.get("path")
    if not isinstance(path, str) or not path:
        return _typed_failure(
            "read_document requires a non-empty path.",
            "invalid_tool_arguments",
        )
    try:
        max_chars = min(max(int(arguments.get("max_chars", 8000)), 1), 20000)
    except (TypeError, ValueError):
        return _typed_failure(
            "read_document max_chars must be an integer.",
            "invalid_tool_arguments",
        )
    try:
        temp_workspace = await _prepare_temp_workspace(
            agent_id,
            tenant_id=tenant_id,
            paths=[path],
        )
    except Exception as exc:
        return _typed_failure(
            f"Document could not be materialized: {type(exc).__name__}.",
            "document_materialize_failed",
            retryable=True,
        )
    try:
        result = await _read_document_result(
            temp_workspace.root,
            path,
            max_chars=max_chars,
            tenant_id=None,
        )
    except Exception as exc:
        return _typed_failure(
            f"Document extraction failed: {type(exc).__name__}.",
            "document_read_failed",
            retryable=True,
        )
    finally:
        temp_workspace.cleanup()
    if result.ok:
        return _typed_success(
            result.content,
            evidence_refs=(_workspace_artifact_ref(agent_id, path),),
        )
    return _typed_failure(
        result.content,
        result.error_code or "document_read_failed",
        retryable=result.retryable,
    )


# ─── Format Conversion Tools ────────────────────────────────────


def _validate_converted_artifact(path: Path, kind: str) -> bool:
    try:
        if not path.is_file() or path.stat().st_size <= 0:
            return False
        if kind == "pdf":
            data = path.read_bytes()
            return data.startswith(b"%PDF-") and b"%%EOF" in data[-2048:]
    except OSError:
        return False
    import zipfile

    try:
        with zipfile.ZipFile(path) as archive:
            if archive.testzip() is not None:
                return False
            names = set(archive.namelist())
    except (OSError, zipfile.BadZipFile):
        return False
    required = {
        "xlsx": {"[Content_Types].xml", "xl/workbook.xml"},
        "docx": {"[Content_Types].xml", "word/document.xml"},
        "pptx": {"[Content_Types].xml", "ppt/presentation.xml"},
    }[kind]
    return required <= names


async def _convert_file_outcome(
    agent_id: uuid.UUID,
    ws: Path,
    arguments: dict,
    *,
    tool_name: str,
) -> ToolExecutionOutcome:
    source_path = arguments.get("source_path")
    target_path = arguments.get("target_path")
    if not isinstance(source_path, str) or not source_path or not isinstance(target_path, str) or not target_path:
        return _typed_failure(
            f"{tool_name} requires source_path and target_path.",
            "invalid_tool_arguments",
        )
    try:
        source = _resolve_tool_source_path(ws, source_path)
        target = _resolve_tool_target_path(ws, target_path)
    except ValueError as exc:
        return _typed_failure(str(exc), "workspace_path_invalid")
    if not source.is_file():
        return _typed_failure(
            f"Source file not found: {source_path}",
            "conversion_source_not_found",
        )
    if source.resolve() == target.resolve():
        return _typed_failure(
            "Conversion source and target must be different files.",
            "invalid_tool_arguments",
        )

    converter_by_name = {
        "convert_csv_to_xlsx": (_convert_csv_to_xlsx, "xlsx"),
        "convert_html_to_pdf": (_convert_html_to_pdf, "pdf"),
        "convert_html_to_pptx": (_convert_html_to_pptx, "pptx"),
        "convert_markdown_to_docx": (_convert_markdown_to_docx, "docx"),
        "convert_markdown_to_pdf": (_convert_markdown_to_pdf, "pdf"),
    }
    converter, kind = converter_by_name[tool_name]
    try:
        previous = target.read_bytes() if target.is_file() else None
        if target.exists() and not target.is_file():
            return _typed_failure(
                f"Conversion target is not a file: {target_path}",
                "conversion_target_invalid",
            )
        target.unlink(missing_ok=True)
    except OSError as exc:
        return _typed_failure(
            f"Conversion target could not be prepared: {type(exc).__name__}.",
            "conversion_target_unavailable",
        )
    try:
        await converter(agent_id, ws, arguments)
        valid = _validate_converted_artifact(target, kind)
    except Exception as exc:
        valid = False
        logger.exception("[Conversion] Typed conversion failed: {}", tool_name)
        failure_class = type(exc).__name__
    else:
        failure_class = None
    if not valid:
        try:
            target.unlink(missing_ok=True)
            if previous is not None:
                target.parent.mkdir(parents=True, exist_ok=True)
                target.write_bytes(previous)
        except OSError as exc:
            return _typed_failure(
                f"Conversion failed and temporary rollback failed after {type(exc).__name__}.",
                "conversion_rollback_failed",
            )
        return _typed_failure(
            (
                f"{tool_name} did not produce a valid {kind.upper()} artifact"
                + (f" ({failure_class})." if failure_class else ".")
            ),
            "conversion_artifact_invalid",
        )
    return _typed_success(
        f"Converted {source_path} to {target_path}.",
        artifact_refs=(_workspace_artifact_ref(agent_id, target_path),),
    )

async def _convert_csv_to_xlsx(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    source_path = arguments.get("source_path")
    target_path = arguments.get("target_path")
    if not source_path or not target_path:
        return "❌ Missing 'source_path' or 'target_path'."
    try:
        src_file = _resolve_tool_source_path(ws, source_path)
        tgt_file = _resolve_tool_target_path(ws, target_path)
    except ValueError as exc:
        return str(exc)
    if not src_file.exists(): return f"❌ Source file not found: {source_path}"

    try:
        import csv
        from openpyxl import Workbook

        text = src_file.read_text(encoding="utf-8-sig")
        lines = [line.strip() for line in text.splitlines() if line.strip()][:10]
        candidates = [",", "，", ";", "\t", "|"]
        delimiter = ","
        if lines:
            scores = {candidate: sum(line.count(candidate) for line in lines) for candidate in candidates}
            if any(scores.values()):
                delimiter = max(scores, key=scores.get)

        wb = Workbook()
        ws_sheet = wb.active
        with src_file.open("r", encoding="utf-8-sig", newline="") as f:
            reader = csv.reader(f, delimiter=delimiter)
            for row in reader:
                values = list(row)
                while values and not str(values[-1] or "").strip():
                    values.pop()
                if values:
                    ws_sheet.append(values)

        tgt_file.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(tgt_file))
        return f"✅ Successfully converted CSV to Excel: {target_path}"
    except Exception as e:
        logger.exception(f"Convert CSV to XLSX failed: {e}")
        return f"❌ Conversion failed: {e}"

async def _convert_html_to_pdf(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    source_path = arguments.get("source_path")
    target_path = arguments.get("target_path")
    if not source_path or not target_path:
        return "❌ Missing 'source_path' or 'target_path'."
    try:
        src_file = _resolve_tool_source_path(ws, source_path)
        tgt_file = _resolve_tool_target_path(ws, target_path)
    except ValueError as exc:
        return str(exc)
    if not src_file.exists():
        return f"❌ Source file not found: {source_path}"

    return await convert_html_file_to_pdf(src_file, tgt_file, str(target_path), arguments)


async def _convert_html_to_pptx(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    source_path = arguments.get("source_path")
    target_path = arguments.get("target_path")
    if not source_path or not target_path:
        return "❌ Missing paths."
    try:
        src_file = _resolve_tool_source_path(ws, source_path)
        tgt_file = _resolve_tool_target_path(ws, target_path)
    except ValueError as exc:
        return str(exc)
    if not src_file.exists():
        return "❌ Source file not found."

    return await convert_html_file_to_pptx(src_file, tgt_file, str(target_path), ws, arguments)

async def _convert_markdown_to_docx(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    source_path = arguments.get("source_path")
    target_path = arguments.get("target_path")
    if not source_path or not target_path: return "❌ Missing paths."
    try:
        src_file = _resolve_tool_source_path(ws, source_path)
        tgt_file = _resolve_tool_target_path(ws, target_path)
    except ValueError as exc:
        return str(exc)
    if not src_file.exists(): return "❌ Source file not found."

    try:
        from docx import Document
        md_text = src_file.read_text(encoding="utf-8")
        doc = Document()

        def flush_paragraph(lines: list[str]) -> None:
            text = " ".join(line.strip() for line in lines if line.strip()).strip()
            if text:
                doc.add_paragraph(text)

        paragraph_lines: list[str] = []
        lines = md_text.splitlines()
        i = 0
        while i < len(lines):
            line = lines[i].rstrip()
            stripped = line.strip()

            if not stripped:
                flush_paragraph(paragraph_lines)
                paragraph_lines = []
                i += 1
                continue

            heading_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
            if heading_match:
                flush_paragraph(paragraph_lines)
                paragraph_lines = []
                level = min(len(heading_match.group(1)), 6)
                doc.add_heading(heading_match.group(2).strip(), level=level)
                i += 1
                continue

            bullet_match = re.match(r"^[-*+]\s+(.*)$", stripped)
            ordered_match = re.match(r"^\d+\.\s+(.*)$", stripped)
            if bullet_match or ordered_match:
                flush_paragraph(paragraph_lines)
                paragraph_lines = []
                text = (bullet_match or ordered_match).group(1).strip()
                if text:
                    doc.add_paragraph(text, style="List Bullet" if bullet_match else "List Number")
                i += 1
                continue

            if "|" in stripped:
                table_lines: list[str] = []
                flush_paragraph(paragraph_lines)
                paragraph_lines = []
                while i < len(lines) and "|" in lines[i]:
                    candidate = lines[i].strip()
                    if candidate:
                        table_lines.append(candidate)
                    i += 1
                data_rows = []
                for raw in table_lines:
                    cells = [cell.strip() for cell in raw.strip("|").split("|")]
                    if cells and all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells):
                        continue
                    if any(cell for cell in cells):
                        data_rows.append(cells)
                if data_rows:
                    table = doc.add_table(rows=len(data_rows), cols=max(len(row) for row in data_rows))
                    table.style = "Table Grid"
                    for row_idx, row in enumerate(data_rows):
                        for col_idx, cell in enumerate(row):
                            table.cell(row_idx, col_idx).text = cell
                continue

            paragraph_lines.append(stripped)
            i += 1

        flush_paragraph(paragraph_lines)

        tgt_file.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(tgt_file))
        return f"✅ Successfully converted Markdown to Word: {target_path}"
    except Exception as e:
        logger.exception(f"Convert MD to Docx failed: {e}")
        return f"❌ Conversion failed: {e}"

async def _convert_markdown_to_pdf(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    source_path = arguments.get("source_path")
    target_path = arguments.get("target_path")
    if not source_path or not target_path: return "❌ Missing paths."
    try:
        src_file = _resolve_tool_source_path(ws, source_path)
        tgt_file = _resolve_tool_target_path(ws, target_path)
    except ValueError as exc:
        return str(exc)
    if not src_file.exists(): return "❌ Source file not found."

    try:
        from weasyprint import HTML

        md_text = src_file.read_text(encoding="utf-8")

        def escape_html(text: str) -> str:
            return (
                text.replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
                .replace('"', "&quot;")
            )

        def render_inline(text: str) -> str:
            text = escape_html(text)
            text = re.sub(r"\*\*\*(.*?)\*\*\*", r"<strong><em>\1</em></strong>", text)
            text = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", text)
            text = re.sub(r"__(.*?)__", r"<strong>\1</strong>", text)
            text = re.sub(r"\*(.*?)\*", r"<em>\1</em>", text)
            text = re.sub(r"_(.*?)_", r"<em>\1</em>", text)
            text = re.sub(r"`([^`]+)`", r"<code>\1</code>", text)
            text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r'<a href="\2">\1</a>', text)
            return text

        def is_table_separator(line: str) -> bool:
            cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
            return bool(cells) and all(re.fullmatch(r":?-{3,}:?", cell or "") for cell in cells)

        html_parts: list[str] = []
        lines = md_text.splitlines()
        in_list = False
        i = 0
        while i < len(lines):
            raw_line = lines[i]
            line = raw_line.rstrip()
            stripped = line.strip()
            if not stripped:
                if in_list:
                    html_parts.append("</ul>")
                    in_list = False
                i += 1
                continue

            heading_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
            if heading_match:
                if in_list:
                    html_parts.append("</ul>")
                    in_list = False
                level = len(heading_match.group(1))
                html_parts.append(f"<h{level}>{render_inline(heading_match.group(2).strip())}</h{level}>")
                i += 1
                continue

            bullet_match = re.match(r"^[-*+]\s+(.*)$", stripped)
            if bullet_match:
                if not in_list:
                    html_parts.append("<ul>")
                    in_list = True
                html_parts.append(f"<li>{render_inline(bullet_match.group(1).strip())}</li>")
                i += 1
                continue

            if "|" in stripped and i + 1 < len(lines) and is_table_separator(lines[i + 1].strip()):
                if in_list:
                    html_parts.append("</ul>")
                    in_list = False
                header_cells = [render_inline(cell.strip()) for cell in stripped.strip("|").split("|")]
                table_rows: list[list[str]] = []
                i += 2
                while i < len(lines) and "|" in lines[i].strip():
                    row = [render_inline(cell.strip()) for cell in lines[i].strip().strip("|").split("|")]
                    table_rows.append(row)
                    i += 1
                html_parts.append("<table><thead><tr>" + "".join(f"<th>{cell}</th>" for cell in header_cells) + "</tr></thead><tbody>")
                html_parts.extend(
                    "<tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>"
                    for row in table_rows
                )
                html_parts.append("</tbody></table>")
                continue

            if in_list:
                html_parts.append("</ul>")
                in_list = False
            html_parts.append(f"<p>{render_inline(stripped)}</p>")
            i += 1

        if in_list:
            html_parts.append("</ul>")

        html_text = "\n".join(html_parts)

        full_html = (
            "<html><head><meta charset='utf-8'><style>"
            "body{font-family:'WenQuanYi Micro Hei','Noto Sans CJK SC',sans-serif;line-height:1.65;padding:2em;color:#111827;}"
            "h1,h2,h3{line-height:1.25;margin:1.2em 0 .55em;}"
            "p{margin:.55em 0;}"
            "table{width:100%;border-collapse:collapse;margin:1em 0;font-size:12px;}"
            "th,td{border:1px solid #d8dee9;padding:7px 9px;text-align:left;vertical-align:top;}"
            "th{background:#f3f4f6;font-weight:700;}"
            "code{background:#f3f4f6;padding:1px 4px;border-radius:4px;}"
            "a{color:#2563eb;text-decoration:none;}"
            "</style></head><body>"
            f"{html_text}"
            "</body></html>"
        )

        tgt_file.parent.mkdir(parents=True, exist_ok=True)
        HTML(string=full_html, base_url=str(ws.resolve())).write_pdf(str(tgt_file))
        return f"✅ Successfully converted Markdown to PDF: {target_path}"
    except Exception as e:
        logger.exception(f"Convert MD to PDF failed: {e}")
        return f"❌ Conversion failed: {e}"


def _write_file(ws: Path, rel_path: str, content: str, tenant_id: str | None = None) -> str:
    # Protect legacy DB-backed tasks.json from direct writes
    if rel_path.strip("/") == "tasks.json":
        return "tasks.json is a legacy read-only snapshot. Use the task APIs/UI to manage tasks."

    if _is_enterprise_info_path(rel_path):
        return "enterprise_info is shared company context and is read-only for agents. Ask an admin to update it."

    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        if not sub:
            return "Write failed: please provide a file path under enterprise_info/, e.g. enterprise_info/knowledge_base/report.md"
        file_path = (enterprise_root / sub).resolve()
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    try:
        file_path.parent.mkdir(parents=True, exist_ok=True)
        file_path.write_text(content, encoding="utf-8")
        return f"✅ Written to {rel_path} ({len(content)} chars)"
    except Exception as e:
        return f"Write failed: {e}"


def _delete_file(ws: Path, rel_path: str) -> str:
    protected = {"tasks.json", "soul.md"}
    if rel_path.strip("/") in protected:
        return f"{rel_path} cannot be deleted (protected)"
    if _is_enterprise_info_path(rel_path):
        return "enterprise_info is shared company context and is read-only for agents. Ask an admin to update it."

    file_path = (ws / rel_path).resolve()
    if not str(file_path).startswith(str(ws.resolve())):
        return "Access denied for this path"
    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        if file_path.is_dir():
            import shutil
            shutil.rmtree(file_path)
            return f"✅ Deleted directory {rel_path}"
        else:
            file_path.unlink()
            return f"✅ Deleted {rel_path}"
    except Exception as e:
        return f"Delete failed: {e}"


def _edit_file(ws: Path, rel_path: str, old_string: str, new_string: str, replace_all: bool = False, tenant_id: str | None = None) -> str:
    """Perform surgical string replacement in a file.

    Args:
        ws: Workspace root path
        rel_path: Relative file path
        old_string: Exact text to find and replace
        new_string: Replacement text
        replace_all: Replace all occurrences if True
        tenant_id: Optional tenant ID for enterprise_info

    Returns:
        Success message or error
    """
    if _is_enterprise_info_path(rel_path):
        return "enterprise_info is shared company context and is read-only for agents. Ask an admin to update it."

    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        file_path = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not file_path.exists():
        return f"File not found: {rel_path}"

    if not file_path.is_file():
        return f"Not a file: {rel_path}"

    try:
        content = file_path.read_text(encoding="utf-8")

        if old_string not in content:
            return f"❌ 'old_string' not found in {rel_path}. Please check the exact text including whitespace and newlines."

        if replace_all:
            new_content = content.replace(old_string, new_string)
            count = content.count(old_string)
        else:
            # Ensure uniqueness for single replacement
            count = content.count(old_string)
            if count > 1:
                return f"❌ 'old_string' appears {count} times in {rel_path}. Use replace_all=true or provide more context to make the match unique."
            new_content = content.replace(old_string, new_string, 1)
            count = 1

        file_path.write_text(new_content, encoding="utf-8")
        return f"✅ Replaced {count} occurrence(s) in {rel_path}"

    except Exception as e:
        return f"Edit failed: {e}"


def _search_files(ws: Path, pattern: str, path: str = ".", file_pattern: str = "*", ignore_case: bool = False, tenant_id: str | None = None) -> str:
    """Search for content patterns across files using regex.

    Args:
        ws: Workspace root path
        pattern: Regex pattern to search for
        path: Directory to search in (relative to workspace root)
        file_pattern: File pattern to match (glob)
        ignore_case: Case-insensitive search
        tenant_id: Optional tenant ID for enterprise_info

    Returns:
        Matching lines with file paths and line numbers
    """
    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if path and path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = path[len("enterprise_info"):].lstrip("/")
        search_path = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(search_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
        ws_for_relative = enterprise_root
    else:
        search_path = (ws / path).resolve() if path and path != "." else ws
        if not str(search_path).startswith(str(ws.resolve())):
            return "Access denied for this path"
        ws_for_relative = ws

    if not search_path.exists():
        return f"Directory not found: {path}"

    flags = re.IGNORECASE if ignore_case else 0

    try:
        regex = re.compile(pattern, flags)
    except re.error as e:
        return f"Invalid regex pattern: {e}"

    results = []
    total_matches = 0
    files_searched = 0

    # Use rglob for recursive search
    for file_path in search_path.rglob(file_pattern):
        if not file_path.is_file():
            continue
        # Skip hidden files and common binary/extensions
        if file_path.name.startswith("."):
            continue
        suffix = file_path.suffix.lower()
        if suffix in {".pyc", ".pyo", ".so", ".dll", ".exe", ".bin", ".png", ".jpg", ".jpeg", ".gif", ".zip", ".tar", ".gz"}:
            continue

        files_searched += 1
        try:
            content = file_path.read_text(encoding="utf-8", errors="ignore")
            for i, line in enumerate(content.splitlines(), 1):
                if regex.search(line):
                    rel_path = file_path.relative_to(ws_for_relative)
                    # Truncate long lines
                    display_line = line.strip()[:100]
                    results.append(f"{rel_path}:{i}: {display_line}")
                    total_matches += 1
                    if len(results) >= 50:  # Limit results per query
                        break
        except Exception:
            continue

        if len(results) >= 50:
            break

    if not results:
        return f"No matches found for pattern '{pattern}' in {files_searched} file(s)"

    # Warn the LLM if results were capped so it knows to refine the search.
    truncated = total_matches > len(results)
    truncation_note = f" (showing first {len(results)} of {total_matches}+ — refine pattern or path for more)" if truncated else ""
    header = f"🔍 Found {total_matches}+ match(es) in {files_searched} file(s) for pattern '{pattern}'{truncation_note}:\n"
    return header + "\n".join(results)


def _find_files(ws: Path, pattern: str, path: str = ".", tenant_id: str | None = None) -> str:
    """Find files matching glob patterns.

    Args:
        ws: Workspace root path
        pattern: Glob pattern to match files
        path: Base directory for search (relative to workspace root)
        tenant_id: Optional tenant ID for enterprise_info

    Returns:
        List of matching files with sizes
    """
    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if path and path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = path[len("enterprise_info"):].lstrip("/")
        search_path = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(search_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
        ws_for_relative = enterprise_root
    else:
        search_path = (ws / path).resolve() if path and path != "." else ws
        if not str(search_path).startswith(str(ws.resolve())):
            return "Access denied for this path"
        ws_for_relative = ws

    if not search_path.exists():
        return f"Directory not found: {path}"

    try:
        matches = list(search_path.glob(pattern))
    except Exception as e:
        return f"Invalid glob pattern: {e}"

    if not matches:
        return f"No files matching pattern: {pattern}"

    # Sort by modification time (most recent first)
    matches.sort(key=lambda x: x.stat().st_mtime if x.exists() else 0, reverse=True)

    results = []
    dir_count = 0
    file_count = 0

    for m in matches[:100]:  # Limit to 100 results
        rel_path = m.relative_to(ws_for_relative)
        if m.is_dir():
            dir_count += 1
            results.append(f"📁 {rel_path}/")
        else:
            file_count += 1
            try:
                size = m.stat().st_size
                size_str = f"{size//1024}KB" if size > 1024 else f"{size}B"
                results.append(f"📄 {rel_path} ({size_str})")
            except Exception:
                results.append(f"📄 {rel_path}")

    header = f"📂 Found {len(matches)} item(s) ({dir_count} dirs, {file_count} files) matching '{pattern}':\n"
    return header + "\n".join(results)


async def _manage_tasks(
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    ws: Path,
    args: dict,
) -> str:
    """Create / update / delete tasks in DB and sync to workspace."""
    from app.models.task import TaskLog
    from datetime import datetime, timezone

    action = args["action"]
    title = args["title"]

    async with async_session() as db:
        if action == "create":
            task_type = args.get("task_type", "todo")
            task = Task(
                agent_id=agent_id,
                title=title,
                description=args.get("description"),
                type=task_type,
                priority=args.get("priority", "medium"),
                created_by=user_id,
                status="pending",
                supervision_target_name=args.get("supervision_target_name"),
                supervision_channel=args.get("supervision_channel", "feishu"),
                remind_schedule=args.get("remind_schedule"),
            )
            db.add(task)
            await db.commit()
            await db.refresh(task)

            if task_type == "todo":
                # Trigger auto-execution for todo tasks
                import asyncio
                from app.services.task_executor import execute_task
                asyncio.create_task(execute_task(task.id, agent_id))
                await _sync_tasks_to_file(agent_id, ws)
                return f"✅ Task created: {title} — auto-execution started"
            else:
                # Supervision task — reminder engine will pick it up
                target = args.get('supervision_target_name', 'someone')
                schedule = args.get('remind_schedule', 'not set')
                await _sync_tasks_to_file(agent_id, ws)
                return f"✅ Supervision task created: '{title}' — will remind {target} on schedule ({schedule})"

        elif action == "update_status":
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id, Task.title.ilike(f"%{title}%"))
            )
            task = result.scalars().first()
            if not task:
                return f"No task found matching '{title}'"
            old = task.status
            task.status = args["status"]
            if args["status"] == "done":
                task.completed_at = datetime.now(timezone.utc)
            await db.commit()
            await _sync_tasks_to_file(agent_id, ws)
            return f"✅ Updated '{task.title}' from {old} to {args['status']}"

        elif action == "delete":
            from sqlalchemy import delete as sa_delete
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id, Task.title.ilike(f"%{title}%"))
            )
            task = result.scalars().first()
            if not task:
                return f"No task found matching '{title}'"
            task_title = task.title
            await db.execute(sa_delete(TaskLog).where(TaskLog.task_id == task.id))
            await db.delete(task)
            await db.commit()
            await _sync_tasks_to_file(agent_id, ws)
            return f"✅ Task deleted: {task_title}"

        return f"Unknown action: {action}"


def _json_tool_result(payload: dict) -> str:
    return json.dumps(payload, ensure_ascii=False)


def _provider_type_value(provider_type: Any) -> str | None:
    if provider_type is None:
        return None
    return getattr(provider_type, "value", provider_type)


def _normalize_roster_provider_type(provider_type: Any) -> str | None:
    value = _provider_type_value(provider_type)
    if value is None:
        return None
    normalized = str(value).strip().lower()
    if not normalized:
        return None
    if normalized == "microsoft_teams":
        return "teams"
    return normalized


@dataclass(frozen=True)
class RosterHumanTarget:
    source_agent: AgentModel
    member: OrgMember
    provider: IdentityProvider | None
    provider_type: str | None
    platform_user: UserModel | None


def _member_has_provider_identity(member: OrgMember) -> bool:
    return bool(
        (getattr(member, "external_id", None) or "").strip()
        or (getattr(member, "open_id", None) or "").strip()
    )


def _provider_identity_condition(provider_user_id: str):
    return or_(
        OrgMember.external_id == provider_user_id,
        OrgMember.open_id == provider_user_id,
        OrgMember.unionid == provider_user_id,
    )


async def _resolve_roster_human_target(
    db,
    agent_id: uuid.UUID,
    *,
    target_member_id: str | None = None,
    platform_user_id: str | None = None,
    provider_user_id: str | None = None,
    member_name: str | None = None,
    provider_type: str | None = None,
    require_platform_user: bool = False,
    require_provider_identity: bool = False,
) -> tuple[RosterHumanTarget | None, str | None]:
    source_result = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
    source_agent = source_result.scalar_one_or_none()
    if not source_agent:
        return None, "❌ Source agent was not found."

    target_member_id_raw = (target_member_id or "").strip()
    platform_user_id_raw = (platform_user_id or "").strip()
    provider_user_id_raw = (provider_user_id or "").strip()
    member_name_raw = (member_name or "").strip()
    requested_provider_type = _normalize_roster_provider_type(provider_type)

    lookup_kind = ""
    conditions = [OrgMember.tenant_id == source_agent.tenant_id]
    if target_member_id_raw:
        lookup_kind = "target_member_id"
        try:
            member_id = uuid.UUID(target_member_id_raw)
        except ValueError:
            return None, "❌ Invalid target_member_id. Use query_directory to get a valid target_member_id."
        conditions.append(OrgMember.id == member_id)
    elif platform_user_id_raw:
        lookup_kind = "platform_user_id"
        try:
            user_id = uuid.UUID(platform_user_id_raw)
        except ValueError:
            return None, "❌ Invalid platform_user_id. Use query_directory to get a valid platform_user_id."
        conditions.append(OrgMember.user_id == user_id)
        require_platform_user = True
    elif provider_user_id_raw:
        lookup_kind = "provider_user_id"
        conditions.append(_provider_identity_condition(provider_user_id_raw))
        require_provider_identity = True
    elif member_name_raw:
        lookup_kind = "member_name"
        conditions.append(OrgMember.name == member_name_raw)
    else:
        return None, "❌ Please provide target_member_id, platform_user_id, provider_user_id, or member_name."

    result = await db.execute(
        select(OrgMember, IdentityProvider)
        .outerjoin(IdentityProvider, OrgMember.provider_id == IdentityProvider.id)
        .where(*conditions)
        .order_by(OrgMember.name.asc(), OrgMember.synced_at.asc())
        .limit(20)
    )
    rows = result.all()
    if not rows:
        return None, "❌ Human recipient not found. Use query_directory to find an available human target."

    candidates: list[RosterHumanTarget] = []
    blocked_reason: str | None = None
    for member, provider in rows:
        authorized_custom_human = False
        if getattr(source_agent, "access_mode", None) == "custom":
            authorized_custom_human = await agent_directory.is_custom_human_authorized(
                db,
                source=source_agent,
                member=member,
            )
        visibility = evaluate_roster_human_visibility(
            source_agent,
            member,
            authorized_custom_human=authorized_custom_human,
        )
        if not visibility.visible:
            blocked_reason = blocked_reason or "not_visible"
            continue
        if not visibility.can_contact:
            blocked_reason = blocked_reason or visibility.unavailable_reason or "not_contactable"
            continue

        member_provider_type = _normalize_roster_provider_type(getattr(provider, "provider_type", None))
        if requested_provider_type and member_provider_type != requested_provider_type:
            blocked_reason = blocked_reason or "provider_type_mismatch"
            continue
        if require_provider_identity:
            if not _member_has_provider_identity(member):
                blocked_reason = blocked_reason or "missing_provider_identity"
                continue
            if not member_provider_type:
                blocked_reason = blocked_reason or "missing_provider_type"
                continue

        platform_user = None
        if getattr(member, "user_id", None):
            user_result = await db.execute(select(UserModel).where(UserModel.id == member.user_id))
            platform_user = user_result.scalar_one_or_none()
            if platform_user and platform_user.tenant_id != source_agent.tenant_id:
                platform_user = None
            if platform_user and not getattr(platform_user, "is_active", False):
                platform_user = None
        if require_platform_user and not platform_user:
            blocked_reason = blocked_reason or "missing_platform_user"
            continue

        candidates.append(RosterHumanTarget(
            source_agent=source_agent,
            member=member,
            provider=provider,
            provider_type=member_provider_type,
            platform_user=platform_user,
        ))

    if not candidates:
        if requested_provider_type and blocked_reason == "provider_type_mismatch":
            return None, f"❌ Human recipient was found, but not in {requested_provider_type} channel."
        return None, f"❌ Human recipient is not contactable ({blocked_reason or 'restricted'}). Use query_directory to choose an available person."
    if len(candidates) > 1:
        if lookup_kind == "member_name":
            return None, "❌ Multiple human recipients match this member_name. Use query_directory and retry with target_member_id."
        return None, "❌ Multiple human recipients match this identifier. Use query_directory and retry with target_member_id."

    return candidates[0], None


def _query_text_match_rank(member: dict, query: str) -> int:
    if not query:
        return 4
    q = query.casefold()
    display_name = (member.get("display_name") or "").casefold()
    if display_name == q:
        return 0
    if display_name.startswith(q):
        return 1
    if q in display_name:
        return 2
    return 3


def _roster_sort_key(member: dict, query: str) -> tuple:
    return agent_directory.roster_sort_key(member, query)


def _department_name(member: OrgMember, department: OrgDepartment | None) -> str | None:
    return agent_directory.department_name(member, department)


def _format_roster_agent(source_agent: AgentModel, target_agent: AgentModel) -> dict | None:
    return agent_directory.format_roster_agent(source_agent, target_agent)


def _format_roster_human(
    source_agent: AgentModel,
    member: OrgMember,
    provider: IdentityProvider | None,
    department: OrgDepartment | None,
    platform_user: UserModel | None = None,
) -> dict | None:
    return agent_directory.format_roster_human(source_agent, member, provider, department, platform_user)


async def _query_directory_payload(agent_id: uuid.UUID, args: dict) -> dict:
    """Return the Directory business payload before display serialization."""
    query = (args.get("query") or "").strip()
    target_member_id_raw = (args.get("target_member_id") or "").strip()
    member_type = (args.get("member_type") or "all").strip().lower()
    include_uncontactable = bool(args.get("include_uncontactable", False))

    try:
        limit = int(args.get("limit", 20))
    except (TypeError, ValueError):
        return {
            "ok": False,
            "error": {"code": "invalid_limit", "message": "limit must be between 1 and 50"},
        }
    try:
        offset = int(args.get("offset", 0))
    except (TypeError, ValueError):
        return {
            "ok": False,
            "error": {"code": "invalid_offset", "message": "offset must be greater than or equal to 0"},
        }

    if member_type not in {"all", "agent", "human"}:
        return {
            "ok": False,
            "error": {"code": "invalid_member_type", "message": "member_type must be all, agent, or human"},
        }
    if limit < 1 or limit > 50:
        return {
            "ok": False,
            "error": {"code": "invalid_limit", "message": "limit must be between 1 and 50"},
        }
    if offset < 0:
        return {
            "ok": False,
            "error": {"code": "invalid_offset", "message": "offset must be greater than or equal to 0"},
        }
    target_member_id = None
    if target_member_id_raw:
        try:
            target_member_id = uuid.UUID(target_member_id_raw)
        except ValueError:
            return {
                "ok": False,
                "error": {"code": "invalid_target_member_id", "message": "target_member_id must be a valid UUID"},
            }
        if member_type == "agent":
            return {
                "ok": False,
                "error": {
                    "code": "invalid_member_type",
                    "message": "target_member_id can only be used with member_type human or all",
                },
            }

    try:
        async with async_session() as db:
            result = await agent_directory.query_agent_directory(
                db,
                source_agent_id=agent_id,
                query=query,
                target_member_id=target_member_id,
                member_type=member_type,
                include_uncontactable=include_uncontactable,
                limit=limit,
                offset=offset,
                max_limit=50,
            )
        return result
    except agent_directory.DirectoryQueryError as e:
        return {
            "ok": False,
            "error": {"code": e.code, "message": e.message},
        }
    except Exception as e:
        logger.exception(f"[Directory] query_directory failed: agent={agent_id}")
        return {
            "ok": False,
            "error": {"code": "query_directory_failed", "message": f"query_directory failed: {type(e).__name__}"},
        }


async def _query_directory_outcome(
    agent_id: uuid.UUID,
    args: dict,
) -> ToolExecutionOutcome:
    payload = await _query_directory_payload(agent_id, args)
    summary = _json_tool_result(payload)
    if payload.get("ok") is True:
        return _typed_success(summary)
    error = payload.get("error") if isinstance(payload.get("error"), Mapping) else {}
    return _typed_failure(
        summary,
        str(error.get("code") or "query_directory_failed"),
        retryable=error.get("code") == "query_directory_failed",
    )


async def _query_directory(agent_id: uuid.UUID, args: dict) -> str:
    """Legacy display adapter for non-Durable callers."""
    return _json_tool_result(await _query_directory_payload(agent_id, args))


async def _send_feishu_message(agent_id: uuid.UUID, args: dict) -> str:
    """Send a Feishu message to a person in the agent's relationship list."""
    target_member_id = (args.get("target_member_id") or "").strip()
    member_name = (args.get("member_name") or "").strip()
    direct_user_id = (args.get("user_id") or "").strip()
    message_text = (args.get("message") or "").strip()

    if not message_text:
        return "❌ Please provide message content"
    if (member_name or direct_user_id) and not target_member_id:
        return (
            "❌ send_feishu_message is a legacy shortcut and no longer accepts member_name or user_id. "
            "Call query_directory(member_type=\"human\", query=\"...\") first, then retry with "
            "send_channel_message(target_member_id=\"...\", channel=\"feishu\", message=\"...\")."
        )
    if not target_member_id:
        return "❌ Please provide target_member_id from query_directory, or use send_channel_message for Feishu."

    return await _send_channel_message(
        agent_id,
        {
            "target_member_id": target_member_id,
            "message": message_text,
            "channel": "feishu",
        },
    )


async def _send_feishu_message_to_member_outcome(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: OrgMember,
) -> ToolExecutionOutcome:
    """Send through Feishu and classify the structured provider response."""
    from app.services.feishu_service import FeishuAPIError, feishu_service

    try:
        async with async_session() as db:
            config_result = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "feishu",
                )
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return _typed_failure(
                    "This Agent has no Feishu channel configured.",
                    "feishu_channel_not_configured",
                )

            feishu_user_id = (target_member.external_id or "").strip()
            if not feishu_user_id:
                return _typed_failure(
                    f"{member_name} has no linked Feishu user_id.",
                    "feishu_recipient_not_linked",
                )

            try:
                response = await feishu_service.send_message(
                    config.app_id,
                    config.app_secret,
                    receive_id=feishu_user_id,
                    msg_type="text",
                    content=json.dumps({"text": message_text}, ensure_ascii=False),
                    receive_id_type="user_id",
                )
            except FeishuAPIError as exc:
                if exc.code is not None or (
                    exc.http_status is not None and exc.http_status < 500
                ):
                    return _typed_failure(
                        f"Feishu rejected the message: {exc.user_message}",
                        "feishu_message_rejected",
                    )
                return _typed_unknown(
                    "Feishu message outcome is unknown; reconcile before retrying.",
                    "feishu_message_outcome_unknown",
                )
            except Exception:
                return _typed_unknown(
                    "Feishu message outcome is unknown; reconcile before retrying.",
                    "feishu_message_outcome_unknown",
                )

            if not isinstance(response, Mapping):
                return _typed_unknown(
                    "Feishu returned an unreadable response; reconcile before retrying.",
                    "feishu_response_invalid",
                )
            if response.get("code") != 0:
                return _typed_failure(
                    f"Feishu rejected the message: {response.get('msg') or 'unknown error'} "
                    f"(code {response.get('code')}).",
                    "feishu_message_rejected",
                )

            # Provider success is the execution fact. Conversation-history
            # persistence is best-effort product synchronization and cannot
            # turn a confirmed send into unknown or trigger a re-send.
            try:
                agent_result = await db.execute(
                    select(AgentModel).where(AgentModel.id == agent_id)
                )
                agent = agent_result.scalar_one_or_none()
                platform_user = await get_platform_user_by_org_member(
                    db=db,
                    org_member=target_member,
                    agent_tenant_id=agent.tenant_id if agent else None,
                )
                session = await find_or_create_channel_session(
                    db=db,
                    agent_id=agent_id,
                    user_id=platform_user.id,
                    external_conv_id=f"feishu_p2p_{feishu_user_id}",
                    source_channel="feishu",
                    first_message_title=f"[Agent → {member_name or feishu_user_id}]",
                )
                db.add(
                    ChatMessage(
                        agent_id=agent_id,
                        user_id=platform_user.id,
                        role="assistant",
                        content=message_text,
                        conversation_id=str(session.id),
                    )
                )
                session.last_message_at = datetime.now(timezone.utc)
                await db.commit()
            except Exception as history_error:
                logger.error(
                    "[Feishu] Confirmed send but failed to sync history: {}",
                    type(history_error).__name__,
                )

            return _typed_success(f"Successfully sent message to {member_name}.")
    except Exception as exc:
        logger.exception("[Feishu] Message setup failed")
        return _typed_failure(
            f"Feishu message could not be prepared: {type(exc).__name__}.",
            "feishu_message_setup_failed",
        )


async def _send_feishu_message_to_member(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: OrgMember,
) -> str:
    """Legacy display adapter; Durable Runtime uses the typed provider helper."""
    outcome = await _send_feishu_message_to_member_outcome(
        agent_id,
        member_name,
        message_text,
        target_member,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Feishu message did not return a summary.",
    )


async def _send_channel_message(agent_id: uuid.UUID, args: dict) -> str:
    """Send message via a resolved human target's configured external channel."""
    target_member_id = (args.get("target_member_id") or "").strip()
    provider_user_id = (args.get("provider_user_id") or "").strip()
    member_name = (args.get("member_name") or "").strip()
    message_text = (args.get("message") or "").strip()
    target_channel = _normalize_roster_provider_type(args.get("channel"))

    if not message_text:
        return "❌ Please provide message content"
    if (provider_user_id or member_name) and not target_member_id:
        return (
            "❌ provider_user_id and member_name are no longer supported for send_channel_message. "
            "Call query_directory(member_type=\"human\", query=\"...\") first, then retry with target_member_id."
        )
    if not target_member_id:
        return "❌ Please provide target_member_id from query_directory."

    try:
        async with async_session() as db:
            target, error = await _resolve_roster_human_target(
                db,
                agent_id,
                target_member_id=target_member_id,
                provider_type=target_channel,
            )
            if error:
                return error

            target_member = target.member
            display_name = target_member.name or target_member_id
            provider_type = target.provider_type
            if not provider_type:
                if target.platform_user and not target_channel:
                    logger.info(
                        "[ChannelMessage] %s is a platform user; rerouting send_channel_message -> send_platform_message",
                        display_name,
                    )
                    return await _send_platform_message(
                        agent_id,
                        {
                            "target_member_id": str(target_member.id),
                            "message": message_text,
                        },
                    )
                if (target_member.external_id or target_member.open_id) and not target_channel:
                    provider_type = "feishu"
                else:
                    return (
                        f"❌ {display_name} has no linked channel. "
                        "If they are a platform user, use send_platform_message instead."
                    )

            logger.info(f"[ChannelMessage] Sending to {display_name} via {provider_type}")

            if provider_type == "feishu":
                return await _send_feishu_message_to_member(agent_id, display_name, message_text, target_member)
            elif provider_type == "dingtalk":
                return await _send_dingtalk_message(agent_id, display_name, message_text, target_member)
            elif provider_type == "wecom":
                return await _send_wecom_message(agent_id, display_name, message_text, target_member)
            elif provider_type == "slack":
                return await _send_slack_message(agent_id, display_name, message_text, target_member)
            elif provider_type == "teams":
                return await _send_teams_channel_message(agent_id, display_name, message_text, target_member)
            elif provider_type == "wechat":
                return await _send_wechat_channel_message(agent_id, display_name, message_text, target_member)
            else:
                return f"❌ Unsupported channel type: {provider_type}"

    except Exception as e:
        logger.exception("[ChannelMessage] Error")
        return f"❌ Channel message error: {str(e)[:200]}"


async def _send_channel_message_outcome(
    agent_id: uuid.UUID,
    args: dict,
) -> ToolExecutionOutcome | str:
    """Typed channel dispatch for providers with structured execution facts.

    Providers that have not yet been migrated deliberately return their legacy
    string so Durable Runtime rejects the result as untyped.
    """
    target_member_id = (args.get("target_member_id") or "").strip()
    provider_user_id = (args.get("provider_user_id") or "").strip()
    member_name = (args.get("member_name") or "").strip()
    message_text = (args.get("message") or "").strip()
    target_channel = _normalize_roster_provider_type(args.get("channel"))
    if not message_text or not target_member_id:
        return _typed_failure(
            "send_channel_message requires target_member_id and message.",
            "invalid_tool_arguments",
        )
    if provider_user_id or member_name:
        return _typed_failure(
            "send_channel_message accepts stable target_member_id, not provider_user_id/member_name.",
            "invalid_tool_arguments",
        )
    try:
        async with async_session() as db:
            target, error = await _resolve_roster_human_target(
                db,
                agent_id,
                target_member_id=target_member_id,
                provider_type=target_channel,
            )
    except Exception as exc:
        return _typed_failure(
            f"Channel recipient could not be resolved: {type(exc).__name__}.",
            "channel_recipient_resolution_failed",
        )
    if error:
        return _typed_failure(error, "channel_recipient_invalid")
    target_member = target.member
    display_name = target_member.name or target_member_id
    provider_type = target.provider_type
    if not provider_type:
        if target.platform_user and not target_channel:
            return await _send_platform_message_outcome(
                agent_id,
                {
                    "target_member_id": str(target_member.id),
                    "message": message_text,
                },
            )
        if (target_member.external_id or target_member.open_id) and not target_channel:
            provider_type = "feishu"
        else:
            return _typed_failure(
                f"{display_name} has no linked external channel.",
                "channel_recipient_unreachable",
            )
    if provider_type == "feishu":
        return await _send_feishu_message_to_member_outcome(
            agent_id,
            display_name,
            message_text,
            target_member,
        )
    if provider_type == "dingtalk":
        return await _send_dingtalk_message_outcome(
            agent_id, display_name, message_text, target_member
        )
    if provider_type == "wecom":
        return await _send_wecom_message_outcome(
            agent_id, display_name, message_text, target_member
        )
    if provider_type == "slack":
        return await _send_slack_message(
            agent_id, display_name, message_text, target_member
        )
    if provider_type == "teams":
        return await _send_teams_channel_message(
            agent_id, display_name, message_text, target_member
        )
    if provider_type == "wechat":
        return await _send_wechat_channel_message(
            agent_id, display_name, message_text, target_member
        )
    return _typed_failure(
        f"Unsupported channel type: {provider_type}",
        "channel_provider_unsupported",
    )


async def _sync_proactive_channel_history(
    db,
    *,
    agent_id: uuid.UUID,
    target_member: "OrgMember",
    member_name: str,
    message_text: str,
    source_channel: str,
    external_user_id: str,
) -> None:
    """Best-effort product sync after the provider confirmed a proactive send."""
    try:
        agent_result = await db.execute(
            select(AgentModel).where(AgentModel.id == agent_id)
        )
        agent = agent_result.scalar_one_or_none()
        platform_user = await get_platform_user_by_org_member(
            db=db,
            org_member=target_member,
            agent_tenant_id=agent.tenant_id if agent else None,
        )
        session = await find_or_create_channel_session(
            db=db,
            agent_id=agent_id,
            user_id=platform_user.id,
            external_conv_id=f"{source_channel}_p2p_{external_user_id}",
            source_channel=source_channel,
            first_message_title=message_text[:30],
        )
        db.add(
            ChatMessage(
                agent_id=agent_id,
                user_id=platform_user.id,
                role="assistant",
                content=message_text,
                conversation_id=str(session.id),
            )
        )
        session.last_message_at = datetime.now(timezone.utc)
        await db.commit()
        logger.info(
            "[{}] Proactive message saved to session {}",
            source_channel,
            session.id,
        )
    except Exception as exc:
        # Provider success is authoritative. A local history-sync failure must
        # never downgrade it or authorize another external send.
        logger.error(
            "[{}] Confirmed send to {} but failed to sync history: {}",
            source_channel,
            member_name,
            type(exc).__name__,
        )


async def _send_dingtalk_message_outcome(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> ToolExecutionOutcome:
    """Send through DingTalk and preserve a typed external-write outcome."""
    from app.services.dingtalk_service import send_dingtalk_message

    try:
        async with async_session() as db:
            config_result = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "dingtalk",
                    ChannelConfig.is_configured.is_(True),
                )
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return _typed_failure(
                    "This Agent has no DingTalk channel configured.",
                    "dingtalk_channel_not_configured",
                )

            user_id = (target_member.external_id or "").strip()
            if not user_id:
                user_id = (target_member.unionid or target_member.open_id or "").strip()
                if not user_id:
                    return _typed_failure(
                        f"{member_name} has no linked DingTalk user_id.",
                        "dingtalk_recipient_not_linked",
                    )

            logger.info(f"[DingTalk] Sending to user_id: {user_id}")
            provider_agent_id = (
                str((config.extra_config or {}).get("agent_id") or "").strip()
                or None
            )
            try:
                result = await send_dingtalk_message(
                    app_id=config.app_id,
                    app_secret=config.app_secret,
                    user_id=user_id,
                    message=message_text,
                    agent_id=provider_agent_id,
                )
            except Exception:
                return _typed_unknown(
                    "DingTalk message outcome is unknown; reconcile before retrying.",
                    "dingtalk_message_outcome_unknown",
                )

            if not isinstance(result, Mapping) or result.get("errcode") in {None, -1}:
                return _typed_unknown(
                    "DingTalk message outcome is unknown; reconcile before retrying.",
                    "dingtalk_message_outcome_unknown",
                )
            if result.get("errcode") != 0:
                return _typed_failure(
                    f"DingTalk rejected the message: {result.get('errmsg') or 'unknown error'} "
                    f"(code {result.get('errcode')}).",
                    "dingtalk_message_rejected",
                )

            await _sync_proactive_channel_history(
                db,
                agent_id=agent_id,
                target_member=target_member,
                member_name=member_name,
                message_text=message_text,
                source_channel="dingtalk",
                external_user_id=user_id,
            )
            return _typed_success(f"Successfully sent message to {member_name} via DingTalk.")
    except Exception as exc:
        logger.exception("[DingTalk] Message setup failed")
        return _typed_failure(
            f"DingTalk message could not be prepared: {type(exc).__name__}.",
            "dingtalk_message_setup_failed",
        )


async def _send_dingtalk_message(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> str:
    """Legacy display adapter; Durable Runtime uses the typed provider helper."""
    outcome = await _send_dingtalk_message_outcome(
        agent_id,
        member_name,
        message_text,
        target_member,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="DingTalk message did not return a summary.",
    )


async def _send_wecom_message_outcome(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> ToolExecutionOutcome:
    """Send through WeCom and preserve a typed external-write outcome."""
    from app.services.wecom_service import send_wecom_message

    try:
        async with async_session() as db:
            config_result = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "wecom",
                    ChannelConfig.is_configured.is_(True),
                )
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return _typed_failure(
                    "This Agent has no WeCom channel configured.",
                    "wecom_channel_not_configured",
                )

            user_id = (target_member.external_id or "").strip()
            if not user_id:
                user_id = (target_member.open_id or "").strip()
                if not user_id:
                    return _typed_failure(
                        f"{member_name} has no linked WeCom user_id.",
                        "wecom_recipient_not_linked",
                    )

            provider_agent_id = str(
                (config.extra_config or {}).get("wecom_agent_id") or ""
            ).strip()
            if not provider_agent_id:
                return _typed_failure(
                    "This Agent's WeCom channel has no application AgentID.",
                    "wecom_agent_id_missing",
                )

            logger.info(f"[WeCom] Sending to user_id: {user_id}")
            try:
                result = await send_wecom_message(
                    config.app_id,
                    config.app_secret,
                    user_id,
                    message_text,
                    agent_id=provider_agent_id,
                )
            except Exception:
                return _typed_unknown(
                    "WeCom message outcome is unknown; reconcile before retrying.",
                    "wecom_message_outcome_unknown",
                )

            if not isinstance(result, Mapping) or result.get("errcode") in {None, -1}:
                return _typed_unknown(
                    "WeCom message outcome is unknown; reconcile before retrying.",
                    "wecom_message_outcome_unknown",
                )
            if result.get("errcode") != 0:
                return _typed_failure(
                    f"WeCom rejected the message: {result.get('errmsg') or 'unknown error'} "
                    f"(code {result.get('errcode')}).",
                    "wecom_message_rejected",
                )

            await _sync_proactive_channel_history(
                db,
                agent_id=agent_id,
                target_member=target_member,
                member_name=member_name,
                message_text=message_text,
                source_channel="wecom",
                external_user_id=user_id,
            )
            return _typed_success(f"Successfully sent message to {member_name} via WeCom.")
    except Exception as exc:
        logger.exception("[WeCom] Message setup failed")
        return _typed_failure(
            f"WeCom message could not be prepared: {type(exc).__name__}.",
            "wecom_message_setup_failed",
        )


async def _send_wecom_message(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> str:
    """Legacy display adapter; Durable Runtime uses the typed provider helper."""
    outcome = await _send_wecom_message_outcome(
        agent_id,
        member_name,
        message_text,
        target_member,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="WeCom message did not return a summary.",
    )

async def _send_slack_message(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> str:
    """Send proactive Slack DM via conversations.open + chat.postMessage."""
    import httpx

    from app.api.slack import _send_slack_messages

    try:
        async with async_session() as db:
            config_result = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "slack",
                    ChannelConfig.is_configured == True,
                )
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return "❌ This agent has no Slack channel configured"

            user_id = (target_member.external_id or "").strip()
            if not user_id:
                return f"❌ {member_name} has no Slack user_id"

            bot_token = (config.app_secret or "").strip()
            if not bot_token:
                return "❌ Slack bot token is missing"

            async with httpx.AsyncClient(timeout=10) as client:
                open_resp = await client.post(
                    "https://slack.com/api/conversations.open",
                    headers={"Authorization": f"Bearer {bot_token}", "Content-Type": "application/json"},
                    json={"users": user_id},
                )
                data = open_resp.json()
                if open_resp.status_code >= 400 or not data.get("ok"):
                    err = data.get("error") or open_resp.text[:200]
                    return f"❌ Slack conversations.open failed: {err}"
                channel_id = (((data.get("channel") or {})).get("id") or "").strip()

            if not channel_id:
                return f"❌ Slack DM channel unavailable for {member_name}"

            await _send_slack_messages(bot_token, channel_id, message_text)

            try:
                agent_r = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
                agent_obj = agent_r.scalar_one_or_none()
                platform_user = await get_platform_user_by_org_member(
                    db=db,
                    org_member=target_member,
                    agent_tenant_id=agent_obj.tenant_id if agent_obj else None,
                )
                conv_id = f"slack_{channel_id}"
                sess = await find_or_create_channel_session(
                    db=db,
                    agent_id=agent_id,
                    user_id=platform_user.id,
                    external_conv_id=conv_id,
                    source_channel="slack",
                    first_message_title=message_text[:30],
                )
                db.add(ChatMessage(
                    agent_id=agent_id,
                    user_id=platform_user.id,
                    role="assistant",
                    content=message_text,
                    conversation_id=str(sess.id),
                ))
                sess.last_message_at = datetime.now(timezone.utc)
                await db.commit()
                logger.info(f"[Slack] Proactive message saved to session {sess.id}")
            except Exception as ex:
                logger.error(f"[Slack] Failed to save proactive message to session: {ex}")

            return f"✅ Message sent to {member_name} via Slack"
    except Exception as e:
        logger.exception("[Slack] Error")
        return f"❌ Slack message error: {str(e)[:200]}"


async def _send_teams_channel_message(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> str:
    """Send proactive Teams message using the latest known conversation context."""
    from app.api.teams import _send_teams_message

    try:
        async with async_session() as db:
            config_result = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "microsoft_teams",
                    ChannelConfig.is_configured == True,
                )
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return "❌ This agent has no Teams channel configured"

            service_url = str((config.extra_config or {}).get("service_url") or "").strip()
            if not service_url:
                return "❌ Teams proactive send requires an existing inbound conversation to capture service_url"

            agent_r = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent_obj = agent_r.scalar_one_or_none()
            platform_user = await get_platform_user_by_org_member(
                db=db,
                org_member=target_member,
                agent_tenant_id=agent_obj.tenant_id if agent_obj else None,
            )

            session_result = await db.execute(
                select(ChatSession)
                .where(
                    ChatSession.agent_id == agent_id,
                    ChatSession.user_id == platform_user.id,
                    ChatSession.source_channel == "microsoft_teams",
                    ChatSession.is_group == False,
                )
                .order_by(ChatSession.last_message_at.desc(), ChatSession.created_at.desc())
                .limit(1)
            )
            session = session_result.scalar_one_or_none()
            conversation_id = str(session.external_conv_id or "").strip() if session else ""
            if not conversation_id:
                return f"❌ Teams proactive send to {member_name} requires them to message the bot first"

            await _send_teams_message(
                config,
                conversation_id,
                {
                    "type": "message",
                    "text": message_text,
                    "conversation": {"id": conversation_id},
                },
            )

            db.add(ChatMessage(
                agent_id=agent_id,
                user_id=platform_user.id,
                role="assistant",
                content=message_text,
                conversation_id=str(session.id),
            ))
            session.last_message_at = datetime.now(timezone.utc)
            await db.commit()
            logger.info(f"[Teams] Proactive message saved to session {session.id}")
            return f"✅ Message sent to {member_name} via Teams"
    except Exception as e:
        logger.exception("[Teams] Error")
        return f"❌ Teams message error: {str(e)[:200]}"


async def _send_wechat_channel_message(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> str:
    """Send proactive WeChat message using the latest cached context_token."""
    from app.services.wechat_channel import (
        WECHAT_ILINK_BASE_URL,
        get_wechat_context_entry,
        send_wechat_text_message,
    )

    try:
        async with async_session() as db:
            config_result = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "wechat",
                    ChannelConfig.is_configured == True,
                )
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return "❌ This agent has no WeChat channel configured"

            user_id = (target_member.external_id or "").strip()
            if not user_id:
                return f"❌ {member_name} has no WeChat user_id"

            ctx_entry = get_wechat_context_entry(config.extra_config, from_user_id=user_id)
            context_token = str((ctx_entry or {}).get("context_token") or "").strip()
            conv_id = str((ctx_entry or {}).get("conv_id") or f"wechat_{user_id}").strip()
            if not context_token:
                return f"❌ WeChat proactive send to {member_name} requires them to message the bot first"

            token = str((config.extra_config or {}).get("bot_token") or "").strip()
            base_url = str((config.extra_config or {}).get("baseurl") or WECHAT_ILINK_BASE_URL).strip()
            route_tag = str((config.extra_config or {}).get("route_tag") or "").strip() or None
            if not token:
                return "❌ WeChat bot token is missing"

            await send_wechat_text_message(
                token=token,
                base_url=base_url,
                to_user_id=user_id,
                context_token=context_token,
                text=message_text,
                route_tag=route_tag,
            )

            agent_r = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent_obj = agent_r.scalar_one_or_none()
            platform_user = await get_platform_user_by_org_member(
                db=db,
                org_member=target_member,
                agent_tenant_id=agent_obj.tenant_id if agent_obj else None,
            )
            sess = await find_or_create_channel_session(
                db=db,
                agent_id=agent_id,
                user_id=platform_user.id,
                external_conv_id=conv_id,
                source_channel="wechat",
                first_message_title=message_text[:30],
            )
            db.add(ChatMessage(
                agent_id=agent_id,
                user_id=platform_user.id,
                role="assistant",
                content=message_text,
                conversation_id=str(sess.id),
            ))
            sess.last_message_at = datetime.now(timezone.utc)
            await db.commit()
            logger.info(f"[WeChat] Proactive message saved to session {sess.id}")
            return f"✅ Message sent to {member_name} via WeChat"
    except Exception as e:
        logger.exception("[WeChat] Error")
        return f"❌ WeChat message error: {str(e)[:200]}"


async def _send_platform_message_outcome(
    agent_id: uuid.UUID,
    args: dict,
) -> ToolExecutionOutcome:
    """Persist a first-party message and expose its transaction outcome."""
    target_member_id = (args.get("target_member_id") or "").strip()
    platform_user_id = (args.get("platform_user_id") or "").strip()
    username = (args.get("username") or "").strip()
    message_text = (args.get("message") or "").strip()
    if username and not target_member_id and not platform_user_id:
        return _typed_failure(
            "username is no longer supported; call query_directory and use target_member_id.",
            "invalid_tool_arguments",
        )
    if not message_text or (not target_member_id and not platform_user_id):
        return _typed_failure(
            "send_platform_message requires message and target_member_id or platform_user_id.",
            "invalid_tool_arguments",
        )
    commit_started = False
    try:
        async with async_session() as db:
            target, error = await _resolve_roster_human_target(
                db,
                agent_id,
                target_member_id=target_member_id,
                platform_user_id=platform_user_id,
                member_name=None,
                require_platform_user=True,
            )
            if error:
                return _typed_failure(error, "platform_recipient_invalid")
            target_user = target.platform_user
            from app.services.chat_session_service import ensure_primary_platform_session

            session = await ensure_primary_platform_session(
                db,
                agent_id,
                target_user.id,
            )
            db.add(
                ChatMessage(
                    agent_id=agent_id,
                    user_id=target_user.id,
                    role="assistant",
                    content=message_text,
                    conversation_id=str(session.id),
                )
            )
            session.last_message_at = datetime.now(timezone.utc)
            try:
                from app.api.websocket import maybe_mark_session_read_for_active_viewer

                await maybe_mark_session_read_for_active_viewer(
                    db,
                    agent_id=agent_id,
                    session_id=str(session.id),
                    user_id=target_user.id,
                )
            except Exception:
                pass
            commit_started = True
            await db.commit()
    except Exception as exc:
        if commit_started:
            return _typed_unknown(
                "Platform message persistence outcome is unknown; reconcile before retrying.",
                "platform_message_outcome_unknown",
            )
        return _typed_failure(
            f"Platform message could not be prepared: {type(exc).__name__}.",
            "platform_message_failed",
        )

    # Push is a best-effort delivery optimization after the durable message
    # exists. Its failure must not cause the durable write to be repeated.
    try:
        from app.api.websocket import manager as ws_manager

        await ws_manager.send_to_user(
            str(agent_id),
            str(target_user.id),
            {
                "type": "trigger_notification",
                "content": message_text,
                "triggers": ["web_message"],
                "session_id": str(session.id),
            },
        )
    except Exception:
        pass
    display = target_user.display_name or target_user.username
    return _typed_success(
        f"Message sent to {display} on the web platform and saved to chat history."
    )


async def _send_platform_message(agent_id: uuid.UUID, args: dict) -> str:
    """Legacy display adapter; Durable Runtime uses the typed transaction helper."""
    outcome = await _send_platform_message_outcome(agent_id, args)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Platform message did not return a summary.",
    )


async def _resolve_a2a_target_by_id(
    db,
    source_agent: AgentModel,
    target_agent_id: str,
) -> tuple[AgentModel | None, str | None]:
    try:
        target_id = uuid.UUID((target_agent_id or "").strip())
    except (TypeError, ValueError):
        return None, "❌ Invalid target_agent_id. Use query_directory to get a valid target_agent_id."

    if target_id == source_agent.id:
        return None, "❌ You cannot send a message to yourself."

    target_result = await db.execute(select(AgentModel).where(AgentModel.id == target_id))
    target = target_result.scalar_one_or_none()
    if not target:
        return None, "❌ Target agent not found. Use query_directory to find an available digital employee."
    if target.tenant_id != source_agent.tenant_id:
        return None, "❌ Target agent is outside your tenant and cannot be contacted."

    authorized_custom_target = False
    if getattr(target, "access_mode", None) == "custom":
        authorized_custom_target = await agent_directory.is_custom_agent_target_authorized(
            db,
            source_agent_id=source_agent.id,
            target_agent_id=target.id,
        )
    visibility = evaluate_roster_agent_visibility(
        source_agent,
        target,
        authorized_custom_target=authorized_custom_target,
    )
    if not visibility.visible:
        return None, "❌ Target agent is not visible to you. Use query_directory to choose a visible digital employee."
    if not visibility.can_contact:
        reason = visibility.unavailable_reason or "target_not_contactable"
        return None, f"❌ Target agent is currently unavailable ({reason})."

    return target, None


async def _send_file_to_agent_outcome(
    from_agent_id: uuid.UUID,
    args: dict,
) -> ToolExecutionOutcome:
    """Copy a file and inbox note, then treat ancillary history as best effort."""
    target_agent_id = args.get("target_agent_id")
    rel_path = args.get("file_path")
    legacy_agent_name = args.get("agent_name", "")
    delivery_note = args.get("message", "")
    if (
        not isinstance(target_agent_id, str)
        or not isinstance(rel_path, str)
        or not isinstance(legacy_agent_name, str)
        or not isinstance(delivery_note, str)
    ):
        return _typed_failure(
            "send_file_to_agent arguments must be strings.",
            "invalid_tool_arguments",
        )
    target_agent_id = target_agent_id.strip()
    legacy_agent_name = legacy_agent_name.strip()
    rel_path = rel_path.strip()
    delivery_note = delivery_note.strip()

    if legacy_agent_name and not target_agent_id:
        return _typed_failure(
            "send_file_to_agent accepts stable target_agent_id, not agent_name.",
            "invalid_tool_arguments",
        )
    if not target_agent_id or not rel_path:
        return _typed_failure(
            "send_file_to_agent requires target_agent_id and file_path.",
            "invalid_tool_arguments",
        )

    storage = get_storage_backend()
    source_key = normalize_storage_key(f"{from_agent_id}/{rel_path}")
    try:
        if not await storage.is_file(source_key):
            return _typed_failure(
                f"Source file not found: {rel_path}",
                "workspace_file_not_found",
            )
        source_entry = await storage.stat(source_key)
    except Exception as exc:
        return _typed_failure(
            f"Source file could not be read: {type(exc).__name__}.",
            "workspace_read_failed",
        )

    # File size limit (50 MB)
    MAX_FILE_SIZE = 50 * 1024 * 1024
    file_size = source_entry.size
    if file_size > MAX_FILE_SIZE:
        size_mb = file_size / (1024 * 1024)
        return _typed_failure(
            f"File too large ({size_mb:.1f} MB). Maximum allowed is 50 MB.",
            "agent_file_too_large",
        )
    try:
        source_bytes = await storage.read_bytes(source_key)
    except Exception as exc:
        return _typed_failure(
            f"Source file could not be read: {type(exc).__name__}.",
            "workspace_read_failed",
        )
    source_name = Path(rel_path).name

    mutation_started = False
    try:
        from app.services.activity_logger import log_activity

        async with async_session() as db:
            src_result = await db.execute(select(AgentModel).where(AgentModel.id == from_agent_id))
            source_agent = src_result.scalar_one_or_none()
            if not source_agent:
                return _typed_failure(
                    "Source Agent not found.",
                    "source_agent_not_found",
                )
            source_agent_name = source_agent.name if source_agent else "Unknown agent"
            source_creator_id = source_agent.creator_id if source_agent else from_agent_id

            target_agent, target_error = await _resolve_a2a_target_by_id(db, source_agent, target_agent_id)
            if target_error:
                return _typed_failure(
                    target_error,
                    "agent_file_recipient_invalid",
                )

            target_name = target_agent.name
            target_id = target_agent.id

        ts = datetime.now(timezone.utc)
        stamp = ts.strftime("%Y%m%d_%H%M%S_%f")
        delivered_name = source_name
        target_rel_path = f"workspace/inbox/files/{delivered_name}"
        target_key = normalize_storage_key(f"{target_id}/{target_rel_path}")
        collision = 0
        while await storage.exists(target_key):
            collision += 1
            delivered_name = f"{stamp}_{collision}_{source_name}"
            target_rel_path = f"workspace/inbox/files/{delivered_name}"
            target_key = normalize_storage_key(f"{target_id}/{target_rel_path}")

        mutation_started = True
        await storage.write_bytes(target_key, source_bytes)

        sender_short = str(from_agent_id)[:8]
        note_rel_path = f"workspace/inbox/{stamp}_{sender_short}_file_delivery.md"
        note_key = normalize_storage_key(f"{target_id}/{note_rel_path}")
        note_lines = [
            f"# File delivery from {source_agent_name}",
            "",
            f"- Time (UTC): {ts.isoformat()}",
            f"- Sender: {source_agent_name}",
            f"- Source path: {rel_path}",
            f"- Delivered file: {target_rel_path}",
            "",
        ]
        if delivery_note:
            note_lines.append("## Note")
            note_lines.append(delivery_note)
            note_lines.append("")
        note_lines.append("## Action")
        note_lines.append(f"- Read the file via `read_file(path=\"{target_rel_path}\")`")
        await storage.write_text(note_key, "\n".join(note_lines), encoding="utf-8")

        try:
            from app.models.audit import AuditLog

            async with async_session() as db:
                db.add(AuditLog(
                    agent_id=from_agent_id,
                    action="collaboration:file_send",
                    details={
                        "to_agent": str(target_id),
                        "to_agent_name": target_name,
                        "source_file": rel_path,
                        "delivered_file": target_rel_path,
                    },
                ))
                db.add(AuditLog(
                    agent_id=target_id,
                    action="collaboration:file_receive",
                    details={
                        "from_agent": str(from_agent_id),
                        "from_agent_name": source_agent_name,
                        "source_file": rel_path,
                        "delivered_file": target_rel_path,
                    },
                ))
                await db.commit()
        except Exception as exc:
            logger.error(
                "[A2A-File] Confirmed delivery but audit sync failed: {}",
                type(exc).__name__,
            )

        try:
            await log_activity(
                from_agent_id,
                "agent_file_sent",
                f"Sent file to {target_name}",
                detail={"target_agent": target_name, "source_file": rel_path, "delivered_file": target_rel_path},
            )
            await log_activity(
                target_id,
                "agent_file_received",
                f"Received file from {source_agent_name}",
                detail={"source_agent": source_agent_name, "source_file": rel_path, "delivered_file": target_rel_path},
            )
        except Exception:
            pass

        # ── Inject file-delivery message into A2A chat session ──
        # This ensures the target agent sees the file delivery in its
        # conversation context when send_message_to_agent is called next.
        logger.info(
            "[A2A-File] Injecting file delivery message: from=%s to=%s file=%s",
            source_name,
            target_name,
            delivered_name,
        )
        try:
            from app.models.audit import ChatMessage
            from app.models.chat_session import ChatSession
            from app.models.participant import Participant
            async with async_session() as db2:
                # Find or create A2A session (same ordering as send_message_to_agent)
                session_agent_id = min(from_agent_id, target_id, key=str)
                session_peer_id = max(from_agent_id, target_id, key=str)
                sess_r = await db2.execute(
                    select(ChatSession).where(
                        ChatSession.agent_id == session_agent_id,
                        ChatSession.peer_agent_id == session_peer_id,
                        ChatSession.source_channel == "agent",
                    )
                )
                chat_session = sess_r.scalar_one_or_none()
                if not chat_session:
                    src_part_r = await db2.execute(
                        select(Participant).where(Participant.type == "agent", Participant.ref_id == from_agent_id)
                    )
                    src_participant = src_part_r.scalar_one_or_none()
                    chat_session = ChatSession(
                        agent_id=session_agent_id,
                        user_id=source_creator_id,
                        title=f"{source_name} ↔ {target_name}",
                        source_channel="agent",
                        participant_id=src_participant.id if src_participant else None,
                        peer_agent_id=session_peer_id,
                    )
                    db2.add(chat_session)
                    await db2.flush()

                file_msg_content = (
                    f"[File delivery from {source_name}]\n"
                    f"{source_name} sent you a file: {delivered_name}\n"
                    f"File path: {target_rel_path}\n"
                    f"Use read_file(path=\"{target_rel_path}\") to inspect it."
                )
                if delivery_note:
                    file_msg_content += f"\nNote: {delivery_note}"

                # Resolve sender participant for proper attribution
                src_part_r2 = await db2.execute(
                    select(Participant).where(Participant.type == "agent", Participant.ref_id == from_agent_id)
                )
                src_part2 = src_part_r2.scalar_one_or_none()

                db2.add(ChatMessage(
                    agent_id=session_agent_id,
                    user_id=source_creator_id,
                    role="user",
                    content=file_msg_content,
                    conversation_id=str(chat_session.id),
                    participant_id=src_part2.id if src_part2 else None,
                ))
                chat_session.last_message_at = ts
                await db2.commit()
                logger.info(
                    "[A2A-File] Injected file delivery message into session %s for %s",
                    chat_session.id,
                    target_name,
                )
        except Exception as e:
            logger.error(f"[A2A-File] FAILED to inject file delivery message: {e}")

        return _typed_success(
            f"File sent to {target_name}.\n"
            f"- Delivered to: {target_rel_path}\n"
            f"- Inbox note: {note_rel_path}"
        )
    except Exception as exc:
        if mutation_started:
            return _typed_unknown(
                "Agent file delivery outcome is unknown; reconcile before retrying.",
                "agent_file_outcome_unknown",
            )
        return _typed_failure(
            f"Agent file delivery failed before dispatch: {type(exc).__name__}.",
            "agent_file_send_failed",
        )


async def _send_file_to_agent(from_agent_id: uuid.UUID, args: dict) -> str:
    """Legacy display adapter for the typed Agent file transfer."""
    outcome = await _send_file_to_agent_outcome(from_agent_id, args)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Agent file delivery returned no summary.",
    )


async def _send_message_to_agent(
    from_agent_id: uuid.UUID,
    args: dict,
    user_id: uuid.UUID | None = None,
    origin_session_id: str | None = None,
) -> str:
    """Fail closed when a caller bypasses the Runtime tool-step service.

    The schema remains in ``agent_tools`` because models still call this tool,
    but execution must be intercepted by ``RuntimeA2AService`` where the source
    Run, tool receipt, target Run or Gateway message, and callback are durable.
    """
    del from_agent_id, args, user_id, origin_session_id
    return (
        "❌ send_message_to_agent requires a durable Agent Runtime Run; "
        "the message was not sent."
    )




# Plaza Tools — Agent Square social feed
# ═══════════════════════════════════════════════════════

async def _plaza_get_new_posts(agent_id: uuid.UUID, arguments: dict) -> str:
    """Get recent posts from the Agent Plaza, scoped to agent's tenant."""
    from app.models.plaza import PlazaPost, PlazaComment
    from app.models.agent import Agent as AgentModel
    from sqlalchemy import desc

    limit = min(arguments.get("limit", 10), 20)

    try:
        async with async_session() as db:
            # Resolve agent's tenant_id
            ar = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = ar.scalar_one_or_none()
            if not agent:
                return "Error: Agent not found."
            if agent.is_system:
                return "System agents cannot access Plaza."

            if (getattr(agent, "access_mode", None) or "company") != "company":
                return "Only company-wide agents can access Plaza."

            tenant_id = agent.tenant_id if agent else None

            q = select(PlazaPost).order_by(desc(PlazaPost.created_at)).limit(limit)
            if tenant_id:
                q = q.where(PlazaPost.tenant_id == tenant_id)
            result = await db.execute(q)
            posts = result.scalars().all()

            if not posts:
                return "📭 No posts in the plaza yet. Be the first to share something!"

            output = []
            for p in posts:
                # Load comments
                cr = await db.execute(
                    select(PlazaComment).where(PlazaComment.post_id == p.id).order_by(PlazaComment.created_at).limit(5)
                )
                comments = cr.scalars().all()
                icon = "🤖" if p.author_type == "agent" else "👤"
                time_str = p.created_at.strftime("%m-%d %H:%M") if p.created_at else ""
                post_text = f"{icon} **{p.author_name}** ({time_str}) [post_id: {p.id}]\n{p.content}\n❤️ {p.likes_count}  💬 {p.comments_count}"
                if comments:
                    for c in comments:
                        c_icon = "🤖" if c.author_type == "agent" else "👤"
                        post_text += f"\n  └─ {c_icon} {c.author_name}: {c.content}"
                output.append(post_text)

            return "🏛️ Agent Plaza — Recent Posts:\n\n" + "\n\n---\n\n".join(output)

    except Exception as e:
        return f"❌ Failed to load plaza posts: {str(e)[:200]}"


async def _plaza_create_post(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new post in the Agent Plaza.

    System agents (is_system=True) are intentionally excluded from Plaza to
    keep the social feed clean — the OKR Agent communicates through Chat and
    reports, not through Plaza posts.
    """
    from app.models.plaza import PlazaPost
    from app.models.agent import Agent as AgentModel

    content = arguments.get("content", "").strip()
    if not content:
        return "Error: Post content cannot be empty."
    if len(content) > 500:
        content = content[:500]

    try:
        async with async_session() as db:
            # Get agent and check is_system
            ar = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = ar.scalar_one_or_none()
            if not agent:
                return "Error: Agent not found."

            # System agents (e.g. OKR Agent) must not post to Plaza
            if agent.is_system:
                return (
                    "System agents are not allowed to post to Plaza. "
                    "Use send_platform_message to communicate with users directly."
                )

            if (getattr(agent, "access_mode", None) or "company") != "company":
                return "Only company-wide agents are allowed to post to Plaza."
            post = PlazaPost(
                author_id=agent_id,
                author_type="agent",
                author_name=agent.name,
                content=content,
                tenant_id=agent.tenant_id,
            )
            db.add(post)
            await db.flush()  # get post.id

            # Extract @mentions
            try:
                import re
                mentions = re.findall(r'@(\S+)', content)
                if mentions:
                    from app.services.notification_service import send_notification
                    a_q = select(AgentModel).where(AgentModel.id != agent_id)
                    if agent.tenant_id:
                        a_q = a_q.where(AgentModel.tenant_id == agent.tenant_id)
                    a_map = {a.name.lower(): a for a in (await db.execute(a_q)).scalars().all()}
                    notified = set()
                    for m in mentions:
                        ma = a_map.get(m.lower())
                        if ma and ma.id not in notified:
                            notified.add(ma.id)
                            await send_notification(
                                db, agent_id=ma.id,
                                type="mention",
                                title=f"{agent.name} mentioned you in a plaza post",
                                body=content[:150],
                                link=f"/plaza?post={post.id}",
                                ref_id=post.id,
                                sender_name=agent.name,
                            )
            except Exception:
                pass

            await db.commit()
            await db.refresh(post)
            return f"Post published! (ID: {post.id})"

    except Exception as e:
        return f"Failed to create post: {str(e)[:200]}"


async def _plaza_add_comment(agent_id: uuid.UUID, arguments: dict) -> str:
    """Add a comment to a plaza post."""
    from app.models.plaza import PlazaPost, PlazaComment
    from app.models.agent import Agent as AgentModel

    post_id = arguments.get("post_id", "")
    content = arguments.get("content", "").strip()
    if not content:
        return "Error: Comment content cannot be empty."
    if len(content) > 300:
        content = content[:300]

    try:
        pid = uuid.UUID(str(post_id))
    except Exception:
        return "Error: Invalid post_id format."

    try:
        async with async_session() as db:
            # Verify post exists
            pr = await db.execute(select(PlazaPost).where(PlazaPost.id == pid))
            post = pr.scalar_one_or_none()
            if not post:
                return "Error: Post not found."

            # Get agent name
            ar = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = ar.scalar_one_or_none()
            if not agent:
                return "Error: Agent not found."
            if agent.is_system:
                return "System agents are not allowed to comment on Plaza posts."

            if (getattr(agent, "access_mode", None) or "company") != "company":
                return "Only company-wide agents are allowed to comment on Plaza posts."

            comment = PlazaComment(
                post_id=pid,
                author_id=agent_id,
                author_type="agent",
                author_name=agent.name,
                content=content,
            )
            db.add(comment)
            post.comments_count = (post.comments_count or 0) + 1

            # Notify post author (if not self)
            if post.author_id != agent_id:
                try:
                    from app.services.notification_service import send_notification
                    if post.author_type == "agent":
                        await send_notification(
                            db, agent_id=post.author_id,
                            type="plaza_reply",
                            title=f"{agent.name} commented on your post",
                            body=content[:150],
                            link=f"/plaza?post={pid}",
                            ref_id=pid,
                            sender_name=agent.name,
                        )
                        # Also notify human creator
                        pa = (await db.execute(select(AgentModel).where(AgentModel.id == post.author_id))).scalar_one_or_none()
                        if pa and pa.creator_id:
                            await send_notification(
                                db, user_id=pa.creator_id,
                                type="plaza_comment",
                                title=f"{agent.name} commented on {pa.name}'s post",
                                body=content[:100],
                                link=f"/plaza?post={pid}",
                                ref_id=pid,
                                sender_name=agent.name,
                            )
                    elif post.author_type == "human":
                        await send_notification(
                            db, user_id=post.author_id,
                            type="plaza_reply",
                            title=f"{agent.name} commented on your post",
                            body=content[:150],
                            link=f"/plaza?post={pid}",
                            ref_id=pid,
                            sender_name=agent.name,
                        )
                except Exception:
                    pass

            # Notify other agents who commented on this post
            try:
                from app.services.notification_service import send_notification
                other_crs = await db.execute(
                    select(PlazaComment.author_id, PlazaComment.author_type)
                    .where(PlazaComment.post_id == pid)
                    .distinct()
                )
                notified = {post.author_id, agent_id}
                for row in other_crs.fetchall():
                    cid, ctype = row
                    if cid in notified:
                        continue
                    notified.add(cid)
                    if ctype == "agent":
                        await send_notification(
                            db, agent_id=cid,
                            type="plaza_reply",
                            title=f"{agent.name} also commented on a post you commented on",
                            body=content[:150],
                            link=f"/plaza?post={pid}",
                            ref_id=pid,
                            sender_name=agent.name,
                        )
            except Exception:
                pass

            # Extract @mentions
            try:
                import re
                mentions = re.findall(r'@(\S+)', content)
                if mentions:
                    from app.services.notification_service import send_notification
                    from app.models.user import User
                    # Load agents in tenant
                    a_q = select(AgentModel).where(AgentModel.id != agent_id)
                    if agent.tenant_id:
                        a_q = a_q.where(AgentModel.tenant_id == agent.tenant_id)
                    a_map = {a.name.lower(): a for a in (await db.execute(a_q)).scalars().all()}
                    notified_m = set()
                    for m in mentions:
                        ma = a_map.get(m.lower())
                        if ma and ma.id not in notified_m:
                            notified_m.add(ma.id)
                            await send_notification(
                                db, agent_id=ma.id,
                                type="mention",
                                title=f"{agent.name} mentioned you in a comment",
                                body=content[:150],
                                link=f"/plaza?post={pid}",
                                ref_id=pid,
                                sender_name=agent.name,
                            )
            except Exception:
                pass

            await db.commit()
            return f"Comment added to post by {post.author_name}."

    except Exception as e:
        return f"Failed to add comment: {str(e)[:200]}"


# ─── Code Execution ─────────────────────────────────────────────

# Dangerous patterns to block (for legacy fallback)
_DANGEROUS_BASH_ALWAYS = [
    "rm -rf /", "rm -rf ~", "sudo ", "mkfs", "dd if=",
    ":(){ :", "chmod 777 /", "chown ", "shutdown", "reboot",
]

_DANGEROUS_BASH_NETWORK = [
    "curl ", "wget ", "nc ", "ncat ", "ssh ", "scp ",
]

_DANGEROUS_PYTHON_IMPORTS_ALWAYS = [
    "shutil.rmtree", "os.system", "os.popen",
    "os.exec", "os.spawn",
]

_DANGEROUS_PYTHON_IMPORTS_NETWORK = [
    "socket", "http.client", "urllib.request", "requests",
    "ftplib", "smtplib", "telnetlib", "ctypes",
]

_DANGEROUS_NODE_ALWAYS = [
    "fs.rmSync", "fs.rmdirSync", "process.exit",
]

_DANGEROUS_NODE_NETWORK = [
    "require('http')", "require('https')", "require('net')",
]


def _check_code_safety(language: str, code: str, allow_network: bool = False) -> str | None:
    """Check code for dangerous patterns. Returns error message if unsafe, None if ok."""
    code_lower = code.lower()

    if language == "bash":
        for pattern in _DANGEROUS_BASH_ALWAYS:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: dangerous command detected ({pattern.strip()})"
        if not allow_network:
            for pattern in _DANGEROUS_BASH_NETWORK:
                if pattern.lower() in code_lower:
                    return f"❌ Blocked: network command not allowed ({pattern.strip()})"
        if "../../" in code:
            return "❌ Blocked: directory traversal not allowed"

    elif language == "python":
        for pattern in _DANGEROUS_PYTHON_IMPORTS_ALWAYS:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: unsafe operation detected ({pattern})"
        if not allow_network:
            for pattern in _DANGEROUS_PYTHON_IMPORTS_NETWORK:
                if pattern.lower() in code_lower:
                    return f"❌ Blocked: network operation not allowed ({pattern})"

    elif language == "node":
        for pattern in _DANGEROUS_NODE_ALWAYS:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: unsafe operation detected ({pattern})"
        if not allow_network:
            for pattern in _DANGEROUS_NODE_NETWORK:
                if pattern.lower() in code_lower:
                    return f"❌ Blocked: network operation not allowed ({pattern})"

    return None


async def _execute_code_outcome(
    agent_id: Optional[uuid.UUID],
    ws: Path,
    arguments: dict,
    *,
    tool_name: str = "execute_code",
    on_output=None,
) -> ToolExecutionOutcome:
    """Execute code using the configured sandbox backend.

    Args:
        agent_id: The agent's UUID (used to fetch per-agent tool config).
        ws: Agent workspace root path.
        arguments: Tool call arguments (language, code, timeout).
        tool_name: The originating tool name — either 'execute_code' (local)
                   or 'execute_code_e2b' (cloud).  Used to look up the
                   correct per-agent tool config entry in the database.
    """
    language = arguments.get("language", "python")
    code = arguments.get("code", "")
    requested_timeout = arguments.get("timeout", 30)

    if not isinstance(code, str) or not code.strip():
        return _typed_failure("No code provided.", "invalid_tool_arguments")

    if language not in ("python", "bash", "node"):
        return _typed_failure(
            f"Unsupported language: {language}. Use python, bash, or node.",
            "invalid_tool_arguments",
        )
    try:
        requested_timeout = int(requested_timeout)
    except (TypeError, ValueError):
        return _typed_failure(
            "execute_code timeout must be an integer.",
            "invalid_tool_arguments",
        )
    if requested_timeout <= 0:
        return _typed_failure(
            "execute_code timeout must be positive.",
            "invalid_tool_arguments",
        )

    # Working directory is the agent's root directory (must be absolute).
    # This allows code to access skills/, workspace/, memory/ etc. directly.
    work_dir = ws.resolve()
    work_dir.mkdir(parents=True, exist_ok=True)

    # For E2B tool: do NOT fall back to local subprocess on error —
    # the user explicitly chose cloud execution.
    is_e2b_tool = (tool_name == "execute_code_e2b")

    fallback_config = None
    execution_started = False
    try:
        # Import here to avoid circular imports
        from app.config import get_sandbox_config
        from app.services.sandbox.config import SandboxConfig
        from app.services.sandbox.registry import get_sandbox_backend

        tool_config = await _get_tool_config(agent_id, tool_name)

        if is_e2b_tool:
            # The explicit E2B tool is available only with its own complete
            # local configuration. Never inherit the platform/local sandbox
            # fallback because that would silently execute code elsewhere.
            if not isinstance(tool_config, dict):
                return _typed_failure(
                    "E2B sandbox credentials are not configured.",
                    "sandbox_configuration_missing",
                )
            if tool_config.get("sandbox_type") != "e2b":
                return _typed_failure(
                    "execute_code_e2b requires sandbox_type=e2b.",
                    "sandbox_configuration_invalid",
                )
            api_key = tool_config.get("api_key")
            if not isinstance(api_key, str) or not api_key.strip():
                return _typed_failure(
                    "E2B sandbox credentials are not configured.",
                    "sandbox_configuration_missing",
                )
            try:
                default_timeout = int(tool_config.get("default_timeout", 30))
                max_timeout = int(tool_config.get("max_timeout", 60))
            except (TypeError, ValueError):
                return _typed_failure(
                    "E2B timeout configuration must be numeric.",
                    "sandbox_configuration_invalid",
                )
            sandbox_config = SandboxConfig(
                type="e2b",
                api_key=api_key.strip(),
                default_timeout=default_timeout,
                max_timeout=max_timeout,
            )
        else:
            # The default execute_code tool retains the established platform
            # fallback behavior; it is a distinct explicit tool contract.
            fallback_config = get_sandbox_config()
            if tool_config:
                sandbox_config = SandboxConfig.from_dict(
                    tool_config,
                    fallback_config,
                )
            else:
                sandbox_config = fallback_config
                logger.info(
                    "[Sandbox] No per-agent config found for '{}', using fallback",
                    tool_name,
                )

        # Clamp timeout by configured max_timeout (default 60s, up to 3600s)
        timeout = min(requested_timeout, sandbox_config.max_timeout)

        backend = get_sandbox_backend(sandbox_config)
        if is_e2b_tool:
            if getattr(backend, "name", None) != "e2b":
                return _typed_failure(
                    "E2B configuration resolved to a non-E2B backend.",
                    "sandbox_configuration_invalid",
                )
            # Load the optional SDK/client class before marking remote dispatch.
            # This is a deterministic local check, not a Provider health ping.
            try:
                getattr(backend, "client")
            except Exception as exc:
                return _typed_failure(
                    f"E2B backend could not start: {type(exc).__name__}.",
                    "sandbox_provider_unavailable",
                )
        logger.info(f"[Sandbox] Executing code with backend: {backend.__class__.__name__} (tool={tool_name}, timeout={timeout}s)")
        execution_started = True
        result = await backend.execute(
            code=code,
            language=language,
            timeout=timeout,
            work_dir=str(work_dir),
            on_output=on_output,
            agent_id=agent_id,
        )

        try:
            summary = backend._format_result(result)
        except Exception:
            summary = (
                "Code executed successfully."
                if result.success and result.exit_code == 0
                else f"Code execution failed with exit code {result.exit_code}."
            )
        if result.success and result.exit_code == 0:
            return _typed_success(summary)
        return _typed_failure(
            summary,
            "sandbox_execution_failed",
        )

    except ValueError as e:
        if execution_started:
            return _typed_unknown(
                "Sandbox execution outcome is unknown after ValueError; reconcile before retrying.",
                "sandbox_execution_outcome_unknown",
            )
        # Sandbox disabled or misconfigured
        if is_e2b_tool:
            # Do not silently fall back — surface the config error to the user
            return _typed_failure(
                f"E2B sandbox configuration error: {str(e)[:300]}",
                "sandbox_configuration_invalid",
            )
        if fallback_config is None:
            return _typed_failure(
                f"Sandbox configuration error: {str(e)[:300]}",
                "sandbox_configuration_invalid",
            )
        logger.warning(f"[Sandbox] Config issue, falling back to legacy subprocess: {e}")
        return await _execute_code_legacy_outcome(
            ws,
            arguments,
            allow_network=fallback_config.allow_network,
            max_timeout=fallback_config.max_timeout,
            on_output=on_output,
        )

    except Exception as e:
        logger.exception(f"[Sandbox] Execution failed for agent {agent_id} (tool={tool_name})")
        # Once backend.execute was entered, it may have run code or emitted
        # network/workspace side effects. Never start a second backend as a
        # fallback when that outcome is unprovable.
        if execution_started:
            return _typed_unknown(
                f"Sandbox execution outcome is unknown after {type(e).__name__}; reconcile before retrying.",
                "sandbox_execution_outcome_unknown",
            )
        return _typed_failure(
            f"Sandbox execution could not start: {type(e).__name__}.",
            "sandbox_execution_failed",
        )


async def _execute_code(
    agent_id: Optional[uuid.UUID],
    ws: Path,
    arguments: dict,
    *,
    tool_name: str = "execute_code",
    on_output=None,
) -> str:
    outcome = await _execute_code_outcome(
        agent_id,
        ws,
        arguments,
        tool_name=tool_name,
        on_output=on_output,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Code execution returned no summary.",
    )


async def _execute_code_legacy_outcome(
    ws: Path,
    arguments: dict,
    allow_network: bool = False,
    max_timeout: int = 60,
    on_output=None,
) -> ToolExecutionOutcome:
    """Legacy subprocess-based code execution (fallback)."""
    import asyncio

    language = arguments.get("language", "python")
    code = arguments.get("code", "")
    try:
        timeout = min(int(arguments.get("timeout", 30)), max_timeout)
    except (TypeError, ValueError):
        return _typed_failure(
            "execute_code timeout must be an integer.",
            "invalid_tool_arguments",
        )
    if timeout <= 0:
        return _typed_failure(
            "execute_code timeout must be positive.",
            "invalid_tool_arguments",
        )

    if not isinstance(code, str) or not code.strip():
        return _typed_failure("No code provided.", "invalid_tool_arguments")

    if language not in ("python", "bash", "node"):
        return _typed_failure(
            f"Unsupported language: {language}. Use python, bash, or node.",
            "invalid_tool_arguments",
        )

    # Security check
    safety_error = _check_code_safety(language, code, allow_network)
    if safety_error:
        return _typed_failure(safety_error, "sandbox_code_blocked")

    # Working directory is the agent's root directory (must be absolute)
    # This allows code to access skills/, workspace/, memory/ etc. directly
    work_dir = ws.resolve()
    work_dir.mkdir(parents=True, exist_ok=True)

    # Determine command and file extension
    if language == "python":
        ext = ".py"
        cmd_prefix = ["python3"]
    elif language == "bash":
        ext = ".sh"
        cmd_prefix = ["bash"]
    elif language == "node":
        ext = ".js"
        cmd_prefix = ["node"]
    else:
        return _typed_failure(
            f"Unsupported language: {language}.",
            "invalid_tool_arguments",
        )

    # Write code to a temp file inside workspace
    script_path = work_dir / f"_exec_tmp{ext}"
    proc = None
    try:
        script_path.write_text(code, encoding="utf-8")

        # Inherit parent environment but override HOME to workspace
        safe_env = dict(os.environ)
        safe_env["HOME"] = str(work_dir)
        safe_env["PYTHONDONTWRITEBYTECODE"] = "1"

        proc = await asyncio.create_subprocess_exec(
            *cmd_prefix, str(script_path),
            cwd=str(work_dir),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            env=safe_env,
        )

        stdout_data = bytearray()
        stderr_data = bytearray()

        async def read_stream(stream, out, label="stdout"):
            capture_limit = MAX_EXEC_STDERR_CAPTURE_BYTES if label == "stderr" else MAX_EXEC_STDOUT_CAPTURE_BYTES
            while True:
                chunk = await stream.read(4096)
                if not chunk:
                    break
                remaining = capture_limit - len(out)
                if remaining > 0:
                    out.extend(chunk[:remaining])
                # Real-time streaming: push each chunk to the WebSocket
                if on_output:
                    try:
                        text = chunk.decode("utf-8", errors="replace")
                        await on_output(text, label)
                    except Exception:
                        pass

        task1 = asyncio.create_task(read_stream(proc.stdout, stdout_data, "stdout"))
        task2 = asyncio.create_task(read_stream(proc.stderr, stderr_data, "stderr"))

        is_timeout = False
        try:
            await asyncio.wait_for(proc.wait(), timeout=timeout)
        except asyncio.TimeoutError:
            proc.kill()
            is_timeout = True

        await asyncio.gather(task1, task2)
        stdout = bytes(stdout_data)
        stderr = bytes(stderr_data)

        stdout_str = stdout.decode("utf-8", errors="replace")[:10000] if stdout else ""
        stderr_str = stderr.decode("utf-8", errors="replace")[:5000] if stderr else ""

        result_parts = []
        if stdout_str.strip():
            result_parts.append(f"📤 Output:\n{stdout_str}")
        if stderr_str.strip():
            result_parts.append(f"⚠️ Stderr:\n{stderr_str}")

        if is_timeout:
            result_parts.append(f"❌ Code execution timed out after {timeout}s. If you expect this code to take longer, try calling the tool again with a higher 'timeout' parameter (up to 3600s).")
            return _typed_failure(
                "\n\n".join(result_parts),
                "sandbox_execution_timeout",
            )

        if proc.returncode != 0:
            result_parts.append(f"Exit code: {proc.returncode}")
            return _typed_failure(
                "\n\n".join(result_parts),
                "sandbox_execution_failed",
            )

        if not result_parts:
            return _typed_success("Code executed successfully (no output).")

        return _typed_success("\n\n".join(result_parts))

    except Exception as e:
        if proc is not None:
            try:
                if proc.returncode is None:
                    proc.kill()
                    await proc.wait()
            except Exception:
                pass
            return _typed_unknown(
                f"Local code execution outcome is unknown after {type(e).__name__}.",
                "sandbox_execution_outcome_unknown",
            )
        return _typed_failure(
            f"Execution could not start: {type(e).__name__}.",
            "sandbox_execution_failed",
        )
    finally:
        # Clean up temp script
        try:
            script_path.unlink(missing_ok=True)
        except Exception:
            pass


async def _execute_code_legacy(
    ws: Path,
    arguments: dict,
    allow_network: bool = False,
    max_timeout: int = 60,
    on_output=None,
) -> str:
    outcome = await _execute_code_legacy_outcome(
        ws,
        arguments,
        allow_network=allow_network,
        max_timeout=max_timeout,
        on_output=on_output,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Code execution returned no summary.",
    )


# ─── Resource Discovery Executors ───────────────────────────────

async def _discover_resources_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "discover_resources requires query.",
            "invalid_tool_arguments",
        )
    try:
        max_results = int(arguments.get("max_results", 5))
    except (TypeError, ValueError):
        return _typed_failure(
            "discover_resources max_results must be an integer.",
            "invalid_tool_arguments",
        )
    if max_results < 1:
        return _typed_failure(
            "discover_resources max_results must be positive.",
            "invalid_tool_arguments",
        )
    from app.services.resource_discovery import search_registries_outcome

    return await search_registries_outcome(
        query.strip(),
        min(max_results, 10),
        agent_id=agent_id,
    )


async def _discover_resources(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter for typed resource discovery."""
    outcome = await _discover_resources_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Resource discovery returned no summary.",
    )


async def _import_mcp_server_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Import one MCP server without interpreting provider display strings."""
    server_id = arguments.get("server_id")
    if not isinstance(server_id, str) or not server_id.strip():
        return _typed_failure(
            "import_mcp_server requires server_id.",
            "invalid_tool_arguments",
        )
    raw_config = arguments.get("config", {})
    if raw_config is None:
        raw_config = {}
    if not isinstance(raw_config, dict):
        return _typed_failure(
            "import_mcp_server config must be an object.",
            "invalid_tool_arguments",
        )
    config = dict(raw_config)
    reauthorize = arguments.get("reauthorize", False)
    if not isinstance(reauthorize, bool):
        return _typed_failure(
            "import_mcp_server reauthorize must be a boolean.",
            "invalid_tool_arguments",
        )

    mcp_url = config.pop("mcp_url", None)
    try:
        if mcp_url is not None:
            if not isinstance(mcp_url, str) or not mcp_url.startswith(
                ("http://", "https://")
            ):
                return _typed_failure(
                    "import_mcp_server config.mcp_url must be an HTTP(S) URL.",
                    "invalid_tool_arguments",
                )
            from app.services.resource_discovery import import_mcp_direct_outcome

            server_name = config.pop("server_name", None) or server_id.strip()
            api_key = config.pop("api_key", None)
            return await import_mcp_direct_outcome(
                mcp_url,
                agent_id,
                server_name,
                api_key,
            )

        from app.services.resource_discovery import import_mcp_from_smithery_outcome

        return await import_mcp_from_smithery_outcome(
            server_id.strip(),
            agent_id,
            config or None,
            reauthorize=reauthorize,
        )
    except Exception as exc:
        logger.error(
            "[ResourceDiscovery] MCP import outcome became unknown: {}",
            type(exc).__name__,
        )
        return _typed_unknown(
            "MCP import outcome is unknown; reconcile before retrying.",
            "mcp_import_outcome_unknown",
        )


async def _import_mcp_server(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter for typed MCP import."""
    outcome = await _import_mcp_server_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="MCP import returned no summary.",
    )


# ─── Trigger Management Handlers (Aware Engine) ────────────────────

MAX_TRIGGERS_PER_AGENT = 20
VALID_TRIGGER_TYPES = {"cron", "once", "interval", "poll", "on_message", "webhook"}


async def _handle_set_trigger_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    session_id: str = "",
    user_id: uuid.UUID | None = None,
) -> ToolExecutionOutcome:
    """Create a trigger from validated config and the committed DB fact."""
    from app.models.trigger import AgentTrigger
    from app.models.chat_session import ChatSession

    raw_name = arguments.get("name", "")
    raw_type = arguments.get("type", "")
    raw_reason = arguments.get("reason", "")
    raw_focus_ref = arguments.get("focus_ref", "") or arguments.get(
        "agenda_ref", ""
    )
    if not all(
        isinstance(value, str)
        for value in (raw_name, raw_type, raw_reason, raw_focus_ref)
    ):
        return _typed_failure(
            "set_trigger name, type, reason, and focus_ref must be strings.",
            "invalid_tool_arguments",
        )
    name = raw_name.strip()
    ttype = raw_type.strip()
    raw_config = arguments.get("config")
    if not isinstance(raw_config, dict):
        return _typed_failure(
            "set_trigger config must be an object.",
            "invalid_tool_arguments",
        )
    config = dict(raw_config)
    reason = raw_reason.strip()
    focus_ref = raw_focus_ref.strip()  # agenda_ref is backward compatibility only

    if not name:
        return _typed_failure(
            "set_trigger requires name.",
            "invalid_tool_arguments",
        )
    if ttype not in VALID_TRIGGER_TYPES:
        return _typed_failure(
            f"Invalid trigger type '{ttype}'.",
            "invalid_tool_arguments",
        )
    if not reason:
        return _typed_failure(
            "set_trigger requires reason.",
            "invalid_tool_arguments",
        )

    # Validate type-specific config
    if ttype == "cron":
        expr = config.get("expr", "")
        if not expr:
            return _typed_failure(
                "cron trigger requires config.expr.",
                "invalid_tool_arguments",
            )
        try:
            from croniter import croniter
            croniter(expr)
        except Exception:
            return _typed_failure(
                f"Invalid cron expression: '{expr}'.",
                "invalid_tool_arguments",
            )
    elif ttype == "once":
        if not config.get("at"):
            return _typed_failure(
                "once trigger requires config.at.",
                "invalid_tool_arguments",
            )
    elif ttype == "interval":
        if not config.get("minutes"):
            return _typed_failure(
                "interval trigger requires config.minutes.",
                "invalid_tool_arguments",
            )
    elif ttype == "poll":
        if not config.get("url"):
            return _typed_failure(
                "poll trigger requires config.url.",
                "invalid_tool_arguments",
            )
    elif ttype == "on_message":
        if not config.get("from_agent_name") and not config.get("from_user_name"):
            return _typed_failure(
                "on_message trigger requires from_agent_name or from_user_name.",
                "invalid_tool_arguments",
            )
        # Snapshot the latest message timestamp so we only detect NEW messages after this point
        # This prevents false positives from already-processed messages
        try:
            from app.models.audit import ChatMessage
            from app.models.chat_session import ChatSession
            from sqlalchemy import cast as sa_cast, String as SaString
            async with async_session() as _snap_db:
                _snap_q = select(ChatMessage.created_at).join(
                    ChatSession, ChatMessage.conversation_id == sa_cast(ChatSession.id, SaString)
                ).where(
                    ChatSession.agent_id == agent_id,
                    ChatMessage.created_at.isnot(None),
                ).order_by(ChatMessage.created_at.desc()).limit(1)
                _snap_r = await _snap_db.execute(_snap_q)
                _latest_ts = _snap_r.scalar_one_or_none()
                if _latest_ts:
                    config["_since_ts"] = _latest_ts.isoformat()
        except Exception:
            pass  # Fallback to trigger.created_at in the daemon
    elif ttype == "webhook":
        # Auto-generate a unique token for the webhook URL
        import secrets
        token = secrets.token_urlsafe(8)  # ~11 chars, URL-safe
        config["token"] = token

    if ttype == "webhook":
        try:
            from app.services.platform_service import platform_service

            base = await platform_service.get_public_base_url()
            if not isinstance(base, str) or not base.strip():
                return _typed_failure(
                    "A public base URL is required for webhook triggers.",
                    "trigger_webhook_base_url_missing",
                )
        except Exception as exc:
            return _typed_failure(
                f"Webhook URL could not be prepared: {type(exc).__name__}.",
                "trigger_webhook_setup_failed",
            )

    # Record the session that created this trigger so trigger results can later be routed to
    # the correct destination instead of being broadcast to every live web session.
    if session_id:
        try:
            async with async_session() as _ctx_db:
                _session_result = await _ctx_db.execute(
                    select(ChatSession).where(ChatSession.id == uuid.UUID(session_id))
                )
                origin_session = _session_result.scalar_one_or_none()
                if origin_session:
                    config["_origin_session_id"] = str(origin_session.id)
                    config["_origin_source_channel"] = origin_session.source_channel
                    if origin_session.source_channel == "agent" and origin_session.peer_agent_id:
                        config["_origin_peer_agent_id"] = str(origin_session.peer_agent_id)
                    elif origin_session.source_channel != "trigger":
                        config["_origin_user_id"] = str(origin_session.user_id)
                elif user_id:
                    config["_origin_user_id"] = str(user_id)
        except Exception:
            if user_id:
                config["_origin_user_id"] = str(user_id)

    mutation_started = False
    try:
        async with async_session() as db:
            # Load agent to get per-agent trigger limit
            from app.models.agent import Agent as _AgentModel
            _a_result = await db.execute(select(_AgentModel).where(_AgentModel.id == agent_id))
            _agent_obj = _a_result.scalar_one_or_none()
            agent_max_triggers = (_agent_obj.max_triggers if _agent_obj else None) or MAX_TRIGGERS_PER_AGENT

            # Check max triggers
            from sqlalchemy import func as sa_func
            result = await db.execute(
                select(sa_func.count()).select_from(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.is_enabled == True,
                )
            )
            count = result.scalar() or 0
            if count >= agent_max_triggers:
                return _typed_failure(
                    f"Maximum trigger limit reached ({agent_max_triggers}).",
                    "trigger_limit_reached",
                )

            # Check for duplicate name
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            existing = result.scalar_one_or_none()
            if existing:
                if existing.is_enabled:
                    return _typed_failure(
                        f"Trigger '{name}' already exists and is active.",
                        "trigger_already_exists",
                    )
                else:
                    focus_ref = await ensure_focus_item(
                        agent_id,
                        focus_ref=focus_ref,
                        description=reason,
                        system=False,
                        db=db,
                    )
                    # Re-enable disabled trigger with new config (preserve fire history)
                    # For webhook triggers: reuse the old token so the URL stays stable
                    if ttype == "webhook":
                        old_token = (existing.config or {}).get("token")
                        if old_token:
                            config["token"] = old_token
                    existing.type = ttype
                    existing.config = config
                    existing.reason = reason
                    existing.focus_ref = focus_ref
                    existing.is_enabled = True
                    # Keep fire_count and last_fired_at — they are cumulative stats,
                    # but reset fire_count if it reached max_fires to allow it to run again.
                    if existing.max_fires and existing.fire_count >= existing.max_fires:
                        existing.fire_count = 0
                    mutation_started = True
                    await db.commit()
                    return _typed_success(
                        f"Trigger '{name}' re-enabled with new configuration "
                        f"({ttype}, fired {existing.fire_count} times so far)."
                    )

            focus_ref = await ensure_focus_item(
                agent_id,
                focus_ref=focus_ref,
                description=reason,
                system=False,
                db=db,
            )
            trigger = AgentTrigger(
                agent_id=agent_id,
                name=name,
                type=ttype,
                config=config,
                reason=reason,
                focus_ref=focus_ref,
            )
            # Fix 4: Safety cap for on_message triggers —
            # prevent infinite loops if agent creates broad watchers.
            if ttype == "on_message":
                trigger.max_fires = trigger.max_fires or 100
                if not trigger.expires_at:
                    trigger.expires_at = datetime.now(timezone.utc) + timedelta(days=7)
            db.add(trigger)
            mutation_started = True
            await db.commit()

        # Activity log
        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_created", {
                "name": name, "type": ttype, "reason": reason[:100],
            }, agent_id=agent_id)
        except Exception:
            pass

        # Return webhook URL for webhook triggers
        if ttype == "webhook":
            return _typed_success(
                f"Webhook trigger '{name}' created. Open the Trigger settings "
                "to copy its private webhook URL."
            )

        return _typed_success(
            f"Trigger '{name}' created ({ttype}). It will wake this Agent "
            "with the configured reason when it fires."
        )

    except Exception as exc:
        if mutation_started:
            return _typed_unknown(
                "Trigger creation outcome is unknown; reconcile before retrying.",
                "trigger_create_outcome_unknown",
            )
        return _typed_failure(
            f"Trigger could not be created: {type(exc).__name__}.",
            "trigger_create_failed",
        )


async def _handle_set_trigger(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    session_id: str = "",
    user_id: uuid.UUID | None = None,
) -> str:
    outcome = await _handle_set_trigger_outcome(
        agent_id,
        arguments,
        session_id=session_id,
        user_id=user_id,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Trigger creation returned no summary.",
    )


async def _handle_update_trigger_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Update an existing trigger's config or reason."""
    from app.models.trigger import AgentTrigger

    name = arguments.get("name", "").strip()
    if not name:
        return _typed_failure(
            "Missing required argument 'name'.",
            "invalid_tool_arguments",
        )

    new_config = arguments.get("config")
    new_reason = arguments.get("reason")

    if new_config is None and new_reason is None:
        return _typed_failure(
            "Provide at least one of config or reason to update.",
            "invalid_tool_arguments",
        )

    commit_started = False
    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            trigger = result.scalar_one_or_none()
            if not trigger:
                return _typed_failure(
                    f"Trigger '{name}' not found.",
                    "trigger_not_found",
                )

            changes = []
            if new_config is not None:
                if not isinstance(new_config, dict):
                    return _typed_failure(
                        "config must be an object.",
                        "invalid_tool_arguments",
                    )
                old_config = dict(trigger.config or {})
                protected = {
                    key: value
                    for key, value in old_config.items()
                    if key == "token" or key.startswith("_")
                }
                user_patch = {
                    key: value
                    for key, value in new_config.items()
                    if key != "token" and not key.startswith("_")
                }
                trigger.config = {**old_config, **user_patch, **protected}
                changes.append(f"config fields patched: {sorted(user_patch)}")
            if new_reason is not None:
                if not isinstance(new_reason, str) or not new_reason.strip():
                    return _typed_failure(
                        "reason must be a non-empty string.",
                        "invalid_tool_arguments",
                    )
                trigger.reason = new_reason
                changes.append(f"reason updated")

            commit_started = True
            await db.commit()

        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_updated", {
                "name": name, "changes": "; ".join(changes),
            }, agent_id=agent_id)
        except Exception:
            pass

        return _typed_success(
            f"Trigger '{name}' updated: {'; '.join(changes)}"
        )

    except Exception as e:
        if commit_started:
            return _typed_unknown(
                "Trigger update outcome is unknown; reconcile before retrying.",
                "trigger_update_outcome_unknown",
            )
        return _typed_failure(
            f"Failed to update trigger: {type(e).__name__}.",
            "trigger_update_failed",
        )


async def _handle_update_trigger(agent_id: uuid.UUID, arguments: dict) -> str:
    outcome = await _handle_update_trigger_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Trigger update returned no summary.",
    )


async def _handle_cancel_trigger_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Cancel (disable) a trigger by name."""
    from app.models.trigger import AgentTrigger

    name = arguments.get("name", "").strip()
    if not name:
        return _typed_failure(
            "Missing required argument 'name'.",
            "invalid_tool_arguments",
        )

    commit_started = False
    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            trigger = result.scalar_one_or_none()
            if not trigger:
                return _typed_failure(
                    f"Trigger '{name}' not found.",
                    "trigger_not_found",
                )
            if not trigger.is_enabled:
                return _typed_success(f"Trigger '{name}' is already disabled.")

            trigger.is_enabled = False
            commit_started = True
            await db.commit()

        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_cancelled", {"name": name}, agent_id=agent_id)
        except Exception:
            pass

        return _typed_success(
            f"Trigger '{name}' cancelled. It will no longer fire."
        )

    except Exception as e:
        if commit_started:
            return _typed_unknown(
                "Trigger cancellation outcome is unknown; reconcile before retrying.",
                "trigger_cancel_outcome_unknown",
            )
        return _typed_failure(
            f"Failed to cancel trigger: {type(e).__name__}.",
            "trigger_cancel_failed",
        )


async def _handle_cancel_trigger(agent_id: uuid.UUID, arguments: dict) -> str:
    outcome = await _handle_cancel_trigger_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Trigger cancellation returned no summary.",
    )


async def _handle_list_triggers_outcome(
    agent_id: uuid.UUID,
) -> ToolExecutionOutcome:
    """List all active triggers for the agent."""
    from app.models.trigger import AgentTrigger

    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                ).order_by(AgentTrigger.created_at.desc())
            )
            triggers = result.scalars().all()

        if not triggers:
            return _typed_success("No triggers found. Use set_trigger to create one.")

        lines = ["| Name | Type | Config | Reason | Status | Fires |", "|------|------|--------|--------|--------|-------|"]
        for t in triggers:
            status = "✅ active" if t.is_enabled else "⏸ disabled"
            config_str = str(t.config)[:50]
            reason_str = t.reason[:40] if t.reason else ""
            lines.append(f"| {t.name} | {t.type} | {config_str} | {reason_str} | {status} | {t.fire_count} |")

        return _typed_success("\n".join(lines))

    except Exception as e:
        return _typed_failure(
            f"Failed to list triggers: {type(e).__name__}.",
            "trigger_list_failed",
            retryable=True,
        )


async def _handle_list_triggers(agent_id: uuid.UUID) -> str:
    outcome = await _handle_list_triggers_outcome(agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Trigger listing returned no summary.",
    )


# ─── Image Upload (ImageKit CDN) ────────────────────────────────

def _image_public_http_url(value: object) -> str | None:
    """Validate a provider-fetchable URL without performing network I/O."""
    import ipaddress
    from urllib.parse import urlsplit

    if not isinstance(value, str):
        return None
    candidate = value.strip()
    if not candidate or len(candidate.encode("utf-8")) > 2048:
        return None
    try:
        parsed = urlsplit(candidate)
        hostname = parsed.hostname
        parsed.port
    except ValueError:
        return None
    if (
        parsed.scheme not in {"http", "https"}
        or not hostname
        or parsed.username is not None
        or parsed.password is not None
    ):
        return None
    normalized_host = hostname.lower().rstrip(".")
    if normalized_host in {"localhost", "localhost.localdomain"} or normalized_host.endswith(
        ".local"
    ):
        return None
    try:
        address = ipaddress.ip_address(normalized_host)
    except ValueError:
        return candidate
    return candidate if address.is_global else None

async def _upload_image_outcome(
    agent_id: uuid.UUID,
    ws: Path,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Upload an image to ImageKit CDN and return the public URL.

    Credential resolution order:
    1. Global tool config (admin-set, shared by all agents)
    2. Per-agent tool config override (agent-specific)
    """
    import httpx
    import base64

    file_path = arguments.get("file_path")
    source_url = arguments.get("url")
    file_name = arguments.get("file_name")
    folder = arguments.get("folder", "/clawith")

    if file_path is not None and not isinstance(file_path, str):
        return _typed_failure(
            "file_path must be a workspace-relative string.",
            "invalid_tool_arguments",
        )
    if source_url is not None and not isinstance(source_url, str):
        return _typed_failure(
            "url must be a public HTTP(S) URL.",
            "invalid_tool_arguments",
        )
    normalized_file_path = file_path.strip() if isinstance(file_path, str) else ""
    normalized_source_url = (
        source_url.strip() if isinstance(source_url, str) else ""
    )
    if bool(normalized_file_path) == bool(normalized_source_url):
        return _typed_failure(
            "Provide exactly one of file_path or url.",
            "invalid_tool_arguments",
        )
    if normalized_source_url:
        validated_source_url = _image_public_http_url(normalized_source_url)
        if not validated_source_url:
            return _typed_failure(
                "url must be a public HTTP(S) URL.",
                "invalid_tool_arguments",
            )
        normalized_source_url = validated_source_url
    if file_name is not None and (
        not isinstance(file_name, str) or not file_name.strip()
    ):
        return _typed_failure(
            "file_name must be a non-empty string when provided.",
            "invalid_tool_arguments",
        )
    if not isinstance(folder, str) or not folder.strip():
        return _typed_failure(
            "folder must be a non-empty string.",
            "invalid_tool_arguments",
        )

    # ── Load ImageKit credentials (Agent > Company priority) ──
    private_key = ""
    url_endpoint = ""
    try:
        # Use standard _get_tool_config (Agent > Company, cached, schema-aware decryption)
        config = await _get_tool_config(agent_id, "upload_image") or {}
        private_key = config.get("private_key", "")
        url_endpoint = config.get("url_endpoint", "")
    except Exception as exc:
        logger.error(
            "[UploadImage] Config load error: {}",
            type(exc).__name__,
        )

    if not private_key:
        return _typed_failure(
            "ImageKit Private Key is not configured.",
            "imagekit_credentials_missing",
        )

    # ── Prepare the file ──
    form_data = {}
    file_content = None

    if normalized_file_path:
        # Read from workspace
        full_path = (ws / normalized_file_path).resolve()
        try:
            full_path.relative_to(ws.resolve())
        except ValueError:
            return _typed_failure(
                "Access denied: path is outside the workspace.",
                "workspace_path_invalid",
            )
        if not full_path.exists():
            return _typed_failure(
                f"File not found: {normalized_file_path}",
                "upload_source_not_found",
            )
        if not full_path.is_file():
            return _typed_failure(
                f"Not a file: {normalized_file_path}",
                "upload_source_invalid",
            )

        # Check file size (max 25MB for free plan)
        try:
            file_size = full_path.stat().st_size
        except OSError as exc:
            return _typed_failure(
                f"Image upload source could not be inspected: {type(exc).__name__}.",
                "upload_source_read_failed",
            )
        size_mb = file_size / (1024 * 1024)
        if size_mb > 25:
            return _typed_failure(
                f"File too large ({size_mb:.1f}MB). Maximum is 25MB.",
                "upload_source_too_large",
            )
        try:
            file_content = full_path.read_bytes()
        except OSError as exc:
            return _typed_failure(
                f"Image upload source could not be read: {type(exc).__name__}.",
                "upload_source_read_failed",
            )

        if not file_name:
            file_name = full_path.name
    else:
        # Pass URL directly to ImageKit
        form_data["file"] = normalized_source_url
        if not file_name:
            from urllib.parse import urlparse
            file_name = (
                urlparse(normalized_source_url).path.split("/")[-1]
                or "image.jpg"
            )

    if not file_name:
        file_name = "image.png"

    form_data["fileName"] = file_name
    form_data["folder"] = folder
    form_data["useUniqueFileName"] = "true"

    # ── Upload to ImageKit V2 ──
    auth_string = base64.b64encode(f"{private_key}:".encode()).decode()

    request_started = False
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            if file_content is not None:
                # Binary upload via multipart
                files = {"file": (file_name, file_content)}
                request_started = True
                resp = await client.post(
                    "https://upload.imagekit.io/api/v2/files/upload",
                    headers={"Authorization": f"Basic {auth_string}"},
                    data=form_data,
                    files=files,
                )
            else:
                # URL upload via form data
                request_started = True
                resp = await client.post(
                    "https://upload.imagekit.io/api/v2/files/upload",
                    headers={"Authorization": f"Basic {auth_string}"},
                    data=form_data,
                )

        if resp.status_code in (200, 201):
            try:
                result = resp.json()
            except Exception:
                return _typed_unknown(
                    "ImageKit accepted the request but returned an unreadable response; reconcile before retrying.",
                    "imagekit_response_invalid",
                )
            if not isinstance(result, Mapping):
                return _typed_unknown(
                    "ImageKit returned an invalid success response; reconcile before retrying.",
                    "imagekit_response_invalid",
                )
            cdn_url = result.get("url", "")
            file_id = result.get("fileId", "")
            if not isinstance(cdn_url, str) or not cdn_url or not isinstance(file_id, str) or not file_id:
                return _typed_unknown(
                    "ImageKit success response omitted the stable file reference; reconcile before retrying.",
                    "imagekit_response_incomplete",
                )
            from urllib.parse import urlsplit

            parsed_cdn_url = urlsplit(cdn_url)
            if parsed_cdn_url.scheme != "https" or not parsed_cdn_url.hostname:
                return _typed_unknown(
                    "ImageKit returned an invalid CDN URL; reconcile before retrying.",
                    "imagekit_response_invalid",
                )
            if url_endpoint:
                normalized_endpoint = url_endpoint.rstrip("/")
                if cdn_url != normalized_endpoint and not cdn_url.startswith(
                    normalized_endpoint + "/"
                ):
                    return _typed_unknown(
                        "ImageKit returned a URL outside the configured endpoint; reconcile before retrying.",
                        "imagekit_response_invalid",
                    )
            try:
                size = max(float(result.get("size", 0)), 0)
            except (TypeError, ValueError):
                size = 0
            size_str = f"{size / 1024:.1f}KB" if size < 1024 * 1024 else f"{size / (1024 * 1024):.1f}MB"
            return _typed_success(
                f"Image uploaded successfully!\n\n"
                f"**CDN URL**: {cdn_url}\n"
                f"**File ID**: {file_id}\n"
                f"**Size**: {size_str}\n"
                f"**Name**: {result.get('name', file_name)}",
                result_ref=f"imagekit://{file_id}",
                artifact_refs=(f"imagekit://{file_id}",),
                evidence_refs=(cdn_url,),
            )
        elif 400 <= resp.status_code < 500:
            return _typed_failure(
                f"ImageKit rejected the upload with HTTP {resp.status_code}.",
                "imagekit_upload_rejected",
            )
        return _typed_unknown(
            f"ImageKit returned HTTP {resp.status_code} after the upload was sent; reconcile before retrying.",
            "imagekit_upload_outcome_unknown",
        )

    except httpx.TimeoutException:
        if request_started:
            return _typed_unknown(
                "ImageKit upload timed out after the request was sent; reconcile before retrying.",
                "imagekit_upload_outcome_unknown",
            )
        return _typed_failure(
            "ImageKit connection timed out before the request was sent.",
            "imagekit_connection_timeout",
        )
    except Exception as e:
        if request_started:
            return _typed_unknown(
                f"ImageKit upload outcome is unknown after {type(e).__name__}; reconcile before retrying.",
                "imagekit_upload_outcome_unknown",
            )
        return _typed_failure(
            f"ImageKit upload could not start: {type(e).__name__}.",
            "imagekit_upload_failed",
        )


async def _upload_image(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    outcome = await _upload_image_outcome(agent_id, ws, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Image upload returned no summary.",
    )



# ─── Image Generation (Multi-Provider) ────────────────────────────────────────

class _ImageGenerationBoundaryError(RuntimeError):
    def __init__(self, status: str, error_code: str, summary: str) -> None:
        super().__init__(summary)
        self.status = status
        self.error_code = error_code
        self.summary = summary


def _image_generation_failure(error_code: str, summary: str) -> None:
    raise _ImageGenerationBoundaryError("failed", error_code, summary)


def _image_generation_unknown(error_code: str, summary: str) -> None:
    raise _ImageGenerationBoundaryError("unknown", error_code, summary)


def _generated_image_media_type(image_bytes: bytes) -> str | None:
    if image_bytes.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if image_bytes.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if (
        len(image_bytes) >= 12
        and image_bytes.startswith(b"RIFF")
        and image_bytes[8:12] == b"WEBP"
    ):
        return "image/webp"
    return None


def _validate_generated_image_bytes(image_bytes: object) -> tuple[bytes, str]:
    if not isinstance(image_bytes, bytes) or not image_bytes:
        _image_generation_unknown(
            "image_result_invalid",
            "The image provider returned an empty or invalid image payload; do not regenerate automatically.",
        )
    if len(image_bytes) > _MAX_GENERATED_IMAGE_BYTES:
        _image_generation_unknown(
            "image_result_too_large",
            "The generated image exceeded the 25 MiB safety limit; do not regenerate automatically.",
        )
    media_type = _generated_image_media_type(image_bytes)
    if not media_type:
        _image_generation_unknown(
            "image_result_invalid",
            "The provider result was not a supported PNG, JPEG, or WebP image; do not regenerate automatically.",
        )
    return image_bytes, media_type


def _image_workspace_target(ws: Path, save_path: str) -> Path:
    if (
        not save_path
        or len(save_path.encode("utf-8")) > 1024
        or "\\" in save_path
    ):
        raise ValueError("invalid workspace image path")
    relative_path = Path(save_path)
    if (
        relative_path.is_absolute()
        or ".." in relative_path.parts
        or relative_path.suffix.lower() not in {".png", ".jpg", ".jpeg", ".webp"}
    ):
        raise ValueError("invalid workspace image path")
    workspace_root = ws.resolve()
    target = (workspace_root / relative_path).resolve()
    try:
        target.relative_to(workspace_root)
    except ValueError as exc:
        raise ValueError("invalid workspace image path") from exc
    return target


async def _generate_image_outcome(
    agent_id: uuid.UUID,
    ws: Path,
    arguments: dict,
    provider: str,
) -> ToolExecutionOutcome:
    """Generate once, then settle Provider and Workspace facts explicitly."""
    prompt_value = arguments.get("prompt")
    if not isinstance(prompt_value, str) or not prompt_value.strip():
        return _typed_failure(
            "Image generation requires a non-empty prompt.",
            "invalid_tool_arguments",
        )
    prompt = prompt_value.strip()

    size = arguments.get("size", "1024x1024")
    if not isinstance(size, str) or size not in _IMAGE_GENERATION_SIZES:
        return _typed_failure(
            "Image size is not supported.",
            "invalid_tool_arguments",
        )

    save_path_value = arguments.get("save_path", "")
    if save_path_value is not None and not isinstance(save_path_value, str):
        return _typed_failure(
            "save_path must be a workspace-relative image path.",
            "invalid_tool_arguments",
        )
    save_path = (save_path_value or "").strip()
    if not save_path:
        slug = "_".join(prompt.split()[:4]).lower()
        slug = "".join(
            character
            for character in slug
            if character.isalnum() or character == "_"
        )[:40] or "generated"
        timestamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        save_path = f"workspace/images/{slug}_{timestamp}.png"
    try:
        full_save_path = _image_workspace_target(ws, save_path)
    except ValueError:
        return _typed_failure(
            "save_path must remain inside the workspace and use PNG, JPEG, or WebP.",
            "workspace_path_invalid",
        )

    if provider not in {"siliconflow", "openai", "google", "custom"}:
        return _typed_failure(
            "Unknown image generation provider.",
            "invalid_tool_arguments",
        )
    tool_key = f"generate_image_{provider}"
    try:
        config = await _get_tool_config(agent_id, tool_key) or {}
    except Exception as exc:
        return _typed_failure(
            f"Image provider configuration could not be loaded: {type(exc).__name__}.",
            "image_configuration_unavailable",
        )
    api_key = str(config.get("api_key") or "").strip()
    if not api_key:
        return _typed_failure(
            "Image generation credentials are not configured.",
            "image_credentials_missing",
        )
    model = str(config.get("model") or "").strip()
    base_url = str(config.get("base_url") or "").strip()

    try:
        if provider == "siliconflow":
            image_bytes = await _generate_image_siliconflow(
                api_key,
                model or "black-forest-labs/FLUX.1-schnell",
                base_url or "https://api.siliconflow.cn/v1",
                prompt,
                size,
            )
        elif provider == "openai":
            image_bytes = await _generate_image_openai(
                api_key,
                model or "gpt-image-1",
                base_url or "https://api.openai.com/v1",
                prompt,
                size,
            )
        elif provider == "google":
            image_bytes = await _generate_image_google(
                api_key,
                model or "gemini-2.5-flash-image",
                base_url
                or "https://generativelanguage.googleapis.com/v1beta",
                prompt,
                size,
            )
        else:
            image_bytes = await _generate_image_custom_api(
                api_key=api_key,
                model=model,
                base_url=base_url,
                endpoint_path=config.get("endpoint_path")
                or "/chat/completions",
                request_body_template_json=config.get(
                    "request_body_template_json"
                )
                or "",
                response_image_path=config.get("response_image_path")
                or "choices.0.message.images.0.image_url.url",
                extra_headers_json=config.get("extra_headers_json") or "",
                timeout_seconds=config.get("timeout_seconds") or 120,
                prompt=prompt,
                size=size,
            )
        image_bytes, media_type = _validate_generated_image_bytes(image_bytes)
    except _ImageGenerationBoundaryError as exc:
        if exc.status == "failed":
            return _typed_failure(exc.summary, exc.error_code)
        return _typed_unknown(exc.summary, exc.error_code)
    except Exception as exc:
        logger.warning(
            "[GenerateImage] Unclassified provider boundary error for {}: {}",
            provider,
            type(exc).__name__,
        )
        return _typed_unknown(
            "The image generation outcome is unknown; do not regenerate automatically.",
            "image_generation_outcome_unknown",
        )

    # Provider generation is already dispatched. Any local persistence failure
    # from this point is unknown and must never trigger another generation.
    try:
        full_save_path.parent.mkdir(parents=True, exist_ok=True)
        if _image_workspace_target(ws, save_path) != full_save_path:
            raise OSError("workspace target changed during image generation")
        full_save_path.write_bytes(image_bytes)
    except Exception as exc:
        return _typed_unknown(
            f"The image was generated but could not be saved durably: {type(exc).__name__}.",
            "image_workspace_write_unknown",
            metadata={
                "provider": provider,
                "workspace_path": save_path,
                "content_hash": hashlib.sha256(image_bytes).hexdigest(),
                "artifact_content_hash": hashlib.sha256(
                    image_bytes
                ).hexdigest(),
                "mime_type": media_type,
                "size": len(image_bytes),
            },
        )

    artifact_ref = _workspace_artifact_ref(agent_id, save_path)
    content_hash = hashlib.sha256(image_bytes).hexdigest()
    api_image_path = (
        f"/api/agents/{agent_id}/files/download?path={save_path}"
    )
    return _typed_success(
        f"Image generated and saved to {save_path} using {provider}.\n\n"
        f"![generated image]({api_image_path})",
        result_ref=artifact_ref,
        artifact_refs=(artifact_ref,),
        metadata={
            "provider": provider,
            "operation": "image_generation",
            "workspace_path": save_path,
            "content_hash": content_hash,
            "artifact_content_hash": content_hash,
            "mime_type": media_type,
            "size": len(image_bytes),
        },
    )


async def _generate_image(
    agent_id: uuid.UUID,
    ws: Path,
    arguments: dict,
    provider: str,
) -> str:
    outcome = await _generate_image_outcome(agent_id, ws, arguments, provider)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Image generation returned no summary.",
    )


def _settle_image_provider_status(provider: str, status_code: int) -> None:
    if 200 <= status_code < 300:
        return
    if 400 <= status_code < 500:
        _image_generation_failure(
            "image_provider_rejected",
            f"The {provider} image provider rejected the request with HTTP {status_code}.",
        )
    _image_generation_unknown(
        "image_generation_outcome_unknown",
        f"The {provider} image provider returned HTTP {status_code} after dispatch; do not regenerate automatically.",
    )


def _decode_generated_image_base64(value: object) -> bytes:
    import base64

    if not isinstance(value, str) or not value:
        _image_generation_unknown(
            "image_result_invalid",
            "The image provider returned an invalid base64 image receipt.",
        )
    try:
        return base64.b64decode(value, validate=True)
    except (ValueError, TypeError):
        _image_generation_unknown(
            "image_result_invalid",
            "The image provider returned an invalid base64 image receipt.",
        )


async def _download_generated_image(image_url: object, client: Any) -> bytes:
    validated_url = _image_public_http_url(image_url)
    if not validated_url:
        _image_generation_unknown(
            "image_download_reference_invalid",
            "The image provider returned an invalid download reference.",
        )
    response = await client.get(validated_url, timeout=60)
    if not 200 <= response.status_code < 300:
        _image_generation_unknown(
            "image_download_outcome_unknown",
            f"The generated image download returned HTTP {response.status_code}; do not regenerate automatically.",
        )
    return response.content


async def _generate_image_siliconflow(
    api_key: str, model: str, base_url: str, prompt: str, size: str
) -> bytes:
    """Generate image via SiliconFlow (OpenAI-compatible images.generate API).

    SiliconFlow returns a temporary URL (expires in ~1 hour), so we download
    the image bytes immediately after generation.
    """
    import httpx
    url = f"{base_url.rstrip('/')}/images/generations"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model,
        "prompt": prompt,
        "image_size": size,  # SiliconFlow uses 'image_size' instead of 'size'
        "n": 1,
    }

    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(url, json=payload, headers=headers)
        _settle_image_provider_status("SiliconFlow", resp.status_code)
        try:
            data = resp.json()
        except Exception:
            _image_generation_unknown(
                "image_provider_response_invalid",
                "SiliconFlow returned an unreadable success response.",
            )
        if not isinstance(data, Mapping):
            _image_generation_unknown(
                "image_provider_response_invalid",
                "SiliconFlow returned an invalid success response.",
            )

        # SiliconFlow may return url or b64_json
        results = data.get("data")
        if not isinstance(results, list) or not results or not isinstance(
            results[0], Mapping
        ):
            _image_generation_unknown(
                "image_provider_response_invalid",
                "SiliconFlow success response omitted the image receipt.",
            )
        image_data = results[0]
        image_url = image_data.get("url")
        if image_url:
            return await _download_generated_image(image_url, client)

        b64 = image_data.get("b64_json")
        if b64:
            return _decode_generated_image_base64(b64)

        _image_generation_unknown(
            "image_provider_response_invalid",
            "SiliconFlow success response omitted the image receipt.",
        )


async def _generate_image_openai(
    api_key: str, model: str, base_url: str, prompt: str, size: str
) -> bytes:
    """Generate image via OpenAI GPT Image API.

    Requests b64_json format to avoid dealing with URL expiry.
    """
    import httpx
    url = f"{base_url.rstrip('/')}/images/generations"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model,
        "prompt": prompt,
        "size": size,
        "n": 1,
        "response_format": "b64_json",
    }

    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(url, json=payload, headers=headers)
        _settle_image_provider_status("OpenAI", resp.status_code)
        try:
            data = resp.json()
        except Exception:
            _image_generation_unknown(
                "image_provider_response_invalid",
                "OpenAI returned an unreadable success response.",
            )
        if not isinstance(data, Mapping):
            _image_generation_unknown(
                "image_provider_response_invalid",
                "OpenAI returned an invalid success response.",
            )

        results = data.get("data")
        if not isinstance(results, list) or not results or not isinstance(
            results[0], Mapping
        ):
            _image_generation_unknown(
                "image_provider_response_invalid",
                "OpenAI success response omitted the image receipt.",
            )
        image_data = results[0]
        b64 = image_data.get("b64_json")
        if b64:
            return _decode_generated_image_base64(b64)

        # Fallback: try URL
        image_url = image_data.get("url")
        if image_url:
            return await _download_generated_image(image_url, client)

        _image_generation_unknown(
            "image_provider_response_invalid",
            "OpenAI success response omitted the image receipt.",
        )


def _json_path_get(data: Any, path: str) -> Any:
    """Read a simple dotted JSON path, with numeric list indexes."""
    if not path:
        return None

    current: Any = data
    for raw_part in path.split("."):
        part = raw_part.strip()
        if not part:
            continue
        if isinstance(current, list):
            if not part.isdigit():
                return None
            index = int(part)
            if index >= len(current):
                return None
            current = current[index]
        elif isinstance(current, dict):
            if part not in current:
                return None
            current = current[part]
        else:
            return None
    return current


def _render_json_template(template_json: str, variables: dict[str, str]) -> dict:
    """Parse JSON first, then replace placeholders inside string values.

    This avoids corrupting JSON when a prompt contains quotes, newlines, or
    other characters that need escaping.
    """
    template_text = template_json.strip()
    parse_errors: list[str] = []

    candidates = [template_text]
    normalized_quotes = (
        template_text
        .replace("\u201c", '"')
        .replace("\u201d", '"')
        .replace("\u2018", "'")
        .replace("\u2019", "'")
    )
    if normalized_quotes != template_text:
        candidates.append(normalized_quotes)

    # Users often paste a JSON example copied from a string literal, leaving
    # escaped quotes like { \"model\": \"{model}\" }. Treat that as JSON too.
    for text in list(candidates):
        if '\\"' in text:
            candidates.append(text.replace('\\"', '"'))

    template = None
    for text in candidates:
        try:
            parsed = json.loads(text)
            if isinstance(parsed, str):
                parsed = json.loads(parsed)
            template = parsed
            break
        except Exception as e:
            parse_errors.append(str(e))

    if template is None:
        detail = parse_errors[-1] if parse_errors else "unknown parse error"
        raise ValueError(detail)

    def render(value: Any) -> Any:
        if isinstance(value, str):
            rendered = value
            for key, replacement in variables.items():
                rendered = rendered.replace("{" + key + "}", replacement)
            return rendered
        if isinstance(value, list):
            return [render(item) for item in value]
        if isinstance(value, dict):
            return {key: render(item) for key, item in value.items()}
        return value

    rendered = render(template)
    if not isinstance(rendered, dict):
        raise ValueError("Request body template must be a JSON object.")
    return rendered


def _json_structure_preview(data: Any, depth: int = 0) -> Any:
    if depth > 4:
        return "..."
    if isinstance(data, dict):
        return {k: _json_structure_preview(v, depth + 1) for k, v in list(data.items())[:12]}
    if isinstance(data, list):
        preview = [_json_structure_preview(item, depth + 1) for item in data[:2]]
        if len(data) > 2:
            preview.append(f"... {len(data)} items total")
        return preview
    if isinstance(data, str):
        if data.startswith("data:image"):
            return f"data:image... len={len(data)}"
        if len(data) > 160:
            return data[:160] + "..."
    return data


def _find_first_image_reference(data: Any) -> Any:
    common_paths = [
        "choices.0.message.images.0.image_url.url",
        "choices.0.message.images.0.image_url",
        "data.0.b64_json",
        "data.0.url",
        "output.0.content.0.image_url",
        "output.0.content.0.image_base64",
    ]
    for path in common_paths:
        value = _json_path_get(data, path)
        if value:
            return value

    def walk(value: Any) -> Any:
        if isinstance(value, dict):
            for key in ("url", "b64_json", "image_url", "image_base64"):
                nested = value.get(key)
                if isinstance(nested, str) and nested:
                    return nested
                if isinstance(nested, dict):
                    found = walk(nested)
                    if found:
                        return found
            for nested in value.values():
                found = walk(nested)
                if found:
                    return found
        elif isinstance(value, list):
            for item in value:
                found = walk(item)
                if found:
                    return found
        elif isinstance(value, str) and (
            value.startswith("data:image")
            or value.startswith("http://")
            or value.startswith("https://")
        ):
            return value
        return None

    return walk(data)


async def _custom_image_reference_to_bytes(image_ref: Any, client: Any) -> bytes:
    if isinstance(image_ref, dict):
        image_ref = image_ref.get("url") or image_ref.get("b64_json") or image_ref.get("image_base64")

    if not isinstance(image_ref, str) or not image_ref:
        _image_generation_unknown(
            "image_provider_response_invalid",
            "The custom image response did not contain a usable image receipt.",
        )

    if image_ref.startswith("data:image"):
        metadata, separator, encoded = image_ref.partition(",")
        if not separator or ";base64" not in metadata.lower() or not encoded:
            _image_generation_unknown(
                "image_result_invalid",
                "The custom image data URL was invalid.",
            )
        return _decode_generated_image_base64(encoded)

    if image_ref.startswith("http://") or image_ref.startswith("https://"):
        return await _download_generated_image(image_ref, client)

    return _decode_generated_image_base64(image_ref)


async def _generate_image_custom_api(
    api_key: str,
    model: str,
    base_url: str,
    endpoint_path: str,
    request_body_template_json: str,
    response_image_path: str,
    extra_headers_json: str,
    timeout_seconds: int | str,
    prompt: str,
    size: str,
) -> bytes:
    """Generate image via a configurable gateway API.

    The default request/response shape supports TokenRouter and OpenRouter:
    POST /chat/completions with image/text modalities, image returned in
    choices.0.message.images.0.image_url.url as a data URL.
    """
    import httpx

    if not isinstance(base_url, str) or not base_url.strip():
        _image_generation_failure(
            "image_configuration_invalid",
            "Custom image API base_url is not configured.",
        )
    if not isinstance(model, str) or not model.strip():
        _image_generation_failure(
            "image_configuration_invalid",
            "Custom image API model is not configured.",
        )

    try:
        timeout = int(timeout_seconds or 120)
    except (TypeError, ValueError):
        _image_generation_failure(
            "image_configuration_invalid",
            "Custom image API timeout_seconds must be an integer.",
        )
    if timeout < 1 or timeout > 600:
        _image_generation_failure(
            "image_configuration_invalid",
            "Custom image API timeout_seconds must be between 1 and 600.",
        )
    endpoint = endpoint_path or "/chat/completions"
    if endpoint.startswith("http://") or endpoint.startswith("https://"):
        url = endpoint
    else:
        url = f"{base_url.rstrip('/')}/{endpoint.lstrip('/')}"
    if not _image_public_http_url(url):
        _image_generation_failure(
            "image_configuration_invalid",
            "Custom image API endpoint must be a valid public HTTP(S) URL.",
        )

    variables = {"prompt": prompt, "size": size, "model": model}
    if request_body_template_json.strip():
        try:
            payload = _render_json_template(request_body_template_json, variables)
        except Exception:
            _image_generation_failure(
                "image_configuration_invalid",
                "Custom image request_body_template_json is invalid.",
            )
    else:
        payload = {
            "model": model,
            "messages": [{"role": "user", "content": prompt}],
            "modalities": ["image", "text"],
            "stream": False,
        }

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    if extra_headers_json.strip():
        try:
            extra_headers = json.loads(extra_headers_json)
        except Exception:
            _image_generation_failure(
                "image_configuration_invalid",
                "Custom image extra_headers_json is invalid.",
            )
        if not isinstance(extra_headers, dict):
            _image_generation_failure(
                "image_configuration_invalid",
                "Custom image extra_headers_json must be a JSON object.",
            )
        headers.update({str(k): str(v) for k, v in extra_headers.items() if v is not None})

    async with httpx.AsyncClient(timeout=timeout) as client:
        resp = await client.post(url, json=payload, headers=headers)
        _settle_image_provider_status("custom", resp.status_code)

        try:
            data = resp.json()
        except Exception:
            _image_generation_unknown(
                "image_provider_response_invalid",
                "The custom image API returned an unreadable success response.",
            )
        if not isinstance(data, Mapping):
            _image_generation_unknown(
                "image_provider_response_invalid",
                "The custom image API returned an invalid success response.",
            )

        image_ref = _json_path_get(data, response_image_path) if response_image_path else None
        if not image_ref:
            image_ref = _find_first_image_reference(data)
        if not image_ref:
            _image_generation_unknown(
                "image_provider_response_invalid",
                "The custom image API success response omitted the image receipt.",
            )

        return await _custom_image_reference_to_bytes(image_ref, client)


async def _generate_image_google(
    api_key: str, model: str, base_url: str, prompt: str, size: str
) -> bytes:
    """Generate image via Google Gemini Native Image API (Nano Banana) or Vertex AI.

    Uses the Gemini generateContent endpoint with responseModalities=["IMAGE"].
    Converts WxH size to aspect ratio format (e.g. 1024x1024 -> 1:1).
    Extracts the generated image from inlineData in the response parts.
    """
    import httpx
    url = f"{base_url.rstrip('/')}/models/{model}:generateContent"

    # Convert WxH size to aspect ratio for Gemini API
    # Supported: 1:1, 3:4, 4:3, 9:16, 16:9
    size_to_ratio = {
        "1024x1024": "1:1",
        "768x1024": "3:4",
        "1024x768": "4:3",
        "768x1366": "9:16",
        "1366x768": "16:9",
        "1024x1536": "3:4",
        "1536x1024": "4:3",
    }
    aspect_ratio = size_to_ratio.get(size, "1:1")

    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "responseModalities": ["IMAGE"],
            "imageConfig": {
                "aspectRatio": aspect_ratio,
            },
        },
    }

    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(
            url,
            json=payload,
            headers={
                "Content-Type": "application/json",
                "x-goog-api-key": api_key,
            },
        )
        _settle_image_provider_status("Google", resp.status_code)
        try:
            data = resp.json()
        except Exception:
            _image_generation_unknown(
                "image_provider_response_invalid",
                "Google returned an unreadable success response.",
            )
        if not isinstance(data, Mapping):
            _image_generation_unknown(
                "image_provider_response_invalid",
                "Google returned an invalid success response.",
            )

        # Extract image from response candidates -> content -> parts
        candidates = data.get("candidates", [])
        if not isinstance(candidates, list) or not candidates:
            _image_generation_unknown(
                "image_provider_response_invalid",
                "Google success response omitted the image receipt.",
            )

        first_candidate = candidates[0]
        if not isinstance(first_candidate, Mapping):
            _image_generation_unknown(
                "image_provider_response_invalid",
                "Google returned an invalid image receipt.",
            )
        content = first_candidate.get("content")
        parts = content.get("parts", []) if isinstance(content, Mapping) else []
        if not isinstance(parts, list):
            _image_generation_unknown(
                "image_provider_response_invalid",
                "Google returned an invalid image receipt.",
            )
        for part in parts:
            if not isinstance(part, Mapping):
                continue
            inline_data = part.get("inlineData")
            if isinstance(inline_data, Mapping):
                return _decode_generated_image_base64(inline_data.get("data"))

        _image_generation_unknown(
            "image_provider_response_invalid",
            "Google success response omitted the image receipt.",
        )


# ─── Feishu Helper ────────────────────────────────────────────────────────────

async def _get_feishu_token(agent_id: uuid.UUID) -> tuple[str, str] | None:
    """Get (app_id, app_access_token) for the agent's configured Feishu channel."""
    import httpx
    from app.models.channel_config import ChannelConfig

    async with async_session() as db:
        result = await db.execute(
            select(ChannelConfig).where(
                ChannelConfig.agent_id == agent_id,
                ChannelConfig.channel_type == "feishu",
                ChannelConfig.is_configured == True,
            )
        )
        config = result.scalar_one_or_none()

    if not config or not config.app_id or not config.app_secret:
        return None

    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal",
            json={"app_id": config.app_id, "app_secret": config.app_secret},
        )
        token = resp.json().get("tenant_access_token", "")

    return (config.app_id, token) if token else None


async def _get_agent_calendar_id(token: str) -> tuple[str | None, str | None]:
    """Get (calendar_id, error_msg) for the agent app's primary calendar.

    Returns (calendar_id, None) on success, or (None, human_readable_error) on failure.
    """
    import httpx
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/calendar/v4/calendars/primary",
            headers={"Authorization": f"Bearer {token}"},
        )
    data = resp.json()
    code = data.get("code", -1)
    if code == 0:
        cals = data.get("data", {}).get("calendars", [])
        if cals:
            cal_id = cals[0].get("calendar", {}).get("calendar_id")
            return cal_id, None
        return None, "日历列表为空，请确认应用有 calendar:calendar 权限并已发布新版本"
    if code == 99991672:
        return None, (
            "❌ 飞书日历权限未开通（错误码 99991672）\n\n"
            "请在飞书开放平台为应用 cli_a9257c5136781ceb 开通以下权限并发布新版本：\n"
            "• calendar:calendar:readonly（应用身份权限）\n"
            "• calendar:calendar.event:create（应用身份权限）\n"
            "• calendar:calendar.event:read（用户身份权限）\n"
            "• calendar:calendar.event:update（用户身份权限）\n"
            "• calendar:calendar.event:delete（用户身份权限）\n\n"
            "开通步骤：飞书开放平台 → 权限管理 → 批量导入权限 → 添加以上权限 → 创建版本 → 确认发布"
        )
    return None, f"获取日历 ID 失败：{data.get('msg')} (code {code})"


async def _feishu_resolve_open_id(token: str, email: str) -> str | None:
    """Resolve a user's open_id from their email."""
    import httpx
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/contact/v3/users/batch_get_id",
            json={"emails": [email]},
            headers={"Authorization": f"Bearer {token}"},
            params={"user_id_type": "open_id"},
        )
    data = resp.json()
    if data.get("code") != 0:
        return None
    for u in data.get("data", {}).get("user_list", []):
        oid = u.get("user_id")
        if oid:
            return oid
    return None


def _iso_to_ts(iso_str: str) -> float:
    """Convert ISO 8601 string to Unix timestamp."""
    from datetime import datetime as _dt
    for fmt in ("%Y-%m-%dT%H:%M:%S%z", "%Y-%m-%dT%H:%M:%SZ", "%Y-%m-%dT%H:%M:%S"):
        try:
            if iso_str.endswith("Z"):
                d = _dt.fromisoformat(iso_str.replace("Z", "+00:00"))
            else:
                d = _dt.strptime(iso_str, fmt)
            return d.timestamp()
        except ValueError:
            continue
    raise ValueError(f"Cannot parse datetime: {iso_str!r}")


async def _get_feishu_credentials(agent_id: uuid.UUID) -> tuple[str, str]:
    """Retrieve Feishu app_id and app_secret for an agent.
    1. Try Agent-specific ChannelConfig
    2. Fallback to global settings (.env)
    """
    from app.models.channel_config import ChannelConfig
    from app.config import get_settings

    settings = get_settings()
    app_id = settings.FEISHU_APP_ID
    app_secret = settings.FEISHU_APP_SECRET

    try:
        async with async_session() as db:
            result = await db.execute(
                select(ChannelConfig).where(ChannelConfig.agent_id == agent_id, ChannelConfig.channel_type == "feishu")
            )
            config = result.scalar_one_or_none()
            if config and config.app_id and config.app_secret:
                app_id = config.app_id
                app_secret = config.app_secret
    except Exception:
        pass

    return app_id, app_secret


async def _get_feishu_tenant_doc_url(tenant_token: str, doc_token: str, doc_type: str = "docx") -> str:
    """Build a user-accessible document URL using the tenant's actual domain.

    The API gateway (open.feishu.cn) cannot serve user documents - we must use
    the tenant's own domain (e.g. xxx.feishu.cn or xxx.larksuite.com).
    Falls back to generating a search link if the tenant domain cannot be resolved.

    Args:
        tenant_token: A valid tenant_access_token.
        doc_token:    The document_id (docx) or wiki node token.
        doc_type:     'docx' or 'wiki' - controls the URL path prefix.
    Returns:
        A fully-formed URL string.
    """
    import httpx
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(
                "https://open.feishu.cn/open-apis/tenant/v2/tenant/query",
                headers={"Authorization": f"Bearer {tenant_token}"},
            )
        data = resp.json()
        if data.get("code") == 0:
            domain = data.get("data", {}).get("tenant", {}).get("domain", "")
            if domain:
                return f"https://{domain}/{doc_type}/{doc_token}"
    except Exception:
        pass
    # Fallback: construct a search URL so the user can locate the document
    return f"https://feishu.cn/{doc_type}/{doc_token}"




async def _get_feishu_bitable_url(tenant_token: str, app_token: str, table_id: str = "") -> str:
    """Build a user-accessible Bitable URL using the tenant's actual domain.

    Constructs https://{tenant_domain}/base/{app_token}?table={table_id}
    Falls back to https://feishu.cn/base/{app_token} if domain resolution fails.

    Args:
        tenant_token: A valid tenant_access_token.
        app_token:    The Bitable app token.
        table_id:     Optional table ID to deep-link to a specific sheet.
    Returns:
        A fully-formed URL string.
    """
    import httpx
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(
                "https://open.feishu.cn/open-apis/tenant/v2/tenant/query",
                headers={"Authorization": f"Bearer {tenant_token}"},
            )
        data = resp.json()
        if data.get("code") == 0:
            domain = data.get("data", {}).get("tenant", {}).get("domain", "")
            if domain:
                base_url = f"https://{domain}/base/{app_token}"
                if table_id:
                    base_url += f"?table={table_id}"
                return base_url
    except Exception:
        pass
    # Fallback
    base_url = f"https://feishu.cn/base/{app_token}"
    if table_id:
        base_url += f"?table={table_id}"
    return base_url


def _parse_feishu_url(url: str) -> dict:
    """Parse various Feishu URLs to extract tokens.
    Supports Bitable (table, view) and Docx.
    """
    import re
    result = {}

    # Bitable URL regex: e.g., https://example.feishu.cn/base/{app_token}?table={table_id}&view={view_id}
    base_match = re.search(r'/base/([a-zA-Z0-9_]+)', url)
    if base_match:
        result['app_token'] = base_match.group(1)

    table_match = re.search(r'table=([a-zA-Z0-9_]+)', url)
    if table_match:
        result['table_id'] = table_match.group(1)

    # support URL with /tblxxxxxx
    if not 'table_id' in result:
        tbl_match = re.search(r'/(tbl[a-zA-Z0-9_]+)', url)
        if tbl_match:
            result['table_id'] = tbl_match.group(1)

    view_match = re.search(r'view=([a-zA-Z0-9_]+)', url)
    if view_match:
        result['view_id'] = view_match.group(1)

    # Docx URL regex
    docx_match = re.search(r'/docx/([a-zA-Z0-9_]+)', url)
    if docx_match:
        result['document_token'] = docx_match.group(1)

    # Wiki URL regex
    wiki_match = re.search(r'/wiki/([a-zA-Z0-9_]+)', url)
    if wiki_match:
        result['wiki_token'] = wiki_match.group(1)

    return result


# ─── Feishu Bitable Tools ──────────────────────────────────────────

async def _resolve_bitable_app_token(agent_id: uuid.UUID, parsed_url: dict) -> str | None:
    app_token = parsed_url.get("app_token")
    if app_token:
        return app_token
    wiki_token = parsed_url.get("wiki_token")
    if wiki_token:
        app_id, app_secret = await _get_feishu_credentials(agent_id)
        if app_id and app_secret:
            from app.services.feishu_service import feishu_service
            token = await feishu_service.get_tenant_access_token(app_id, app_secret)
            node_info = await _feishu_wiki_get_node(wiki_token, token)
            if node_info and node_info.get("obj_token"):
                return node_info["obj_token"]
    return None

def _check_feishu_err(resp: dict) -> str | None:
    """Check Feishu API response for errors and return a user-friendly message.

    For permission-related errors, returns detailed step-by-step instructions
    guiding the user through the Feishu UI to grant the bot app access to
    the target document or Bitable.
    """
    code = resp.get("code")
    if code != 0:
        msg = str(resp.get("msg", ""))
        msg_lower = msg.lower()
        # Common Feishu permission/access error codes:
        #   99991663 - no permission to access resource
        #   99991661 - token invalid / access denied
        #   99991668 - permission denied on resource
        #   10006    - no permission
        #   91403    - Forbidden (bitable write without app access)
        #   91404    - bitable record not found (sometimes permission)
        #   1063001  - doc permission denied
        #   1063004  - doc operation forbidden
        _perm_codes = {99991663, 10006, 99991661, 99991668, 91403, 1063001, 1063004}
        _perm_keywords = ("permission", "forbidden", "no access", "access denied", "403")
        is_perm_error = code in _perm_codes or any(kw in msg_lower for kw in _perm_keywords)
        if is_perm_error:
            return (
                f"Failed: Permission denied (code: {code}, msg: {msg}). "
                "The bot app does not have access to this document/Bitable. "
                "CRITICAL INSTRUCTION for the Agent: You MUST explicitly output ALL of the following "
                "step-by-step instructions to the user in your reply. Do NOT omit or summarize any step:\n\n"
                "---\n"
                "**The bot needs to be added to this document/table as an authorized app. "
                "Please follow these steps:**\n\n"
                "1. Open the target document or Bitable in the Feishu client (web or desktop).\n"
                "2. Click the **「...」** menu button in the top-right corner of the page.\n"
                "3. In the dropdown menu, hover over **「更多」** (More) at the bottom.\n"
                "4. In the sub-menu that appears, click **「添加文档应用」** (Add Document App).\n"
                "5. In the search box, type the name of your Feishu bot app (the one bound to this Agent's channel), then click to add it.\n"
                "6. After adding, retry the same operation.\n\n"
                "If you cannot find 「添加文档应用」, it means the document owner may need to enable this option, "
                "or you can try: click **「分享」** (Share) button -> invite the bot app directly.\n"
                "---"
            )
        return f"Failed: API Error {code} - {msg}"
    return None


async def _feishu_credentials_outcome(
    agent_id: uuid.UUID,
) -> tuple[str | None, str | None, ToolExecutionOutcome | None]:
    """Resolve the locally configured Feishu app credentials."""
    try:
        app_id, app_secret = await _get_feishu_credentials(agent_id)
    except Exception as exc:
        return None, None, _feishu_read_exception_outcome(
            "bitable_credentials",
            exc,
        )
    if not (
        isinstance(app_id, str)
        and app_id.strip()
        and isinstance(app_secret, str)
        and app_secret.strip()
    ):
        return None, None, _typed_failure(
            "The Agent has no complete Feishu channel credentials.",
            "feishu_channel_not_configured",
        )
    return app_id, app_secret, None


async def _bitable_target_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
    *,
    require_table: bool,
) -> tuple[
    str | None,
    str | None,
    str | None,
    str | None,
    ToolExecutionOutcome | None,
]:
    """Resolve credentials and stable app/table IDs before Provider dispatch."""
    url = arguments.get("url")
    if not isinstance(url, str) or not url.strip():
        return None, None, None, None, _typed_failure(
            "Bitable tools require a Feishu Bitable or Wiki URL.",
            "invalid_tool_arguments",
        )
    table_argument = arguments.get("table_id")
    if table_argument is not None and not isinstance(table_argument, str):
        return None, None, None, None, _typed_failure(
            "Bitable table_id must be a string.",
            "invalid_tool_arguments",
        )
    parsed = _parse_feishu_url(url.strip())
    try:
        app_token = await _resolve_bitable_app_token(agent_id, parsed)
    except Exception as exc:
        return None, None, None, None, _feishu_read_exception_outcome(
            "bitable_target",
            exc,
        )
    table_id = (
        table_argument.strip()
        if isinstance(table_argument, str) and table_argument.strip()
        else str(parsed.get("table_id") or "")
    )
    if not isinstance(app_token, str) or not app_token.strip():
        return None, None, None, None, _typed_failure(
            "Could not resolve a Bitable app token from the supplied URL.",
            "bitable_app_token_missing",
        )
    if require_table and not table_id:
        return None, None, None, None, _typed_failure(
            "Could not resolve a Bitable table ID from the supplied arguments.",
            "bitable_table_id_missing",
        )
    app_id, app_secret, error = await _feishu_credentials_outcome(agent_id)
    return app_id, app_secret, app_token.strip(), table_id, error


def _bitable_read_data(
    response: object,
    operation: str,
) -> tuple[Mapping | None, ToolExecutionOutcome | None]:
    """Validate the common Bitable read envelope without string inference."""
    if not isinstance(response, Mapping):
        return None, _typed_failure(
            "Feishu Bitable returned an unreadable response.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    if response.get("code") != 0:
        return None, _typed_failure(
            f"Feishu rejected {operation}.",
            f"feishu_{operation}_rejected",
        )
    data = response.get("data")
    if not isinstance(data, Mapping):
        return None, _typed_failure(
            "Feishu Bitable returned an invalid data object.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    return data, None


def _bitable_write_data(
    response: object,
    operation: str,
    *,
    result_ref: str | None = None,
) -> tuple[Mapping | None, ToolExecutionOutcome | None]:
    """Validate a dispatched Bitable write response before trusting receipts."""
    if not isinstance(response, Mapping):
        return None, _typed_unknown(
            f"Feishu {operation} returned no readable receipt; reconcile first.",
            f"feishu_{operation}_outcome_unknown",
            result_ref=result_ref,
        )
    if response.get("code") != 0:
        return None, _typed_failure(
            f"Feishu rejected {operation}.",
            f"feishu_{operation}_rejected",
            result_ref=result_ref,
        )
    data = response.get("data", {})
    if not isinstance(data, Mapping):
        return None, _typed_unknown(
            f"Feishu {operation} returned an invalid receipt; reconcile first.",
            f"feishu_{operation}_receipt_invalid",
            result_ref=result_ref,
        )
    return data, None


async def _bitable_enriched_url(
    app_id: str,
    app_secret: str,
    app_token: str,
    table_id: str = "",
) -> str | None:
    """Best-effort product link enrichment after the Provider fact settles."""
    from app.services.feishu_service import feishu_service

    try:
        tenant_token = await feishu_service.get_tenant_access_token(
            app_id,
            app_secret,
        )
        if not isinstance(tenant_token, str) or not tenant_token:
            return None
        return await _get_feishu_bitable_url(
            tenant_token,
            app_token,
            table_id,
        )
    except Exception:
        return None


async def _bitable_read_outcome(
    tool_name: str,
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Execute one of the three typed Bitable reads."""
    from app.services.feishu_service import feishu_service

    if tool_name == "bitable_query_records":
        filter_info = arguments.get("filter_info", {})
        if not isinstance(filter_info, dict):
            return _typed_failure(
                "bitable_query_records filter_info must be an object.",
                "invalid_tool_arguments",
            )
        max_results_value = arguments.get("max_results", 100)
        if (
            isinstance(max_results_value, bool)
            or not isinstance(max_results_value, int)
            or max_results_value <= 0
        ):
            return _typed_failure(
                "bitable_query_records max_results must be a positive integer.",
                "invalid_tool_arguments",
            )
        max_results = min(max_results_value, 1000)
    else:
        filter_info = {}
        max_results = 0

    require_table = tool_name != "bitable_list_tables"
    (
        app_id,
        app_secret,
        app_token,
        table_id,
        target_error,
    ) = await _bitable_target_outcome(
        agent_id,
        arguments,
        require_table=require_table,
    )
    if (
        target_error is not None
        or app_id is None
        or app_secret is None
        or app_token is None
        or table_id is None
    ):
        return target_error or _typed_failure(
            "Bitable target resolution failed.",
            "bitable_target_invalid",
        )

    if tool_name == "bitable_query_records":
        records: list[Mapping] = []
        page_token: str | None = None
        seen_page_tokens: set[str] = set()
        while len(records) < max_results:
            page_size = min(100, max_results - len(records))
            try:
                response = await feishu_service.bitable_query_records(
                    app_id,
                    app_secret,
                    app_token,
                    table_id,
                    filters=filter_info,
                    page_size=page_size,
                    page_token=page_token,
                )
            except Exception as exc:
                return _feishu_read_exception_outcome(
                    "bitable_query_records",
                    exc,
                )
            data, read_error = _bitable_read_data(
                response,
                "bitable_query_records",
            )
            if read_error is not None or data is None:
                return read_error or _typed_failure(
                    "Bitable query returned no data.",
                    "feishu_bitable_query_records_response_invalid",
                )
            items = data.get("items", [])
            if not isinstance(items, list):
                return _typed_failure(
                    "Feishu Bitable returned an invalid record list.",
                    "feishu_bitable_query_records_response_invalid",
                    retryable=True,
                )
            remaining = max_results - len(records)
            records.extend(
                item
                for item in items[:remaining]
                if isinstance(item, Mapping)
            )
            if len(records) >= max_results or not bool(
                data.get("has_more", False)
            ):
                break
            next_page_token = data.get("page_token")
            if (
                not isinstance(next_page_token, str)
                or not next_page_token
                or next_page_token in seen_page_tokens
            ):
                return _typed_failure(
                    "Feishu Bitable pagination returned no new page token.",
                    "feishu_bitable_pagination_invalid",
                    retryable=True,
                )
            seen_page_tokens.add(next_page_token)
            page_token = next_page_token

        if not records:
            return _typed_success(
                "The Bitable query matched no records.",
                result_ref=f"{app_token}:{table_id}",
            )
        lines = [f"Bitable query returned {len(records)} record(s):"]
        for record in records:
            lines.append(
                f"- record_id={str(record.get('record_id') or '')}; "
                f"fields={json.dumps(record.get('fields', {}), ensure_ascii=False)}"
            )
        return _typed_success(
            "\n".join(lines),
            result_ref=f"{app_token}:{table_id}",
            metadata={"record_count": len(records)},
        )

    try:
        if tool_name == "bitable_list_tables":
            response = await feishu_service.bitable_list_tables(
                app_id,
                app_secret,
                app_token,
            )
        else:
            response = await feishu_service.bitable_list_fields(
                app_id,
                app_secret,
                app_token,
                table_id,
            )
    except Exception as exc:
        return _feishu_read_exception_outcome(tool_name, exc)
    data, read_error = _bitable_read_data(response, tool_name)
    if read_error is not None or data is None:
        return read_error or _typed_failure(
            "Bitable read returned no data.",
            f"feishu_{tool_name}_response_invalid",
        )
    items = data.get("items", [])
    if not isinstance(items, list):
        return _typed_failure(
            "Feishu Bitable returned an invalid items list.",
            f"feishu_{tool_name}_response_invalid",
            retryable=True,
        )
    link = await _bitable_enriched_url(
        app_id,
        app_secret,
        app_token,
        table_id,
    )
    if not items:
        summary = (
            "The Bitable app has no tables."
            if tool_name == "bitable_list_tables"
            else "The Bitable table has no fields."
        )
    elif tool_name == "bitable_list_tables":
        lines = [f"Bitable contains {len(items)} table(s):"]
        for item in items:
            if isinstance(item, Mapping):
                lines.append(
                    f"- {str(item.get('name') or '(untitled)')} "
                    f"(table_id={str(item.get('table_id') or '')})"
                )
        summary = "\n".join(lines)
    else:
        lines = [f"Bitable table contains {len(items)} field(s):"]
        for item in items:
            if isinstance(item, Mapping):
                lines.append(
                    f"- {str(item.get('field_name') or '(unnamed)')} "
                    f"(field_id={str(item.get('field_id') or '')}, "
                    f"type={str(item.get('type') or '')})"
                )
        summary = "\n".join(lines)
    if link:
        summary += f"\nBitable URL: {link}"
    return _typed_success(
        summary,
        result_ref=f"{app_token}:{table_id}" if table_id else app_token,
        metadata={"item_count": len(items)},
    )


async def _bitable_write_outcome(
    tool_name: str,
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Execute one Bitable write exactly once and require a stable receipt."""
    from app.services.feishu_service import feishu_service

    if tool_name == "bitable_create_app":
        name = arguments.get("name")
        folder_token = arguments.get("folder_token", "")
        if not isinstance(name, str) or not name.strip():
            return _typed_failure(
                "bitable_create_app requires name.",
                "invalid_tool_arguments",
            )
        if not isinstance(folder_token, str):
            return _typed_failure(
                "bitable_create_app folder_token must be a string.",
                "invalid_tool_arguments",
            )
        app_id, app_secret, credential_error = (
            await _feishu_credentials_outcome(agent_id)
        )
        if (
            credential_error is not None
            or app_id is None
            or app_secret is None
        ):
            return credential_error or _typed_failure(
                "Bitable credentials are unavailable.",
                "feishu_channel_not_configured",
            )
        try:
            response = await feishu_service.bitable_create_app(
                app_id,
                app_secret,
                name.strip(),
                folder_token.strip(),
            )
        except Exception as exc:
            return _feishu_write_exception_outcome(
                "bitable_create_app",
                exc,
            )
        data, write_error = _bitable_write_data(
            response,
            "bitable_create_app",
        )
        if write_error is not None or data is None:
            return write_error or _typed_unknown(
                "Bitable app creation returned no receipt; reconcile first.",
                "feishu_bitable_create_app_receipt_missing",
            )
        app = data.get("app")
        app_token = (
            str(app.get("app_token") or "")
            if isinstance(app, Mapping)
            else ""
        )
        if not app_token:
            return _typed_unknown(
                "Feishu accepted Bitable app creation but returned no app token; "
                "reconcile before any retry.",
                "feishu_bitable_create_app_receipt_missing",
            )
        link = await _bitable_enriched_url(
            app_id,
            app_secret,
            app_token,
        )
        summary = f"Created Bitable app {app_token}."
        if link:
            summary += f"\nBitable URL: {link}"
        return _typed_success(summary, result_ref=app_token)

    fields = arguments.get("fields")
    if tool_name in {"bitable_create_record", "bitable_update_record"} and not isinstance(fields, dict):
        return _typed_failure(
            f"{tool_name} fields must be an object.",
            "invalid_tool_arguments",
        )
    record_id_value = arguments.get("record_id")
    if tool_name in {"bitable_update_record", "bitable_delete_record"} and (
        not isinstance(record_id_value, str) or not record_id_value.strip()
    ):
        return _typed_failure(
            f"{tool_name} requires record_id.",
            "invalid_tool_arguments",
        )
    (
        app_id,
        app_secret,
        app_token,
        table_id,
        target_error,
    ) = await _bitable_target_outcome(
        agent_id,
        arguments,
        require_table=True,
    )
    if (
        target_error is not None
        or app_id is None
        or app_secret is None
        or app_token is None
        or table_id is None
    ):
        return target_error or _typed_failure(
            "Bitable target resolution failed.",
            "bitable_target_invalid",
        )
    requested_record_id = (
        record_id_value.strip()
        if isinstance(record_id_value, str)
        else None
    )
    try:
        if tool_name == "bitable_create_record":
            response = await feishu_service.bitable_create_record(
                app_id,
                app_secret,
                app_token,
                table_id,
                fields,
            )
        elif tool_name == "bitable_update_record":
            response = await feishu_service.bitable_update_record(
                app_id,
                app_secret,
                app_token,
                table_id,
                requested_record_id,
                fields,
            )
        else:
            response = await feishu_service.bitable_delete_record(
                app_id,
                app_secret,
                app_token,
                table_id,
                requested_record_id,
            )
    except Exception as exc:
        return _feishu_write_exception_outcome(
            tool_name,
            exc,
            result_ref=requested_record_id,
        )
    data, write_error = _bitable_write_data(
        response,
        tool_name,
        result_ref=requested_record_id,
    )
    if write_error is not None or data is None:
        return write_error or _typed_unknown(
            f"{tool_name} returned no receipt; reconcile first.",
            f"feishu_{tool_name}_receipt_missing",
            result_ref=requested_record_id,
        )

    if tool_name == "bitable_delete_record":
        receipt = requested_record_id or ""
    else:
        record_data = data.get("record")
        receipt = (
            str(record_data.get("record_id") or "")
            if isinstance(record_data, Mapping)
            else ""
        )
        if not receipt:
            return _typed_unknown(
                f"Feishu accepted {tool_name} but returned no record ID; "
                "reconcile before any retry.",
                f"feishu_{tool_name}_receipt_missing",
                result_ref=requested_record_id,
            )
        if (
            tool_name == "bitable_update_record"
            and receipt != requested_record_id
        ):
            return _typed_unknown(
                "Feishu returned a different record ID for the update; "
                "reconcile before any retry.",
                "feishu_bitable_update_record_receipt_mismatch",
                result_ref=requested_record_id,
                metadata={"returned_record_id": receipt},
            )

    link = await _bitable_enriched_url(
        app_id,
        app_secret,
        app_token,
        table_id,
    )
    summary = f"{tool_name} succeeded for record {receipt}."
    if link:
        summary += f"\nBitable URL: {link}"
    return _typed_success(
        summary,
        result_ref=receipt,
        metadata={"app_token": app_token, "table_id": table_id},
    )

async def _bitable_list_tables(agent_id: uuid.UUID, arguments: dict) -> str:
    """List all tables in a Feishu Bitable app."""
    url = arguments.get("url", "")
    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    if not app_token:
        return "Failed: Could not extract Bitable app_token from the URL (also could not resolve wiki_token)."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_list_tables(app_id, app_secret, app_token)
        err = _check_feishu_err(resp)
        if err: return err

        tables = resp.get("data", {}).get("items", [])
        if not tables:
            return "OK: No tables found in this Bitable."
        lines = [f"- {t.get('name')} (ID: {t.get('table_id')})" for t in tables]
        # Provide a user-accessible link so the user can open the Bitable directly
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        bitable_url = await _get_feishu_bitable_url(tenant_token, app_token)
        return "OK: Tables in this Bitable:\n" + "\n".join(lines) + f"\n\n🔗 多维表格链接: {bitable_url}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


async def _bitable_create_app(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new Feishu Bitable (多维表格) app.

    Calls the Bitable v1 apps API: POST /open-apis/bitable/v1/apps
    The API response includes a user-accessible URL with the tenant's own domain.
    """
    name = arguments.get("name", "").strip()
    if not name:
        return "Failed: Missing required argument 'name' — please provide a name for the new Bitable."

    folder_token = arguments.get("folder_token", "").strip()

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_create_app(app_id, app_secret, name, folder_token)
        err = _check_feishu_err(resp)
        if err:
            return err

        # API response structure: data.app.{app_token, name, url, default_table_id, folder_token}
        app_info = resp.get("data", {}).get("app", {})
        app_token = app_info.get("app_token", "")
        bitable_url = app_info.get("url", "")
        default_table_id = app_info.get("default_table_id", "")
        if not app_token:
            return f"Failed: Bitable created but could not extract app_token from response: {resp}"

        # Fallback URL resolution if the API didn't return one
        if not bitable_url:
            tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
            bitable_url = await _get_feishu_bitable_url(tenant_token, app_token)

        result = (
            f"OK: Bitable created successfully!\n"
            f"Name: {name}\n"
            f"App Token: {app_token}\n"
            f"URL: {bitable_url}"
        )
        if default_table_id:
            result += f"\nDefault Table ID: {default_table_id}"
        return result
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


async def _bitable_list_fields(agent_id: uuid.UUID, arguments: dict) -> str:
    """List all fields (columns) in a specific Bitable table."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")

    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")

    if not app_token:
        return "Failed: Could not extract Bitable app_token from the URL."
    if not table_id:
        return "Failed: table_id is required. Provide it as a parameter or include it in the URL."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_list_fields(app_id, app_secret, app_token, table_id)
        err = _check_feishu_err(resp)
        if err: return err

        fields = resp.get("data", {}).get("items", [])
        if not fields:
            return "OK: No fields found in this table."
        lines = [f"- {f.get('field_name')} (type: {f.get('type')}, ID: {f.get('field_id')})" for f in fields]
        return "OK: Fields in this table:\n" + "\n".join(lines)
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _bitable_query_records(agent_id: uuid.UUID, arguments: dict) -> str:
    """Query records (rows) from a Bitable table, with optional FQL filter."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    filter_info = arguments.get("filter_info", "")
    max_results = arguments.get("max_results", 100)

    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")

    if not app_token or not table_id:
        return "Failed: Could not resolve app_token or table_id from the provided parameters/URL."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        import json
        filters_dict = {}
        if isinstance(filter_info, dict):
            filters_dict = filter_info
        elif isinstance(filter_info, str) and filter_info.strip():
            try:
                filters_dict = json.loads(filter_info)
            except json.JSONDecodeError:
                pass

        resp = await feishu_service.bitable_query_records(app_id, app_secret, app_token, table_id, filters_dict)
        err = _check_feishu_err(resp)
        if err:
            return err

        records = resp.get("data", {}).get("items", [])
        if not records:
            return "OK: No matching records found."

        lines = []
        for r in records[:max_results]:
            lines.append(f"Record {r.get('record_id')}: {json.dumps(r.get('fields', {}), ensure_ascii=False)}")
        return "OK: Query results:\n" + "\n".join(lines)
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _bitable_create_record(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new record (row) in a Bitable table."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    fields_value = arguments.get("fields", {})

    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")

    if not app_token or not table_id:
        return "Failed: Could not resolve app_token or table_id from the provided parameters/URL."

    import json
    if isinstance(fields_value, dict):
        fields = dict(fields_value)
    elif isinstance(fields_value, str):
        try:
            fields = json.loads(fields_value)
        except json.JSONDecodeError:
            return "Failed: The 'fields' parameter is not valid JSON."
        if not isinstance(fields, dict):
            return "Failed: The 'fields' parameter must be an object."
    else:
        return "Failed: The 'fields' parameter must be an object."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_create_record(app_id, app_secret, app_token, table_id, fields)
        err = _check_feishu_err(resp)
        if err:
            return err

        record = resp.get("data", {}).get("record", {})
        # Provide a user-accessible link so they can verify the new row in the table
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        bitable_url = await _get_feishu_bitable_url(tenant_token, app_token, table_id)
        return (
            f"OK: Record created. Record ID: {record.get('record_id')}\n"
            f"Fields: {json.dumps(record.get('fields', {}), ensure_ascii=False)}\n"
            f"🔗 多维表格链接: {bitable_url}"
        )
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _bitable_update_record(agent_id: uuid.UUID, arguments: dict) -> str:
    """Update an existing record in a Bitable table by record_id."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    record_id = arguments.get("record_id", "")
    fields_value = arguments.get("fields", {})

    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")

    if not app_token or not table_id or not record_id:
        return "Failed: Missing required parameters. Need app_token (from URL), table_id, and record_id."

    import json
    if isinstance(fields_value, dict):
        fields = dict(fields_value)
    elif isinstance(fields_value, str):
        try:
            fields = json.loads(fields_value)
        except json.JSONDecodeError:
            return "Failed: The 'fields' parameter is not valid JSON."
        if not isinstance(fields, dict):
            return "Failed: The 'fields' parameter must be an object."
    else:
        return "Failed: The 'fields' parameter must be an object."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_update_record(app_id, app_secret, app_token, table_id, record_id, fields)
        err = _check_feishu_err(resp)
        if err:
            return err

        record = resp.get("data", {}).get("record", {})
        # Provide a user-accessible link so they can verify the updated row
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        bitable_url = await _get_feishu_bitable_url(tenant_token, app_token, table_id)
        return (
            f"OK: Record updated. Record ID: {record.get('record_id')}\n"
            f"Fields: {json.dumps(record.get('fields', {}), ensure_ascii=False)}\n"
            f"🔗 多维表格链接: {bitable_url}"
        )
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _bitable_delete_record(agent_id: uuid.UUID, arguments: dict) -> str:
    """Delete a record from a Bitable table by record_id."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    record_id = arguments.get("record_id", "")

    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")

    if not app_token or not table_id or not record_id:
        return "Failed: Missing required parameters. Need app_token (from URL), table_id, and record_id."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_delete_record(app_id, app_secret, app_token, table_id, record_id)
        err = _check_feishu_err(resp)
        if err: return err

        # Provide a user-accessible link so they can verify the deletion
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        bitable_url = await _get_feishu_bitable_url(tenant_token, app_token, table_id)
        return f"OK: Record {record_id} deleted successfully.\n🔗 多维表格链接: {bitable_url}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


# ─── Feishu Document Tools ──────────────────────────────────────────

async def _resolve_docx_document_token(agent_id: uuid.UUID, parsed_url: dict) -> str | None:
    doc_token = parsed_url.get("document_token")
    if doc_token:
        return doc_token
    wiki_token = parsed_url.get("wiki_token")
    if wiki_token:
        app_id, app_secret = await _get_feishu_credentials(agent_id)
        if app_id and app_secret:
            from app.services.feishu_service import feishu_service
            token = await feishu_service.get_tenant_access_token(app_id, app_secret)
            node_info = await _feishu_wiki_get_node(wiki_token, token)
            if node_info and node_info.get("obj_token"):
                return node_info["obj_token"]
    return None

async def _feishu_read_doc(agent_id: uuid.UUID, arguments: dict) -> str:
    """Read full text content of a Feishu Docx."""
    url = arguments.get("url", "")
    parsed = _parse_feishu_url(url)
    doc_token = await _resolve_docx_document_token(agent_id, parsed)
    if not doc_token:
        return "Failed: Could not extract Document token from the URL."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.read_feishu_doc(app_id, app_secret, doc_token)
        err = _check_feishu_err(resp)
        if err: return err

        content = resp.get("data", {}).get("content", "")
        if not content:
            return "OK: Document is empty or content is unavailable."
        return f"OK: Document Content:\n{content}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _feishu_create_doc(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new blank Feishu Docx."""
    title = arguments.get("title", "Untitled Document")
    folder_token = arguments.get("folder_token", "")

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.create_feishu_doc(app_id, app_secret, folder_token or None, title)
        err = _check_feishu_err(resp)
        if err: return err

        doc = resp.get("data", {}).get("document", {})
        doc_id = doc.get("document_id")
        # Get the tenant's actual domain (open.feishu.cn is the API gateway, not for users)
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        url = await _get_feishu_tenant_doc_url(tenant_token, doc_id)
        return f"OK: Document created perfectly. Document ID: {doc_id}\nURL: {url}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _feishu_append_doc(agent_id: uuid.UUID, arguments: dict) -> str:
    """Append text to the bottom of a Feishu Docx."""
    url = arguments.get("url", "")
    content = arguments.get("content", "")
    if not content:
        return "Failed: Content to append cannot be empty."

    parsed = _parse_feishu_url(url)
    doc_token = await _resolve_docx_document_token(agent_id, parsed)
    if not doc_token:
        return "Failed: Could not extract Document token from the URL."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    try:
        # Feishu uses the document_id as the root block_id to append entirely to the document
        resp = await feishu_service.append_feishu_doc(app_id, app_secret, doc_token, content)
        err = _check_feishu_err(resp)
        if err: return err

        return "OK: Content appended successfully to the end of the document."
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

# ─── Feishu Wiki Tools ───────────────────────────────────────────────────────

async def _feishu_wiki_get_node(token_str: str, auth_token: str) -> dict | None:
    """Call wiki get_node API to resolve a wiki node token → {obj_token, space_id, has_child, title}.
    Returns None if the token is not a wiki node."""
    import httpx
    async with httpx.AsyncClient(timeout=5) as client:
        r = await client.get(
            "https://open.feishu.cn/open-apis/wiki/v2/spaces/get_node",
            headers={"Authorization": f"Bearer {auth_token}"},
            params={"token": token_str, "obj_type": "wiki"},
        )
    d = r.json()
    if d.get("code") != 0:
        return None
    node = d.get("data", {}).get("node", {})
    return {
        "obj_token": node.get("obj_token", ""),
        "space_id": node.get("origin_space_id", node.get("space_id", "")),
        "has_child": node.get("has_child", False),
        "title": node.get("title", ""),
        "node_token": node.get("node_token", token_str),
    }


def _feishu_error_is_known_rejection(exc: Exception) -> bool:
    """Return whether Feishu conclusively rejected a provider request."""
    from app.services.feishu_service import FeishuAPIError

    if not isinstance(exc, FeishuAPIError):
        return False
    return (
        exc.code not in {None, 0}
        or (
            exc.http_status is not None
            and 400 <= exc.http_status < 500
        )
    )


def _feishu_read_exception_outcome(
    operation: str,
    exc: Exception,
) -> ToolExecutionOutcome:
    """Classify a Feishu read without converting display text into facts."""
    import httpx
    from app.services.feishu_service import FeishuAPIError

    if _feishu_error_is_known_rejection(exc):
        return _typed_failure(
            f"Feishu rejected {operation}.",
            f"feishu_{operation}_rejected",
        )
    retryable = isinstance(
        exc,
        (asyncio.TimeoutError, httpx.TransportError, FeishuAPIError),
    )
    return _typed_failure(
        f"Feishu {operation} failed before a durable result was read: "
        f"{type(exc).__name__}.",
        f"feishu_{operation}_failed",
        retryable=retryable,
    )


def _feishu_write_exception_outcome(
    operation: str,
    exc: Exception,
    *,
    result_ref: str | None = None,
    metadata: dict | None = None,
) -> ToolExecutionOutcome:
    """Classify a Feishu write after its business request was dispatched."""
    if _feishu_error_is_known_rejection(exc):
        return _typed_failure(
            f"Feishu rejected {operation}.",
            f"feishu_{operation}_rejected",
            result_ref=result_ref,
            metadata=metadata,
        )
    return _typed_unknown(
        f"Feishu {operation} may have taken effect; reconcile before any retry.",
        f"feishu_{operation}_outcome_unknown",
        result_ref=result_ref,
        metadata=metadata,
    )


async def _feishu_access_token_outcome(
    agent_id: uuid.UUID,
) -> tuple[str | None, ToolExecutionOutcome | None]:
    """Resolve execution credentials locally, then obtain a provider token."""
    from app.services.feishu_service import feishu_service

    try:
        app_id, app_secret = await _get_feishu_credentials(agent_id)
    except Exception as exc:
        return None, _feishu_read_exception_outcome("credentials", exc)
    if not (
        isinstance(app_id, str)
        and app_id.strip()
        and isinstance(app_secret, str)
        and app_secret.strip()
    ):
        return None, _typed_failure(
            "The Agent has no complete Feishu channel credentials.",
            "feishu_channel_not_configured",
        )
    try:
        token = await feishu_service.get_tenant_access_token(
            app_id,
            app_secret,
        )
    except Exception as exc:
        return None, _feishu_read_exception_outcome("token", exc)
    if not isinstance(token, str) or not token.strip():
        return None, _typed_failure(
            "Feishu did not return a tenant access token.",
            "feishu_token_rejected",
        )
    return token, None


async def _feishu_calendar_context_outcome(
    agent_id: uuid.UUID,
) -> tuple[str | None, str | None, ToolExecutionOutcome | None]:
    """Resolve the Bot primary calendar before dispatching an event operation."""
    token, error = await _feishu_access_token_outcome(agent_id)
    if error is not None or token is None:
        return None, None, error
    try:
        calendar_id, calendar_error = await _get_agent_calendar_id(token)
    except Exception as exc:
        return None, None, _feishu_read_exception_outcome(
            "calendar_primary",
            exc,
        )
    if not isinstance(calendar_id, str) or not calendar_id.strip():
        return None, None, _typed_failure(
            calendar_error or "Feishu Bot primary calendar is unavailable.",
            "feishu_calendar_unavailable",
        )
    return token, calendar_id, None


async def _feishu_wiki_list_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """List Wiki children with provider pagination and a fixed depth bound."""
    import httpx
    from app.services.feishu_service import feishu_service

    node_token = arguments.get("node_token")
    if not isinstance(node_token, str) or not node_token.strip():
        return _typed_failure(
            "feishu_wiki_list requires node_token.",
            "invalid_tool_arguments",
        )
    node_token = node_token.strip()
    recursive = bool(arguments.get("recursive", False))

    token, error = await _feishu_access_token_outcome(agent_id)
    if error is not None or token is None:
        return error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    try:
        node_info = await _feishu_wiki_get_node(node_token, token)
    except Exception as exc:
        return _feishu_read_exception_outcome("wiki_node", exc)
    if not isinstance(node_info, Mapping):
        return _typed_failure(
            f"Feishu rejected or could not resolve Wiki node {node_token}.",
            "feishu_wiki_node_rejected",
        )
    space_id = node_info.get("space_id")
    if not isinstance(space_id, str) or not space_id.strip():
        return _typed_failure(
            f"Wiki node {node_token} has no stable space ID.",
            "feishu_wiki_space_missing",
        )

    pages: list[dict] = []
    visited_parents: set[str] = set()

    async with httpx.AsyncClient(timeout=15) as client:
        async def collect(parent_token: str, depth: int) -> ToolExecutionOutcome | None:
            if parent_token in visited_parents:
                return None
            visited_parents.add(parent_token)
            provider_page_token: str | None = None
            seen_page_tokens: set[str] = set()
            children: list[dict] = []

            while True:
                params: dict[str, object] = {
                    "parent_node_token": parent_token,
                    "page_size": 50,
                }
                if provider_page_token is not None:
                    params["page_token"] = provider_page_token
                try:
                    response = await client.get(
                        f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/nodes",
                        headers={"Authorization": f"Bearer {token}"},
                        params=params,
                    )
                    data = feishu_service._parse_api_response(
                        response,
                        stage="wiki_list",
                    )
                except Exception as exc:
                    return _feishu_read_exception_outcome("wiki_list", exc)

                body = data.get("data")
                if not isinstance(body, Mapping):
                    return _typed_failure(
                        "Feishu Wiki returned an invalid data object.",
                        "feishu_wiki_response_invalid",
                        retryable=True,
                    )
                items = body.get("items", [])
                if not isinstance(items, list):
                    return _typed_failure(
                        "Feishu Wiki returned an invalid items list.",
                        "feishu_wiki_response_invalid",
                        retryable=True,
                    )
                for item in items:
                    if not isinstance(item, Mapping):
                        continue
                    child_token = str(item.get("node_token") or "")
                    entry = {
                        "title": str(item.get("title") or "(untitled)"),
                        "node_token": child_token,
                        "obj_token": str(item.get("obj_token") or ""),
                        "has_child": bool(item.get("has_child", False)),
                        "depth": depth,
                    }
                    pages.append(entry)
                    children.append(entry)

                if not bool(body.get("has_more", False)):
                    break
                next_page_token = body.get("page_token")
                if (
                    not isinstance(next_page_token, str)
                    or not next_page_token
                    or next_page_token in seen_page_tokens
                ):
                    return _typed_failure(
                        "Feishu Wiki pagination did not provide a new page token.",
                        "feishu_wiki_pagination_invalid",
                        retryable=True,
                    )
                seen_page_tokens.add(next_page_token)
                provider_page_token = next_page_token

            if recursive and depth < 2:
                for child in children:
                    child_token = child["node_token"]
                    if child["has_child"] and child_token:
                        child_error = await collect(child_token, depth + 1)
                        if child_error is not None:
                            return child_error
            return None

        collection_error = await collect(node_token, 0)

    if collection_error is not None:
        return collection_error
    if not pages:
        return _typed_success(
            f"Wiki node {node_token} has no child pages.",
            result_ref=node_token,
        )
    lines = [
        f"Wiki node {node_token} has {len(pages)} child page(s) in space {space_id}:"
    ]
    for page_entry in pages:
        indent = "  " * int(page_entry["depth"])
        lines.append(
            f"{indent}- {page_entry['title']} "
            f"(node_token={page_entry['node_token']}, "
            f"obj_token={page_entry['obj_token']})"
        )
    return _typed_success(
        "\n".join(lines),
        result_ref=node_token,
        metadata={"space_id": space_id, "page_count": len(pages)},
    )


def _feishu_doc_read_data(
    response: object,
    operation: str,
) -> tuple[Mapping | None, ToolExecutionOutcome | None]:
    """Validate a Feishu Doc read envelope returned by a service adapter."""
    if not isinstance(response, Mapping):
        return None, _typed_failure(
            f"Feishu {operation} returned an unreadable response.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    if response.get("code") != 0:
        return None, _typed_failure(
            f"Feishu rejected {operation}.",
            f"feishu_{operation}_rejected",
        )
    data = response.get("data")
    if not isinstance(data, Mapping):
        return None, _typed_failure(
            f"Feishu {operation} returned an invalid data object.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    return data, None


def _feishu_doc_write_data(
    response: object,
    operation: str,
    *,
    result_ref: str | None = None,
) -> tuple[Mapping | None, ToolExecutionOutcome | None]:
    """Validate a Feishu Doc/Drive write receipt without inferring from text."""
    if not isinstance(response, Mapping):
        return None, _typed_unknown(
            f"Feishu {operation} returned no readable receipt; reconcile first.",
            f"feishu_{operation}_outcome_unknown",
            result_ref=result_ref,
        )
    if response.get("code") != 0:
        return None, _typed_failure(
            f"Feishu rejected {operation}.",
            f"feishu_{operation}_rejected",
            result_ref=result_ref,
        )
    data = response.get("data", {})
    if not isinstance(data, Mapping):
        return None, _typed_unknown(
            f"Feishu {operation} returned an invalid receipt; reconcile first.",
            f"feishu_{operation}_receipt_invalid",
            result_ref=result_ref,
        )
    return data, None


async def _feishu_doc_search_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Search documents with bounded pagination and stable document tokens."""
    import httpx
    from app.services.feishu_service import feishu_service

    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "feishu_doc_search requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()

    count_value = arguments.get("count", 10)
    offset_value = arguments.get("offset", 0)
    if (
        isinstance(count_value, bool)
        or not isinstance(count_value, int)
        or isinstance(offset_value, bool)
        or not isinstance(offset_value, int)
    ):
        return _typed_failure(
            "feishu_doc_search count and offset must be integers.",
            "invalid_tool_arguments",
        )
    count = max(1, min(count_value, 50))
    offset = max(0, offset_value)

    docs_types = arguments.get("docs_types", [])
    if docs_types is None:
        docs_types = []
    valid_doc_types = {
        "doc",
        "docx",
        "sheet",
        "bitable",
        "file",
        "folder",
        "mindnote",
        "slides",
    }
    if not isinstance(docs_types, list) or any(
        not isinstance(value, str) or value not in valid_doc_types
        for value in docs_types
    ):
        return _typed_failure(
            "feishu_doc_search docs_types must contain supported file types.",
            "invalid_tool_arguments",
        )

    token, token_error = await _feishu_access_token_outcome(agent_id)
    if token_error is not None or token is None:
        return token_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )

    payload: dict[str, object] = {
        "search_key": query,
        "count": count,
        "offset": offset,
    }
    if docs_types:
        payload["docs_types"] = docs_types
    try:
        async with httpx.AsyncClient(timeout=20) as client:
            response = await client.post(
                "https://open.feishu.cn/open-apis/suite/docs-api/search/object",
                headers={
                    "Authorization": f"Bearer {token}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )
        parsed = feishu_service._parse_api_response(
            response,
            stage="doc_search",
        )
    except Exception as exc:
        return _feishu_read_exception_outcome("doc_search", exc)

    data, data_error = _feishu_doc_read_data(parsed, "doc_search")
    if data_error is not None or data is None:
        return data_error or _typed_failure(
            "Feishu document search returned no data.",
            "feishu_doc_search_response_invalid",
            retryable=True,
        )
    entities = data.get("docs_entities", [])
    if not isinstance(entities, list):
        return _typed_failure(
            "Feishu document search returned an invalid result list.",
            "feishu_doc_search_response_invalid",
            retryable=True,
        )

    normalized: list[dict[str, str]] = []
    for item in entities:
        if not isinstance(item, Mapping):
            return _typed_failure(
                "Feishu document search returned an invalid result item.",
                "feishu_doc_search_response_invalid",
                retryable=True,
            )
        docs_token = item.get("docs_token")
        if not isinstance(docs_token, str) or not docs_token.strip():
            return _typed_failure(
                "Feishu document search omitted a stable docs_token.",
                "feishu_doc_search_receipt_missing",
                retryable=True,
            )
        normalized.append(
            {
                "title": str(item.get("title") or "(untitled)"),
                "docs_type": str(item.get("docs_type") or "unknown"),
                "docs_token": docs_token.strip(),
                "owner_id": str(item.get("owner_id") or ""),
            }
        )

    if not normalized:
        return _typed_success(
            f'No Feishu documents matched "{query}".',
        )
    lines = [
        f'Feishu document search returned {len(normalized)} result(s) for "{query}":'
    ]
    for index, item in enumerate(normalized, start=offset + 1):
        lines.append(
            f"{index}. {item['title']} "
            f"(docs_type={item['docs_type']}, "
            f"docs_token={item['docs_token']}, owner_id={item['owner_id']})"
        )
    return _typed_success("\n".join(lines))


async def _feishu_doc_read_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Read an explicitly supplied ordinary Docx token without Wiki guessing."""
    from app.services.feishu_service import feishu_service

    document_token = arguments.get("document_token")
    if not isinstance(document_token, str) or not document_token.strip():
        return _typed_failure(
            "feishu_doc_read requires an explicit document_token.",
            "invalid_tool_arguments",
        )
    document_token = document_token.strip()

    max_chars_value = arguments.get("max_chars", 6000)
    if isinstance(max_chars_value, bool) or not isinstance(max_chars_value, int):
        return _typed_failure(
            "feishu_doc_read max_chars must be an integer.",
            "invalid_tool_arguments",
        )
    max_chars = max(1, min(max_chars_value, 20000))

    app_id, app_secret, credential_error = await _feishu_credentials_outcome(
        agent_id
    )
    if credential_error is not None or app_id is None or app_secret is None:
        return credential_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    try:
        response = await feishu_service.read_feishu_doc(
            app_id,
            app_secret,
            document_token,
        )
    except Exception as exc:
        return _feishu_read_exception_outcome("doc_read", exc)

    data, data_error = _feishu_doc_read_data(response, "doc_read")
    if data_error is not None or data is None:
        return data_error or _typed_failure(
            "Feishu document read returned no data.",
            "feishu_doc_read_response_invalid",
            retryable=True,
        )
    content = data.get("content")
    if not isinstance(content, str):
        return _typed_failure(
            "Feishu document read returned invalid text content.",
            "feishu_doc_read_response_invalid",
            retryable=True,
        )

    bounded_content = content[:max_chars]
    if not bounded_content:
        summary = f"Feishu document {document_token} is empty."
    else:
        summary = f"Feishu document {document_token}:\n\n{bounded_content}"
        if len(content) > max_chars:
            summary += f"\n\n[truncated to {max_chars} characters]"
    return _typed_success(summary, result_ref=document_token)


async def _feishu_doc_create_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Create one ordinary Docx and require its stable document receipt."""
    from app.services.feishu_service import feishu_service

    if any(
        legacy_name in arguments
        for legacy_name in ("wiki_space_id", "parent_node_token")
    ):
        return _typed_failure(
            "feishu_doc_create no longer accepts Wiki placement arguments; use the canonical Wiki tools instead.",
            "legacy_tool_arguments_unsupported",
        )

    title = arguments.get("title")
    if not isinstance(title, str) or not title.strip():
        return _typed_failure(
            "feishu_doc_create requires title.",
            "invalid_tool_arguments",
        )
    title = title.strip()
    folder_token = arguments.get("folder_token", "")
    if folder_token is None:
        folder_token = ""
    if not isinstance(folder_token, str):
        return _typed_failure(
            "feishu_doc_create folder_token must be a string.",
            "invalid_tool_arguments",
        )
    folder_token = folder_token.strip()

    app_id, app_secret, credential_error = await _feishu_credentials_outcome(
        agent_id
    )
    if credential_error is not None or app_id is None or app_secret is None:
        return credential_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    try:
        response = await feishu_service.create_feishu_doc(
            app_id,
            app_secret,
            folder_token or None,
            title,
        )
    except Exception as exc:
        return _feishu_write_exception_outcome("doc_create", exc)

    data, data_error = _feishu_doc_write_data(response, "doc_create")
    if data_error is not None or data is None:
        return data_error or _typed_unknown(
            "Feishu document creation returned no receipt; reconcile first.",
            "feishu_doc_create_outcome_unknown",
        )
    document = data.get("document")
    document_id = document.get("document_id") if isinstance(document, Mapping) else None
    if not isinstance(document_id, str) or not document_id.strip():
        return _typed_unknown(
            "Feishu document creation omitted document_id; reconcile first.",
            "feishu_doc_create_receipt_missing",
        )
    document_id = document_id.strip()

    document_url: str | None = None
    try:
        token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        document_url = await _get_feishu_tenant_doc_url(token, document_id)
    except Exception:
        pass
    summary = f"Created Feishu Docx {document_id} with title {title}."
    if document_url:
        summary += f" URL: {document_url}"
    return _typed_success(summary, result_ref=document_id)


async def _feishu_doc_append_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Append once after a read-only root-block preflight."""
    import httpx
    from app.services.feishu_service import feishu_service

    document_token = arguments.get("document_token")
    content = arguments.get("content")
    if not isinstance(document_token, str) or not document_token.strip():
        return _typed_failure(
            "feishu_doc_append requires an explicit document_token.",
            "invalid_tool_arguments",
        )
    if not isinstance(content, str) or not content.strip():
        return _typed_failure(
            "feishu_doc_append requires non-empty content.",
            "invalid_tool_arguments",
        )
    document_token = document_token.strip()
    content = content.strip()

    token, token_error = await _feishu_access_token_outcome(agent_id)
    if token_error is not None or token is None:
        return token_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    headers = {"Authorization": f"Bearer {token}"}
    try:
        async with httpx.AsyncClient(timeout=20) as client:
            metadata_response = await client.get(
                "https://open.feishu.cn/open-apis/docx/v1/documents/"
                f"{document_token}",
                headers=headers,
            )
        metadata_payload = feishu_service._parse_api_response(
            metadata_response,
            stage="doc_append_preflight",
        )
    except Exception as exc:
        return _feishu_read_exception_outcome("doc_append_preflight", exc)

    metadata_data = metadata_payload.get("data")
    document = (
        metadata_data.get("document")
        if isinstance(metadata_data, Mapping)
        else None
    )
    body = document.get("body") if isinstance(document, Mapping) else None
    body_block_id = body.get("block_id") if isinstance(body, Mapping) else None
    if not isinstance(body_block_id, str) or not body_block_id.strip():
        return _typed_failure(
            "Feishu document append preflight omitted the body block ID.",
            "feishu_doc_append_preflight_invalid",
            retryable=True,
        )
    body_block_id = body_block_id.strip()
    children = _markdown_to_feishu_blocks(content)

    try:
        async with httpx.AsyncClient(timeout=20) as client:
            append_response = await client.post(
                "https://open.feishu.cn/open-apis/docx/v1/documents/"
                f"{document_token}/blocks/{body_block_id}/children",
                json={"children": children},
                headers=headers,
            )
        append_payload = feishu_service._parse_api_response(
            append_response,
            stage="doc_append",
        )
    except Exception as exc:
        return _feishu_write_exception_outcome(
            "doc_append",
            exc,
            result_ref=document_token,
        )

    append_data, append_error = _feishu_doc_write_data(
        append_payload,
        "doc_append",
        result_ref=document_token,
    )
    if append_error is not None or append_data is None:
        return append_error or _typed_unknown(
            "Feishu document append returned no receipt; reconcile first.",
            "feishu_doc_append_outcome_unknown",
            result_ref=document_token,
        )
    receipt_children = append_data.get("children")
    revision = append_data.get("document_revision_id")
    first_child = (
        receipt_children[0]
        if isinstance(receipt_children, list) and receipt_children
        else None
    )
    block_id = first_child.get("block_id") if isinstance(first_child, Mapping) else None
    revision_valid = (
        isinstance(revision, (int, str))
        and not isinstance(revision, bool)
        and bool(str(revision).strip())
    )
    if not isinstance(block_id, str) or not block_id.strip() or not revision_valid:
        return _typed_unknown(
            "Feishu document append omitted block or revision receipt; reconcile first.",
            "feishu_doc_append_receipt_missing",
            result_ref=document_token,
        )
    block_id = block_id.strip()

    document_url: str | None = None
    try:
        document_url = await _get_feishu_tenant_doc_url(token, document_token)
    except Exception:
        pass
    summary = (
        f"Appended {len(children)} block(s) to Feishu document {document_token}; "
        f"block_id={block_id}, revision={revision}."
    )
    if document_url:
        summary += f" URL: {document_url}"
    return _typed_success(summary, result_ref=block_id)


async def _feishu_drive_share_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Settle each collaborator mutation independently and stop on uncertainty."""
    import httpx
    from app.services.feishu_service import feishu_service

    document_token = arguments.get("document_token")
    action = arguments.get("action")
    doc_type = arguments.get("doc_type", "docx")
    permission = arguments.get("permission", "edit")
    if not isinstance(document_token, str) or not document_token.strip():
        return _typed_failure(
            "feishu_drive_share requires document_token.",
            "invalid_tool_arguments",
        )
    if not isinstance(action, str) or action not in {"add", "remove", "list"}:
        return _typed_failure(
            "feishu_drive_share action must be add, remove, or list.",
            "invalid_tool_arguments",
        )
    valid_doc_types = {
        "docx",
        "bitable",
        "sheet",
        "doc",
        "folder",
        "mindnote",
        "slides",
    }
    if not isinstance(doc_type, str) or doc_type not in valid_doc_types:
        return _typed_failure(
            "feishu_drive_share doc_type is unsupported.",
            "invalid_tool_arguments",
        )
    if not isinstance(permission, str) or permission not in {
        "view",
        "edit",
        "full_access",
    }:
        return _typed_failure(
            "feishu_drive_share permission is unsupported.",
            "invalid_tool_arguments",
        )
    document_token = document_token.strip()

    member_names = arguments.get("member_names", [])
    member_open_ids = arguments.get("member_open_ids", [])
    if member_names is None:
        member_names = []
    if member_open_ids is None:
        member_open_ids = []
    if not isinstance(member_names, list) or any(
        not isinstance(value, str) or not value.strip()
        for value in member_names
    ):
        return _typed_failure(
            "feishu_drive_share member_names must contain non-empty strings.",
            "invalid_tool_arguments",
        )
    if not isinstance(member_open_ids, list) or any(
        not isinstance(value, str) or not value.strip()
        for value in member_open_ids
    ):
        return _typed_failure(
            "feishu_drive_share member_open_ids must contain non-empty strings.",
            "invalid_tool_arguments",
        )
    if member_names:
        return _typed_failure(
            "Name lookup is not part of the typed Doc/Drive adapter; provide member_open_ids.",
            "feishu_drive_share_member_lookup_unsupported",
        )
    normalized_member_ids = [value.strip() for value in member_open_ids]
    if len(set(normalized_member_ids)) != len(normalized_member_ids):
        return _typed_failure(
            "feishu_drive_share member_open_ids must not contain duplicates.",
            "invalid_tool_arguments",
        )
    if action in {"add", "remove"} and not normalized_member_ids:
        return _typed_failure(
            "feishu_drive_share add/remove requires member_open_ids.",
            "invalid_tool_arguments",
        )

    token, token_error = await _feishu_access_token_outcome(agent_id)
    if token_error is not None or token is None:
        return token_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    headers = {"Authorization": f"Bearer {token}"}
    base_url = (
        "https://open.feishu.cn/open-apis/drive/v1/permissions/"
        f"{document_token}/members"
    )

    if action == "list":
        try:
            async with httpx.AsyncClient(timeout=15) as client:
                response = await client.get(
                    base_url,
                    params={"type": doc_type},
                    headers=headers,
                )
            payload = feishu_service._parse_api_response(
                response,
                stage="drive_share_list",
            )
        except Exception as exc:
            if _feishu_error_is_known_rejection(exc):
                return _typed_failure(
                    "Feishu rejected drive_share_list.",
                    "feishu_drive_share_list_rejected",
                    result_ref=document_token,
                )
            return _typed_failure(
                f"Feishu drive_share_list failed: {type(exc).__name__}.",
                "feishu_drive_share_list_failed",
                result_ref=document_token,
            )
        data = payload.get("data")
        items = data.get("items", []) if isinstance(data, Mapping) else None
        if not isinstance(items, list) or any(
            not isinstance(item, Mapping) for item in items
        ):
            return _typed_failure(
                "Feishu drive collaborator list returned an invalid payload.",
                "feishu_drive_share_list_response_invalid",
                result_ref=document_token,
            )
        lines = [
            f"Feishu file {document_token} has {len(items)} collaborator(s)."
        ]
        for item in items:
            lines.append(
                f"- member_id={item.get('member_id', '')}, "
                f"member_type={item.get('member_type', '')}, "
                f"permission={item.get('perm', '')}"
            )
        return _typed_success("\n".join(lines), result_ref=document_token)

    confirmed: list[str] = []
    async with httpx.AsyncClient(timeout=15) as client:
        for member_id in normalized_member_ids:
            try:
                if action == "add":
                    response = await client.post(
                        base_url,
                        params={"type": doc_type},
                        json={
                            "member_type": "openid",
                            "member_id": member_id,
                            "perm": permission,
                        },
                        headers=headers,
                    )
                else:
                    response = await client.delete(
                        f"{base_url}/{member_id}",
                        params={"type": doc_type, "member_type": "openid"},
                        headers=headers,
                    )
                payload = feishu_service._parse_api_response(
                    response,
                    stage=f"drive_share_{action}",
                )
            except Exception as exc:
                prefix = (
                    f"Confirmed members: {', '.join(confirmed)}. "
                    if confirmed
                    else ""
                )
                if _feishu_error_is_known_rejection(exc):
                    return _typed_failure(
                        f"{prefix}Feishu rejected member {member_id}.",
                        f"feishu_drive_share_{action}_rejected",
                        result_ref=document_token,
                    )
                return _typed_unknown(
                    f"{prefix}Outcome for member {member_id} is unknown; "
                    "later members were not dispatched.",
                    f"feishu_drive_share_{action}_outcome_unknown",
                    result_ref=document_token,
                )

            if action == "add":
                data = payload.get("data")
                member = data.get("member") if isinstance(data, Mapping) else None
                receipt_member_id = (
                    member.get("member_id")
                    if isinstance(member, Mapping)
                    else None
                )
                receipt_member_type = (
                    member.get("member_type")
                    if isinstance(member, Mapping)
                    else None
                )
                receipt_permission = (
                    member.get("perm")
                    if isinstance(member, Mapping)
                    else None
                )
                if (
                    receipt_member_id != member_id
                    or receipt_member_type != "openid"
                    or receipt_permission != permission
                ):
                    prefix = (
                        f"Confirmed members: {', '.join(confirmed)}. "
                        if confirmed
                        else ""
                    )
                    return _typed_unknown(
                        f"{prefix}Feishu omitted the receipt for member {member_id}; "
                        "later members were not dispatched.",
                        "feishu_drive_share_member_receipt_missing",
                        result_ref=document_token,
                    )
            confirmed.append(member_id)

    document_url: str | None = None
    try:
        document_url = await _get_feishu_tenant_doc_url(
            token,
            document_token,
            doc_type=doc_type,
        )
    except Exception:
        pass
    summary = (
        f"Feishu drive share {action} confirmed for document {document_token}: "
        f"{', '.join(confirmed)}."
    )
    if document_url:
        summary += f" URL: {document_url}"
    return _typed_success(summary, result_ref=document_token)


async def _feishu_drive_delete_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Delete exactly once and require the folder task receipt when applicable."""
    import httpx
    from app.services.feishu_service import feishu_service

    file_token = arguments.get("file_token")
    file_type = arguments.get("file_type")
    valid_types = {
        "file",
        "docx",
        "bitable",
        "folder",
        "doc",
        "sheet",
        "mindnote",
        "shortcut",
        "slides",
    }
    if not isinstance(file_token, str) or not file_token.strip():
        return _typed_failure(
            "feishu_drive_delete requires file_token.",
            "invalid_tool_arguments",
        )
    if not isinstance(file_type, str) or file_type not in valid_types:
        return _typed_failure(
            "feishu_drive_delete file_type is unsupported.",
            "invalid_tool_arguments",
        )
    file_token = file_token.strip()

    token, token_error = await _feishu_access_token_outcome(agent_id)
    if token_error is not None or token is None:
        return token_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    try:
        async with httpx.AsyncClient(timeout=15) as client:
            response = await client.delete(
                "https://open.feishu.cn/open-apis/drive/v1/files/"
                f"{file_token}",
                params={"type": file_type},
                headers={"Authorization": f"Bearer {token}"},
            )
        payload = feishu_service._parse_api_response(
            response,
            stage="drive_delete",
        )
    except Exception as exc:
        return _feishu_write_exception_outcome(
            "drive_delete",
            exc,
            result_ref=file_token,
        )

    data, data_error = _feishu_doc_write_data(
        payload,
        "drive_delete",
        result_ref=file_token,
    )
    if data_error is not None or data is None:
        return data_error or _typed_unknown(
            "Feishu drive delete returned no receipt; reconcile first.",
            "feishu_drive_delete_outcome_unknown",
            result_ref=file_token,
        )
    if file_type == "folder":
        task_id = data.get("task_id")
        if not isinstance(task_id, str) or not task_id.strip():
            return _typed_unknown(
                f"Folder delete for {file_token} omitted task_id; reconcile first.",
                "feishu_drive_delete_task_receipt_missing",
                result_ref=file_token,
            )
        task_id = task_id.strip()
        return _typed_success(
            f"Feishu folder {file_token} delete task accepted as {task_id}.",
            result_ref=task_id,
        )

    file_url: str | None = None
    try:
        file_url = await _get_feishu_tenant_doc_url(
            token,
            file_token,
            doc_type=file_type,
        )
    except Exception:
        pass
    summary = f"Feishu {file_type} {file_token} was moved to the recycle bin."
    if file_url:
        summary += f" Previous URL: {file_url}"
    return _typed_success(summary, result_ref=file_token)


async def _feishu_doc_search(agent_id: uuid.UUID, arguments: dict) -> str:
    """Search Feishu documents by keyword using the official document search API."""
    import httpx

    query = (arguments.get("query") or arguments.get("search_key") or "").strip()
    if not query:
        return "❌ Missing required argument 'query'"

    count = max(1, min(int(arguments.get("count", 10)), 50))
    offset = max(0, int(arguments.get("offset", 0)))
    docs_types = arguments.get("docs_types") or []
    if docs_types and not isinstance(docs_types, list):
        return "❌ 'docs_types' must be an array of strings."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."

    from app.services.feishu_service import feishu_service

    token = await feishu_service.get_tenant_access_token(app_id, app_secret)
    payload: dict[str, object] = {
        "search_key": query,
        "count": count,
        "offset": offset,
    }
    if docs_types:
        payload["docs_types"] = docs_types

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/suite/docs-api/search/object",
            headers={
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/json",
            },
            json=payload,
        )

    data = resp.json()
    err = _check_feishu_err(data)
    if err:
        return err

    result = data.get("data", {})
    entities = result.get("docs_entities", []) or []
    total = result.get("total", len(entities))
    has_more = bool(result.get("has_more", False))
    if not entities:
        return (
            f"🔎 未找到与 `{query}` 匹配的飞书文档。"
            "\n可以尝试："
            "\n1. 缩短关键词"
            "\n2. 换同义词"
            "\n3. 指定 docs_types 过滤，例如 ['docx'] 或 ['bitable']"
        )

    lines = [
        f"🔎 飞书文档搜索结果：关键词 `{query}`",
        f"返回 {len(entities)} 条，total={total}，offset={offset}，has_more={str(has_more).lower()}",
        "",
    ]
    for idx, item in enumerate(entities, start=offset + 1):
        title = item.get("title") or "(无标题)"
        docs_token = item.get("docs_token") or ""
        docs_type = item.get("docs_type") or "unknown"
        owner_id = item.get("owner_id") or ""
        lines.append(
            f"{idx}. **{title}**\n"
            f"   - docs_type: `{docs_type}`\n"
            f"   - docs_token: `{docs_token}`\n"
            f"   - owner_id: `{owner_id}`"
        )

    lines.append("")
    lines.append("💡 后续操作建议：")
    lines.append("- 读取普通文档/知识库页：`feishu_doc_read(document_token=\"...\")`")
    lines.append("- 管理权限：`feishu_drive_share(document_token=\"...\", doc_type=\"...\", action=\"list|add|remove\")`")
    lines.append("- 删除文件：`feishu_drive_delete(file_token=\"...\", file_type=\"...\")`")
    if has_more:
        lines.append(f"- 下一页：`feishu_doc_search(query=\"{query}\", offset={offset + len(entities)}, count={count})`")

    return "\n".join(lines)


async def _feishu_wiki_list(agent_id: uuid.UUID, arguments: dict) -> str:
    """List sub-pages of a Feishu Wiki node, optionally recursive."""
    import httpx

    node_token = (arguments.get("node_token") or "").strip()
    recursive = bool(arguments.get("recursive", False))

    if not node_token:
        return "❌ Missing required argument 'node_token'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)
    headers = {"Authorization": f"Bearer {token}"}

    # Resolve node → space_id
    node_info = await _feishu_wiki_get_node(node_token, token)
    if not node_info:
        return (
            f"❌ 无法解析 Wiki 节点 `{node_token}`。\n"
            "请确认 token 来自飞书知识库 URL（https://xxx.feishu.cn/wiki/NodeToken），"
            "而非普通文档 URL。"
        )

    space_id = node_info["space_id"]
    if not space_id:
        return f"❌ 无法获取知识库 space_id，请检查 token 是否正确。"

    async def _list_children(parent_token: str, depth: int) -> list[dict]:
        """Return flat list of {title, node_token, obj_token, has_child, depth}."""
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/nodes",
                headers=headers,
                params={"parent_node_token": parent_token, "page_size": 50},
            )
        data = resp.json()
        if data.get("code") != 0:
            return []
        items = data.get("data", {}).get("items", [])
        result = []
        for item in items:
            entry = {
                "title": item.get("title", "(无标题)"),
                "node_token": item.get("node_token", ""),
                "obj_token": item.get("obj_token", ""),
                "has_child": item.get("has_child", False),
                "depth": depth,
            }
            result.append(entry)
            if recursive and entry["has_child"] and depth < 2:
                children = await _list_children(entry["node_token"], depth + 1)
                result.extend(children)
        return result

    pages = await _list_children(node_token, 0)
    if not pages:
        return f"📂 Wiki 页面 `{node_token}` 下没有子页面。"

    lines = [f"📂 Wiki 页面 `{node_token}` 的子页面（共 {len(pages)} 个）：\nspace_id: `{space_id}`\n"]
    for p in pages:
        indent = "  " * p["depth"]
        child_hint = " _(有子页面)_" if p["has_child"] else ""
        lines.append(
            f"{indent}• **{p['title']}**{child_hint}\n"
            f"{indent}  node_token: `{p['node_token']}`\n"
            f"{indent}  obj_token: `{p['obj_token']}`"
        )
    lines.append(
        "\n💡 用 `feishu_doc_read(document_token=\"<node_token>\")` 读取每个子页面的内容。"
        "\n   对有子页面的条目，再次调用 `feishu_wiki_list(node_token=\"...\")` 继续展开。"
    )
    return "\n".join(lines)


async def _feishu_doc_read(agent_id: uuid.UUID, arguments: dict) -> str:
    document_token = arguments.get("document_token", "").strip()
    if not document_token:
        url = arguments.get("url", "")
        parsed = _parse_feishu_url(url)
        document_token = parsed.get("document_token", parsed.get("wiki_token", ""))
        
    if not document_token:
        return "Failed: Missing required argument 'document_token'"
    max_chars = min(int(arguments.get("max_chars", 6000)), 20000)

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
    
    read_token = document_token
    wiki_hint = ""
    node_info = await _feishu_wiki_get_node(document_token, tenant_token)
    if node_info and node_info.get("obj_token"):
        read_token = node_info["obj_token"]
        if node_info.get("has_child"):
            wiki_hint = (
                "\n\n> 💡 这是一个 Wiki 目录页，它有多个子页面。"
                "使用 `feishu_wiki_list` 工具（传入相同的 node_token）可以查看所有子页面列表。"
            )

    try:
        resp = await feishu_service.read_feishu_doc(app_id, app_secret, read_token)
        err = _check_feishu_err(resp)
        if err: return err
        
        content = resp.get("data", {}).get("content", "")
        if not content:
            return f"📄 Document '{document_token}' is empty.{wiki_hint}"

        truncated = ""
        if len(content) > max_chars:
            content = content[:max_chars]
            truncated = f"\n\n_(Truncated to {max_chars} chars)_"

        return f"📄 **Document content** (`{document_token}`):\n\n{content}{truncated}{wiki_hint}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


async def _feishu_doc_create(agent_id: uuid.UUID, arguments: dict) -> str:
    title = arguments.get("title", "").strip()
    if not title:
        return "Failed: Missing required argument 'title'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    folder_token = (arguments.get("folder_token") or "").strip()
    wiki_space_id = (arguments.get("wiki_space_id") or "").strip()
    parent_node_token = (arguments.get("parent_node_token") or "").strip()

    from app.services.feishu_service import feishu_service
    tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    try:
        import httpx

        # ── Smart fallback: if folder_token is actually a wiki node token,
        #    auto-redirect to wiki creation branch. This handles LLMs that
        #    pass the wiki node token via the old folder_token param.
        if folder_token and not wiki_space_id and not parent_node_token:
            probe = await _feishu_wiki_get_node(folder_token, tenant_token)
            if probe and probe.get("space_id"):
                wiki_space_id = probe["space_id"]
                parent_node_token = probe.get("node_token", folder_token)
                folder_token = ""  # Don't use as Drive folder

        # ── Wiki branch: create as a wiki node ──────────────────────────
        # If parent_node_token is given but wiki_space_id is not,
        # resolve space_id from the parent node automatically.
        if parent_node_token and not wiki_space_id:
            node_info = await _feishu_wiki_get_node(parent_node_token, tenant_token)
            if node_info and node_info.get("space_id"):
                wiki_space_id = node_info["space_id"]

        if wiki_space_id:
            body: dict = {
                "obj_type": "docx",
                "node_type": "origin",  # Required by Feishu Wiki API: "origin" = new entity
                "title": title,
            }
            if parent_node_token:
                body["parent_node_token"] = parent_node_token

            import logging
            _wiki_log = logging.getLogger("feishu_wiki_create")
            _wiki_log.info(f"Creating wiki node in space={wiki_space_id}, body={body}")

            async with httpx.AsyncClient(timeout=15) as client:
                resp = await client.post(
                    f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{wiki_space_id}/nodes",
                    json=body,
                    headers={"Authorization": f"Bearer {tenant_token}"},
                )
            result = resp.json()
            _wiki_log.info(f"Wiki create response: code={result.get('code')}, msg={result.get('msg')}")
            err = _check_feishu_err(result)
            if err:
                return err

            node = result.get("data", {}).get("node", {})
            # obj_token is the underlying docx token used by feishu_doc_append
            doc_token = node.get("obj_token", "")
            node_token = node.get("node_token", "")
            # Wiki docs are accessed via /wiki/{node_token}, not /docx/{obj_token}
            doc_url = await _get_feishu_tenant_doc_url(tenant_token, node_token, doc_type="wiki")

            return (
                f"✅ 知识库文档创建成功！\n"
                f"标题：{title}\n"
                f"文档 Token（用于 feishu_doc_append）：{doc_token}\n"
                f"Wiki Node Token：{node_token}\n"
                f"🔗 访问链接：{doc_url}\n"
                f"下一步：调用 feishu_doc_append(document_token=\"{doc_token}\", content=\"...\") 写入正文内容。"
            )

        # ── Regular Drive branch (original behavior) ─────────────────────
        resp = await feishu_service.create_feishu_doc(app_id, app_secret, folder_token, title)
        err = _check_feishu_err(resp)
        if err: return err
        
        doc = resp.get("data", {}).get("document", {})
        doc_token = doc.get("document_id", "")
        doc_url = await _get_feishu_tenant_doc_url(tenant_token, doc_token)
        
        # Auto-share with the Feishu sender so they can access the document.
        # channel_feishu_sender_open_id is a module-level ContextVar defined in this file;
        # no import needed — it is already in scope.
        share_note = ""
        try:
            sender_open_id = channel_feishu_sender_open_id.get(None)
            if sender_open_id and doc_token:
                async with httpx.AsyncClient(timeout=10) as client:
                    share_resp = await client.post(
                        f"https://open.feishu.cn/open-apis/drive/v1/permissions/{doc_token}/members",
                        params={"type": "docx"},
                        json={
                            "member_type": "openid",
                            "member_id": sender_open_id,
                            "perm": "full_access",
                        },
                        headers={"Authorization": f"Bearer {tenant_token}"},
                    )
                sr = share_resp.json()
                if sr.get("code") == 0:
                    share_note = "\n✅ 已自动为你开通访问权限。"
                else:
                    share_note = f"\n⚠️ 自动授权失败（{sr.get('code')}），你可能需要手动在飞书前端搜索此文件。"
        except Exception as _e:
            share_note = f"\n⚠️ 自动授权异常: {_e}"

        return (
            f"✅ 文档创建成功！{share_note}\n"
            f"标题：{title}\n"
            f"Token：{doc_token}\n"
            f"🔗 访问链接：{doc_url}\n"
            f"下一步：调用 feishu_doc_append(document_token=\"{doc_token}\", content=\"...\") 写入正文内容。"
        )
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


def _parse_inline_markdown(text: str) -> list[dict]:
    """Parse inline markdown (bold, italic, strikethrough) into Feishu text_run elements.
    Note: inline `code` is deliberately NOT rendered as inline_code style because
    Feishu's API rejects inline_code inside heading blocks (field validation error).
    Instead, backtick-wrapped text is returned as plain text.
    Empty text_element_style dicts are intentionally omitted to avoid API validation errors.
    """
    import re as _re

    def _make_run(content: str, style: dict | None = None) -> dict:
        run: dict = {"content": content}
        if style:
            run["text_element_style"] = style
        return {"text_run": run}

    elements = []
    # Only handle **bold**, *italic*, ~~strikethrough~~; backticks become plain text
    pattern = r'(\*\*(.+?)\*\*|\*(.+?)\*|~~(.+?)~~|`(.+?)`)'
    pos = 0
    for m in _re.finditer(pattern, text):
        if m.start() > pos:
            elements.append(_make_run(text[pos:m.start()]))
        raw = m.group(0)
        if raw.startswith("**"):
            elements.append(_make_run(m.group(2), {"bold": True}))
        elif raw.startswith("~~"):
            elements.append(_make_run(m.group(4), {"strikethrough": True}))
        elif raw.startswith("`"):
            # Render as plain text to avoid inline_code validation issues in headings
            elements.append(_make_run(m.group(5)))
        else:
            elements.append(_make_run(m.group(3), {"italic": True}))
        pos = m.end()
    if pos < len(text):
        elements.append(_make_run(text[pos:]))
    if not elements:
        elements.append(_make_run(text or " "))
    return elements


def _markdown_to_feishu_blocks(markdown: str) -> list[dict]:
    """Convert Markdown text to Feishu docx v1 block list.

    Supported:
      # / ## / ### / ####  → heading1-4 (block_type 3-6)
      - / * / + text       → bullet      (block_type 12)
      1. text              → ordered     (block_type 13)
      > text               → quote       (block_type 15)
      --- / ***            → divider     (block_type 22)
      ``` ... ```          → code block  (block_type 14)
      plain text           → text        (block_type 2)
      inline **bold** *italic* `code` ~~strike~~  → text_element_style
    """
    import re as _re

    _HEADING_BLOCK = {1: (3, "heading1"), 2: (4, "heading2"),
                      3: (5, "heading3"), 4: (6, "heading4")}

    def _text_block(bt: int, key: str, line: str) -> dict:
        # Omit "style" entirely to avoid Feishu field validation errors on empty style dicts
        return {
            "block_type": bt,
            key: {"elements": _parse_inline_markdown(line)},
        }

    blocks: list[dict] = []
    lines = markdown.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]

        # ── Code fence ──────────────────────────────────────────────────────
        if line.strip().startswith("```"):
            lang = line.strip()[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            blocks.append({
                "block_type": 14,
                "code": {
                    "elements": [{"text_run": {"content": "\n".join(code_lines)}}],
                    "style": {"language": 1 if not lang else
                              {"python": 49, "javascript": 22, "js": 22,
                               "typescript": 56, "ts": 56, "bash": 4, "sh": 4,
                               "sql": 53, "java": 21, "go": 17, "rust": 51,
                               "json": 25, "yaml": 60, "html": 19, "css": 10,
                               }.get(lang.lower(), 1)},
                },
            })
            i += 1
            continue

        # ── Divider ──────────────────────────────────────────────────────────
        if _re.fullmatch(r'[-*_]{3,}', line.strip()):
            # NOTE: block_type 22 (Feishu native divider) is rejected by the batch children
            # creation API with error 99992402 (field validation failed).  Render as a plain
            # text block containing a visual em-dash separator instead — always accepted.
            blocks.append({
                "block_type": 2,
                "text": {"elements": [{"text_run": {"content": "\u2500" * 24}}]},
            })
            i += 1
            continue

        # ── Headings ─────────────────────────────────────────────────────────
        hm = _re.match(r'^(#{1,4})\s+(.*)', line)
        if hm:
            level = min(len(hm.group(1)), 4)
            bt, key = _HEADING_BLOCK[level]
            blocks.append(_text_block(bt, key, hm.group(2)))
            i += 1
            continue

        # ── Bullet list ──────────────────────────────────────────────────────
        if _re.match(r'^[\-\*\+]\s+', line):
            text = _re.sub(r'^[\-\*\+]\s+', '', line)
            blocks.append(_text_block(12, "bullet", text))
            i += 1
            continue

        # ── Ordered list ─────────────────────────────────────────────────────
        if _re.match(r'^\d+\.\s+', line):
            text = _re.sub(r'^\d+\.\s+', '', line)
            blocks.append(_text_block(13, "ordered", text))
            i += 1
            continue

        # ── Blockquote ───────────────────────────────────────────────────────
        if line.startswith("> "):
            blocks.append(_text_block(15, "quote", line[2:]))
            i += 1
            continue

        # ── Empty line → empty text block ────────────────────────────────────
        if line.strip() == "":
            blocks.append({
                "block_type": 2,
                "text": {"elements": [{"text_run": {"content": " "}}]},
            })
            i += 1
            continue

        # ── Markdown table separator line (|---|---| ) → skip ───────────────
        if _re.match(r'^\|[\s\-:]+(\|[\s\-:]+)*\|?\s*$', line.strip()):
            i += 1
            continue

        # ── Markdown table row → plain text ──────────────────────────────────
        if line.strip().startswith("|") and line.strip().endswith("|"):
            # Strip pipe separators and render each cell as plain text
            cells = [c.strip() for c in line.strip().strip("|").split("|")]
            cell_text = "  |  ".join(c for c in cells if c)
            blocks.append(_text_block(2, "text", cell_text))
            i += 1
            continue

        # ── Plain text (with inline formatting) ──────────────────────────────
        blocks.append(_text_block(2, "text", line))
        i += 1

    return blocks


async def _feishu_doc_append(agent_id: uuid.UUID, arguments: dict) -> str:
    document_token = arguments.get("document_token", "").strip()
    if not document_token:
        url = arguments.get("url", "")
        parsed = _parse_feishu_url(url)
        document_token = parsed.get("document_token", parsed.get("wiki_token", ""))
        
    content = arguments.get("content", "").strip()
    if not document_token:
        return "Failed: Missing required argument 'document_token'"
    if not content:
        return "Failed: Missing required argument 'content'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    # For wiki node tokens, use the obj_token for the docx API
    node_info = await _feishu_wiki_get_node(document_token, tenant_token)
    docx_token = node_info["obj_token"] if (node_info and node_info.get("obj_token")) else document_token

    try:
        import httpx
        async with httpx.AsyncClient(timeout=20) as client:
            meta_resp = (await client.get(
                f"https://open.feishu.cn/open-apis/docx/v1/documents/{docx_token}",
                headers={"Authorization": f"Bearer {tenant_token}"},
            )).json()
            err = _check_feishu_err(meta_resp)
            if err: return err

            body_block_id = (
                meta_resp.get("data", {}).get("document", {}).get("body", {}).get("block_id")
                or docx_token
            )

            children = _markdown_to_feishu_blocks(content)

            result = (await client.post(
                f"https://open.feishu.cn/open-apis/docx/v1/documents/{docx_token}/blocks/{body_block_id}/children",
                # Do NOT pass index: -1.  Omitting the field lets Feishu default to
                # append-at-end, which is always valid.  Passing -1 explicitly can
                # trigger error 1770001 (invalid param) with certain block type mixes.
                json={"children": children},
                headers={"Authorization": f"Bearer {tenant_token}"},
            )).json()

            err = _check_feishu_err(result)
            if err: return err

        doc_url = await _get_feishu_tenant_doc_url(tenant_token, docx_token)
        return (
            f"✅ 已写入 {len(children)} 个段落到文档。\n"
            f"🔗 文档直链（原文发给用户，勿修改）：{doc_url}"
        )
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


# ─── Feishu Drive Share (All File Types) ────────────────────────────────────────

async def _feishu_drive_share(agent_id: uuid.UUID, arguments: dict) -> str:
    """Manage Feishu drive file collaborators.
    Automatically handles both regular docs/files (Drive permissions API)
    and Wiki node documents (Wiki space members API).
    """
    import httpx

    document_token = (arguments.get("document_token") or "").strip()
    doc_type = (arguments.get("doc_type") or "docx").strip()
    action = (arguments.get("action") or "list").strip()
    permission = (arguments.get("permission") or "edit").strip()

    if not document_token:
        return "❌ Missing required argument 'document_token'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)
    headers = {"Authorization": f"Bearer {token}"}

    # ── Detect if this is a Wiki node token ─────────────────────────────────
    node_info = await _feishu_wiki_get_node(document_token, token)
    is_wiki = node_info is not None
    space_id = node_info.get("space_id", "") if node_info else ""
    obj_token = node_info.get("obj_token", "") if node_info else ""

    # Permission level mapping: Feishu API uses "view" / "edit" / "full_access"
    api_perm = {"view": "view", "edit": "edit", "full_access": "full_access"}.get(permission, "edit")
    # Wiki space role mapping: only "admin" / "member" are valid roles
    wiki_role = "admin" if api_perm in ("edit", "full_access") else "member"

    # ── LIST collaborators ────────────────────────────────────────────────────
    if action == "list":
        use_token = obj_token if (is_wiki and obj_token) else document_token
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                f"https://open.feishu.cn/open-apis/drive/v1/permissions/{use_token}/members",
                params={"type": doc_type},
                headers=headers,
            )
        data = resp.json()
        if data.get("code") != 0:
            _c = data.get("code")
            if _c == 1063003 and is_wiki:
                return (
                    f"ℹ️ 文档 `{document_token}` 是知识库页面，其权限由知识库空间统一管理。\n"
                    "知识库空间 ID：`" + space_id + "`\n"
                    "请直接在飞书知识库中管理成员权限。"
                )
            if _c in (99991672, 99991668):
                return (
                    f"❌ 权限不足（code {_c}）\n"
                    "需要在飞书开放平台开通：\n"
                    "• drive:drive（云文档权限管理）"
                )
            return f"❌ 获取协作者列表失败：{data.get('msg')} (code {_c})"

        members = data.get("data", {}).get("items", [])
        if not members:
            return f"📄 文档 `{document_token}` 当前没有其他协作者。"

        lines = [f"📄 文档 `{document_token}` 的协作者列表（共 {len(members)} 人）：\n"]
        for m in members:
            perm = m.get("perm", "")
            member_type = m.get("member_type", "")
            member_id = m.get("member_id", "")
            _type_label = {"openid": "用户", "openchat": "群组", "opendepartmentid": "部门"}.get(member_type, member_type)
            lines.append(f"• {_type_label} `{member_id}` | 权限: **{perm}**")
        return "\n".join(lines)

    # ── ADD / REMOVE collaborators ─────────────────────────────────────────────
    member_names: list[str] = list(arguments.get("member_names") or [])
    member_open_ids: list[str] = list(arguments.get("member_open_ids") or [])

    if not member_names and not member_open_ids:
        return "❌ 请提供 member_names（姓名列表）或 member_open_ids（open_id 列表）"

    # Resolve names → open_ids
    resolved: list[tuple[str, str]] = []  # (display_name, open_id)
    for name in member_names:
        open_id = await _feishu_open_id_for_visible_name(agent_id, name)
        resolved.append((name, open_id or ""))

    for oid in member_open_ids:
        if oid:
            resolved.append((oid, oid))

    results = []
    async with httpx.AsyncClient(timeout=15) as client:
        for display, oid in resolved:
            if not oid:
                results.append(f"❌ 无法找到「{display}」的 open_id，跳过")
                continue

            if action == "add":
                # ── Wiki node: use wiki space members API ──────────────────
                if is_wiki and space_id:
                    resp = await client.post(
                        f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/members",
                        json={"member_type": "openid", "member_id": oid, "member_role": wiki_role},
                        headers=headers,
                    )
                    d = resp.json()
                    _c = d.get("code")
                    if _c == 0:
                        results.append(f"✅ 已将「{display}」加入知识库空间（角色：{wiki_role}）")
                    elif _c == 131008:
                        results.append(f"ℹ️ 「{display}」已经是知识库成员，无需重复添加")
                    elif _c == 131101:
                        # Public wiki space — everyone already has access
                        results.append(
                            f"ℹ️ 这是一个**公开知识库**，所有人已可访问。\n"
                            f"「{display}」无需单独添加权限。"
                        )
                    else:
                        results.append(f"❌ 添加「{display}」到知识库失败：{d.get('msg')} (code {_c})")
                    continue

                # ── Regular docx: use Drive permissions API ────────────────
                body = {
                    "member_type": "openid",
                    "member_id": oid,
                    "perm": api_perm,
                }
                resp = await client.post(
                    f"https://open.feishu.cn/open-apis/drive/v1/permissions/{document_token}/members",
                    json=body,
                    headers=headers,
                    params={"type": doc_type},
                )
                d = resp.json()
                if d.get("code") == 0:
                    results.append(f"✅ 已将「{display}」添加为**{permission}**权限协作者")
                else:
                    _c = d.get("code")
                    if _c == 99992402:
                        # Feishu platform policy: you cannot add yourself as a collaborator via API.
                        # Permissions must be granted by others, or set manually in the UI.
                        results.append(
                            f"⚠️ 飞书平台安全限制：无法通过 API 为自己添加协作权限。\n"
                            f"请手动操作：打开文档 → 右上角「分享」→ 添加自己并设置权限。"
                        )
                    elif _c in (99991672, 99991668):
                        return (
                            f"❌ 权限不足（code {_c}）\n"
                            "需要在飞书开放平台开通：\n"
                            "• drive:drive（云文档权限管理）"
                        )
                    else:
                        results.append(f"❌ 添加「{display}」失败：{d.get('msg')} (code {_c})")

            elif action == "remove":
                if is_wiki and space_id:
                    resp = await client.delete(
                        f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/members/{oid}",
                        headers=headers,
                        params={"member_type": "openid"},
                    )
                    d = resp.json()
                    if d.get("code") == 0:
                        results.append(f"✅ 已将「{display}」从知识库移除")
                    else:
                        results.append(f"❌ 移除「{display}」失败：{d.get('msg')} (code {d.get('code')})")
                    continue

                resp = await client.delete(
                    f"https://open.feishu.cn/open-apis/drive/v1/permissions/{document_token}/members/{oid}",
                    headers=headers,
                    params={"type": doc_type, "member_type": "openid"},
                )
                d = resp.json()
                if d.get("code") == 0:
                    results.append(f"✅ 已移除「{display}」的协作权限")
                else:
                    results.append(f"❌ 移除「{display}」失败：{d.get('msg')} (code {d.get('code')})")

    return "\n".join(results) if results else "没有需要处理的成员"


# ─── Feishu Drive Delete ──────────────────────────────────────────────────────

async def _feishu_drive_delete(agent_id: uuid.UUID, arguments: dict) -> str:
    """Delete a file or folder from Feishu Drive (cloud space).
    The file is moved to the recycle bin, not permanently deleted.
    For folders, the deletion is asynchronous and returns a task_id.
    """
    import httpx

    file_token = (arguments.get("file_token") or "").strip()
    file_type = (arguments.get("file_type") or "").strip()

    if not file_token:
        return "❌ Missing required argument 'file_token'"
    if not file_type:
        return "❌ Missing required argument 'file_type'. Valid values: file, docx, bitable, folder, doc, sheet, mindnote, shortcut, slides"

    valid_types = {"file", "docx", "bitable", "folder", "doc", "sheet", "mindnote", "shortcut", "slides"}
    if file_type not in valid_types:
        return f"❌ Invalid file_type '{file_type}'. Valid values: {', '.join(sorted(valid_types))}"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    # Type label mapping for user-friendly output
    type_labels = {
        "file": "文件", "docx": "文档", "bitable": "多维表格",
        "folder": "文件夹", "doc": "旧版文档", "sheet": "电子表格",
        "mindnote": "思维笔记", "shortcut": "快捷方式", "slides": "幻灯片",
    }
    type_label = type_labels.get(file_type, file_type)

    try:
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.delete(
                f"https://open.feishu.cn/open-apis/drive/v1/files/{file_token}",
                params={"type": file_type},
                headers={"Authorization": f"Bearer {token}"},
            )
        data = resp.json()
        code = data.get("code", -1)

        if code == 0:
            # Folder deletion returns a task_id for async tracking
            task_id = data.get("data", {}).get("task_id")
            if task_id:
                return (
                    f"✅ 已提交{type_label}删除任务（异步执行中）。\n"
                    f"📋 任务 ID: `{task_id}`\n"
                    f"文件夹删除为异步操作，文件会被移至回收站。"
                )
            return f"✅ {type_label} `{file_token}` 已删除（移至回收站）。"

        # Error handling with specific codes
        msg = data.get("msg", "Unknown error")
        if code == 1061003:
            return f"❌ 未找到文件 `{file_token}`。请确认文件 token 和类型是否正确。"
        elif code == 1061004:
            return (
                f"❌ 权限不足（code {code}）\n"
                "需要满足以下条件之一：\n"
                "• 文件所有者 + 父文件夹编辑权限\n"
                "• 父文件夹的所有者或 full_access 权限\n"
                "同时需要在飞书开放平台开通：drive:drive 或 space:document:delete"
            )
        elif code == 1061007:
            return f"❌ 文件 `{file_token}` 已被删除。"
        elif code == 1061045:
            return f"⚠️ 接口频率限制，请稍后重试。（每秒最多 5 次）"
        else:
            return f"❌ 删除{type_label}失败：{msg} (code {code})"

    except Exception as e:
        return f"❌ 删除文件异常: {str(e)[:300]}"


# ─── Feishu Calendar Tools ────────────────────────────────────────────────────

async def _feishu_calendar_list_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Read Bot-calendar events; freebusy remains best-effort context only."""
    import httpx
    from app.services.feishu_service import feishu_service

    try:
        max_results = max(1, min(int(arguments.get("max_results", 20)), 100))
    except (TypeError, ValueError):
        return _typed_failure(
            "feishu_calendar_list max_results must be an integer.",
            "invalid_tool_arguments",
        )

    now = datetime.now(timezone.utc)
    start_value = arguments.get("start_time")
    end_value = arguments.get("end_time")
    if start_value is not None and not isinstance(start_value, str):
        return _typed_failure(
            "feishu_calendar_list start_time must be an ISO 8601 string.",
            "invalid_tool_arguments",
        )
    if end_value is not None and not isinstance(end_value, str):
        return _typed_failure(
            "feishu_calendar_list end_time must be an ISO 8601 string.",
            "invalid_tool_arguments",
        )
    try:
        start_epoch = (
            _iso_to_ts(start_value)
            if start_value
            else now.timestamp()
        )
        end_epoch = (
            _iso_to_ts(end_value)
            if end_value
            else (now + timedelta(days=7)).timestamp()
        )
    except ValueError:
        return _typed_failure(
            "feishu_calendar_list requires valid ISO 8601 times.",
            "invalid_tool_arguments",
        )
    if end_epoch <= start_epoch:
        return _typed_failure(
            "feishu_calendar_list end_time must be after start_time.",
            "invalid_tool_arguments",
        )

    token, calendar_id, error = await _feishu_calendar_context_outcome(agent_id)
    if error is not None or token is None or calendar_id is None:
        return error or _typed_failure(
            "Feishu Bot primary calendar is unavailable.",
            "feishu_calendar_unavailable",
        )

    freebusy_status = "not_requested"
    sender_open_id = channel_feishu_sender_open_id.get(None)
    async with httpx.AsyncClient(timeout=20) as client:
        if sender_open_id:
            try:
                freebusy_response = await client.post(
                    "https://open.feishu.cn/open-apis/calendar/v4/freebusy/list",
                    headers={"Authorization": f"Bearer {token}"},
                    params={"user_id_type": "open_id"},
                    json={
                        "time_min": datetime.fromtimestamp(
                            start_epoch,
                            tz=timezone.utc,
                        ).isoformat(),
                        "time_max": datetime.fromtimestamp(
                            end_epoch,
                            tz=timezone.utc,
                        ).isoformat(),
                        "user_id": sender_open_id,
                    },
                )
                feishu_service._parse_api_response(
                    freebusy_response,
                    stage="calendar_freebusy",
                )
                freebusy_status = "succeeded"
            except Exception:
                # Freebusy is supplemental context. It must never replace or
                # mask the Bot-calendar execution fact.
                freebusy_status = "failed"

        try:
            response = await client.get(
                f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{calendar_id}/events",
                headers={"Authorization": f"Bearer {token}"},
                params={
                    "start_time": str(int(start_epoch)),
                    "end_time": str(int(end_epoch)),
                },
            )
            data = feishu_service._parse_api_response(
                response,
                stage="calendar_list",
            )
        except Exception as exc:
            return _feishu_read_exception_outcome("calendar_list", exc)

    body = data.get("data")
    if not isinstance(body, Mapping):
        return _typed_failure(
            "Feishu calendar returned an invalid data object.",
            "feishu_calendar_response_invalid",
            retryable=True,
        )
    items = body.get("items", [])
    if not isinstance(items, list):
        return _typed_failure(
            "Feishu calendar returned an invalid event list.",
            "feishu_calendar_response_invalid",
            retryable=True,
        )
    selected = [item for item in items if isinstance(item, Mapping)][
        :max_results
    ]
    if not selected:
        return _typed_success(
            "The Feishu Bot calendar has no events in the requested range.",
            result_ref=calendar_id,
            metadata={
                "calendar_id": calendar_id,
                "event_count": 0,
                "freebusy_status": freebusy_status,
            },
        )
    lines = [f"Feishu Bot calendar returned {len(selected)} event(s):"]
    for item in selected:
        event_id = str(item.get("event_id") or "")
        summary = str(item.get("summary") or "(untitled)")
        lines.append(f"- {summary} (event_id={event_id})")
    return _typed_success(
        "\n".join(lines),
        result_ref=calendar_id,
        metadata={
            "calendar_id": calendar_id,
            "event_count": len(selected),
            "freebusy_status": freebusy_status,
        },
    )


async def _feishu_calendar_create_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Create one event, then record each attendee write independently."""
    import httpx
    from app.services.feishu_service import feishu_service

    if any(
        legacy_name in arguments
        for legacy_name in ("attendee_open_ids", "attendee_emails")
    ):
        return _typed_failure(
            "feishu_calendar_create no longer accepts direct attendee IDs or emails; use attendee_names.",
            "legacy_tool_arguments_unsupported",
        )

    required: dict[str, str] = {}
    for field in ("summary", "start_time", "end_time"):
        value = arguments.get(field)
        if not isinstance(value, str) or not value.strip():
            return _typed_failure(
                f"feishu_calendar_create requires {field}.",
                "invalid_tool_arguments",
            )
        required[field] = value.strip()
    timezone_name = arguments.get("timezone", "Asia/Shanghai")
    if not isinstance(timezone_name, str) or not timezone_name.strip():
        return _typed_failure(
            "feishu_calendar_create timezone must be a non-empty string.",
            "invalid_tool_arguments",
        )
    try:
        start_epoch = _iso_to_ts(required["start_time"])
        end_epoch = _iso_to_ts(required["end_time"])
    except ValueError:
        return _typed_failure(
            "feishu_calendar_create requires valid ISO 8601 times.",
            "invalid_tool_arguments",
        )
    if end_epoch <= start_epoch:
        return _typed_failure(
            "feishu_calendar_create end_time must be after start_time.",
            "invalid_tool_arguments",
        )

    attendee_names = arguments.get("attendee_names", []) or []
    if not isinstance(attendee_names, list) or any(
        not isinstance(name, str) for name in attendee_names
    ):
        return _typed_failure(
            "feishu_calendar_create attendee_names must be an array of strings.",
            "invalid_tool_arguments",
        )
    attendee_open_ids: list[str] = []
    unresolved_names: list[str] = []
    for attendee_name in attendee_names[:20]:
        name = attendee_name.strip()
        if not name:
            continue
        try:
            open_id = await _feishu_open_id_for_visible_name(agent_id, name)
        except Exception as exc:
            return _feishu_read_exception_outcome("attendee_lookup", exc)
        if open_id is None:
            unresolved_names.append(name)
            continue
        if open_id not in attendee_open_ids:
            attendee_open_ids.append(open_id)
    if unresolved_names:
        return _typed_failure(
            "Could not resolve requested Feishu attendee(s): "
            + ", ".join(unresolved_names),
            "feishu_attendee_not_found",
        )
    sender_open_id = channel_feishu_sender_open_id.get(None)
    if sender_open_id and sender_open_id not in attendee_open_ids:
        attendee_open_ids.append(sender_open_id)

    token, calendar_id, error = await _feishu_calendar_context_outcome(agent_id)
    if error is not None or token is None or calendar_id is None:
        return error or _typed_failure(
            "Feishu Bot primary calendar is unavailable.",
            "feishu_calendar_unavailable",
        )
    body: dict[str, object] = {
        "summary": required["summary"],
        "start_time": {
            "timestamp": str(int(start_epoch)),
            "timezone": timezone_name.strip(),
        },
        "end_time": {
            "timestamp": str(int(end_epoch)),
            "timezone": timezone_name.strip(),
        },
    }
    description = arguments.get("description")
    location = arguments.get("location")
    if isinstance(description, str) and description:
        body["description"] = description
    if isinstance(location, str) and location:
        body["location"] = {"name": location}

    async with httpx.AsyncClient(timeout=20) as client:
        try:
            response = await client.post(
                f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{calendar_id}/events",
                json=body,
                headers={"Authorization": f"Bearer {token}"},
            )
            data = feishu_service._parse_api_response(
                response,
                stage="calendar_create",
            )
        except Exception as exc:
            return _feishu_write_exception_outcome(
                "calendar_create",
                exc,
            )

        event_data = data.get("data")
        event = (
            event_data.get("event")
            if isinstance(event_data, Mapping)
            else None
        )
        event_id = (
            str(event.get("event_id") or "")
            if isinstance(event, Mapping)
            else ""
        )
        if not event_id:
            return _typed_unknown(
                "Feishu accepted calendar_create but returned no event ID; "
                "reconcile before any retry.",
                "feishu_calendar_create_receipt_missing",
            )

        invited: list[str] = []
        for attendee_open_id in attendee_open_ids:
            receipt_metadata = {
                "calendar_id": calendar_id,
                "event_id": event_id,
                "attendee_receipt_count": len(invited),
            }
            try:
                attendee_response = await client.post(
                    f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{calendar_id}/events/{event_id}/attendees",
                    json={
                        "attendees": [
                            {"type": "user", "user_id": attendee_open_id}
                        ]
                    },
                    headers={"Authorization": f"Bearer {token}"},
                    params={"user_id_type": "open_id"},
                )
                feishu_service._parse_api_response(
                    attendee_response,
                    stage="calendar_attendee_create",
                )
            except Exception as exc:
                return _feishu_write_exception_outcome(
                    "calendar_attendee_create",
                    exc,
                    result_ref=event_id,
                    metadata=receipt_metadata,
                )
            invited.append(attendee_open_id)

    return _typed_success(
        f"Created Feishu event {event_id} and confirmed "
        f"{len(invited)} attendee invitation(s).",
        result_ref=event_id,
        metadata={
            "calendar_id": calendar_id,
            "event_id": event_id,
            "attendee_receipt_count": len(invited),
        },
    )


async def _feishu_calendar_mutation_outcome(
    tool_name: str,
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Update or delete one event on the Bot primary calendar."""
    import httpx
    from app.services.feishu_service import feishu_service

    event_id = arguments.get("event_id")
    if not isinstance(event_id, str) or not event_id.strip():
        return _typed_failure(
            f"{tool_name} requires event_id.",
            "invalid_tool_arguments",
        )
    event_id = event_id.strip()

    patch: dict[str, object] | None = None
    if tool_name == "feishu_calendar_update":
        patch = {}
        timezone_name = arguments.get("timezone", "Asia/Shanghai")
        if not isinstance(timezone_name, str) or not timezone_name.strip():
            return _typed_failure(
                "feishu_calendar_update timezone must be a non-empty string.",
                "invalid_tool_arguments",
            )
        for field in ("summary", "description"):
            value = arguments.get(field)
            if isinstance(value, str) and value:
                patch[field] = value
        location = arguments.get("location")
        if isinstance(location, str) and location:
            patch["location"] = {"name": location}
        for field in ("start_time", "end_time"):
            value = arguments.get(field)
            if value is None:
                continue
            if not isinstance(value, str) or not value.strip():
                return _typed_failure(
                    f"feishu_calendar_update {field} must be an ISO 8601 string.",
                    "invalid_tool_arguments",
                )
            try:
                timestamp = _iso_to_ts(value)
            except ValueError:
                return _typed_failure(
                    f"feishu_calendar_update {field} is not valid ISO 8601.",
                    "invalid_tool_arguments",
                )
            patch[field] = {
                "timestamp": str(int(timestamp)),
                "timezone": timezone_name.strip(),
            }
        if not patch:
            return _typed_failure(
                "feishu_calendar_update requires at least one changed field.",
                "invalid_tool_arguments",
            )

    token, calendar_id, error = await _feishu_calendar_context_outcome(agent_id)
    if error is not None or token is None or calendar_id is None:
        return error or _typed_failure(
            "Feishu Bot primary calendar is unavailable.",
            "feishu_calendar_unavailable",
        )
    operation = (
        "calendar_update"
        if tool_name == "feishu_calendar_update"
        else "calendar_delete"
    )
    async with httpx.AsyncClient(timeout=20) as client:
        try:
            url = (
                "https://open.feishu.cn/open-apis/calendar/v4/calendars/"
                f"{calendar_id}/events/{event_id}"
            )
            if patch is not None:
                response = await client.patch(
                    url,
                    json=patch,
                    headers={"Authorization": f"Bearer {token}"},
                )
            else:
                response = await client.delete(
                    url,
                    headers={"Authorization": f"Bearer {token}"},
                )
            feishu_service._parse_api_response(
                response,
                stage=operation,
            )
        except Exception as exc:
            return _feishu_write_exception_outcome(
                operation,
                exc,
                result_ref=event_id,
                metadata={
                    "calendar_id": calendar_id,
                    "event_id": event_id,
                },
            )
    action = "updated" if patch is not None else "deleted"
    return _typed_success(
        f"Feishu event {event_id} was {action}.",
        result_ref=event_id,
        metadata={"calendar_id": calendar_id, "event_id": event_id},
    )


async def _feishu_calendar_list(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx
    import re as _re
    from datetime import timedelta as _td

    user_email = arguments.get("user_email", "").strip()

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    now = datetime.now(timezone.utc)

    def _to_iso(t: str | None, default: datetime) -> str:
        """Return an ISO-8601 string with timezone for freebusy API."""
        if not t:
            return default.strftime("%Y-%m-%dT%H:%M:%S+00:00")
        if _re.fullmatch(r'\d+', t.strip()):
            from datetime import datetime as _dt2
            return _dt2.fromtimestamp(int(t.strip()), tz=timezone.utc).strftime("%Y-%m-%dT%H:%M:%S+00:00")
        return t.strip()

    def _to_unix(t: str | None, default: datetime) -> str:
        """Convert ISO-8601 / Unix string / None to Unix timestamp string."""
        if not t:
            return str(int(default.timestamp()))
        if _re.fullmatch(r'\d+', t.strip()):
            return t.strip()
        try:
            from datetime import datetime as _dt2
            for fmt in ("%Y-%m-%dT%H:%M:%S%z", "%Y-%m-%dT%H:%M:%SZ", "%Y-%m-%dT%H:%M:%S"):
                try:
                    dt = _dt2.strptime(t.strip(), fmt)
                    if dt.tzinfo is None:
                        dt = dt.replace(tzinfo=timezone.utc)
                    return str(int(dt.timestamp()))
                except ValueError:
                    continue
            from dateutil import parser as _dp
            return str(int(_dp.parse(t).timestamp()))
        except Exception:
            return str(int(default.timestamp()))

    start_arg = arguments.get("start_time")
    end_arg = arguments.get("end_time")
    start_ts = _to_unix(start_arg, now)
    end_ts = _to_unix(end_arg, now + _td(days=7))
    start_iso = _to_iso(start_arg, now)
    end_iso = _to_iso(end_arg, now + _td(days=7))

    # ── 1. Query sender's real freebusy from Feishu Calendar ─────────────────
    sender_open_id = channel_feishu_sender_open_id.get(None)
    # Allow explicit override via argument
    if arguments.get("user_open_id"):
        sender_open_id = arguments["user_open_id"]
    elif user_email:
        resolved = await _feishu_resolve_open_id(token, user_email)
        if resolved:
            sender_open_id = resolved

    freebusy_section = ""
    if sender_open_id:
        try:
            async with httpx.AsyncClient(timeout=10) as fb_client:
                fb_resp = await fb_client.post(
                    "https://open.feishu.cn/open-apis/calendar/v4/freebusy/list",
                    headers={"Authorization": f"Bearer {token}"},
                    params={"user_id_type": "open_id"},
                    json={
                        "time_min": start_iso,
                        "time_max": end_iso,
                        "user_id": sender_open_id,
                    },
                )
            fb_data = fb_resp.json()
            if fb_data.get("code") == 0:
                busy_slots = fb_data.get("data", {}).get("freebusy_list", [])
                if busy_slots:
                    from datetime import datetime as _dt2
                    from zoneinfo import ZoneInfo
                    tz_cn = ZoneInfo("Asia/Shanghai")
                    busy_lines = []
                    for slot in sorted(busy_slots, key=lambda x: x.get("start_time", "")):
                        try:
                            s = _dt2.fromisoformat(slot["start_time"]).astimezone(tz_cn).strftime("%H:%M")
                            e = _dt2.fromisoformat(slot["end_time"]).astimezone(tz_cn).strftime("%H:%M")
                            busy_lines.append(f"  🔴 {s}–{e}")
                        except Exception:
                            busy_lines.append(f"  🔴 {slot.get('start_time')}–{slot.get('end_time')}")
                    freebusy_section = f"\n📌 **用户真实日历（忙碌时段）**：\n" + "\n".join(busy_lines)
                else:
                    freebusy_section = "\n📌 **用户真实日历**：该时段全部空闲。"
        except Exception as _fe:
            freebusy_section = f"\n⚠️ Freebusy 查询异常: {_fe}"

    # ── 2. Also list bot's own calendar events ───────────────────────────────
    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        # Return freebusy results even if bot calendar fails
        if freebusy_section:
            return freebusy_section.strip()
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    # Note: page_size is NOT a valid param for this API — omit it entirely
    params: dict = {}
    if start_ts:
        params["start_time"] = start_ts
    if end_ts:
        params["end_time"] = end_ts

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.get(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events",
            headers={"Authorization": f"Bearer {token}"},
            params=params,
        )

    data = resp.json()
    if data.get("code") != 0:
        if freebusy_section:
            return freebusy_section.strip()
        return f"❌ Calendar API error: {data.get('msg')} (code {data.get('code')})"

    items = data.get("data", {}).get("items", [])
    if not items and not freebusy_section:
        return "📅 该时间段内没有日程。"

    lines = []
    if items:
        lines.append(f"📅 Bot 日历共 {len(items)} 个日程：\n")
    for ev in items:
        summary = ev.get("summary", "(no title)")
        start = ev.get("start_time", {}).get("timestamp", "")
        end_t = ev.get("end_time", {}).get("timestamp", "")
        location = ev.get("location", {}).get("name", "")
        event_id = ev.get("event_id", "")
        try:
            from datetime import datetime as _dt
            s = _dt.fromtimestamp(int(start), tz=timezone.utc).strftime("%m-%d %H:%M") if start else "?"
            e = _dt.fromtimestamp(int(end_t), tz=timezone.utc).strftime("%H:%M") if end_t else "?"
        except Exception:
            s, e = start, end_t
        loc_str = f" | 📍{location}" if location else ""
        lines.append(f"- **{summary}** | 🕐{s}–{e}{loc_str}  (ID: `{event_id}`)")

    if freebusy_section:
        lines.append(freebusy_section)

    return "\n".join(lines) if lines else "📅 该时间段内没有日程。"


async def _feishu_calendar_create(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    user_email = arguments.get("user_email", "").strip()
    summary = arguments.get("summary", "").strip()
    start_time = arguments.get("start_time", "").strip()
    end_time = arguments.get("end_time", "").strip()

    for f, v in [("summary", summary), ("start_time", start_time), ("end_time", end_time)]:
        if not v:
            return f"❌ Missing required argument '{f}'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    # Resolve organizer open_id from email — soft failure
    organizer_open_id: str | None = None
    if user_email:
        organizer_open_id = await _feishu_resolve_open_id(token, user_email)
        if not organizer_open_id:
            logger.warning(f"[Feishu Calendar] Could not resolve open_id for '{user_email}', continuing without organizer invite")

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    tz = arguments.get("timezone", "Asia/Shanghai")
    body: dict = {
        "summary": summary,
        "start_time": {"timestamp": str(int(_iso_to_ts(start_time))), "timezone": tz},
        "end_time": {"timestamp": str(int(_iso_to_ts(end_time))), "timezone": tz},
    }
    if arguments.get("description"):
        body["description"] = arguments["description"]
    if arguments.get("location"):
        body["location"] = {"name": arguments["location"]}

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.post(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events",
            json=body,
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to create event: {data.get('msg')} (code {data.get('code')})"

    event_id = data.get("data", {}).get("event", {}).get("event_id", "")

    # Collect all attendee open_ids to invite
    attendee_open_ids: list[str] = []
    attendee_display: list[str] = []  # for summary message

    # 1. Direct open_ids provided by caller
    for oid in (arguments.get("attendee_open_ids") or []):
        if oid and oid not in attendee_open_ids:
            attendee_open_ids.append(oid)
            attendee_display.append(oid)

    # 2. Names → look up via feishu_user_search
    for aname in (arguments.get("attendee_names") or []):
        aname = aname.strip()
        if not aname:
            continue
        _oid = await _feishu_open_id_for_visible_name(agent_id, aname)
        if _oid:
            if _oid not in attendee_open_ids:
                attendee_open_ids.append(_oid)
                attendee_display.append(aname)
        else:
            logger.warning(
                f"[Calendar] Could not resolve attendee '{aname}'"
            )

    # 3. From explicit attendee_emails
    attendee_emails: list[str] = list(arguments.get("attendee_emails") or [])
    if user_email and user_email not in attendee_emails:
        attendee_emails.append(user_email)
    for email in attendee_emails[:20]:
        oid = await _feishu_resolve_open_id(token, email)
        if oid and oid not in attendee_open_ids:
            attendee_open_ids.append(oid)
            attendee_display.append(email)

    # 4. Auto-invite the Feishu message sender (from context var)
    sender_oid = channel_feishu_sender_open_id.get(None)
    if sender_oid and sender_oid not in attendee_open_ids:
        attendee_open_ids.append(sender_oid)

    if attendee_open_ids and event_id:
        async with httpx.AsyncClient(timeout=20) as client:
            for oid in attendee_open_ids:
                await client.post(
                    f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}/attendees",
                    json={"attendees": [{"type": "user", "user_id": oid}]},
                    headers={"Authorization": f"Bearer {token}"},
                    params={"user_id_type": "open_id"},
                )

    att_str = f"\n**参与人**: {', '.join(attendee_display)}" if attendee_display else ""
    invite_note = "\n（已向您发送日历邀请，请在飞书日历中确认）" if attendee_open_ids else ""
    return (
        f"✅ 日历事件已创建！\n"
        f"**标题**: {summary}\n"
        f"**时间**: {start_time} → {end_time}{att_str}\n"
        f"**Event ID**: `{event_id}`{invite_note}"
    )


async def _feishu_calendar_update(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    event_id = arguments.get("event_id", "").strip()
    if not event_id:
        return "❌ 'event_id' is required."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    patch: dict = {}
    tz = arguments.get("timezone", "Asia/Shanghai")
    if arguments.get("summary"):
        patch["summary"] = arguments["summary"]
    if arguments.get("description"):
        patch["description"] = arguments["description"]
    if arguments.get("location"):
        patch["location"] = {"name": arguments["location"]}
    if arguments.get("start_time"):
        patch["start_time"] = {"timestamp": str(int(_iso_to_ts(arguments["start_time"]))), "timezone": tz}
    if arguments.get("end_time"):
        patch["end_time"] = {"timestamp": str(int(_iso_to_ts(arguments["end_time"]))), "timezone": tz}

    if not patch:
        return "ℹ️ No fields to update."

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.patch(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}",
            json=patch,
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to update: {data.get('msg')} (code {data.get('code')})"

    return f"✅ Event `{event_id}` updated. Changed: {', '.join(patch.keys())}."


async def _feishu_calendar_delete(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    event_id = arguments.get("event_id", "").strip()
    if not event_id:
        return "❌ 'event_id' is required."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.delete(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}",
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to delete: {data.get('msg')} (code {data.get('code')})"

    return f"✅ Event `{event_id}` deleted successfully."

# ─── Feishu Approval Tools ───────────────────────────────────────────────────

_FEISHU_APPROVAL_STATUSES = frozenset(
    {"PENDING", "APPROVED", "REJECTED", "CANCELED", "DELETED"}
)
_FEISHU_APPROVAL_SECTIONS = frozenset(
    {"summary", "form", "tasks", "timeline", "comments"}
)
_FEISHU_APPROVAL_SECTION_KEYS = {
    "form": "form",
    "tasks": "task_list",
    "timeline": "timeline",
    "comments": "comment_list",
}


def _feishu_approval_read_response(
    response: object,
    operation: str,
) -> tuple[Mapping | None, ToolExecutionOutcome | None]:
    """Validate one approval read response without losing HTTP status facts."""
    status_code = getattr(response, "status_code", None)
    if not isinstance(status_code, int) or isinstance(status_code, bool):
        return None, _typed_failure(
            f"Feishu {operation} returned no readable HTTP status.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    if status_code == 429 or status_code >= 500:
        return None, _typed_failure(
            f"Feishu {operation} is temporarily unavailable.",
            f"feishu_{operation}_http_retryable",
            retryable=True,
        )
    if 400 <= status_code < 500:
        return None, _typed_failure(
            f"Feishu rejected {operation}.",
            f"feishu_{operation}_http_rejected",
        )
    if not 200 <= status_code < 300:
        return None, _typed_failure(
            f"Feishu {operation} returned an unexpected HTTP status.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    try:
        payload = response.json()
    except Exception:
        return None, _typed_failure(
            f"Feishu {operation} returned unreadable JSON.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    if not isinstance(payload, Mapping):
        return None, _typed_failure(
            f"Feishu {operation} returned an invalid response.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    code = payload.get("code")
    if isinstance(code, bool) or not isinstance(code, int):
        return None, _typed_failure(
            f"Feishu {operation} returned no valid business code.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    if code != 0:
        return None, _typed_failure(
            f"Feishu rejected {operation}.",
            f"feishu_{operation}_rejected",
        )
    data = payload.get("data")
    if not isinstance(data, Mapping):
        return None, _typed_failure(
            f"Feishu {operation} returned an invalid data object.",
            f"feishu_{operation}_response_invalid",
            retryable=True,
        )
    return data, None


def _bounded_feishu_json(payload: Mapping, *, max_bytes: int = 8192) -> str:
    """Keep approval/directory summaries within the Tool Ledger text bound."""
    serialized = json.dumps(payload, ensure_ascii=False, separators=(",", ":"))
    encoded = serialized.encode("utf-8")
    if len(encoded) <= max_bytes:
        return serialized
    preview = encoded[: max_bytes - 128].decode("utf-8", errors="ignore")
    while preview:
        bounded = json.dumps(
            {"truncated": True, "preview": preview},
            ensure_ascii=False,
            separators=(",", ":"),
        )
        if len(bounded.encode("utf-8")) <= max_bytes:
            return bounded
        preview = preview[: (len(preview) * 3) // 4]
    return '{"truncated":true}'


async def _feishu_user_search_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Project tenant-scoped Directory facts without exposing Provider IDs."""
    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "feishu_user_search requires query.",
            "invalid_tool_arguments",
        )
    limit = arguments.get("limit", 20)
    offset = arguments.get("offset", 0)
    if (
        isinstance(limit, bool)
        or not isinstance(limit, int)
        or not 1 <= limit <= 50
        or isinstance(offset, bool)
        or not isinstance(offset, int)
        or offset < 0
    ):
        return _typed_failure(
            "feishu_user_search requires limit 1..50 and offset >= 0.",
            "invalid_tool_arguments",
        )

    payload = await _query_directory_payload(
        agent_id,
        {
            "query": query.strip(),
            "member_type": "human",
            "include_uncontactable": False,
            "limit": limit,
            "offset": offset,
        },
    )
    if payload.get("ok") is not True:
        error = (
            payload.get("error")
            if isinstance(payload.get("error"), Mapping)
            else {}
        )
        error_code = str(error.get("code") or "query_directory_failed")
        return _typed_failure(
            "The tenant directory search could not be completed.",
            error_code,
            retryable=error_code == "query_directory_failed",
        )

    raw_members = payload.get("members", [])
    if not isinstance(raw_members, list):
        return _typed_failure(
            "The tenant directory returned an invalid member list.",
            "query_directory_failed",
            retryable=True,
        )
    members: list[dict[str, object]] = []
    for raw_member in raw_members:
        if not isinstance(raw_member, Mapping):
            continue
        provider = raw_member.get("provider")
        if (
            raw_member.get("member_type") != "human"
            or raw_member.get("can_contact") is not True
            or not isinstance(provider, Mapping)
            or _normalize_roster_provider_type(
                provider.get("provider_type")
            )
            != "feishu"
        ):
            continue
        target_member_id = raw_member.get("target_member_id")
        if not isinstance(target_member_id, str) or not target_member_id.strip():
            continue
        member: dict[str, object] = {
            "target_member_id": target_member_id.strip(),
            "display_name": str(raw_member.get("display_name") or ""),
        }
        title = raw_member.get("title")
        if isinstance(title, str) and title:
            member["title"] = title
        department = raw_member.get("department")
        if isinstance(department, Mapping):
            department_name = department.get("name")
            if isinstance(department_name, str) and department_name:
                member["department"] = {"name": department_name}
        members.append(member)

    has_more = payload.get("has_more", False)
    if not isinstance(has_more, bool):
        return _typed_failure(
            "The tenant directory returned invalid pagination facts.",
            "query_directory_failed",
            retryable=True,
        )
    summary_payload = {
        "query": query.strip(),
        "returned_count": len(members),
        "has_more": has_more,
        "members": members,
    }
    return _typed_success(
        _bounded_feishu_json(summary_payload),
        metadata={
            "returned_count": len(members),
            "has_more": has_more,
            "limit": limit,
            "offset": offset,
        },
    )


async def _feishu_approval_query_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Read one Provider page of approval instance facts."""
    import httpx

    approval_code = arguments.get("approval_code")
    if not isinstance(approval_code, str) or not approval_code.strip():
        return _typed_failure(
            "feishu_approval_query requires approval_code.",
            "invalid_tool_arguments",
        )
    instance_status = arguments.get("instance_status")
    if instance_status is not None and (
        not isinstance(instance_status, str)
        or instance_status not in _FEISHU_APPROVAL_STATUSES
    ):
        return _typed_failure(
            "feishu_approval_query instance_status is invalid.",
            "invalid_tool_arguments",
        )
    page_size = arguments.get("page_size", 20)
    page_token = arguments.get("page_token", "")
    if (
        isinstance(page_size, bool)
        or not isinstance(page_size, int)
        or not 1 <= page_size <= 100
        or not isinstance(page_token, str)
    ):
        return _typed_failure(
            "feishu_approval_query requires page_size 1..100 and a string page_token.",
            "invalid_tool_arguments",
        )

    token, token_error = await _feishu_access_token_outcome(agent_id)
    if token_error is not None or token is None:
        return token_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    body: dict[str, object] = {"approval_code": approval_code.strip()}
    if instance_status:
        body["instance_status"] = instance_status
    params: dict[str, object] = {"page_size": page_size}
    if page_token:
        params["page_token"] = page_token

    try:
        async with httpx.AsyncClient(timeout=20) as client:
            response = await client.post(
                "https://open.feishu.cn/open-apis/approval/v4/instances/query",
                headers={"Authorization": f"Bearer {token}"},
                json=body,
                params=params,
            )
    except Exception as exc:
        return _feishu_read_exception_outcome("approval_query", exc)

    data, response_error = _feishu_approval_read_response(
        response,
        "approval_query",
    )
    if response_error is not None or data is None:
        return response_error or _typed_failure(
            "Feishu approval_query returned no data.",
            "feishu_approval_query_response_invalid",
            retryable=True,
        )
    raw_instances = data.get("instance_list")
    if not isinstance(raw_instances, list):
        return _typed_failure(
            "Feishu approval_query returned an invalid instance list.",
            "feishu_approval_query_response_invalid",
            retryable=True,
        )

    instances: list[dict[str, str]] = []
    for item in raw_instances:
        instance = item.get("instance") if isinstance(item, Mapping) else None
        if not isinstance(instance, Mapping):
            return _typed_failure(
                "Feishu approval_query returned an invalid instance.",
                "feishu_approval_query_response_invalid",
                retryable=True,
            )
        instance_code = instance.get("code")
        if not isinstance(instance_code, str) or not instance_code:
            return _typed_failure(
                "Feishu approval_query returned an instance without a code.",
                "feishu_approval_query_response_invalid",
                retryable=True,
            )
        fact = {"instance_id": instance_code}
        for key in ("status", "title"):
            value = instance.get(key)
            if isinstance(value, str) and value:
                fact[key] = value
        instances.append(fact)

    has_more = data.get("has_more", False)
    returned_page_token = data.get("page_token")
    if not isinstance(has_more, bool) or (
        returned_page_token is not None
        and not isinstance(returned_page_token, str)
    ):
        return _typed_failure(
            "Feishu approval_query returned invalid pagination facts.",
            "feishu_approval_query_response_invalid",
            retryable=True,
        )
    if has_more and not returned_page_token:
        return _typed_failure(
            "Feishu approval_query omitted the next page token.",
            "feishu_approval_query_response_invalid",
            retryable=True,
        )
    summary = _bounded_feishu_json(
        {
            "returned_count": len(instances),
            "instances": instances,
        }
    )
    return _typed_success(
        summary,
        result_ref=approval_code.strip(),
        metadata={
            "instance_count": len(instances),
            "has_more": has_more,
            "page_token": returned_page_token,
        },
    )


async def _feishu_approval_get_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Read a safe instance summary or one explicitly selected section."""
    import httpx
    from urllib.parse import quote

    instance_id = arguments.get("instance_id")
    section = arguments.get("section", "summary")
    offset = arguments.get("offset", 0)
    limit = arguments.get("limit", 20)
    if not isinstance(instance_id, str) or not instance_id.strip():
        return _typed_failure(
            "feishu_approval_get requires instance_id.",
            "invalid_tool_arguments",
        )
    if not isinstance(section, str) or section not in _FEISHU_APPROVAL_SECTIONS:
        return _typed_failure(
            "feishu_approval_get section is invalid.",
            "invalid_tool_arguments",
        )
    if (
        isinstance(offset, bool)
        or not isinstance(offset, int)
        or offset < 0
        or isinstance(limit, bool)
        or not isinstance(limit, int)
        or not 1 <= limit <= 50
    ):
        return _typed_failure(
            "feishu_approval_get requires offset >= 0 and limit 1..50.",
            "invalid_tool_arguments",
        )

    token, token_error = await _feishu_access_token_outcome(agent_id)
    if token_error is not None or token is None:
        return token_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    stable_instance_id = instance_id.strip()
    try:
        async with httpx.AsyncClient(timeout=20) as client:
            response = await client.get(
                "https://open.feishu.cn/open-apis/approval/v4/instances/"
                + quote(stable_instance_id, safe=""),
                headers={"Authorization": f"Bearer {token}"},
            )
    except Exception as exc:
        return _feishu_read_exception_outcome("approval_get", exc)

    data, response_error = _feishu_approval_read_response(
        response,
        "approval_get",
    )
    if response_error is not None or data is None:
        return response_error or _typed_failure(
            "Feishu approval_get returned no data.",
            "feishu_approval_get_response_invalid",
            retryable=True,
        )

    if section == "summary":
        summary_fields: dict[str, object] = {}
        for key in (
            "approval_name",
            "approval_code",
            "status",
            "serial_number",
            "title",
            "start_time",
            "end_time",
        ):
            value = data.get(key)
            if isinstance(value, (str, int, float, bool)) and not isinstance(
                value,
                complex,
            ):
                summary_fields[key] = value
        return _typed_success(
            _bounded_feishu_json(
                {
                    "instance_id": stable_instance_id,
                    "summary": summary_fields,
                }
            ),
            result_ref=stable_instance_id,
            metadata={"section": "summary"},
        )

    provider_key = _FEISHU_APPROVAL_SECTION_KEYS[section]
    raw_section = data.get(provider_key, [])
    if section == "form" and isinstance(raw_section, str):
        try:
            raw_section = json.loads(raw_section)
        except (TypeError, ValueError):
            return _typed_failure(
                "Feishu approval_get returned an invalid form section.",
                "feishu_approval_get_response_invalid",
                retryable=True,
            )
    if not isinstance(raw_section, list):
        return _typed_failure(
            f"Feishu approval_get returned an invalid {section} section.",
            "feishu_approval_get_response_invalid",
            retryable=True,
        )
    selected = raw_section[offset : offset + limit]
    next_offset = offset + len(selected)
    has_more = next_offset < len(raw_section)
    return _typed_success(
        _bounded_feishu_json(
            {
                "instance_id": stable_instance_id,
                "section": section,
                "offset": offset,
                "returned_count": len(selected),
                "items": selected,
            }
        ),
        result_ref=stable_instance_id,
        metadata={
            "section": section,
            "offset": offset,
            "returned_count": len(selected),
            "has_more": has_more,
            "next_offset": next_offset if has_more else None,
        },
    )


async def _feishu_approval_create_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Hidden external-write adapter retained behind the future confirmation gate."""
    import httpx

    approval_code = arguments.get("approval_code")
    target_member_id = arguments.get("target_member_id")
    form_data = arguments.get("form_data")
    if not (
        isinstance(approval_code, str)
        and approval_code.strip()
        and isinstance(target_member_id, str)
        and target_member_id.strip()
        and isinstance(form_data, str)
        and form_data.strip()
    ):
        return _typed_failure(
            "feishu_approval_create requires approval_code, target_member_id, and form_data.",
            "invalid_tool_arguments",
        )
    try:
        parsed_form = json.loads(form_data)
    except (TypeError, ValueError):
        return _typed_failure(
            "feishu_approval_create form_data must be a JSON array.",
            "invalid_tool_arguments",
        )
    if not isinstance(parsed_form, list):
        return _typed_failure(
            "feishu_approval_create form_data must be a JSON array.",
            "invalid_tool_arguments",
        )

    try:
        async with async_session() as db:
            target, target_error = await _resolve_roster_human_target(
                db,
                agent_id,
                target_member_id=target_member_id.strip(),
                provider_type="feishu",
                require_provider_identity=True,
            )
    except Exception as exc:
        return _typed_failure(
            f"Feishu approval target resolution failed: {type(exc).__name__}.",
            "feishu_approval_target_resolution_failed",
        )
    if target is None or target_error is not None:
        return _typed_failure(
            "The requested Feishu approval applicant is unavailable.",
            "feishu_approval_target_unavailable",
        )
    if _normalize_roster_provider_type(target.provider_type) != "feishu":
        return _typed_failure(
            "The approval applicant is not a Feishu member.",
            "feishu_approval_target_provider_mismatch",
        )
    provider_user_id = str(
        getattr(target.member, "external_id", "") or ""
    ).strip()
    if not provider_user_id:
        return _typed_failure(
            "The approval applicant has no Feishu user_id.",
            "feishu_approval_target_identity_missing",
        )

    token, token_error = await _feishu_access_token_outcome(agent_id)
    if token_error is not None or token is None:
        return token_error or _typed_failure(
            "Feishu credentials are unavailable.",
            "feishu_channel_not_configured",
        )
    try:
        async with httpx.AsyncClient(timeout=20) as client:
            response = await client.post(
                "https://open.feishu.cn/open-apis/approval/v4/instances",
                headers={"Authorization": f"Bearer {token}"},
                json={
                    "approval_code": approval_code.strip(),
                    "user_id": provider_user_id,
                    "form": form_data,
                },
            )
    except Exception as exc:
        return _feishu_write_exception_outcome(
            "approval_create",
            exc,
        )

    status_code = getattr(response, "status_code", None)
    if not isinstance(status_code, int) or isinstance(status_code, bool):
        return _typed_unknown(
            "Feishu approval_create returned no readable HTTP receipt; reconcile before retrying.",
            "feishu_approval_create_outcome_unknown",
        )
    if status_code == 429 or status_code >= 500:
        return _typed_unknown(
            "Feishu approval_create may have taken effect; reconcile before retrying.",
            "feishu_approval_create_outcome_unknown",
        )
    if 400 <= status_code < 500:
        return _typed_failure(
            "Feishu rejected approval_create.",
            "feishu_approval_create_rejected",
        )
    try:
        payload = response.json()
    except Exception:
        return _typed_unknown(
            "Feishu approval_create returned an unreadable receipt; reconcile before retrying.",
            "feishu_approval_create_outcome_unknown",
        )
    if not isinstance(payload, Mapping):
        return _typed_unknown(
            "Feishu approval_create returned an invalid receipt; reconcile before retrying.",
            "feishu_approval_create_outcome_unknown",
        )
    code = payload.get("code")
    if isinstance(code, bool) or not isinstance(code, int):
        return _typed_unknown(
            "Feishu approval_create returned no business receipt; reconcile before retrying.",
            "feishu_approval_create_outcome_unknown",
        )
    if code != 0:
        return _typed_failure(
            "Feishu rejected approval_create.",
            "feishu_approval_create_rejected",
        )
    data = payload.get("data")
    instance_code = (
        str(data.get("instance_code") or "").strip()
        if isinstance(data, Mapping)
        else ""
    )
    if not instance_code:
        return _typed_unknown(
            "Feishu accepted approval_create but returned no instance receipt; reconcile before retrying.",
            "feishu_approval_create_receipt_missing",
        )
    return _typed_success(
        f"Feishu approval instance {instance_code} was created.",
        result_ref=instance_code,
    )


async def _feishu_approval_create(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter; Durable Runtime keeps this write hidden."""
    outcome = await _feishu_approval_create_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Feishu approval creation returned no summary.",
    )


async def _feishu_approval_query(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter for the typed approval page read."""
    outcome = await _feishu_approval_query_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Feishu approval query returned no summary.",
    )


async def _feishu_approval_get(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter for the typed approval instance read."""
    outcome = await _feishu_approval_get_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Feishu approval read returned no summary.",
    )


# ─── Feishu User Search ───────────────────────────────────────────────────────

async def _feishu_user_search(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter for the stable-ID directory projection."""
    canonical_arguments = dict(arguments)
    if "query" not in canonical_arguments and isinstance(
        canonical_arguments.get("name"),
        str,
    ):
        canonical_arguments["query"] = canonical_arguments["name"]
    outcome = await _feishu_user_search_outcome(
        agent_id,
        canonical_arguments,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Feishu user search returned no summary.",
    )


_NATIVE_FEISHU_USER_SEARCH_ADAPTER = _feishu_user_search


async def _feishu_open_id_for_visible_name(
    agent_id: uuid.UUID,
    name: str,
) -> str | None:
    """Resolve one visible Feishu human privately for legacy attendee APIs."""
    normalized_name = name.strip()
    if not normalized_name:
        return None
    # Calendar F1 tests and old extension points replace the legacy adapter.
    # Preserve that narrow injection seam without making raw IDs part of the
    # production user-search result again.
    if _feishu_user_search is not _NATIVE_FEISHU_USER_SEARCH_ADAPTER:
        legacy_result = await _feishu_user_search(
            agent_id,
            {"name": normalized_name, "query": normalized_name},
        )
        match = re.search(
            r"open_id:\s*`(ou_[A-Za-z0-9]+)`",
            str(legacy_result),
        )
        return match.group(1) if match is not None else None

    payload = await _query_directory_payload(
        agent_id,
        {
            "query": normalized_name,
            "member_type": "human",
            "include_uncontactable": False,
            "limit": 20,
            "offset": 0,
        },
    )
    raw_members = payload.get("members") if payload.get("ok") is True else None
    if not isinstance(raw_members, list):
        return None
    exact_open_ids: list[str] = []
    for member in raw_members:
        provider = member.get("provider") if isinstance(member, Mapping) else None
        if (
            not isinstance(member, Mapping)
            or member.get("member_type") != "human"
            or member.get("can_contact") is not True
            or str(member.get("display_name") or "").casefold()
            != normalized_name.casefold()
            or not isinstance(provider, Mapping)
            or _normalize_roster_provider_type(provider.get("provider_type"))
            != "feishu"
        ):
            continue
        open_id = provider.get("open_id")
        if isinstance(open_id, str) and open_id and open_id not in exact_open_ids:
            exact_open_ids.append(open_id)
    return exact_open_ids[0] if len(exact_open_ids) == 1 else None


async def _feishu_contacts_refresh(agent_id: uuid.UUID) -> None:
    """Force-clear the local contacts cache so next search re-fetches from API."""
    import pathlib as _pl
    _cache_file = _pl.Path("/data/workspaces") / str(agent_id) / "feishu_contacts_cache.json"
    try:
        if _cache_file.exists():
            _cache_file.unlink()
    except Exception:
        pass


# ─── Email Tool Helpers ─────────────────────────────────────

async def _get_email_config(agent_id: uuid.UUID) -> dict:
    """Retrieve per-agent email config from the send_email tool's AgentTool config."""
    from app.models.tool import Tool, AgentTool

    async with async_session() as db:
        # Find the send_email tool
        r = await db.execute(select(Tool).where(Tool.name == "send_email"))
        tool = r.scalar_one_or_none()
        if not tool:
            return {}

        # Get per-agent config
        at_r = await db.execute(
            select(AgentTool).where(
                AgentTool.agent_id == agent_id,
                AgentTool.tool_id == tool.id,
            )
        )
        at = at_r.scalar_one_or_none()
        agent_config = (at.config or {}) if at else {}
        merged = {**(tool.config or {}), **agent_config}
        return _decrypt_sensitive_fields(merged, tool.config_schema)


def _resolve_local_email_configuration(
    config: object,
) -> tuple[dict | None, frozenset[str]]:
    """Resolve Email presets and local protocol readiness without provider I/O."""
    from app.services import email_service

    if not isinstance(config, Mapping):
        return None, frozenset()
    try:
        resolved = email_service.resolve_config(dict(config))
    except (TypeError, ValueError):
        return None, frozenset()

    address = resolved.get("email_address")
    password = resolved.get("auth_code")
    if not (
        isinstance(address, str)
        and address.strip()
        and isinstance(password, str)
        and password.strip()
    ):
        return None, frozenset()

    def endpoint_ready(host_key: str, port_key: str) -> bool:
        host = resolved.get(host_key)
        port = resolved.get(port_key)
        return (
            isinstance(host, str)
            and bool(host.strip())
            and isinstance(port, int)
            and not isinstance(port, bool)
            and 1 <= port <= 65535
        )

    protocols: set[str] = set()
    if endpoint_ready("imap_host", "imap_port"):
        protocols.add("imap")
    if endpoint_ready("smtp_host", "smtp_port"):
        protocols.add("smtp")
    return resolved, frozenset(protocols)


class _EmailIMAPRejected(RuntimeError):
    def __init__(self, stage: str) -> None:
        self.stage = stage
        super().__init__(stage)


class _EmailIMAPMalformed(RuntimeError):
    def __init__(self, stage: str) -> None:
        self.stage = stage
        super().__init__(stage)


def _checked_email_imap_status(response: object, stage: str) -> object:
    if not isinstance(response, (tuple, list)) or len(response) != 2:
        raise _EmailIMAPMalformed(stage)
    status, payload = response
    if isinstance(status, bytes):
        try:
            status = status.decode("ascii")
        except UnicodeDecodeError as exc:
            raise _EmailIMAPMalformed(stage) from exc
    if not isinstance(status, str):
        raise _EmailIMAPMalformed(stage)
    if status.upper() != "OK":
        raise _EmailIMAPRejected(stage)
    return payload


async def _read_emails_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Read IMAP messages using explicit status facts at every provider stage."""
    import socket

    from app.services import email_service

    limit = arguments.get("limit", 10)
    if isinstance(limit, bool) or not isinstance(limit, int) or not 1 <= limit <= 30:
        return _typed_failure(
            "read_emails limit must be an integer from 1 through 30.",
            "invalid_tool_arguments",
        )
    folder = arguments.get("folder", "INBOX")
    if not isinstance(folder, str) or not folder.strip():
        return _typed_failure(
            "read_emails folder must be a non-empty string.",
            "invalid_tool_arguments",
        )
    folder = folder.strip()
    search = arguments.get("search")
    if search is not None and (
        not isinstance(search, str) or not search.strip()
    ):
        return _typed_failure(
            "read_emails search must be a non-empty string when supplied.",
            "invalid_tool_arguments",
        )
    search_criteria = search.strip() if isinstance(search, str) else "ALL"

    try:
        stored_config = await _get_email_config(agent_id)
    except Exception as exc:
        return _typed_failure(
            f"Email configuration could not be read: {type(exc).__name__}.",
            "email_configuration_unavailable",
        )
    config, protocols = _resolve_local_email_configuration(stored_config)
    if config is None or "imap" not in protocols:
        return _typed_failure(
            "read_emails requires complete local Email and IMAP configuration.",
            "email_imap_not_configured",
        )

    def read_mailbox() -> list[dict[str, str]]:
        with email_service.force_ipv4():
            ssl_context = email_service.ssl.create_default_context()
            with email_service.imaplib.IMAP4_SSL(
                config["imap_host"],
                config["imap_port"],
                ssl_context=ssl_context,
            ) as mailbox:
                login_payload = _checked_email_imap_status(
                    mailbox.login(
                        config["email_address"],
                        config["auth_code"],
                    ),
                    "login",
                )
                if not isinstance(login_payload, (tuple, list)):
                    raise _EmailIMAPMalformed("login")

                select_payload = _checked_email_imap_status(
                    mailbox.select(folder, readonly=True),
                    "select",
                )
                if not isinstance(select_payload, (tuple, list)):
                    raise _EmailIMAPMalformed("select")

                search_payload = _checked_email_imap_status(
                    mailbox.search(None, search_criteria),
                    "search",
                )
                if not isinstance(search_payload, (tuple, list)) or not search_payload:
                    raise _EmailIMAPMalformed("search")
                packed_ids = search_payload[0]
                if not isinstance(packed_ids, (bytes, str)):
                    raise _EmailIMAPMalformed("search")
                message_ids = packed_ids.split()
                if not message_ids:
                    return []

                selected_ids = list(reversed(message_ids[-limit:]))
                messages: list[dict[str, str]] = []
                for message_number in selected_ids:
                    fetch_payload = _checked_email_imap_status(
                        mailbox.fetch(message_number, "(RFC822)"),
                        "fetch",
                    )
                    if not isinstance(fetch_payload, (tuple, list)):
                        raise _EmailIMAPMalformed("fetch")
                    raw_message: bytes | None = None
                    for item in fetch_payload:
                        if (
                            isinstance(item, (tuple, list))
                            and len(item) >= 2
                            and isinstance(item[1], bytes)
                        ):
                            raw_message = item[1]
                            break
                    if raw_message is None:
                        raise _EmailIMAPMalformed("fetch")
                    try:
                        parsed = email_service.email_lib.message_from_bytes(
                            raw_message
                        )
                        body = email_service._extract_body(parsed)
                        if len(body) > 500:
                            body = body[:500] + "..."
                        messages.append(
                            {
                                "from": email_service._decode_header_value(
                                    parsed.get("From", "")
                                ),
                                "subject": email_service._decode_header_value(
                                    parsed.get("Subject", "(No subject)")
                                ),
                                "date": str(parsed.get("Date", "")),
                                "message_id": str(
                                    parsed.get("Message-ID", "")
                                ),
                                "body": body,
                            }
                        )
                    except Exception as exc:
                        raise _EmailIMAPMalformed("message_parse") from exc
                return messages

    try:
        messages = await asyncio.to_thread(read_mailbox)
    except _EmailIMAPRejected as exc:
        return _typed_failure(
            f"IMAP rejected the {exc.stage} operation.",
            f"email_imap_{exc.stage}_rejected",
        )
    except _EmailIMAPMalformed as exc:
        return _typed_failure(
            f"IMAP returned a malformed {exc.stage} response.",
            f"email_imap_{exc.stage}_response_invalid",
            retryable=True,
        )
    except email_service.imaplib.IMAP4.abort:
        return _typed_failure(
            "IMAP disconnected before the read completed.",
            "email_imap_transport_failed",
            retryable=True,
        )
    except email_service.imaplib.IMAP4.error as exc:
        error_text = str(exc).upper()
        if "AUTH" in error_text or "LOGIN" in error_text:
            return _typed_failure(
                "IMAP authentication failed.",
                "email_imap_authentication_failed",
            )
        return _typed_failure(
            "IMAP rejected the mailbox read.",
            "email_imap_rejected",
        )
    except (socket.timeout, ConnectionError, OSError) as exc:
        return _typed_failure(
            f"IMAP transport failed before the read completed: "
            f"{type(exc).__name__}.",
            "email_imap_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"IMAP read failed before a reliable result was parsed: "
            f"{type(exc).__name__}.",
            "email_imap_read_failed",
            retryable=True,
        )

    if not messages:
        return _typed_success(
            f"No emails found in {folder}.",
            result_ref=folder,
        )
    lines = [f"{len(messages)} email(s) from {folder}:"]
    for message in messages:
        lines.extend(
            [
                "---",
                f"From: {message['from']}",
                f"Subject: {message['subject']}",
                f"Date: {message['date']}",
                f"Message-ID: {message['message_id']}",
                f"Body:\n{message['body']}",
            ]
        )
    return _typed_success("\n".join(lines), result_ref=folder)


class _EmailSMTPDispatchError(RuntimeError):
    """Preserve whether SMTP DATA may have started without exposing secrets."""

    def __init__(self, cause: Exception, *, data_started: bool) -> None:
        self.cause = cause
        self.data_started = data_started
        super().__init__(type(cause).__name__)


class _EmailOriginalNotFound(RuntimeError):
    pass


class _EmailOriginalSenderInvalid(RuntimeError):
    pass


def _email_recipient_list(value: object) -> list[str] | None:
    if not isinstance(value, str) or not value.strip():
        return None
    recipients: list[str] = []
    seen: set[str] = set()
    for candidate in value.split(","):
        recipient = candidate.strip()
        if not recipient or "\r" in recipient or "\n" in recipient:
            return None
        key = recipient.casefold()
        if key not in seen:
            seen.add(key)
            recipients.append(recipient)
    return recipients or None


async def _email_attachment_payloads(
    agent_id: uuid.UUID,
    attachments: object,
) -> tuple[list[tuple[str, bytes]] | None, ToolExecutionOutcome | None]:
    """Read every attachment before opening SMTP so preflight is atomic."""
    if attachments is None:
        return [], None
    if not isinstance(attachments, list) or any(
        not isinstance(path, str) or not path.strip()
        for path in attachments
    ):
        return None, _typed_failure(
            "send_email attachments must be a list of non-empty paths.",
            "invalid_tool_arguments",
        )

    storage = get_storage_backend()
    tenant_id = await _get_agent_tenant_id(agent_id)
    workspace = _agent_workspace_root(agent_id)
    payloads: list[tuple[str, bytes]] = []
    total_bytes = 0
    for path in attachments:
        try:
            storage_key, normalized_path, _ = _tool_storage_key(
                agent_id,
                path,
                tenant_id,
            )
            file_bytes: bytes | None = None
            if await storage.exists(storage_key) and await storage.is_file(
                storage_key
            ):
                file_bytes = await storage.read_bytes(storage_key)
            if file_bytes is None:
                local_path = _resolve_tool_source_path(
                    workspace,
                    path,
                    tenant_id,
                )
                if local_path.exists() and local_path.is_file():
                    file_bytes = await asyncio.to_thread(local_path.read_bytes)
            if file_bytes is None:
                return None, _typed_failure(
                    "An email attachment was not found.",
                    "email_attachment_not_found",
                )
            if len(file_bytes) > TOOL_MATERIALIZE_MAX_FILE_BYTES:
                return None, _typed_failure(
                    "An email attachment exceeds the per-file size limit.",
                    "email_attachment_too_large",
                )
            total_bytes += len(file_bytes)
            if total_bytes > TOOL_MATERIALIZE_MAX_TOTAL_BYTES:
                return None, _typed_failure(
                    "Email attachments exceed the total size limit.",
                    "email_attachments_too_large",
                )
            payloads.append((Path(normalized_path).name, file_bytes))
        except (TypeError, ValueError):
            return None, _typed_failure(
                "An email attachment path is invalid.",
                "email_attachment_path_invalid",
            )
        except Exception as exc:
            return None, _typed_failure(
                "Email attachment preflight failed: "
                f"{type(exc).__name__}.",
                "email_attachment_preflight_failed",
                retryable=True,
            )
    return payloads, None


def _email_reply_source(
    config: dict,
    *,
    message_id: str,
    folder: str,
) -> tuple[str, str]:
    from app.services import email_service

    with email_service.force_ipv4():
        ssl_context = email_service.ssl.create_default_context()
        with email_service.imaplib.IMAP4_SSL(
            config["imap_host"],
            config["imap_port"],
            ssl_context=ssl_context,
        ) as mailbox:
            login_payload = _checked_email_imap_status(
                mailbox.login(
                    config["email_address"],
                    config["auth_code"],
                ),
                "login",
            )
            if not isinstance(login_payload, (tuple, list)):
                raise _EmailIMAPMalformed("login")
            select_payload = _checked_email_imap_status(
                mailbox.select(folder, readonly=True),
                "select",
            )
            if not isinstance(select_payload, (tuple, list)):
                raise _EmailIMAPMalformed("select")
            escaped_message_id = message_id.replace("\\", "\\\\").replace(
                '"',
                '\\"',
            )
            search_payload = _checked_email_imap_status(
                mailbox.search(
                    None,
                    f'HEADER Message-ID "{escaped_message_id}"',
                ),
                "search",
            )
            if not isinstance(search_payload, (tuple, list)) or not search_payload:
                raise _EmailIMAPMalformed("search")
            packed_ids = search_payload[0]
            if not isinstance(packed_ids, (bytes, str)):
                raise _EmailIMAPMalformed("search")
            message_numbers = packed_ids.split()
            if not message_numbers:
                raise _EmailOriginalNotFound
            fetch_payload = _checked_email_imap_status(
                mailbox.fetch(message_numbers[0], "(RFC822)"),
                "fetch",
            )
            if not isinstance(fetch_payload, (tuple, list)):
                raise _EmailIMAPMalformed("fetch")
            raw_message: bytes | None = None
            for item in fetch_payload:
                if (
                    isinstance(item, (tuple, list))
                    and len(item) >= 2
                    and isinstance(item[1], bytes)
                ):
                    raw_message = item[1]
                    break
            if raw_message is None:
                raise _EmailIMAPMalformed("fetch")
            try:
                original = email_service.email_lib.message_from_bytes(raw_message)
                sender = email_service.parseaddr(original.get("From", ""))[1]
                subject = email_service._decode_header_value(
                    original.get("Subject", "")
                )
            except Exception as exc:
                raise _EmailIMAPMalformed("message_parse") from exc
            if not sender or "\r" in sender or "\n" in sender:
                raise _EmailOriginalSenderInvalid
            return sender, subject


def _email_message(
    config: dict,
    *,
    recipients: list[str],
    subject: str,
    body: str,
    message_id: str,
    cc_recipients: list[str] | None = None,
    attachments: list[tuple[str, bytes]] | None = None,
    reply_to_message_id: str | None = None,
):
    from app.services import email_service

    message = email_service.MIMEMultipart()
    message["From"] = config["email_address"]
    message["To"] = ", ".join(recipients)
    message["Subject"] = subject
    message["Message-ID"] = message_id
    message["Date"] = email_service.datetime.now().strftime(
        "%a, %d %b %Y %H:%M:%S %z"
    )
    if cc_recipients:
        message["Cc"] = ", ".join(cc_recipients)
    if reply_to_message_id:
        message["In-Reply-To"] = reply_to_message_id
        message["References"] = reply_to_message_id
    message.attach(email_service.MIMEText(body, "plain", "utf-8"))
    for filename, file_bytes in attachments or []:
        part = email_service.MIMEBase("application", "octet-stream")
        part.set_payload(file_bytes)
        email_service.encoders.encode_base64(part)
        part.add_header(
            "Content-Disposition",
            "attachment",
            filename=filename,
        )
        message.attach(part)
    return message


def _send_email_message(
    config: dict,
    *,
    recipients: list[str],
    message,
) -> Mapping[str, object]:
    """Call ``sendmail`` exactly once and return its recipient receipt."""
    from app.services import email_service

    data_started = False
    try:
        with email_service.force_ipv4():
            if config.get("smtp_ssl", True):
                ssl_context = email_service.ssl.create_default_context()
                with email_service.smtplib.SMTP_SSL(
                    config["smtp_host"],
                    config["smtp_port"],
                    context=ssl_context,
                    timeout=15,
                ) as server:
                    server.login(
                        config["email_address"],
                        config["auth_code"],
                    )
                    data_started = True
                    receipt = server.sendmail(
                        config["email_address"],
                        recipients,
                        message.as_string(),
                    )
            else:
                with email_service.smtplib.SMTP(
                    config["smtp_host"],
                    config["smtp_port"],
                    timeout=15,
                ) as server:
                    server.ehlo()
                    if "starttls" in server.esmtp_features:
                        server.starttls(
                            context=email_service.ssl.create_default_context()
                        )
                        server.ehlo()
                    if (
                        config["email_address"] or config["auth_code"]
                    ) and "auth" in server.esmtp_features:
                        server.login(
                            config["email_address"],
                            config["auth_code"],
                        )
                    data_started = True
                    receipt = server.sendmail(
                        config["email_address"],
                        recipients,
                        message.as_string(),
                    )
    except Exception as exc:
        raise _EmailSMTPDispatchError(
            exc,
            data_started=data_started,
        ) from exc
    if not isinstance(receipt, Mapping):
        raise _EmailSMTPDispatchError(
            TypeError("SMTP sendmail receipt was not a mapping"),
            data_started=True,
        )
    return receipt


def _email_receipt_metadata(
    message_id: str,
    recipients: list[str],
    refused: Mapping[object, object],
) -> tuple[list[str], list[str], dict[str, object]]:
    refused_recipients = [str(recipient) for recipient in refused]
    refused_keys = {recipient.casefold() for recipient in refused_recipients}
    accepted_recipients = [
        recipient
        for recipient in recipients
        if recipient.casefold() not in refused_keys
    ]
    return (
        accepted_recipients,
        refused_recipients,
        {
            "message_id": message_id,
            "accepted_recipients": accepted_recipients,
            "refused_recipients": refused_recipients,
        },
    )


async def _email_write_outcome(
    tool_name: str,
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Return SMTP provider facts without inferring success from display text."""
    import socket

    from app.services import email_service

    body = arguments.get("body")
    if not isinstance(body, str) or not body.strip():
        return _typed_failure(
            f"{tool_name} requires a non-empty body.",
            "invalid_tool_arguments",
        )

    try:
        stored_config = await _get_email_config(agent_id)
    except Exception as exc:
        return _typed_failure(
            f"Email configuration could not be read: {type(exc).__name__}.",
            "email_configuration_unavailable",
        )
    config, protocols = _resolve_local_email_configuration(stored_config)
    required_protocols = {"smtp", "imap"} if tool_name == "reply_email" else {"smtp"}
    if config is None or not required_protocols.issubset(protocols):
        return _typed_failure(
            f"{tool_name} requires complete local Email configuration.",
            "email_not_configured",
        )

    reply_to_message_id: str | None = None
    attachments: list[tuple[str, bytes]] = []
    cc_recipients: list[str] = []
    if tool_name == "send_email":
        recipients = _email_recipient_list(arguments.get("to"))
        subject = arguments.get("subject")
        cc_value = arguments.get("cc")
        if recipients is None:
            return _typed_failure(
                "send_email requires valid recipients.",
                "invalid_tool_arguments",
            )
        if (
            not isinstance(subject, str)
            or not subject.strip()
            or "\r" in subject
            or "\n" in subject
        ):
            return _typed_failure(
                "send_email requires a valid non-empty subject.",
                "invalid_tool_arguments",
            )
        if cc_value is not None:
            parsed_cc = _email_recipient_list(cc_value)
            if parsed_cc is None:
                return _typed_failure(
                    "send_email cc must contain valid recipients.",
                    "invalid_tool_arguments",
                )
            cc_recipients = parsed_cc
        attachments_result, attachment_failure = await _email_attachment_payloads(
            agent_id,
            arguments.get("attachments"),
        )
        if attachment_failure is not None:
            return attachment_failure
        attachments = attachments_result or []
        subject = subject.strip()
        envelope_recipients = list(
            dict.fromkeys([*recipients, *cc_recipients])
        )
    else:
        reply_to_message_id = arguments.get("message_id")
        folder = arguments.get("folder", "INBOX")
        if (
            not isinstance(reply_to_message_id, str)
            or not reply_to_message_id.strip()
            or "\r" in reply_to_message_id
            or "\n" in reply_to_message_id
        ):
            return _typed_failure(
                "reply_email requires a valid message_id.",
                "invalid_tool_arguments",
            )
        if not isinstance(folder, str) or not folder.strip():
            return _typed_failure(
                "reply_email folder must be a non-empty string.",
                "invalid_tool_arguments",
            )
        reply_to_message_id = reply_to_message_id.strip()
        try:
            sender, original_subject = await asyncio.to_thread(
                _email_reply_source,
                config,
                message_id=reply_to_message_id,
                folder=folder.strip(),
            )
        except _EmailOriginalNotFound:
            return _typed_failure(
                "The original email was not found in the requested folder.",
                "email_original_not_found",
            )
        except _EmailOriginalSenderInvalid:
            return _typed_failure(
                "The original email has no valid reply address.",
                "email_original_sender_invalid",
            )
        except _EmailIMAPRejected as exc:
            return _typed_failure(
                f"IMAP rejected the {exc.stage} operation.",
                f"email_imap_{exc.stage}_rejected",
            )
        except _EmailIMAPMalformed as exc:
            return _typed_failure(
                f"IMAP returned a malformed {exc.stage} response.",
                f"email_imap_{exc.stage}_response_invalid",
            )
        except email_service.imaplib.IMAP4.error as exc:
            error_text = str(exc).upper()
            error_code = (
                "email_imap_authentication_failed"
                if "AUTH" in error_text or "LOGIN" in error_text
                else "email_imap_rejected"
            )
            return _typed_failure("IMAP reply preflight failed.", error_code)
        except (socket.timeout, ConnectionError, OSError) as exc:
            return _typed_failure(
                "IMAP reply preflight transport failed: "
                f"{type(exc).__name__}.",
                "email_imap_transport_failed",
                retryable=True,
            )
        except Exception as exc:
            return _typed_failure(
                "IMAP reply preflight failed: "
                f"{type(exc).__name__}.",
                "email_imap_reply_preflight_failed",
            )
        recipients = [sender]
        envelope_recipients = recipients
        normalized_subject = original_subject.strip() or "(No subject)"
        subject = (
            normalized_subject
            if normalized_subject.casefold().startswith("re:")
            else f"Re: {normalized_subject}"
        )

    message_id = email_service.make_msgid()
    message = _email_message(
        config,
        recipients=recipients,
        subject=subject,
        body=body,
        message_id=message_id,
        cc_recipients=cc_recipients,
        attachments=attachments,
        reply_to_message_id=reply_to_message_id,
    )
    try:
        refused = await asyncio.to_thread(
            _send_email_message,
            config,
            recipients=envelope_recipients,
            message=message,
        )
    except _EmailSMTPDispatchError as exc:
        cause = exc.cause
        if isinstance(cause, email_service.smtplib.SMTPRecipientsRefused):
            refused = cause.recipients
            accepted, refused_names, metadata = _email_receipt_metadata(
                message_id,
                envelope_recipients,
                refused,
            )
            return _typed_failure(
                f"SMTP refused all {len(refused_names)} recipient(s).",
                "email_smtp_all_recipients_refused",
                result_ref=message_id,
                metadata=metadata,
            )
        if isinstance(cause, email_service.smtplib.SMTPAuthenticationError):
            return _typed_failure(
                "SMTP authentication failed before message submission.",
                "email_smtp_authentication_failed",
            )
        if exc.data_started:
            return _typed_unknown(
                "SMTP submission outcome is unknown; reconcile by Message-ID "
                "before any retry.",
                "email_smtp_submission_unknown",
                result_ref=message_id,
                metadata={"message_id": message_id},
            )
        retryable = isinstance(cause, (socket.timeout, ConnectionError, OSError))
        return _typed_failure(
            f"SMTP failed before message submission: {type(cause).__name__}.",
            "email_smtp_preflight_failed",
            retryable=retryable,
        )

    accepted, refused_names, metadata = _email_receipt_metadata(
        message_id,
        envelope_recipients,
        refused,
    )
    if not refused_names:
        return _typed_success(
            f"Email accepted for {len(accepted)} recipient(s).",
            result_ref=message_id,
            metadata=metadata,
        )
    if not accepted:
        return _typed_failure(
            f"SMTP refused all {len(refused_names)} recipient(s).",
            "email_smtp_all_recipients_refused",
            result_ref=message_id,
            metadata=metadata,
        )
    return _typed_unknown(
        f"SMTP accepted {len(accepted)} recipient(s) and refused "
        f"{len(refused_names)} recipient(s); reconcile before retrying.",
        "email_smtp_partial_acceptance",
        result_ref=message_id,
        metadata=metadata,
    )


# ── Pages: public HTML hosting ──────────────────────────

async def _publish_page_outcome(
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    ws: Path,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Publish an HTML file as a public page."""
    import secrets
    import re

    path = arguments.get("path", "")
    if not path:
        return _typed_failure(
            "Missing required argument 'path'.",
            "invalid_tool_arguments",
        )

    # Validate file extension
    if not path.lower().endswith((".html", ".htm")):
        return _typed_failure(
            "Only .html and .htm files can be published.",
            "published_page_format_invalid",
        )

    # Resolve via storage backend (supports local FS and S3)
    try:
        storage = get_storage_backend()
        storage_key, normalized_path, is_enterprise = _tool_storage_key(
            agent_id,
            path,
        )
        if is_enterprise:
            return _typed_failure(
                "Shared enterprise files cannot be published as Agent pages.",
                "published_page_source_forbidden",
            )
        source_exists = await storage.exists(storage_key)
        source_is_file = source_exists and await storage.is_file(storage_key)
    except Exception as exc:
        return _typed_failure(
            f"Published page source could not be checked: {type(exc).__name__}.",
            "published_page_source_check_failed",
        )
    if not source_is_file:
        return _typed_failure(
            f"File not found: {path}",
            "published_page_source_not_found",
        )
    path = normalized_path

    # Extract title from HTML
    try:
        content = await storage.read_text(storage_key, encoding="utf-8", errors="replace")
        title_match = re.search(r"<title[^>]*>(.*?)</title>", content, re.IGNORECASE | re.DOTALL)
        title = title_match.group(1).strip()[:200] if title_match else Path(path).stem
    except Exception:
        title = Path(path).stem

    # Generate short_id
    short_id = secrets.token_urlsafe(6)[:8]  # 8-char URL-safe string

    # Look up tenant_id
    tenant_id = None
    try:
        from app.models.agent import Agent as _AgModel
        async with async_session() as _db:
            _r = await _db.execute(select(_AgModel.tenant_id).where(_AgModel.id == agent_id))
            tenant_id = _r.scalar_one_or_none()
    except Exception:
        pass

    # Create record
    from app.models.published_page import PublishedPage
    commit_started = False
    try:
        async with async_session() as db:
            page = PublishedPage(
                short_id=short_id,
                agent_id=agent_id,
                user_id=user_id,
                tenant_id=tenant_id,
                source_path=path,
                title=title,
            )
            db.add(page)
            commit_started = True
            await db.commit()
    except Exception as e:
        if commit_started:
            return _typed_unknown(
                "Published page commit outcome is unknown; reconcile before retrying.",
                "published_page_outcome_unknown",
            )
        return _typed_failure(
            f"Page could not be prepared for publishing: {type(e).__name__}.",
            "published_page_failed",
        )

    # Build public URL from the same settings loader used by the app. Reading
    # os.environ directly misses values that come from the local .env file.
    try:
        from app.config import get_settings as _get_publish_settings
        public_base = (_get_publish_settings().PUBLIC_BASE_URL or os.environ.get("PUBLIC_BASE_URL", "")).rstrip("/")
    except Exception:
        public_base = os.environ.get("PUBLIC_BASE_URL", "").rstrip("/")
    public_base_error = False
    if public_base:
        validated_base, validation_error = await _validate_public_http_url(
            public_base
        )
        if validation_error or not validated_base:
            public_base = ""
            public_base_error = True
        else:
            public_base = validated_base.rstrip("/")
    if not public_base:
        # Relative path works inside the same deployment; include a note so
        # the user can configure PUBLIC_BASE_URL for a fully-qualified link.
        url = f"/p/{short_id}"
        url_note = (
            "\n\n> Note: PUBLIC_BASE_URL is not configured with a public URL on this server. "
            "The link above is a relative path — prepend your server's domain "
            "to get the full URL. Set PUBLIC_BASE_URL in your .env to have "
            "the agent generate complete links automatically."
        )
        if public_base_error:
            url_note += " The configured value failed the public-URL safety check."
    else:
        url = f"{public_base}/p/{short_id}"
        url_note = ""

    evidence_refs = (url,) if url.startswith(("http://", "https://")) else ()
    return _typed_success(
        f"Published successfully!\n\n"
        f"Public URL: {url}\n"
        f"Title: {title}\n\n"
        f"Anyone can access this page without logging in.{url_note}",
        result_ref=f"published-page://{short_id}",
        artifact_refs=(f"published-page://{short_id}",),
        evidence_refs=evidence_refs,
    )


async def _publish_page(
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    ws: Path,
    arguments: dict,
) -> str:
    outcome = await _publish_page_outcome(agent_id, user_id, ws, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Page publishing returned no summary.",
    )


async def _list_published_pages_outcome(
    agent_id: uuid.UUID,
) -> ToolExecutionOutcome:
    """List all published pages for this agent."""
    from app.models.published_page import PublishedPage
    try:
        from app.config import get_settings as _get_publish_settings
        public_base = (_get_publish_settings().PUBLIC_BASE_URL or os.environ.get("PUBLIC_BASE_URL", "")).rstrip("/")
    except Exception:
        public_base = os.environ.get("PUBLIC_BASE_URL", "").rstrip("/")

    try:
        async with async_session() as db:
            result = await db.execute(
                select(PublishedPage)
                .where(PublishedPage.agent_id == agent_id)
                .order_by(PublishedPage.created_at.desc())
            )
            pages = result.scalars().all()

        if not pages:
            return _typed_success("No published pages yet.")

        lines = [f"Published pages ({len(pages)} total):\n"]
        for p in pages:
            url = f"{public_base}/p/{p.short_id}" if public_base else f"/p/{p.short_id}"
            lines.append(f"- {p.title or 'Untitled'}")
            lines.append(f"  URL: {url}")
            lines.append(f"  Source: {p.source_path}")
            lines.append(f"  Views: {p.view_count}")
            lines.append("")
        evidence_refs = tuple(
            f"published-page://{page.short_id}" for page in pages
        )
        return _typed_success(
            "\n".join(lines),
            evidence_refs=evidence_refs,
        )
    except Exception as e:
        return _typed_failure(
            f"Failed to list pages: {type(e).__name__}.",
            "published_page_list_failed",
            retryable=True,
        )


async def _list_published_pages(agent_id: uuid.UUID) -> str:
    outcome = await _list_published_pages_outcome(agent_id)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Published page listing returned no summary.",
    )


# ─── AgentBay Tool Handlers ─────────────────────────────────────

def _agentbay_normalize_image_bytes(data) -> bytes | None:
    """Normalize AgentBay image payloads to raw bytes."""
    import base64 as _base64

    if isinstance(data, str):
        if data.startswith("data:image"):
            data = data.split(",", 1)[1]
        return _base64.b64decode(data)
    if isinstance(data, bytes):
        return data
    return None


def _agentbay_save_image_to_workspace(
    *,
    agent_id: uuid.UUID,
    ws: Path,
    raw_bytes: bytes,
    prefix: str,
    label: str,
) -> str:
    """Save an explicitly requested screenshot under workspace/screenshots/."""
    import time as _time

    rel_path = f"workspace/screenshots/{prefix}-{int(_time.time())}.png"
    screenshot_path = ws / rel_path
    screenshot_path.parent.mkdir(parents=True, exist_ok=True)
    screenshot_path.write_bytes(raw_bytes)
    logger.info(f"[AgentBay] Explicit screenshot saved to workspace: {rel_path}")
    return (
        f"Screenshot saved to `{rel_path}`.\n"
        f"![{label}](/api/agents/{agent_id}/files/download?path={rel_path})"
    )

def _agentbay_result_field(
    result: object,
    field: str,
    default: object = None,
) -> object:
    if isinstance(result, Mapping):
        return result.get(field, default)
    return getattr(result, field, default)


def _agentbay_json_summary(label: str, value: object) -> str | None:
    try:
        encoded = json.dumps(
            value,
            ensure_ascii=False,
            allow_nan=False,
            sort_keys=True,
        )
    except (TypeError, ValueError):
        return None
    return f"{label}: {encoded}"


def _agentbay_read_failure(*, malformed: bool = False) -> ToolExecutionOutcome:
    if malformed:
        return _typed_failure(
            "AgentBay returned an invalid read payload; retry the read.",
            "agentbay_read_payload_invalid",
            retryable=True,
        )
    return _typed_failure(
        "AgentBay rejected the read request.",
        "agentbay_read_rejected",
        retryable=False,
    )


def _agentbay_decode_image(value: object) -> tuple[bytes, str] | None:
    import base64
    from io import BytesIO

    from PIL import Image, ImageFile

    raw: bytes
    if isinstance(value, bytes):
        raw = value
    elif isinstance(value, str):
        encoded = value.strip()
        if encoded.startswith("data:"):
            _, separator, encoded = encoded.partition(",")
            if not separator:
                return None
        try:
            raw = base64.b64decode("".join(encoded.split()), validate=True)
        except (ValueError, TypeError):
            return None
    else:
        return None
    if not raw or len(raw) > _MAX_GENERATED_IMAGE_BYTES:
        return None
    allow_truncated = ImageFile.LOAD_TRUNCATED_IMAGES
    try:
        ImageFile.LOAD_TRUNCATED_IMAGES = True
        with Image.open(BytesIO(raw)) as image:
            image.load()
            if image.width <= 0 or image.height <= 0:
                return None
            image_format = str(image.format or "").upper()
    except Exception:
        return None
    finally:
        ImageFile.LOAD_TRUNCATED_IMAGES = allow_truncated
    mime_type = {
        "PNG": "image/png",
        "JPEG": "image/jpeg",
        "JPG": "image/jpeg",
        "WEBP": "image/webp",
    }.get(image_format)
    if mime_type is None:
        return None
    return raw, mime_type


def _agentbay_crop_image(
    raw: bytes,
    *,
    x: int,
    y: int,
    width: int,
    height: int,
) -> bytes | None:
    from io import BytesIO

    from PIL import Image, ImageFile

    allow_truncated = ImageFile.LOAD_TRUNCATED_IMAGES
    try:
        ImageFile.LOAD_TRUNCATED_IMAGES = True
        with Image.open(BytesIO(raw)) as image:
            image.load()
            if (
                x < 0
                or y < 0
                or width <= 0
                or height <= 0
                or x + width > image.width
                or y + height > image.height
            ):
                return None
            if (x, y, width, height) == (0, 0, image.width, image.height):
                return raw
            cropped = image.crop((x, y, x + width, y + height))
            output = BytesIO()
            cropped.save(output, format="PNG")
            return output.getvalue()
    except Exception:
        return None
    finally:
        ImageFile.LOAD_TRUNCATED_IMAGES = allow_truncated


def _agentbay_screenshot_outcome(
    tool_name: str,
    result: object,
    arguments: Mapping[str, Any],
) -> ToolExecutionOutcome:
    success = _agentbay_result_field(result, "success")
    if success is False:
        return _agentbay_read_failure()
    if success is not True:
        return _agentbay_read_failure(malformed=True)
    field = "screenshot" if tool_name == "agentbay_browser_screenshot" else "data"
    decoded = _agentbay_decode_image(_agentbay_result_field(result, field))
    if decoded is None:
        return _agentbay_read_failure(malformed=True)
    raw, mime_type = decoded
    if tool_name == "agentbay_computer_precision_screenshot":
        coordinates = tuple(arguments.get(name) for name in ("x", "y", "width", "height"))
        if any(
            not isinstance(value, int) or isinstance(value, bool)
            for value in coordinates
        ):
            return _typed_failure(
                "Precision screenshot coordinates must be integers.",
                "invalid_tool_arguments",
            )
        cropped = _agentbay_crop_image(
            raw,
            x=coordinates[0],
            y=coordinates[1],
            width=coordinates[2],
            height=coordinates[3],
        )
        if cropped is None:
            return _agentbay_read_failure(malformed=True)
        raw = cropped
        mime_type = "image/png"
    return _typed_success(
        "AgentBay screenshot captured for internal vision.",
        metadata={
            "provider": "agentbay",
            "operation": tool_name,
            "content_hash": hashlib.sha256(raw).hexdigest(),
            "mime_type": mime_type,
            "size": len(raw),
        },
        private_binary=raw,
    )


def _agentbay_structured_read_outcome(
    tool_name: str,
    result: object,
) -> ToolExecutionOutcome:
    success = _agentbay_result_field(result, "success")
    if success is False:
        return _agentbay_read_failure()
    if success is not True:
        return _agentbay_read_failure(malformed=True)

    summary: str | None = None
    if tool_name == "agentbay_browser_extract":
        value = _agentbay_result_field(result, "data")
        summary = _agentbay_json_summary("AgentBay extracted data", value)
    elif tool_name == "agentbay_browser_observe":
        value = _agentbay_result_field(result, "elements")
        if isinstance(value, list):
            summary = _agentbay_json_summary("AgentBay observed elements", value)
    elif tool_name == "agentbay_code_read_file":
        value = _agentbay_result_field(result, "content")
        if isinstance(value, str):
            summary = f"AgentBay file content:\n{value}"
    elif tool_name == "agentbay_computer_get_screen_size":
        value = _agentbay_result_field(result, "data")
        if (
            isinstance(value, Mapping)
            and isinstance(value.get("width"), int)
            and not isinstance(value.get("width"), bool)
            and isinstance(value.get("height"), int)
            and not isinstance(value.get("height"), bool)
            and value["width"] > 0
            and value["height"] > 0
        ):
            summary = _agentbay_json_summary("AgentBay screen size", dict(value))
    elif tool_name == "agentbay_computer_get_installed_apps":
        value = _agentbay_result_field(result, "apps")
        if isinstance(value, list):
            summary = _agentbay_json_summary("AgentBay installed apps", value)
    elif tool_name == "agentbay_computer_get_cursor_position":
        value = _agentbay_result_field(result, "data")
        if (
            isinstance(value, Mapping)
            and isinstance(value.get("x"), int)
            and not isinstance(value.get("x"), bool)
            and isinstance(value.get("y"), int)
            and not isinstance(value.get("y"), bool)
        ):
            summary = _agentbay_json_summary(
                "AgentBay cursor position",
                dict(value),
            )
    elif tool_name == "agentbay_computer_get_active_window":
        value = _agentbay_result_field(result, "window")
        if isinstance(value, Mapping):
            summary = _agentbay_json_summary(
                "AgentBay active window",
                dict(value),
            )
    elif tool_name == "agentbay_computer_list_windows":
        value = _agentbay_result_field(result, "windows")
        if isinstance(value, list):
            summary = _agentbay_json_summary("AgentBay windows", value)
    elif tool_name == "agentbay_computer_list_visible_apps":
        value = _agentbay_result_field(result, "apps")
        if isinstance(value, list):
            summary = _agentbay_json_summary("AgentBay visible apps", value)
    if summary is None:
        return _agentbay_read_failure(malformed=True)
    return _typed_success(
        summary,
        metadata={"provider": "agentbay", "operation": tool_name},
    )


async def _agentbay_read_outcome(
    tool_name: str,
    agent_id: uuid.UUID,
    arguments: Mapping[str, Any],
    *,
    session_id: str,
) -> ToolExecutionOutcome:
    from app.services.agentbay_client import get_agentbay_client_for_agent

    if tool_name in {
        "agentbay_browser_extract",
        "agentbay_browser_observe",
    }:
        instruction = arguments.get("instruction")
        selector = arguments.get("selector", "")
        if (
            not isinstance(instruction, str)
            or not instruction.strip()
            or not isinstance(selector, str)
        ):
            return _typed_failure(
                "Browser read requires instruction and an optional string selector.",
                "invalid_tool_arguments",
            )
    if tool_name == "agentbay_code_read_file":
        remote_path = arguments.get("remote_path")
        if not isinstance(remote_path, str) or not remote_path.strip():
            return _typed_failure(
                "Code file read requires remote_path.",
                "invalid_tool_arguments",
            )
    if tool_name == "agentbay_computer_precision_screenshot":
        coordinates = tuple(
            arguments.get(name) for name in ("x", "y", "width", "height")
        )
        if (
            any(
                not isinstance(value, int) or isinstance(value, bool)
                for value in coordinates
            )
            or coordinates[0] < 0
            or coordinates[1] < 0
            or coordinates[2] <= 0
            or coordinates[3] <= 0
        ):
            return _typed_failure(
                "Precision screenshot requires non-negative x/y and positive integer width/height.",
                "invalid_tool_arguments",
            )
    if tool_name == "agentbay_computer_get_installed_apps" and any(
        not isinstance(arguments.get(name, default), bool)
        for name, default in (
            ("start_menu", True),
            ("desktop", True),
            ("ignore_system_apps", True),
        )
    ):
        return _typed_failure(
            "Installed-app read options must be booleans.",
            "invalid_tool_arguments",
        )
    if tool_name == "agentbay_computer_list_windows":
        timeout_ms = arguments.get("timeout_ms", 3000)
        if (
            not isinstance(timeout_ms, int)
            or isinstance(timeout_ms, bool)
            or timeout_ms <= 0
        ):
            return _typed_failure(
                "Window-list timeout_ms must be a positive integer.",
                "invalid_tool_arguments",
            )

    image_type = (
        "browser"
        if tool_name.startswith("agentbay_browser_")
        else "code"
        if tool_name.startswith("agentbay_code_")
        else "computer"
    )
    try:
        client = await get_agentbay_client_for_agent(
            agent_id,
            image_type,
            session_id=session_id,
            run_id=agentbay_run_scope_id.get().strip(),
        )
    except Exception:
        return _typed_unknown(
            "AgentBay session creation or restore outcome is unknown; do not retry automatically.",
            "agentbay_session_outcome_unknown",
        )

    try:
        if tool_name == "agentbay_browser_screenshot":
            result = await client.browser_screenshot()
        elif tool_name == "agentbay_browser_extract":
            instruction = cast(str, arguments.get("instruction"))
            selector = cast(str, arguments.get("selector", ""))
            result = await client.browser_extract(instruction, selector)
        elif tool_name == "agentbay_browser_observe":
            instruction = cast(str, arguments.get("instruction"))
            selector = cast(str, arguments.get("selector", ""))
            result = await client.browser_observe(instruction, selector)
        elif tool_name == "agentbay_code_read_file":
            remote_path = cast(str, arguments.get("remote_path"))
            result = await client.code_read_file(remote_path)
        elif tool_name in {
            "agentbay_computer_screenshot",
            "agentbay_computer_precision_screenshot",
        }:
            result = await client.computer_screenshot()
        elif tool_name == "agentbay_computer_get_screen_size":
            result = await client.computer_get_screen_size()
        elif tool_name == "agentbay_computer_get_installed_apps":
            options = tuple(
                arguments.get(name, default)
                for name, default in (
                    ("start_menu", True),
                    ("desktop", True),
                    ("ignore_system_apps", True),
                )
            )
            result = await client.computer_get_installed_apps(
                start_menu=options[0],
                desktop=options[1],
                ignore_system_apps=options[2],
            )
        elif tool_name == "agentbay_computer_get_cursor_position":
            result = await client.computer_get_cursor_position()
        elif tool_name == "agentbay_computer_get_active_window":
            result = await client.computer_get_active_window()
        elif tool_name == "agentbay_computer_list_windows":
            timeout_ms = cast(int, arguments.get("timeout_ms", 3000))
            result = await client.computer_list_windows(timeout_ms=timeout_ms)
        elif tool_name == "agentbay_computer_list_visible_apps":
            result = await client.computer_list_visible_apps()
        else:  # pragma: no cover - guarded by the fixed A1 workset
            return _typed_failure(
                "AgentBay read is not part of the typed A1 workset.",
                "unsupported_tool",
            )
    except Exception:
        return _typed_failure(
            "AgentBay read failed before a valid result was received; retry the read.",
            "agentbay_read_transport_failed",
            retryable=True,
        )

    if tool_name in {
        "agentbay_browser_screenshot",
        "agentbay_computer_screenshot",
        "agentbay_computer_precision_screenshot",
    }:
        return _agentbay_screenshot_outcome(tool_name, result, arguments)
    return _agentbay_structured_read_outcome(tool_name, result)


async def _agentbay_browser_navigate(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """AgentBay browser navigation.

    After navigating, always captures an internal screenshot for LLM vision.
    The screenshot is held in memory and consumed by vision_inject.py in the
    same request cycle; it is not persisted to the user's workspace.
    """
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    url = arguments.get("url", "")
    wait_for = arguments.get("wait_for", "")

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id, run_id=_run_id)
        # Always request a screenshot for navigation so the model can observe the result
        result = await client.browser_navigate(url, wait_for=wait_for, screenshot=True)

        # Build text parts from the navigation result
        parts = [f"✅ 已访问: {url}"]
        if result.get("title"):
            parts.append(f"标题: {result['title']}")
        if result.get("content"):
            content = result["content"][:3000]
            parts.append(f"内容:\n{content}")
        logger.info(f"[AgentBay] Browser navigate result: {result.get('title')}")

        screenshot_data = result.get("screenshot")
        if screenshot_data:
            raw_bytes = _agentbay_normalize_image_bytes(screenshot_data)

            if raw_bytes:
                # Store in memory only — vision_inject.py will consume it.
                from app.services.vision_inject import store_temp_screenshot
                img_id = store_temp_screenshot(raw_bytes)
                parts.append(
                    f"Internal screenshot captured for analysis. [ImageID: {img_id}]\n"
                    f"NOTE: This screenshot is for LLM vision only and is not saved to the user's workspace."
                )
                logger.info(f"[AgentBay] Browser navigate screenshot stored in memory (id={img_id})")

        return "\n\n".join(parts)

    except RuntimeError as e:
        return f"❌ {str(e)}。请先在 Agent 设置中配置 AgentBay 通道。"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser navigate failed for agent {agent_id}")
        return f"❌ AgentBay 浏览器访问失败: {str(e)[:200]}"


async def _agentbay_browser_screenshot(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Take a screenshot of the CURRENT browser page without navigating.

    Correct way to observe the result of a click, type, or form submit — never
    call browser_navigate again just to screenshot, that refreshes the page.

    The image is held in the process-level memory cache and consumed once by
    the LLM vision pipeline — no disk write, nothing shown in the user's file
    manager or chat history.
    """
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id, run_id=_run_id)
        result = await client.browser_screenshot()

        screenshot_data = result.get("screenshot")
        if not screenshot_data:
            return "❌ 截图失败：未返回图像数据"

        raw_bytes = _agentbay_normalize_image_bytes(screenshot_data)
        if raw_bytes is None:
            return "❌ 截图失败：未知数据格式"

        # Store in memory only — vision_inject.py will consume it for LLM vision
        from app.services.vision_inject import store_temp_screenshot
        img_id = store_temp_screenshot(raw_bytes)
        logger.info(f"[AgentBay] Browser screenshot stored in memory (id={img_id})")
        return (
            f"Internal screenshot captured for analysis. [ImageID: {img_id}]\n"
            f"NOTE: This screenshot is for LLM vision only and is not saved to the user's workspace."
        )

    except RuntimeError as e:
        return f"❌ {str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser screenshot failed for agent {agent_id}")
        return f"❌ 截图失败: {str(e)[:200]}"


async def _agentbay_browser_save_screenshot(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Save the current AgentBay browser screenshot to workspace/screenshots/."""
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id, run_id=_run_id)
        result = await client.browser_screenshot()
        raw_bytes = _agentbay_normalize_image_bytes(result.get("screenshot"))
        if raw_bytes is None:
            return "❌ 截图保存失败：未返回可保存的图像数据"
        return _agentbay_save_image_to_workspace(
            agent_id=agent_id,
            ws=ws,
            raw_bytes=raw_bytes,
            prefix="browser-screenshot",
            label="Browser Screenshot",
        )
    except RuntimeError as e:
        return f"❌ {str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser save screenshot failed for agent {agent_id}")
        return f"❌ 截图保存失败: {str(e)[:200]}"


async def _agentbay_browser_click(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """AgentBay 浏览器点击。"""
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    selector = arguments.get("selector", "")

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id, run_id=_run_id)
        await client.browser_click(selector)
        return f"✅ 已点击元素: {selector}"
    except RuntimeError as e:
        return f"❌ {str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser click failed")
        return f"❌ 点击失败: {str(e)[:200]}"


async def _agentbay_browser_type(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """AgentBay 浏览器输入。"""
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    selector = arguments.get("selector", "")
    text = arguments.get("text", "")

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id, run_id=_run_id)
        await client.browser_type(selector, text)
        return f"✅ 已在 {selector} 输入文本"
    except RuntimeError as e:
        return f"❌ {str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser type failed")
        return f"❌ 输入失败: {str(e)[:200]}"


async def _agentbay_code_execute(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """在 AgentBay 代码空间执行代码。"""
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    language = arguments.get("language", "python")
    code = arguments.get("code", "")
    timeout = arguments.get("timeout", 30)

    if not code.strip():
        return "❌ 请提供要执行的代码"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "code", session_id=_session_id, run_id=_run_id)
        result = await client.code_execute(language, code, timeout)

        # 格式化返回结果
        parts = [f"✅ 代码执行完成 ({language})"]
        if result.get("stdout"):
            parts.append(f"📤 输出:\n{result['stdout']}")
        if result.get("stderr"):
            parts.append(f"⚠️ 错误输出:\n{result['stderr']}")
        if result.get("exit_code") != 0:
            parts.append(f"退出码: {result['exit_code']}")

        return "\n\n".join(parts)

    except RuntimeError as e:
        return f"❌ {str(e)}。请先在 Agent 设置中配置 AgentBay 通道。"
    except Exception as e:
        logger.exception(f"[AgentBay] Code execution failed for agent {agent_id}")
        return f"❌ 代码执行失败: {str(e)[:200]}"


async def _agentbay_code_write_file(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Write a text file in the AgentBay Code Sandbox."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    remote_path = arguments.get("remote_path") or arguments.get("path") or ""
    content = arguments.get("content")
    mode = arguments.get("mode", "overwrite")

    if not remote_path.strip():
        return "Missing required argument 'remote_path'"
    if content is None:
        return "Missing required argument 'content'"
    if mode not in ("overwrite", "append"):
        return "Invalid mode. Use 'overwrite' or 'append'."

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "code", session_id=_session_id, run_id=_run_id)
        result = await asyncio.to_thread(
            client._session.file_system.write_file,
            remote_path,
            str(content),
            mode,
        )
        if result.success:
            byte_count = len(str(content).encode("utf-8"))
            return f"File written in AgentBay Code Sandbox: {remote_path} ({byte_count} bytes, mode={mode})"
        return f"Write failed: {result.error_message}"
    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Code write file failed for agent {agent_id}")
        return f"Write file failed: {str(e)[:200]}"


async def _agentbay_code_read_file(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Read a text file from the AgentBay Code Sandbox."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    remote_path = arguments.get("remote_path") or arguments.get("path") or ""
    if not remote_path.strip():
        return "Missing required argument 'remote_path'"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "code", session_id=_session_id, run_id=_run_id)
        result = await asyncio.to_thread(
            client._session.file_system.read_file,
            remote_path,
        )
        if result.success:
            content = getattr(result, "content", "") or ""
            return f"File read from AgentBay Code Sandbox: {remote_path}\n\n{content[:12000]}"
        return f"Read failed: {result.error_message}"
    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Code read file failed for agent {agent_id}")
        return f"Read file failed: {str(e)[:200]}"


async def _agentbay_code_edit_file(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Edit a text file in the AgentBay Code Sandbox."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    remote_path = arguments.get("remote_path") or arguments.get("path") or ""
    edits = arguments.get("edits")
    dry_run = bool(arguments.get("dry_run", False))

    if not remote_path.strip():
        return "Missing required argument 'remote_path'"
    if not isinstance(edits, list) or not edits:
        return "Missing required argument 'edits'"

    normalized_edits = []
    for edit in edits:
        if not isinstance(edit, dict):
            return "Each edit must be an object with oldText and newText."
        old_text = edit.get("oldText")
        new_text = edit.get("newText")
        if old_text is None or new_text is None:
            return "Each edit must include oldText and newText."
        normalized_edits.append({"oldText": str(old_text), "newText": str(new_text)})

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "code", session_id=_session_id, run_id=_run_id)
        result = await asyncio.to_thread(
            client._session.file_system.edit_file,
            remote_path,
            normalized_edits,
            dry_run,
        )
        if result.success:
            action = "Previewed edits for" if dry_run else "Edited"
            return f"{action} AgentBay Code Sandbox file: {remote_path} ({len(normalized_edits)} replacement(s))"
        return f"Edit failed: {result.error_message}"
    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Code edit file failed for agent {agent_id}")
        return f"Edit file failed: {str(e)[:200]}"


async def _handle_email_tool(tool_name: str, agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Dispatch email tool calls to the email_service module."""
    from app.services.email_service import send_email, read_emails, reply_email

    config = await _get_email_config(agent_id)
    if not config.get("email_address") or not config.get("auth_code"):
        return (
            "❌ Email not configured for this agent.\n\n"
            "Please go to Agent → Tools → Send Email → Config to set up your email:\n"
            "1. Select your email provider\n"
            "2. Enter your email address\n"
            "3. Enter your authorization code (not your login password)"
        )

    try:
        if tool_name == "send_email":
            return await send_email(
                config=config,
                to=arguments.get("to", ""),
                subject=arguments.get("subject", ""),
                body=arguments.get("body", ""),
                cc=arguments.get("cc"),
                attachments=arguments.get("attachments"),
                workspace_path=ws,
                agent_id=agent_id,
            )
        elif tool_name == "read_emails":
            return await read_emails(
                config=config,
                limit=arguments.get("limit", 10),
                search=arguments.get("search"),
                folder=arguments.get("folder", "INBOX"),
            )
        elif tool_name == "reply_email":
            return await reply_email(
                config=config,
                message_id=arguments.get("message_id", ""),
                body=arguments.get("body", ""),
                folder=arguments.get("folder", "INBOX"),
            )
        else:
            return f"❌ Unknown email tool: {tool_name}"
    except Exception as e:
        return f"❌ Email tool error: {str(e)[:200]}"


# ─── Skill Management Tools ────────────────────────────────────


async def _search_clawhub_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Search ClawHub using its decoded JSON response as the read fact."""
    query = arguments.get("query")
    if not isinstance(query, str) or not query.strip():
        return _typed_failure(
            "search_clawhub requires query.",
            "invalid_tool_arguments",
        )
    query = query.strip()

    # Resolve tenant ClawHub API key
    from app.api.skills import _clawhub_search_endpoint, _fetch_clawhub_json, _get_clawhub_key
    tenant_id = await _get_agent_tenant_id(agent_id)
    api_key = await _get_clawhub_key(tenant_id)

    try:
        data, _ = await _fetch_clawhub_json(
            _clawhub_search_endpoint,
            api_key=api_key,
            params={"q": query},
        )
    except Exception as e:
        status_code = getattr(e, "status_code", None)
        return _typed_failure(
            f"ClawHub search failed: {type(e).__name__}.",
            "clawhub_search_failed",
            retryable=(
                status_code in {408, 429}
                or isinstance(status_code, int)
                and status_code >= 500
            ),
        )

    if not isinstance(data, Mapping):
        return _typed_failure(
            "ClawHub returned an unreadable search response.",
            "clawhub_response_invalid",
            retryable=True,
        )

    results = data.get("results", [])
    if not isinstance(results, list):
        return _typed_failure(
            "ClawHub returned an invalid results collection.",
            "clawhub_response_invalid",
            retryable=True,
        )
    if any(not isinstance(result, Mapping) for result in results):
        return _typed_failure(
            "ClawHub returned an invalid Skill entry.",
            "clawhub_response_invalid",
            retryable=True,
        )
    if not results:
        return _typed_success(f"No skills found matching '{query}'.")

    lines = [f"Found {len(results)} skill(s) matching '{query}':\n"]
    for r in results:
        name = r.get("displayName") or r.get("slug", "?")
        slug = r.get("slug", "")
        summary = (r.get("summary") or "")[:120]
        updated = ""
        if r.get("updatedAt"):
            from datetime import datetime
            try:
                dt = datetime.fromtimestamp(r["updatedAt"] / 1000)
                updated = f" | Updated: {dt.strftime('%Y-%m-%d')}"
            except Exception:
                pass
        lines.append(f"• **{name}** (`{slug}`){updated}")
        if summary:
            lines.append(f"  {summary}")
    lines.append("\nTo install a skill, use: install_skill(source=\"<slug>\")")
    return _typed_success("\n".join(lines))


async def _search_clawhub(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter for typed ClawHub search."""
    outcome = await _search_clawhub_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="ClawHub search returned no summary.",
    )


async def _install_skill_outcome(
    agent_id: uuid.UUID,
    ws: Path,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Fetch, validate, and write one Skill package into a temp workspace."""
    source = arguments.get("source")
    if not isinstance(source, str) or not source.strip():
        return _typed_failure(
            "install_skill requires source.",
            "invalid_tool_arguments",
        )
    source = source.strip()

    is_url = source.startswith("http://") or source.startswith("https://")
    base = ws  # agent workspace dir (skills/ lives under workspace/)

    try:
        if is_url:
            # ── GitHub URL path ──
            from app.api.skills import _parse_github_url, _fetch_github_directory, _get_github_token

            parsed = _parse_github_url(source)
            if not parsed:
                return _typed_failure(
                    "Invalid GitHub Skill URL.",
                    "invalid_tool_arguments",
                )

            owner, repo, branch, path = parsed["owner"], parsed["repo"], parsed["branch"], parsed["path"]
            tenant_id = await _get_agent_tenant_id(agent_id)
            token = await _get_github_token(tenant_id)
            files = await _fetch_github_directory(owner, repo, path, branch, token)
            if not files:
                return _typed_failure(
                    "No files found at the specified GitHub URL.",
                    "skill_source_not_found",
                )

            folder_name = path.rstrip("/").split("/")[-1] if path else repo
        else:
            # ── ClawHub slug path ──
            slug = source
            from app.api.skills import _fetch_clawhub_skill_archive, _fetch_clawhub_skill_meta, _get_clawhub_key

            # 1. Fetch metadata from ClawHub (with tenant API key)
            tenant_id = await _get_agent_tenant_id(agent_id)
            api_key = await _get_clawhub_key(tenant_id)
            try:
                _meta, meta_base = await _fetch_clawhub_skill_meta(slug, api_key=api_key)
            except Exception as e:
                return _typed_failure(
                    f"ClawHub Skill lookup failed: {type(e).__name__}.",
                    "skill_source_lookup_failed",
                )

            # 2. Fetch files from the ClawHub archive
            files, _ = await _fetch_clawhub_skill_archive(slug, api_key=api_key, preferred_base=meta_base)
            if not files:
                return _typed_failure(
                    f"No files found for Skill '{slug}'.",
                    "skill_source_not_found",
                )

            folder_name = slug

        if (
            not isinstance(folder_name, str)
            or not folder_name.strip()
            or Path(folder_name).name != folder_name
            or folder_name in {".", ".."}
        ):
            return _typed_failure(
                "Skill source resolved to an invalid folder name.",
                "skill_package_invalid",
            )
        if not any(
            isinstance(file, Mapping)
            and str(file.get("path") or "").upper() == "SKILL.MD"
            for file in files
        ):
            return _typed_failure(
                "Skill package does not contain a root SKILL.md.",
                "skill_package_invalid",
            )

        # 3. Write files to the temporary Agent workspace. Durable sync is
        # performed only after this function returns a typed success.
        skill_dir = base / "skills" / folder_name
        skill_dir.mkdir(parents=True, exist_ok=True)
        skill_root = skill_dir.resolve()

        written = []
        for file in files:
            if not isinstance(file, Mapping):
                return _typed_failure(
                    "Skill package contains an invalid file entry.",
                    "skill_package_invalid",
                )
            rel_path = file.get("path")
            content = file.get("content")
            if not isinstance(rel_path, str) or not isinstance(content, str):
                return _typed_failure(
                    "Skill package contains an invalid file entry.",
                    "skill_package_invalid",
                )
            file_path = (skill_root / rel_path).resolve()
            if not file_path.is_relative_to(skill_root):
                return _typed_failure(
                    "Skill package contains an unsafe file path.",
                    "skill_package_invalid",
                )
            file_path.parent.mkdir(parents=True, exist_ok=True)
            file_path.write_text(content, encoding="utf-8")
            written.append(rel_path)

        refs = tuple(
            _workspace_artifact_ref(
                agent_id,
                f"skills/{folder_name}/{rel_path}",
            )
            for rel_path in written
        )
        shown = ", ".join(written[:20])
        if len(written) > 20:
            shown += f", ... and {len(written) - 20} more"
        return _typed_success(
            f"Skill '{folder_name}' installed ({len(written)} files).\n\n"
            f"Files: {shown}",
            artifact_refs=refs,
        )

    except Exception as e:
        return _typed_failure(
            f"Skill installation failed: {type(e).__name__}.",
            "skill_install_failed",
        )


async def _install_skill(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Legacy display adapter for typed Skill installation."""
    outcome = await _install_skill_outcome(agent_id, ws, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Skill installation returned no summary.",
    )


# ─── AgentBay: Browser Extract & Observe ────────────────────────────────

async def _agentbay_browser_extract(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Extract structured data from current browser page."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    instruction = arguments.get("instruction", "")
    selector = arguments.get("selector", "")

    if not instruction.strip():
        return "Missing required argument 'instruction'"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id, run_id=_run_id)
        result = await client.browser_extract(instruction, selector=selector)

        if result.get("success"):
            import json
            data = result.get("data", {})
            data_str = json.dumps(data, ensure_ascii=False, indent=2) if isinstance(data, (dict, list)) else str(data)
            return f"Extraction successful:\n\n{data_str[:5000]}"
        else:
            return f"Extraction failed: {result}"

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Browser extract failed for agent {agent_id}")
        return f"Browser extract failed: {str(e)[:200]}"


async def _agentbay_browser_observe(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Observe the current browser page state."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    instruction = arguments.get("instruction", "")
    selector = arguments.get("selector", "")

    if not instruction.strip():
        return "Missing required argument 'instruction'"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id, run_id=_run_id)
        result = await client.browser_observe(instruction, selector=selector)

        if result.get("success"):
            import json
            elements = result.get("elements", [])
            if not elements:
                return "No interactive elements found matching your instruction."
            elements_str = json.dumps(elements, ensure_ascii=False, indent=2)
            return f"Found {len(elements)} interactive element(s):\n\n{elements_str[:5000]}"
        else:
            return f"Observation failed: {result}"

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Browser observe failed for agent {agent_id}")
        return f"Browser observe failed: {str(e)[:200]}"


# ─── AgentBay: Command (Shell) ──────────────────────────────────────────

async def _agentbay_browser_login(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Perform an automated login using AgentBay's built-in login skill.

    Supports complex login flows including CAPTCHAs, OTP inputs,
    and multi-step authentication via AgentBay's AI-driven capability.
    """
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    url = arguments.get("url", "")
    login_config = arguments.get("login_config", "")

    if not url.strip():
        return "Missing required argument 'url'"
    if not login_config.strip():
        return "Missing required argument 'login_config' (JSON string with api_key + skill_id)"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id, run_id=_run_id)
        result = await client.browser_login(url, login_config)

        if result.get("success"):
            return f"Login completed successfully. {result.get('message', '')}"
        else:
            return f"Login failed: {result.get('message', 'Unknown error')}"

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Browser login failed for agent {agent_id}")
        return f"Login failed: {str(e)[:200]}"


async def _agentbay_command_exec(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Execute a shell command in the AgentBay environment."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    command = arguments.get("command", "")
    timeout_ms = arguments.get("timeout_ms", 50000)
    cwd = arguments.get("cwd", "")

    if not command.strip():
        return "Missing required argument 'command'"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "code", session_id=_session_id, run_id=_run_id)
        result = await client.command_exec(command, timeout_ms=timeout_ms, cwd=cwd)

        parts = []
        if result.get("success"):
            parts.append(f"Command executed successfully (exit code: {result.get('exit_code', 0)})")
        else:
            parts.append(f"Command failed (exit code: {result.get('exit_code', -1)})")

        if result.get("stdout"):
            parts.append(f"stdout:\n{result['stdout'][:3000]}")
        if result.get("stderr"):
            parts.append(f"stderr:\n{result['stderr'][:1000]}")
        if result.get("error_message"):
            parts.append(f"Error: {result['error_message']}")

        return "\n\n".join(parts)

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Command exec failed for agent {agent_id}")
        return f"Command execution failed: {str(e)[:200]}"


# ─── AgentBay: Computer Use Handlers ────────────────────────────────────

def _agentbay_extract_screen_dimensions(screen_data) -> tuple[int | None, int | None, str]:
    """Return width/height/dpi text from AgentBay get_screen_size payload."""
    if not isinstance(screen_data, dict):
        return None, None, ""
    width = screen_data.get("width")
    height = screen_data.get("height")
    dpi = screen_data.get("dpiScalingFactor")
    try:
        width = int(width) if width is not None else None
        height = int(height) if height is not None else None
    except (TypeError, ValueError):
        width, height = None, None
    parts = []
    if width and height:
        parts.append(f"width={width}, height={height}")
    if dpi is not None:
        parts.append(f"dpiScalingFactor={dpi}")
    return width, height, ", ".join(parts)


async def _agentbay_get_screen_metadata(client) -> tuple[int | None, int | None, str]:
    try:
        size_result = await client.computer_get_screen_size()
        if size_result.get("success"):
            return _agentbay_extract_screen_dimensions(size_result.get("data"))
    except Exception as e:
        logger.debug(f"[AgentBay] Could not fetch computer screen size: {e}")
    return None, None, ""


def _agentbay_image_dimensions(raw_bytes: bytes) -> tuple[int | None, int | None]:
    try:
        from io import BytesIO
        from PIL import Image

        with Image.open(BytesIO(raw_bytes)) as img:
            return img.width, img.height
    except Exception:
        return None, None


def _agentbay_crop_image_bytes(
    raw_bytes: bytes,
    *,
    x: int,
    y: int,
    width: int,
    height: int,
) -> tuple[bytes, tuple[int, int, int, int], int] | None:
    try:
        from io import BytesIO
        from PIL import Image

        with Image.open(BytesIO(raw_bytes)) as img:
            img_width, img_height = img.width, img.height
            left = max(0, min(int(x), img_width - 1))
            top = max(0, min(int(y), img_height - 1))
            right = max(left + 1, min(left + int(width), img_width))
            bottom = max(top + 1, min(top + int(height), img_height))
            cropped = img.crop((left, top, right, bottom))

            # Enlarge precision crops before vision injection so small controls
            # occupy more pixels without changing the absolute coordinate labels.
            max_side = max(cropped.width, cropped.height)
            scale = 1
            if max_side <= 260:
                scale = 3
            elif max_side <= 520:
                scale = 2
            if scale > 1:
                cropped = cropped.resize((cropped.width * scale, cropped.height * scale), Image.Resampling.LANCZOS)

            buf = BytesIO()
            cropped.save(buf, format="PNG")
            return buf.getvalue(), (left, top, right - left, bottom - top), scale
    except Exception as e:
        logger.debug(f"[AgentBay] Could not crop desktop screenshot: {e}")
        return None


def _agentbay_expand_precision_crop(
    x: int,
    y: int,
    width: int,
    height: int,
    *,
    min_width: int = 360,
    min_height: int = 240,
) -> tuple[int, int, int, int]:
    """Expand small requested crops so near-miss targeting still shows context."""
    width = max(1, int(width))
    height = max(1, int(height))
    expanded_width = max(width, min_width)
    expanded_height = max(height, min_height)
    center_x = int(x) + width / 2
    center_y = int(y) + height / 2
    expanded_x = int(round(center_x - expanded_width / 2))
    expanded_y = int(round(center_y - expanded_height / 2))
    return expanded_x, expanded_y, expanded_width, expanded_height


def _agentbay_desktop_coordinate_note(
    screen_note: str,
    image_width: int | None = None,
    image_height: int | None = None,
    crop: tuple[int, int, int, int] | None = None,
) -> str:
    parts = []
    if screen_note:
        parts.append(f"Cloud Desktop coordinate system for mouse tools: {screen_note}.")
    if image_width and image_height:
        parts.append(f"Latest screenshot pixel size: width={image_width}, height={image_height}.")
    if crop:
        x, y, width, height = crop
        parts.append(
            f"Precision crop shown to vision: absolute origin=({x}, {y}), size={width}x{height}. "
            "Grid labels in the crop are absolute Cloud Desktop coordinates, not crop-local coordinates."
        )
    if parts:
        parts.append(
            "The injected analysis image includes a coordinate grid; use the grid labels to choose the center of the target. "
            "Before clicking dialog buttons, text buttons, tabs, menus, checkboxes, close buttons, small controls, "
            "or any target whose center is not unambiguous, take a precision screenshot around that target area. "
            "For popup dismissal, prefer agentbay_computer_dismiss_dialog before coordinate clicking. "
            "Use absolute desktop pixels from the top-left corner (0, 0); do not use the size of the right-side preview panel."
        )
    return "\n".join(parts)


def _agentbay_normalize_text(value) -> str:
    import re

    return re.sub(r"[^a-z0-9]+", "", str(value or "").lower())


def _agentbay_app_field(app: dict, *keys: str) -> str:
    for key in keys:
        value = app.get(key)
        if value:
            return str(value)
    return ""


def _agentbay_format_apps(apps: list, limit: int = 40) -> str:
    import json

    if not apps:
        return "[]"
    compact_apps = []
    for app in apps[:limit]:
        if isinstance(app, dict):
            compact_apps.append(
                {
                    key: app.get(key)
                    for key in ("name", "start_cmd", "startCmd", "work_directory", "workDirectory", "stop_cmd", "stopCmd")
                    if app.get(key)
                }
            )
        else:
            compact_apps.append(str(app))
    rendered = json.dumps(compact_apps, ensure_ascii=False, indent=2)
    if len(apps) > limit:
        rendered += f"\n... {len(apps) - limit} more app(s) omitted"
    return rendered[:5000]


def _agentbay_find_installed_app_match(query: str, apps: list) -> tuple[dict | None, float]:
    from difflib import SequenceMatcher

    query_norm = _agentbay_normalize_text(query.split()[0] if query else query)
    if not query_norm:
        return None, 0.0

    best_app = None
    best_score = 0.0
    for app in apps:
        if not isinstance(app, dict):
            continue
        fields = [
            _agentbay_app_field(app, "name"),
            _agentbay_app_field(app, "start_cmd", "startCmd"),
            _agentbay_app_field(app, "work_directory", "workDirectory"),
        ]
        for field in fields:
            field_norm = _agentbay_normalize_text(field)
            if not field_norm:
                continue
            if query_norm == field_norm:
                score = 1.0
            elif query_norm in field_norm or field_norm in query_norm:
                score = 0.9
            else:
                score = SequenceMatcher(None, query_norm, field_norm).ratio()
            if score > best_score:
                best_app, best_score = app, score

    return best_app, best_score


def _agentbay_uncertain_start_error(error_message: str) -> bool:
    text = (error_message or "").lower()
    return "may have launched" in text or "no processes found" in text


async def _agentbay_visible_apps_note(client) -> str:
    try:
        visible = await client.computer_list_visible_apps()
        if visible.get("success"):
            apps = visible.get("apps", [])
            return f"Visible applications after the launch attempt ({len(apps)}):\n{_agentbay_format_apps(apps, limit=20)}"
        return f"Could not verify visible applications: {visible.get('error_message', 'Unknown error')}"
    except Exception as e:
        logger.debug(f"[AgentBay] Could not list visible apps after start_app: {e}")
        return f"Could not verify visible applications: {str(e)[:200]}"


async def _agentbay_computer_screenshot(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Take a screenshot of the AgentBay cloud desktop.

    The image is held in the process-level memory cache for LLM vision analysis
    only — no disk write, nothing shown in the user's file manager or chat
    history.
    """
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    focus_x = arguments.get("focus_x")
    focus_y = arguments.get("focus_y")
    focus_width = arguments.get("focus_width")
    focus_height = arguments.get("focus_height")

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_screenshot()

        if not (result.get("success") and result.get("data")):
            return f"Screenshot failed: {result.get('error_message', 'Unknown error')}"

        raw_data = result["data"]

        raw_bytes = _agentbay_normalize_image_bytes(raw_data)
        if raw_bytes is None:
            return "Screenshot captured but data format is unrecognised."

        crop_bounds: tuple[int, int, int, int] | None = None
        crop_scale = 1
        analysis_bytes = raw_bytes
        if (
            focus_x is not None
            and focus_y is not None
            and focus_width is not None
            and focus_height is not None
        ):
            try:
                crop_result = _agentbay_crop_image_bytes(
                    raw_bytes,
                    x=int(round(float(focus_x))),
                    y=int(round(float(focus_y))),
                    width=int(round(float(focus_width))),
                    height=int(round(float(focus_height))),
                )
                if crop_result:
                    analysis_bytes, crop_bounds, crop_scale = crop_result
            except (TypeError, ValueError):
                crop_bounds = None

        # Store in memory only — vision_inject.py will consume it for LLM vision
        from app.services.vision_inject import store_temp_screenshot
        grid_options = {}
        if crop_bounds:
            crop_x, crop_y, crop_width, crop_height = crop_bounds
            grid_options = {
                "origin_x": crop_x,
                "origin_y": crop_y,
                "minor_step": 10,
                "major_step": 50,
                "pixel_scale": crop_scale,
            }
        img_id = store_temp_screenshot(analysis_bytes, grid_options=grid_options)
        logger.info(f"[AgentBay] Desktop screenshot stored in memory (id={img_id})")
        screen_width, screen_height, screen_note = await _agentbay_get_screen_metadata(client)
        image_width, image_height = _agentbay_image_dimensions(raw_bytes)
        coordinate_note = _agentbay_desktop_coordinate_note(
            screen_note,
            image_width or screen_width,
            image_height or screen_height,
            crop=crop_bounds,
        )
        return (
            f"Internal desktop screenshot captured for analysis. [ImageID: {img_id}]\n"
            f"{coordinate_note}\n"
            "TARGETING NOTE: Before clicking dialog buttons, text buttons, tabs, menus, checkboxes, "
            "close buttons, small controls, or any target whose center is not unambiguous, call "
            "agentbay_computer_precision_screenshot around the target and click from that enlarged crop.\n"
            f"NOTE: This screenshot is for LLM vision only and is not saved to the user's workspace."
        )

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Computer screenshot failed for agent {agent_id}")
        return f"Desktop screenshot failed: {str(e)[:200]}"


async def _agentbay_computer_save_screenshot(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Save the current AgentBay cloud desktop screenshot to workspace/screenshots/."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_screenshot()
        if not (result.get("success") and result.get("data")):
            return f"Screenshot save failed: {result.get('error_message', 'Unknown error')}"
        raw_bytes = _agentbay_normalize_image_bytes(result.get("data"))
        if raw_bytes is None:
            return "Screenshot save failed: captured data format is unrecognised."
        screen_width, screen_height, screen_note = await _agentbay_get_screen_metadata(client)
        image_width, image_height = _agentbay_image_dimensions(raw_bytes)
        coordinate_note = _agentbay_desktop_coordinate_note(
            screen_note,
            image_width or screen_width,
            image_height or screen_height,
        )
        saved = _agentbay_save_image_to_workspace(
            agent_id=agent_id,
            ws=ws,
            raw_bytes=raw_bytes,
            prefix="desktop-screenshot",
            label="Desktop Screenshot",
        )
        return f"{saved}\n{coordinate_note}"
    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Computer save screenshot failed for agent {agent_id}")
        return f"Desktop screenshot save failed: {str(e)[:200]}"


async def _agentbay_computer_precision_screenshot(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Take an enlarged precision crop for desktop controls."""
    aliases = {
        "focus_x": "x",
        "focus_y": "y",
        "focus_width": "width",
        "focus_height": "height",
    }
    for alias, canonical in aliases.items():
        if arguments.get(canonical) is None and arguments.get(alias) is not None:
            arguments[canonical] = arguments.get(alias)

    required = ("x", "y", "width", "height")
    missing = [key for key in required if arguments.get(key) is None]
    if missing:
        return (
            f"Missing required precision crop argument(s): {', '.join(missing)}. "
            "Use x, y, width, height for the absolute desktop crop rectangle."
        )

    try:
        requested_x = int(round(float(arguments["x"])))
        requested_y = int(round(float(arguments["y"])))
        requested_width = int(round(float(arguments["width"])))
        requested_height = int(round(float(arguments["height"])))
    except (TypeError, ValueError):
        return (
            "Precision crop failed: x, y, width, and height must be numeric absolute desktop pixels. "
            f"Got x={arguments.get('x')!r}, y={arguments.get('y')!r}, "
            f"width={arguments.get('width')!r}, height={arguments.get('height')!r}."
        )

    expanded_x, expanded_y, expanded_width, expanded_height = _agentbay_expand_precision_crop(
        requested_x,
        requested_y,
        requested_width,
        requested_height,
    )

    precision_args = dict(arguments)
    precision_args["focus_x"] = expanded_x
    precision_args["focus_y"] = expanded_y
    precision_args["focus_width"] = expanded_width
    precision_args["focus_height"] = expanded_height
    result = await _agentbay_computer_screenshot(agent_id, ws, precision_args)
    expansion_note = ""
    if (
        expanded_x,
        expanded_y,
        expanded_width,
        expanded_height,
    ) != (requested_x, requested_y, requested_width, requested_height):
        expansion_note = (
            f"Requested crop ({requested_x}, {requested_y}, {requested_width}x{requested_height}) "
            f"was expanded for context to ({expanded_x}, {expanded_y}, {expanded_width}x{expanded_height}). "
        )
    return (
        "Precision desktop crop captured for accurate targeting. "
        f"{expansion_note}"
        "Use the absolute coordinate labels in this enlarged crop for the next click; click the visual center "
        "of the target and do not reuse a guessed coordinate from the full screenshot.\n"
        f"{result}"
    )


async def _agentbay_computer_click(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Click the mouse at specific coordinates on the desktop."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    x = arguments.get("x", 0)
    y = arguments.get("y", 0)
    button = arguments.get("button", "left")

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        try:
            x = int(round(float(x)))
            y = int(round(float(y)))
        except (TypeError, ValueError):
            return f"Click failed: x and y must be numeric desktop pixel coordinates, got x={x!r}, y={y!r}."

        screen_width, screen_height, screen_note = await _agentbay_get_screen_metadata(client)
        if screen_width and screen_height and not (0 <= x < screen_width and 0 <= y < screen_height):
            return (
                f"Click refused: ({x}, {y}) is outside the Cloud Desktop coordinate system "
                f"({screen_note}). Use coordinates from the latest full desktop screenshot."
            )
        result = await client.computer_click(x, y, button=button)
        if result.get("success"):
            note = f" within {screen_note}" if screen_note else ""
            return (
                f"Clicked at ({x}, {y}) with {button} button{note}. "
                f"This only confirms the mouse event was sent; call agentbay_computer_screenshot to verify the UI changed."
            )
        note = f" Coordinate system: {screen_note}." if screen_note else ""
        return f"Click failed at ({x}, {y}).{note}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer click failed")
        return f"Click failed: {str(e)[:200]}"


async def _agentbay_computer_input_text(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Type text at the current cursor position."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    text = arguments.get("text", "")
    if not text:
        return "Missing required argument 'text'"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_input_text(text)
        if result.get("success"):
            return f"Typed text: {text[:100]}"
        return f"Text input failed"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer input_text failed")
        return f"Text input failed: {str(e)[:200]}"


async def _agentbay_computer_press_keys(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Press keyboard keys or shortcuts."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    keys = arguments.get("keys", [])
    hold = arguments.get("hold", False)

    if not keys:
        return "Missing required argument 'keys'"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_press_keys(keys, hold=hold)
        key_str = "+".join(keys)
        if result.get("success"):
            return f"Pressed keys: {key_str}" + (" (held)" if hold else "")
        return f"Key press failed: {key_str}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer press_keys failed")
        return f"Key press failed: {str(e)[:200]}"


async def _agentbay_computer_scroll(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Scroll the screen at a specific position."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    x = arguments.get("x", 0)
    y = arguments.get("y", 0)
    direction = arguments.get("direction", "down")
    amount = arguments.get("amount", 1)

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_scroll(x, y, direction=direction, amount=amount)
        if result.get("success"):
            return f"Scrolled {direction} by {amount} step(s) at ({x}, {y})"
        return f"Scroll failed"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer scroll failed")
        return f"Scroll failed: {str(e)[:200]}"


async def _agentbay_computer_move_mouse(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Move mouse to coordinates without clicking."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    x = arguments.get("x", 0)
    y = arguments.get("y", 0)

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_move_mouse(x, y)
        if result.get("success"):
            return f"Mouse moved to ({x}, {y})"
        return f"Mouse move failed"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer move_mouse failed")
        return f"Mouse move failed: {str(e)[:200]}"


async def _agentbay_computer_drag_mouse(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Drag mouse from one position to another."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    from_x = arguments.get("from_x", 0)
    from_y = arguments.get("from_y", 0)
    to_x = arguments.get("to_x", 0)
    to_y = arguments.get("to_y", 0)
    button = arguments.get("button", "left")

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_drag_mouse(from_x, from_y, to_x, to_y, button=button)
        if result.get("success"):
            return f"Dragged from ({from_x}, {from_y}) to ({to_x}, {to_y})"
        return f"Drag failed"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer drag_mouse failed")
        return f"Drag failed: {str(e)[:200]}"


async def _agentbay_computer_get_screen_size(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Get the screen resolution."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_get_screen_size()
        if result.get("success"):
            import json
            data = result.get("data")
            data_str = json.dumps(data, ensure_ascii=False) if isinstance(data, (dict, list)) else str(data)
            return f"Screen size: {data_str}"
        return f"Failed to get screen size: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer get_screen_size failed")
        return f"Get screen size failed: {str(e)[:200]}"


async def _agentbay_computer_start_app(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Start an application on the desktop."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    cmd = arguments.get("cmd", "")
    work_dir = arguments.get("work_dir", "")

    if not cmd.strip():
        return "Missing required argument 'cmd'"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_start_app(cmd, work_dir=work_dir)
        if result.get("success"):
            # result.data may contain non-serializable objects (e.g. Process),
            # so convert to string safely instead of json.dumps()
            data = result.get("data")
            if data is not None:
                try:
                    import json
                    data_str = json.dumps(data, ensure_ascii=False, indent=2) if isinstance(data, (dict, list, str, int, float, bool)) else str(data)
                except (TypeError, ValueError):
                    data_str = str(data)
            else:
                data_str = ""
            return f"Application started: {cmd}" + (f"\n\n{data_str[:1000]}" if data_str else "")

        # A launch has already been dispatched. Do not guess another command or
        # perform a second start when the Provider result is failed/unknown.
        return (
            "Failed to start application: "
            f"{result.get('error_message', 'Unknown error')}"
        )
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer start_app failed")
        return f"Start application failed: {str(e)[:200]}"


async def _agentbay_computer_get_installed_apps(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """List installed desktop applications and launch commands."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    start_menu = arguments.get("start_menu", True)
    desktop = arguments.get("desktop", True)
    ignore_system_apps = arguments.get("ignore_system_apps", True)

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_get_installed_apps(
            start_menu=bool(start_menu),
            desktop=bool(desktop),
            ignore_system_apps=bool(ignore_system_apps),
        )
        if result.get("success"):
            apps = result.get("apps", [])
            if not apps:
                return "No installed applications found."
            return (
                f"Installed applications ({len(apps)}). Use the returned start_cmd exactly with "
                f"agentbay_computer_start_app; do not guess app launch commands.\n\n"
                f"{_agentbay_format_apps(apps, limit=80)}"
            )
        return f"Failed to get installed applications: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer get_installed_apps failed")
        return f"Get installed applications failed: {str(e)[:200]}"


async def _agentbay_computer_get_cursor_position(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Get current cursor position."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_get_cursor_position()
        if result.get("success"):
            import json
            data = result.get("data")
            data_str = json.dumps(data, ensure_ascii=False) if isinstance(data, (dict, list)) else str(data)
            return f"Cursor position: {data_str}"
        return f"Failed to get cursor position: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer get_cursor_position failed")
        return f"Get cursor position failed: {str(e)[:200]}"


async def _agentbay_computer_get_active_window(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Get info about the currently active window."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_get_active_window()
        if result.get("success"):
            import json
            window = result.get("window")
            window_str = json.dumps(window, ensure_ascii=False, indent=2) if isinstance(window, dict) else str(window)
            return f"Active window:\n\n{window_str}"
        return f"Failed to get active window: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer get_active_window failed")
        return f"Get active window failed: {str(e)[:200]}"


async def _agentbay_computer_activate_window(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Activate (bring to front) a window by its ID."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    window_id = arguments.get("window_id")
    if window_id is None:
        return "Missing required argument 'window_id'"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_activate_window(int(window_id))
        if result.get("success"):
            return f"Window {window_id} activated (brought to front)"
        return f"Failed to activate window {window_id}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer activate_window failed")
        return f"Activate window failed: {str(e)[:200]}"


async def _agentbay_computer_list_windows(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """List OS-level root windows with IDs and geometry."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    timeout_ms = arguments.get("timeout_ms", 3000)

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_list_windows(timeout_ms=int(timeout_ms))
        if result.get("success"):
            import json
            windows = result.get("windows", [])
            if not windows:
                return "No root windows found."
            windows_str = json.dumps(windows, ensure_ascii=False, indent=2)
            return (
                f"OS-level root desktop windows ({len(windows)}). These window_id values refer to whole "
                f"application windows. Use them for activation, or for closing only when the user explicitly "
                f"asked to close/quit an entire desktop window or app. Do NOT use these IDs for in-app popups, "
                f"modals, embedded marketplace/store panels, browser/app tabs, document tabs, or software-internal "
                f"dialogs; close those with the app UI, Escape, Ctrl+W, or agentbay_computer_dismiss_dialog.\n\n"
                f"{windows_str[:5000]}"
            )
        return f"Failed to list windows: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer list_windows failed")
        return f"List windows failed: {str(e)[:200]}"


async def _agentbay_computer_close_window(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Close an entire OS-level root desktop window/application by explicit ID."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    window_id = arguments.get("window_id")
    title = str(arguments.get("title") or "").strip()

    if window_id is None:
        if not title:
            return (
                "Missing required argument `window_id`. Only use agentbay_computer_close_window when the user "
                "explicitly wants to close or quit an entire OS-level desktop window/application. If the target "
                "is an in-app popup, modal, embedded marketplace/store panel, browser/app tab, document tab, "
                "or software-internal dialog, use app UI controls, Escape, Ctrl+W, or "
                "agentbay_computer_dismiss_dialog instead."
            )

        try:
            _session_id, _run_id = _agentbay_scope_ids(arguments)
            client = await get_agentbay_client_for_agent(
                agent_id,
                "computer",
                session_id=_session_id,
                run_id=_run_id,
            )
            windows_result = await client.computer_list_windows()
            if not windows_result.get("success"):
                return f"Failed to list windows before closing: {windows_result.get('error_message', 'Unknown error')}"

            from difflib import SequenceMatcher
            import json

            title_norm = _agentbay_normalize_text(title)
            candidates: list[dict] = []
            for window in windows_result.get("windows", []):
                if not isinstance(window, dict):
                    continue
                candidate = str(window.get("title") or window.get("window_title") or "")
                candidate_norm = _agentbay_normalize_text(candidate)
                if not candidate_norm:
                    continue
                if title_norm in candidate_norm or candidate_norm in title_norm:
                    score = 0.95
                else:
                    score = SequenceMatcher(None, title_norm, candidate_norm).ratio()
                if score >= 0.35:
                    item = dict(window)
                    item["match_score"] = round(score, 3)
                    candidates.append(item)
            candidates.sort(key=lambda item: item.get("match_score", 0), reverse=True)
            return (
                f"Refusing to close by title-only match for `{title}` because it can close the wrong application. "
                f"The candidates below are whole OS-level root windows. Choose a root window_id only if the user "
                f"explicitly wants to close/quit that entire application window. For in-app popups, modals, "
                f"embedded marketplace/store panels, browser/app tabs, document tabs, or software-internal dialogs, "
                f"do not close a root window; use app UI controls, Escape, Ctrl+W, or "
                f"agentbay_computer_dismiss_dialog instead.\n\n"
                f"{json.dumps(candidates[:8], ensure_ascii=False, indent=2)[:3000]}"
            )
        except RuntimeError as e:
            return f"{str(e)}"
        except Exception as e:
            logger.exception(f"[AgentBay] Computer close_window candidate lookup failed")
            return f"Close window requires window_id. Candidate lookup failed: {str(e)[:200]}"

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_close_window(int(window_id))
        if result.get("success"):
            return (
                f"Closed OS-level root desktop window {window_id}; the whole application window may now be gone. "
                f"Call agentbay_computer_screenshot to verify."
            )
        return f"Failed to close window {window_id}: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer close_window failed")
        return f"Close window failed: {str(e)[:200]}"


async def _agentbay_computer_dismiss_dialog(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Safely dismiss the current in-app popup/dialog without closing root windows."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    title = str(arguments.get("title") or "").strip()
    window_id = arguments.get("window_id")

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)

        if window_id is not None:
            return (
                "agentbay_computer_dismiss_dialog does not close root desktop windows. "
                "It only sends Escape to the active in-app popup/dialog. "
                "For in-app tabs, embedded panels, marketplace/store windows, or document tabs, use the app UI "
                "or shortcuts such as Ctrl+W. If the user explicitly wants to close/quit a whole desktop window "
                "or app, call agentbay_computer_close_window with a window_id returned by "
                "agentbay_computer_list_windows."
            )

        esc_result = await client.computer_press_keys(["esc"])
        if esc_result.get("success"):
            title_note = f" Target hint: `{title}`." if title else ""
            return (
                f"Sent Escape to safely dismiss the active in-app popup/dialog.{title_note} "
                f"Call agentbay_computer_screenshot to verify. This tool never closes the root application window; "
                f"if Escape does not affect an in-app tab or embedded panel, use that app's own close control "
                f"or a shortcut such as Ctrl+W instead of root-window close."
            )

        return (
            f"Could not send Escape to dismiss the active popup/dialog: "
            f"{esc_result.get('error_message', 'Unknown error')}. "
            f"Do not use this tool to close root application windows."
        )
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer dismiss_dialog failed")
        return f"Dismiss dialog failed: {str(e)[:200]}"


async def _agentbay_computer_list_visible_apps(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """List currently visible/running applications."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id, _run_id = _agentbay_scope_ids(arguments)
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id, run_id=_run_id)
        result = await client.computer_list_visible_apps()
        if result.get("success"):
            import json
            apps = result.get("apps", [])
            if not apps:
                return "No visible applications running."
            apps_str = json.dumps(apps, ensure_ascii=False, indent=2)
            return f"Visible applications ({len(apps)}):\n\n{apps_str[:3000]}"
        return f"Failed to list applications: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer list_visible_apps failed")
        return f"List applications failed: {str(e)[:200]}"


async def _agentbay_file_transfer(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Transfer a file between workspace and an AgentBay environment, or between two environments.

    Supported transfer directions:
      - workspace  → env:      upload_file(local_workspace_path, remote_path)   [single SDK call]
      - env        → workspace: download_file(remote_path, local_workspace_path) [single SDK call]
      - env A      → env B:    download to /tmp/<uuid>, upload to env B, cleanup /tmp [transparent]

    The 'local' side of the SDK calls is always the Clawith backend server,
    which has access to the agent workspace directory.
    """
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    from_type = arguments.get("from_type", "")
    from_path = arguments.get("from_path", "")
    to_type   = arguments.get("to_type", "")
    to_path   = arguments.get("to_path", "")
    session_id, run_id = _agentbay_scope_ids(arguments)

    if not all([from_type, from_path, to_type, to_path]):
        return "Missing required parameters: from_type, from_path, to_type, to_path"

    # Reject no-op transfers
    if from_type == "workspace" and to_type == "workspace":
        return "Cannot transfer workspace → workspace. Use write_file or workspace tools instead."
    if from_type == to_type and from_type != "workspace":
        return f"Same environment ({from_type}) transfer: use agentbay_command_exec with 'cp' to copy files within the same environment."

    env_types = {"browser", "computer", "code"}

    # ── Helper: resolve and validate a workspace-relative path ──────────────
    def resolve_workspace(rel_path: str) -> tuple[str | None, str]:
        """Return (absolute_local_path_str, error_message). error_message is '' on success."""
        local = (ws / rel_path).resolve()
        if not str(local).startswith(str(ws.resolve())):
            return None, "Permission denied: path must be inside the agent workspace"
        return str(local), ""

    try:
        # ── Case 1: workspace → env ──────────────────────────────────────────
        if from_type == "workspace" and to_type in env_types:
            local_path, err = resolve_workspace(from_path)
            if err:
                return err
            import os
            if not os.path.exists(local_path):
                return f"File not found in workspace: {from_path}"
            client = await get_agentbay_client_for_agent(
                agent_id,
                to_type,
                session_id=session_id,
                run_id=run_id,
            )
            result = await asyncio.to_thread(
                client._session.file_system.upload_file,
                local_path, to_path
            )
            if result.success:
                msg = (
                    f"Transferred workspace/{from_path} → [{to_type}]{to_path} "
                    f"({result.bytes_sent} bytes)"
                )
                # After uploading to the computer desktop directory, notify the GNOME
                # file manager so the file icon appears immediately without manual refresh.
                desktop_dir = "/home/wuying/桌面"
                if to_type == "computer" and to_path.startswith(desktop_dir):
                    try:
                        await asyncio.to_thread(
                            client._session.command.exec,
                            f"DISPLAY=:0 gio info '{to_path}' 2>/dev/null || true"
                        )
                    except Exception:
                        pass  # Non-critical: desktop refresh failure doesn't affect transfer result
                return msg
            return f"Upload failed: {result.error_message}"

        # ── Case 2: env → workspace ──────────────────────────────────────────
        elif from_type in env_types and to_type == "workspace":
            local_path, err = resolve_workspace(to_path)
            if err:
                return err
            import os
            os.makedirs(os.path.dirname(local_path) or ".", exist_ok=True)
            client = await get_agentbay_client_for_agent(
                agent_id,
                from_type,
                session_id=session_id,
                run_id=run_id,
            )
            result = await asyncio.to_thread(
                client._session.file_system.download_file,
                from_path, local_path
            )
            if result.success:
                return (
                    f"Transferred [{from_type}]{from_path} → workspace/{to_path} "
                    f"({result.bytes_received} bytes). "
                    f"File available in workspace at: {to_path}"
                )
            return f"Download failed: {result.error_message}"

        # ── Case 3: env A → env B (transparent /tmp/ intermediary) ──────────
        elif from_type in env_types and to_type in env_types:
            import uuid as _uuid
            import os
            tmp_path = f"/tmp/agentbay_transfer_{_uuid.uuid4().hex}"
            try:
                # Step 1: download from source env to backend /tmp/
                src_client = await get_agentbay_client_for_agent(
                    agent_id,
                    from_type,
                    session_id=session_id,
                    run_id=run_id,
                )
                dl_result = await asyncio.to_thread(
                    src_client._session.file_system.download_file,
                    from_path, tmp_path
                )
                if not dl_result.success:
                    return f"Transfer failed (download from {from_type}): {dl_result.error_message}"

                # Step 2: upload from backend /tmp/ to destination env
                dst_client = await get_agentbay_client_for_agent(
                    agent_id,
                    to_type,
                    session_id=session_id,
                    run_id=run_id,
                )
                ul_result = await asyncio.to_thread(
                    dst_client._session.file_system.upload_file,
                    tmp_path, to_path
                )
                if not ul_result.success:
                    return f"Transfer failed (upload to {to_type}): {ul_result.error_message}"

                return (
                    f"Transferred [{from_type}]{from_path} → [{to_type}]{to_path} "
                    f"({dl_result.bytes_received} bytes)"
                )
            finally:
                # Always clean up the temporary file regardless of success or failure
                try:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)
                except Exception:
                    pass  # Non-critical: ignore cleanup errors

        else:
            return f"Unsupported transfer: {from_type} → {to_type}"

    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] File transfer failed for agent {agent_id}")
        return f"File transfer failed: {str(e)[:200]}"


# ─── OKR Tools ───────────────────────────────────────────────────────────────


async def _get_agent_owner_info(agent_id: uuid.UUID) -> tuple[str, str]:
    """Return (owner_type, owner_id_str) for the calling agent.

    Used by get_my_okr and update_kr_progress to scope queries to the
    correct owner without requiring the caller to pass their own ID.
    """
    from app.database import async_session
    from app.models.agent import Agent
    from sqlalchemy import select as _select

    async with async_session() as db:
        result = await db.execute(_select(Agent).where(Agent.id == agent_id))
        agent = result.scalar_one_or_none()
    if not agent:
        return "agent", str(agent_id)
    return "agent", str(agent_id)


def _compute_okr_period_bounds(frequency: str, length_days: int | None):
    """Return the current OKR period using the tenant's configured cadence."""
    from datetime import date, timedelta

    today = date.today()
    if frequency == "monthly":
        start = today.replace(day=1)
        if today.month == 12:
            end = today.replace(month=12, day=31)
        else:
            end = today.replace(month=today.month + 1, day=1) - timedelta(days=1)
    elif frequency == "custom" and length_days:
        epoch = date(1970, 1, 1)
        days_since_epoch = (today - epoch).days
        period_index = days_since_epoch // length_days
        start = epoch + timedelta(days=period_index * length_days)
        end = start + timedelta(days=length_days - 1)
    else:
        quarter = (today.month - 1) // 3 + 1
        start = date(today.year, (quarter - 1) * 3 + 1, 1)
        if quarter == 4:
            end = date(today.year, 12, 31)
        else:
            end = date(today.year, quarter * 3 + 1, 1) - timedelta(days=1)
    return start, end


def _explicit_okr_period(
    arguments: Mapping[str, object],
) -> tuple[object | None, object | None, str | None]:
    """Parse a caller-supplied OKR range, requiring both dates or neither."""
    from datetime import date

    period_start = arguments.get("period_start")
    period_end = arguments.get("period_end")
    has_start = period_start is not None
    has_end = period_end is not None
    if has_start != has_end:
        return (
            None,
            None,
            "period_start and period_end must be provided together.",
        )
    if not has_start:
        return None, None, None
    if not (
        isinstance(period_start, str)
        and period_start.strip()
        and isinstance(period_end, str)
        and period_end.strip()
    ):
        return None, None, "period_start and period_end must be ISO dates."
    try:
        start = date.fromisoformat(period_start.strip())
        end = date.fromisoformat(period_end.strip())
    except ValueError:
        return None, None, "period_start and period_end must use YYYY-MM-DD."
    if start > end:
        return None, None, "period_start must be on or before period_end."
    return start, end, None


_OKR_KR_STATUSES = frozenset(
    {"on_track", "at_risk", "behind", "completed"}
)
_OKR_OBJECTIVE_STATUSES = frozenset(
    {"draft", "active", "completed", "archived"}
)


def _okr_uuid(value: object, field: str) -> tuple[uuid.UUID | None, ToolExecutionOutcome | None]:
    if not isinstance(value, str) or not value.strip():
        return None, _typed_failure(
            f"{field} must be a non-empty UUID string.",
            "invalid_tool_arguments",
        )
    try:
        return uuid.UUID(value.strip()), None
    except ValueError:
        return None, _typed_failure(
            f"{field} must be a UUID.",
            "invalid_tool_arguments",
        )


def _okr_finite_number(
    value: object,
    field: str,
) -> tuple[float | None, ToolExecutionOutcome | None]:
    if isinstance(value, bool):
        return None, _typed_failure(
            f"{field} must be a finite number.",
            "invalid_tool_arguments",
        )
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None, _typed_failure(
            f"{field} must be a finite number.",
            "invalid_tool_arguments",
        )
    if not math.isfinite(number):
        return None, _typed_failure(
            f"{field} must be a finite number.",
            "invalid_tool_arguments",
        )
    return number, None


def _okr_progress_status(value: float, target: float) -> str:
    if target == 0:
        return "completed" if value >= 0 else "behind"
    ratio = value / target
    if ratio >= 1.0:
        return "completed"
    if ratio >= 0.7:
        return "on_track"
    if ratio >= 0.4:
        return "at_risk"
    return "behind"


async def _require_designated_okr_agent(
    agent_id: uuid.UUID | None,
) -> ToolExecutionOutcome | None:
    if agent_id is None:
        return _typed_failure(
            "This OKR tool requires Agent context.",
            "invalid_tool_arguments",
        )
    if not await _agent_is_designated_okr_agent(agent_id):
        return _typed_failure(
            "Only the tenant's designated OKR Agent may use this tool.",
            "okr_agent_permission_denied",
        )
    return None


def _okr_period_ref(
    tenant_id: object,
    period_start: date,
    period_end: date,
    *,
    owner_id: uuid.UUID | None = None,
) -> str:
    owner_suffix = f"/owner/{owner_id}" if owner_id else ""
    return (
        f"okr://tenant/{tenant_id}/period/{period_start.isoformat()}"
        f"/{period_end.isoformat()}{owner_suffix}"
    )


async def _get_okr_outcome(
    agent_id: uuid.UUID | None,
    arguments: dict,
    *,
    own_only: bool,
) -> ToolExecutionOutcome:
    if agent_id is None:
        return _typed_failure(
            "OKR reads require Agent context.",
            "invalid_tool_arguments",
        )
    explicit_start, explicit_end, period_error = _explicit_okr_period(arguments)
    if period_error:
        return _typed_failure(period_error, "invalid_tool_arguments")

    try:
        from app.models.agent import Agent as AgentModel
        from app.models.okr import OKRKeyResult, OKRObjective, OKRSettings

        async with async_session() as db:
            agent_result = await db.execute(
                select(AgentModel).where(AgentModel.id == agent_id)
            )
            agent = agent_result.scalar_one_or_none()
            if agent is None:
                return _typed_failure(
                    "Agent not found.",
                    "source_agent_not_found",
                )

            settings_result = await db.execute(
                select(OKRSettings).where(
                    OKRSettings.tenant_id == agent.tenant_id
                )
            )
            settings = settings_result.scalar_one_or_none()
            if settings is None or not settings.enabled:
                return _typed_failure(
                    "OKR is not enabled for this organization.",
                    "okr_not_enabled",
                )

            if explicit_start is not None and explicit_end is not None:
                period_start, period_end = explicit_start, explicit_end
            else:
                period_start, period_end = _compute_okr_period_bounds(
                    settings.period_frequency,
                    settings.period_length_days,
                )

            objective_query = select(OKRObjective).where(
                OKRObjective.tenant_id == agent.tenant_id,
                OKRObjective.period_start >= period_start,
                OKRObjective.period_end <= period_end,
                OKRObjective.status != "archived",
            )
            if own_only:
                objective_query = objective_query.where(
                    OKRObjective.owner_type == "agent",
                    OKRObjective.owner_id == agent_id,
                )
            objective_result = await db.execute(
                objective_query.order_by(OKRObjective.created_at)
            )
            objectives = objective_result.scalars().all()
            result_ref = _okr_period_ref(
                agent.tenant_id,
                period_start,
                period_end,
                owner_id=agent_id if own_only else None,
            )
            if not objectives:
                scope = "your" if own_only else "organization"
                return _typed_success(
                    f"No {scope} OKRs found for {period_start.isoformat()} through {period_end.isoformat()}.",
                    result_ref=result_ref,
                    metadata={
                        "period_start": period_start.isoformat(),
                        "period_end": period_end.isoformat(),
                        "objective_count": 0,
                        "kr_count": 0,
                    },
                )

            objective_ids = [objective.id for objective in objectives]
            kr_result = await db.execute(
                select(OKRKeyResult)
                .where(OKRKeyResult.objective_id.in_(objective_ids))
                .order_by(OKRKeyResult.created_at)
            )
            key_results = kr_result.scalars().all()
            key_results_by_objective: dict[str, list] = {}
            for key_result in key_results:
                key_results_by_objective.setdefault(
                    str(key_result.objective_id), []
                ).append(key_result)

            lines = [
                f"OKRs for {period_start.isoformat()} through {period_end.isoformat()}:"
            ]
            for objective in objectives:
                lines.append(
                    f"Objective {objective.id}: {objective.title} [{objective.status}]"
                )
                for key_result in key_results_by_objective.get(
                    str(objective.id), []
                ):
                    lines.append(
                        "  KR "
                        f"{key_result.id}: {key_result.title} — "
                        f"{key_result.current_value}/{key_result.target_value} "
                        f"{key_result.unit or ''} [{key_result.status}]"
                    )
            return _typed_success(
                "\n".join(lines),
                result_ref=result_ref,
                metadata={
                    "period_start": period_start.isoformat(),
                    "period_end": period_end.isoformat(),
                    "objective_count": len(objectives),
                    "kr_count": len(key_results),
                },
            )
    except Exception as exc:
        logger.exception("[OKR] typed OKR read failed")
        return _typed_failure(
            f"OKR read failed: {type(exc).__name__}.",
            "okr_read_failed",
            retryable=True,
        )


async def _get_okr_settings_outcome(
    agent_id: uuid.UUID | None,
) -> ToolExecutionOutcome:
    designated_error = await _require_designated_okr_agent(agent_id)
    if designated_error is not None:
        return designated_error
    assert agent_id is not None

    try:
        from app.models.agent import Agent as AgentModel
        from app.services.okr_scheduler import get_okr_settings_for_agent

        async with async_session() as db:
            agent_result = await db.execute(
                select(AgentModel).where(AgentModel.id == agent_id)
            )
            agent = agent_result.scalar_one_or_none()
            if agent is None:
                return _typed_failure(
                    "Agent not found.",
                    "source_agent_not_found",
                )
        settings = await get_okr_settings_for_agent(agent.tenant_id)
        summary = json.dumps(settings, ensure_ascii=False, sort_keys=True, default=str)
        if len(summary.encode("utf-8")) > 8192:
            summary = summary.encode("utf-8")[:8192].decode(
                "utf-8", errors="ignore"
            )
        return _typed_success(
            summary,
            result_ref=f"okr-settings://tenant/{agent.tenant_id}",
            metadata={"tenant_id": str(agent.tenant_id)},
        )
    except Exception as exc:
        logger.exception("[OKR] typed settings read failed")
        return _typed_failure(
            f"OKR settings read failed: {type(exc).__name__}.",
            "okr_settings_read_failed",
            retryable=True,
        )


async def _okr_transaction_outcome(
    tool_name: str,
    agent_id: uuid.UUID | None,
    user_id: uuid.UUID | None,
    arguments: dict,
) -> ToolExecutionOutcome:
    if tool_name == "get_okr":
        return await _get_okr_outcome(agent_id, arguments, own_only=False)
    if tool_name == "get_my_okr":
        return await _get_okr_outcome(agent_id, arguments, own_only=True)
    if tool_name == "get_okr_settings":
        return await _get_okr_settings_outcome(agent_id)
    if tool_name in {"update_kr_progress", "update_any_kr_progress"}:
        return await _update_kr_progress_outcome(
            agent_id,
            user_id,
            arguments,
            any_owner=tool_name == "update_any_kr_progress",
        )
    if tool_name == "update_kr_content":
        return await _update_kr_content_outcome(
            agent_id,
            user_id,
            arguments,
        )
    if tool_name == "create_objective":
        return await _create_objective_outcome(agent_id, user_id, arguments)
    if tool_name == "create_key_result":
        return await _create_key_result_outcome(agent_id, user_id, arguments)
    if tool_name == "update_objective":
        return await _update_objective_outcome(agent_id, user_id, arguments)
    if tool_name == "upsert_member_daily_report":
        return await _upsert_member_daily_report_outcome(
            agent_id,
            user_id,
            arguments,
        )
    return _typed_failure(
        f"Unsupported OKR transaction tool: {tool_name}.",
        "unsupported_tool",
    )


async def _get_okr(agent_id: uuid.UUID | None, arguments: dict) -> str:
    """Return the full OKR board for the current period as formatted text.

    Includes company-level O+KR and every member's individual O+KR.
    This is a read-only tool available to all agents.
    """
    # Resolve tenant_id from the calling agent
    if not agent_id:
        return "OKR tools require agent context."
    explicit_start, explicit_end, period_error = _explicit_okr_period(
        arguments
    )
    if period_error:
        return period_error

    try:
        from app.database import async_session
        from app.models.agent import Agent
        from app.models.okr import OKRObjective, OKRKeyResult, OKRSettings
        from app.models.org import OrgMember
        from app.models.user import User
        from sqlalchemy import select as _select

        async with async_session() as db:
            # Look up the agent's tenant
            agent_result = await db.execute(_select(Agent).where(Agent.id == agent_id))
            agent = agent_result.scalar_one_or_none()
            if not agent:
                return "Agent not found."

            tenant_id = agent.tenant_id

            # Get OKR settings to determine period
            settings_result = await db.execute(
                _select(OKRSettings).where(OKRSettings.tenant_id == tenant_id)
            )
            settings = settings_result.scalar_one_or_none()

            if not settings or not settings.enabled:
                return "OKR is not enabled for your organization."

            # Compute period bounds
            if explicit_start is not None and explicit_end is not None:
                ps = explicit_start
                pe = explicit_end
            else:
                ps, pe = _compute_okr_period_bounds(
                    settings.period_frequency,
                    settings.period_length_days,
                )

            # Fetch all active objectives
            obj_result = await db.execute(
                _select(OKRObjective).where(
                    OKRObjective.tenant_id == tenant_id,
                    OKRObjective.period_start >= ps,
                    OKRObjective.period_end <= pe,
                    OKRObjective.status != "archived",
                ).order_by(OKRObjective.owner_type, OKRObjective.created_at)
            )
            objectives = obj_result.scalars().all()

            if not objectives:
                return f"No OKRs found for the current period ({ps} – {pe})."

            # Fetch all KRs
            obj_ids = [o.id for o in objectives]
            kr_result = await db.execute(
                _select(OKRKeyResult)
                .where(OKRKeyResult.objective_id.in_(obj_ids))
                .order_by(OKRKeyResult.created_at)
            )
            all_krs = kr_result.scalars().all()

            krs_by_obj: dict = {}
            for kr in all_krs:
                krs_by_obj.setdefault(str(kr.objective_id), []).append(kr)

            # Resolve readable owner names so the OKR Agent can reason about
            # members by display name instead of raw UUIDs.
            user_owner_ids = [
                o.owner_id for o in objectives
                if o.owner_type == "user" and o.owner_id
            ]
            agent_owner_ids = [
                o.owner_id for o in objectives
                if o.owner_type == "agent" and o.owner_id
            ]

            user_names: dict[uuid.UUID, str] = {}
            if user_owner_ids:
                u_result = await db.execute(
                    _select(User.id, User.display_name).where(User.id.in_(user_owner_ids))
                )
                user_names = {
                    row.id: (row.display_name or "")
                    for row in u_result.fetchall()
                }

                unresolved_ids = [oid for oid in user_owner_ids if oid not in user_names]
                if unresolved_ids:
                    m_result = await db.execute(
                        _select(OrgMember.id, OrgMember.name).where(
                            OrgMember.id.in_(unresolved_ids)
                        )
                    )
                    for row in m_result.fetchall():
                        user_names[row.id] = row.name or ""

            agent_names: dict[uuid.UUID, str] = {}
            if agent_owner_ids:
                a_result = await db.execute(
                    _select(Agent.id, Agent.name).where(Agent.id.in_(agent_owner_ids))
                )
                agent_names = {
                    row.id: (row.name or "")
                    for row in a_result.fetchall()
                }

            def _resolve_owner_label(obj: OKRObjective) -> str:
                if obj.owner_type == "company":
                    return "Company"
                if not obj.owner_id:
                    return f"{obj.owner_type}:unassigned"
                if obj.owner_type == "user":
                    return user_names.get(obj.owner_id) or f"user:{obj.owner_id}"
                if obj.owner_type == "agent":
                    return agent_names.get(obj.owner_id) or f"agent:{obj.owner_id}"
                return f"{obj.owner_type}:{obj.owner_id}"

        # Format output
        lines = [f"# OKR Board — {ps} to {pe}\n"]

        company_objs = [o for o in objectives if o.owner_type == "company"]
        member_objs = [o for o in objectives if o.owner_type != "company"]

        if company_objs:
            lines.append("## Company Objectives")
            for o in company_objs:
                krs = krs_by_obj.get(str(o.id), [])
                pct = 0
                if krs:
                    pct = int(sum(min(k.current_value / k.target_value, 1) for k in krs) / len(krs) * 100)
                lines.append(f"\n**O: {o.title}** [{pct}%]  objective_id={o.id}")
                for kr in krs:
                    lines.append(
                        f"  - KR ({kr.status}): {kr.title}  "
                        f"[{kr.current_value}/{kr.target_value} {kr.unit or ''}]  "
                        f" kr_id={kr.id}"
                    )

        if member_objs:
            lines.append("\n## Member Objectives")
            for o in member_objs:
                owner_label = _resolve_owner_label(o)
                krs = krs_by_obj.get(str(o.id), [])
                lines.append(f"\n**{owner_label}** | O: {o.title}  objective_id={o.id}")
                for kr in krs:
                    lines.append(
                        f"  - KR ({kr.status}): {kr.title}  "
                        f"[{kr.current_value}/{kr.target_value} {kr.unit or ''}]  "
                        f" kr_id={kr.id}"
                    )

        return "\n".join(lines)

    except Exception as e:
        logger.exception(f"[OKR] get_okr failed for agent {agent_id}")
        return f"Failed to retrieve OKR data: {str(e)[:200]}"


async def _get_my_okr(agent_id: uuid.UUID | None, arguments: dict) -> str:
    """Return the calling agent's own Objectives and KRs.

    Includes objective_id and kr_id values so the agent can update existing OKRs
    instead of accidentally creating duplicate ones.
    """
    if not agent_id:
        return "OKR tools require agent context."
    explicit_start, explicit_end, period_error = _explicit_okr_period(
        arguments
    )
    if period_error:
        return period_error

    try:
        from app.database import async_session
        from app.models.agent import Agent
        from app.models.okr import OKRObjective, OKRKeyResult, OKRSettings
        from sqlalchemy import select as _select

        async with async_session() as db:
            agent_result = await db.execute(_select(Agent).where(Agent.id == agent_id))
            agent = agent_result.scalar_one_or_none()
            if not agent:
                return "Agent not found."

            settings_result = await db.execute(
                _select(OKRSettings).where(OKRSettings.tenant_id == agent.tenant_id)
            )
            settings = settings_result.scalar_one_or_none()
            if not settings or not settings.enabled:
                return "OKR is not enabled for your organization."

            if explicit_start is not None and explicit_end is not None:
                ps = explicit_start
                pe = explicit_end
            else:
                ps, pe = _compute_okr_period_bounds(
                    settings.period_frequency,
                    settings.period_length_days,
                )

            obj_result = await db.execute(
                _select(OKRObjective).where(
                    OKRObjective.tenant_id == agent.tenant_id,
                    OKRObjective.owner_type == "agent",
                    OKRObjective.owner_id == agent_id,
                    OKRObjective.period_start >= ps,
                    OKRObjective.period_end <= pe,
                    OKRObjective.status != "archived",
                )
            )
            objectives = obj_result.scalars().all()

            if not objectives:
                return (
                    f"You have no OKRs set for the current period ({ps} – {pe}). "
                    "Contact the OKR Agent to set up your Objectives and Key Results."
                )

            obj_ids = [o.id for o in objectives]
            kr_result = await db.execute(
                _select(OKRKeyResult)
                .where(OKRKeyResult.objective_id.in_(obj_ids))
                .order_by(OKRKeyResult.created_at)
            )
            all_krs = kr_result.scalars().all()

            krs_by_obj: dict = {}
            for kr in all_krs:
                krs_by_obj.setdefault(str(kr.objective_id), []).append(kr)

        lines = [
            f"# My OKRs — {ps} to {pe}\n",
            "If you need to revise an existing OKR, reuse the IDs below:",
            "- change Objective title/description/status with update_objective(objective_id=...)",
            "- change KR title/target/unit/focus/status with update_kr_content(kr_id=...)",
            "- change KR numeric progress with update_kr_progress(kr_id=...)",
            "",
        ]
        for o in objectives:
            krs = krs_by_obj.get(str(o.id), [])
            lines.append(f"**O: {o.title}**  objective_id={o.id}")
            if o.description:
                lines.append(f"  {o.description}")
            for kr in krs:
                lines.append(
                    f"  - [{kr.status}] {kr.title}  "
                    f"Progress: {kr.current_value}/{kr.target_value} {kr.unit or ''}  "
                    f"  kr_id={kr.id}"
                )
        return "\n".join(lines)

    except Exception as e:
        logger.exception(f"[OKR] get_my_okr failed for agent {agent_id}")
        return f"Failed to retrieve your OKR: {str(e)[:200]}"


async def _load_okr_request_context(
    db,
    agent_id: uuid.UUID,
    user_id: uuid.UUID | None,
) -> dict:
    from app.models.agent import Agent as AgentModel
    from app.models.okr import OKRSettings
    from app.models.user import User as UserModel

    ag_res = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
    agent = ag_res.scalar_one_or_none()
    requester = None
    if user_id:
        user_res = await db.execute(select(UserModel).where(UserModel.id == user_id))
        requester = user_res.scalar_one_or_none()
    designated_okr_agent_id = None
    if agent:
        settings_res = await db.execute(
            select(OKRSettings.okr_agent_id).where(
                OKRSettings.tenant_id == agent.tenant_id
            )
        )
        designated_okr_agent_id = settings_res.scalar_one_or_none()

    return {
        "agent": agent,
        "tenant_id": getattr(agent, "tenant_id", None),
        "agent_is_system": bool(agent and agent.is_system),
        "agent_is_designated_okr_agent": bool(
            agent and designated_okr_agent_id == agent.id
        ),
        "requester": requester,
        "requester_user_id": user_id,
        "requester_is_admin": bool(requester and requester.role in ("org_admin", "platform_admin")),
    }


def _okr_permission_denied(message: str) -> str:
    return f"Permission denied: {message}"


def _can_access_existing_okr_target(ctx: dict, owner_type: str, owner_id: uuid.UUID | None) -> str | None:
    if ctx.get("agent_is_designated_okr_agent", False):
        if ctx["requester_is_admin"]:
            return None
        if owner_type != "user" or owner_id != ctx["requester_user_id"]:
            return _okr_permission_denied(
                "non-admin requests may only create or modify the requester's own personal OKRs. "
                "Do not create or edit company OKRs or other members' OKRs."
            )
        return None

    if owner_type != "agent" or owner_id != ctx["agent"].id:
        return _okr_permission_denied(
            "you can only create or modify your own agent OKRs."
        )
    return None


def _can_create_okr_target(ctx: dict, owner_type: str, owner_id: uuid.UUID | None) -> str | None:
    if ctx.get("agent_is_designated_okr_agent", False):
        if ctx["requester_is_admin"]:
            return None
        if owner_type != "user" or owner_id != ctx["requester_user_id"]:
            return _okr_permission_denied(
                "non-admin requests may only create the requester's own personal OKRs. "
                "Creating company OKRs or other members' OKRs requires an org admin."
            )
        return None

    if owner_type != "agent" or owner_id != ctx["agent"].id:
        return _okr_permission_denied(
            "you can only create OKRs for yourself."
        )
    return None


async def _update_kr_progress_outcome(
    agent_id: uuid.UUID | None,
    user_id: uuid.UUID | None,
    arguments: dict,
    *,
    any_owner: bool,
) -> ToolExecutionOutcome:
    if agent_id is None:
        return _typed_failure(
            "KR progress updates require Agent context.",
            "invalid_tool_arguments",
        )
    kr_id, argument_error = _okr_uuid(arguments.get("kr_id"), "kr_id")
    if argument_error is not None:
        return argument_error
    value, argument_error = _okr_finite_number(arguments.get("value"), "value")
    if argument_error is not None:
        return argument_error
    status = arguments.get("status")
    if status is not None and status not in _OKR_KR_STATUSES:
        return _typed_failure(
            "status is not a supported KR status.",
            "invalid_tool_arguments",
        )
    note = arguments.get("note")
    if note is not None and not isinstance(note, str):
        return _typed_failure(
            "note must be a string.",
            "invalid_tool_arguments",
        )
    if any_owner:
        designated_error = await _require_designated_okr_agent(agent_id)
        if designated_error is not None:
            return designated_error

    assert kr_id is not None and value is not None
    result_ref = str(kr_id)
    commit_started = False
    metadata: dict[str, object] = {"kr_id": result_ref}
    try:
        from app.models.okr import OKRKeyResult, OKRObjective, OKRProgressLog

        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if ctx.get("agent") is None:
                return _typed_failure(
                    "Agent not found.",
                    "source_agent_not_found",
                    result_ref=result_ref,
                )
            if any_owner:
                ctx = dict(ctx)
                ctx["agent_is_designated_okr_agent"] = True

            result = await db.execute(
                select(OKRKeyResult, OKRObjective)
                .join(
                    OKRObjective,
                    OKRKeyResult.objective_id == OKRObjective.id,
                )
                .where(
                    OKRKeyResult.id == kr_id,
                    OKRObjective.tenant_id == ctx["tenant_id"],
                )
            )
            row = result.first()
            if row is None:
                return _typed_failure(
                    f"Key Result {kr_id} was not found.",
                    "key_result_not_found",
                    result_ref=result_ref,
                )
            key_result, objective = row
            permission_error = _can_access_existing_okr_target(
                ctx,
                objective.owner_type,
                objective.owner_id,
            )
            if permission_error:
                return _typed_failure(
                    permission_error,
                    "okr_permission_denied",
                    result_ref=result_ref,
                )

            previous_value = float(key_result.current_value)
            key_result.current_value = value
            key_result.status = status or _okr_progress_status(
                value,
                float(key_result.target_value),
            )
            key_result.last_updated_at = datetime.now(timezone.utc)
            progress_log_id = uuid.uuid4()
            progress_log = OKRProgressLog(
                id=progress_log_id,
                kr_id=kr_id,
                previous_value=previous_value,
                new_value=value,
                source="okr_agent" if any_owner else "self_report",
                note=note,
            )
            db.add(progress_log)
            metadata.update(
                {
                    "progress_log_id": str(progress_log_id),
                    "previous_value": previous_value,
                    "current_value": value,
                    "target_value": float(key_result.target_value),
                    "status": key_result.status,
                }
            )
            commit_started = True
            await db.commit()
            return _typed_success(
                f"Updated KR {kr_id}: {previous_value} -> {value}; status={key_result.status}.",
                result_ref=result_ref,
                metadata=metadata,
            )
    except Exception as exc:
        logger.exception("[OKR] typed KR progress update failed")
        if commit_started:
            return _typed_unknown(
                "KR progress commit acknowledgement was lost; reconcile before retrying.",
                "kr_progress_outcome_unknown",
                result_ref=result_ref,
                metadata=metadata,
            )
        return _typed_failure(
            f"KR progress update failed: {type(exc).__name__}.",
            "kr_progress_update_failed",
            result_ref=result_ref,
            metadata=metadata,
        )


async def _update_kr_content_outcome(
    agent_id: uuid.UUID | None,
    user_id: uuid.UUID | None,
    arguments: dict,
) -> ToolExecutionOutcome:
    if agent_id is None:
        return _typed_failure(
            "KR content updates require Agent context.",
            "invalid_tool_arguments",
        )
    kr_id, argument_error = _okr_uuid(arguments.get("kr_id"), "kr_id")
    if argument_error is not None:
        return argument_error
    supported_fields = ("title", "target_value", "unit", "focus_ref", "status")
    updates = {
        field: arguments[field]
        for field in supported_fields
        if field in arguments
    }
    if not updates:
        return _typed_failure(
            "At least one KR content field must be provided.",
            "invalid_tool_arguments",
        )
    for field in ("title", "unit", "focus_ref"):
        if field in updates and not isinstance(updates[field], str):
            return _typed_failure(
                f"{field} must be a string.",
                "invalid_tool_arguments",
            )
    if "title" in updates and not updates["title"].strip():
        return _typed_failure(
            "title must be non-empty.",
            "invalid_tool_arguments",
        )
    if "target_value" in updates:
        target_value, argument_error = _okr_finite_number(
            updates["target_value"],
            "target_value",
        )
        if argument_error is not None:
            return argument_error
        updates["target_value"] = target_value
    if "status" in updates and updates["status"] not in _OKR_KR_STATUSES:
        return _typed_failure(
            "status is not a supported KR status.",
            "invalid_tool_arguments",
        )

    assert kr_id is not None
    result_ref = str(kr_id)
    metadata: dict[str, object] = {
        "kr_id": result_ref,
        "changed_fields": sorted(updates),
    }
    commit_started = False
    try:
        from app.models.okr import OKRKeyResult, OKRObjective

        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if ctx.get("agent") is None:
                return _typed_failure(
                    "Agent not found.",
                    "source_agent_not_found",
                    result_ref=result_ref,
                )
            result = await db.execute(
                select(OKRKeyResult, OKRObjective)
                .join(
                    OKRObjective,
                    OKRKeyResult.objective_id == OKRObjective.id,
                )
                .where(
                    OKRKeyResult.id == kr_id,
                    OKRObjective.tenant_id == ctx["tenant_id"],
                )
            )
            row = result.first()
            if row is None:
                return _typed_failure(
                    f"Key Result {kr_id} was not found.",
                    "key_result_not_found",
                    result_ref=result_ref,
                )
            key_result, objective = row
            permission_error = _can_access_existing_okr_target(
                ctx,
                objective.owner_type,
                objective.owner_id,
            )
            if permission_error:
                return _typed_failure(
                    permission_error,
                    "okr_permission_denied",
                    result_ref=result_ref,
                )

            for field, value in updates.items():
                if field in {"title", "unit", "focus_ref"}:
                    value = value.strip()
                    if field != "title" and not value:
                        value = None
                setattr(key_result, field, value)
            metadata.update(
                {
                    "target_value": float(key_result.target_value),
                    "status": key_result.status,
                }
            )
            commit_started = True
            await db.commit()
            return _typed_success(
                f"Updated KR {kr_id}. Changed fields: {', '.join(sorted(updates))}.",
                result_ref=result_ref,
                metadata=metadata,
            )
    except Exception as exc:
        logger.exception("[OKR] typed KR content update failed")
        if commit_started:
            return _typed_unknown(
                "KR content commit acknowledgement was lost; reconcile before retrying.",
                "kr_content_outcome_unknown",
                result_ref=result_ref,
                metadata=metadata,
            )
        return _typed_failure(
            f"KR content update failed: {type(exc).__name__}.",
            "kr_content_update_failed",
            result_ref=result_ref,
            metadata=metadata,
        )


async def _update_kr_progress(agent_id: uuid.UUID | None, user_id: uuid.UUID | None, arguments: dict) -> str:
    """Update a KR's current_value. Only the owning agent may call this.

    Automatically writes an OKRProgressLog entry for history tracking.
    """
    if not agent_id:
        return "OKR tools require agent context."

    kr_id_str = arguments.get("kr_id", "").strip()
    value = arguments.get("value")
    note = arguments.get("note")
    status = arguments.get("status")

    if not kr_id_str:
        return "Missing required argument 'kr_id'. Call get_my_okr first to get your KR IDs."
    if value is None:
        return "Missing required argument 'value'."
    if status is not None and status not in {
        "on_track",
        "at_risk",
        "behind",
        "completed",
    }:
        return "Invalid status for update_kr_progress."

    try:
        kr_id = uuid.UUID(kr_id_str)
    except ValueError:
        return f"Invalid kr_id format: {kr_id_str}"

    try:
        from app.models.okr import OKRObjective, OKRKeyResult, OKRProgressLog
        from sqlalchemy import select as _select
        from datetime import datetime

        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if not ctx["agent"]:
                return "Agent not found."

            result = await db.execute(
                _select(OKRKeyResult, OKRObjective)
                .join(OKRObjective, OKRKeyResult.objective_id == OKRObjective.id)
                .where(
                    OKRKeyResult.id == kr_id,
                    OKRObjective.tenant_id == ctx["tenant_id"],
                )
            )
            row = result.first()
            if not row:
                return f"Key Result {kr_id_str} not found in your organization."

            kr, obj = row
            permission_error = _can_access_existing_okr_target(ctx, obj.owner_type, obj.owner_id)
            if permission_error:
                return permission_error

            prev_value = kr.current_value
            kr.current_value = float(value)
            kr.last_updated_at = datetime.utcnow()

            if status is not None:
                kr.status = status
            else:
                # Auto-determine status based on progress ratio.
                ratio = (
                    kr.current_value / kr.target_value
                    if kr.target_value
                    else 0
                )
                if ratio >= 1.0:
                    kr.status = "completed"
                elif ratio >= 0.7:
                    kr.status = "on_track"
                elif ratio >= 0.4:
                    kr.status = "at_risk"
                else:
                    kr.status = "behind"

            log = OKRProgressLog(
                kr_id=kr_id,
                previous_value=prev_value,
                new_value=float(value),
                source="self_report",
                note=note,
            )
            db.add(log)
            await db.commit()

        return (
            f"KR updated: {kr.title}\n"
            f"  {prev_value} → {value} {kr.unit or ''} (status: {kr.status})"
        )

    except Exception as e:
        logger.exception(f"[OKR] update_kr_progress failed for agent {agent_id}")
        return f"Failed to update KR progress: {str(e)[:200]}"


async def _update_kr_content(agent_id: uuid.UUID | None, user_id: uuid.UUID | None, arguments: dict) -> str:
    """Update metadata/content fields of one of the caller's own KRs."""
    if not agent_id:
        return "OKR tools require agent context."

    kr_id_str = arguments.get("kr_id", "").strip()
    if not kr_id_str:
        return "Missing required argument 'kr_id'. Call get_my_okr first to get your KR IDs."

    try:
        kr_id = uuid.UUID(kr_id_str)
    except ValueError:
        return f"Invalid kr_id format: {kr_id_str}"

    supported_fields = {
        "title": arguments.get("title"),
        "target_value": arguments.get("target_value"),
        "unit": arguments.get("unit"),
        "focus_ref": arguments.get("focus_ref"),
        "status": arguments.get("status"),
    }
    provided_updates = {key: value for key, value in supported_fields.items() if value is not None}
    if not provided_updates:
        return "No KR content fields provided. You can update: title, target_value, unit, focus_ref, status."

    try:
        from app.models.okr import OKRObjective, OKRKeyResult
        from sqlalchemy import select as _select

        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if not ctx["agent"]:
                return "Agent not found."

            result = await db.execute(
                _select(OKRKeyResult, OKRObjective)
                .join(OKRObjective, OKRKeyResult.objective_id == OKRObjective.id)
                .where(
                    OKRKeyResult.id == kr_id,
                    OKRObjective.tenant_id == ctx["tenant_id"],
                )
            )
            row = result.first()
            if not row:
                return f"Key Result {kr_id_str} not found in your organization."

            kr, obj = row
            permission_error = _can_access_existing_okr_target(ctx, obj.owner_type, obj.owner_id)
            if permission_error:
                return permission_error

            changed_fields: list[str] = []
            if "title" in provided_updates:
                kr.title = str(provided_updates["title"]).strip()
                changed_fields.append("title")
            if "target_value" in provided_updates:
                kr.target_value = float(provided_updates["target_value"])
                changed_fields.append("target_value")
            if "unit" in provided_updates:
                kr.unit = str(provided_updates["unit"]).strip() or None
                changed_fields.append("unit")
            if "focus_ref" in provided_updates:
                kr.focus_ref = str(provided_updates["focus_ref"]).strip() or None
                changed_fields.append("focus_ref")
            if "status" in provided_updates:
                kr.status = str(provided_updates["status"]).strip()
                changed_fields.append("status")

            await db.commit()

        return (
            f"KR content updated: {kr.title}\n"
            f"Changed fields: {', '.join(changed_fields)}"
        )

    except Exception as e:
        logger.exception(f"[OKR] update_kr_content failed for agent {agent_id}")
        return f"Failed to update KR content: {str(e)[:200]}"


async def _load_okr_job_agent(
    agent_id: uuid.UUID,
):
    from app.models.agent import Agent as AgentModel

    async with async_session() as db:
        result = await db.execute(
            select(AgentModel).where(AgentModel.id == agent_id)
        )
        return result.scalar_one_or_none()


def _okr_collection_outcome_from_receipt(
    receipt: Mapping,
) -> ToolExecutionOutcome:
    operation_id = receipt.get("operation_id")
    result_ref = (
        f"okr-collection://{operation_id}" if operation_id else None
    )
    metadata = {
        "operation_id": str(operation_id) if operation_id else None,
        "updated_count": int(receipt.get("updated_count", 0)),
        "skipped_count": int(receipt.get("skipped_count", 0)),
        "error_count": int(receipt.get("error_count", 0)),
        "updated_refs": list(receipt.get("updated_refs") or []),
    }
    status = receipt.get("status")
    summary = (
        "OKR focus collection settled: "
        f"updated={metadata['updated_count']}, "
        f"skipped={metadata['skipped_count']}, "
        f"errors={metadata['error_count']}."
    )
    if status == "succeeded":
        return _typed_success(
            summary,
            result_ref=result_ref,
            evidence_refs=tuple(metadata["updated_refs"]),
            metadata=metadata,
        )
    if status == "unknown":
        return _typed_unknown(
            "OKR focus collection commit outcome is unknown; reconcile before retrying.",
            str(
                receipt.get("error_code")
                or "okr_collection_commit_outcome_unknown"
            ),
            result_ref=result_ref,
            metadata=metadata,
        )
    return _typed_failure(
        summary,
        (
            "okr_collection_partial_failure"
            if status == "partial"
            else str(receipt.get("error_code") or "okr_collection_failed")
        ),
        result_ref=result_ref,
        metadata=metadata,
    )


def _okr_report_outcome_from_receipt(
    agent_id: uuid.UUID,
    receipt: Mapping,
) -> ToolExecutionOutcome:
    report_id = receipt.get("report_id")
    result_ref = f"okr-report://{report_id}" if report_id else None
    workspace_path = receipt.get("workspace_path")
    projection_status = str(
        receipt.get("projection_status") or "not_started"
    )
    db_status = str(receipt.get("db_status") or receipt.get("status") or "failed")
    metadata = {
        "operation_id": receipt.get("operation_id"),
        "report_id": str(report_id) if report_id else None,
        "report_type": receipt.get("report_type"),
        "period_start": receipt.get("period_start"),
        "period_end": receipt.get("period_end"),
        "workspace_path": workspace_path,
        "db_status": db_status,
        "projection_status": projection_status,
    }
    report_type = metadata["report_type"] or "OKR"
    summary = (
        f"{report_type} report database status={db_status}; "
        f"workspace projection status={projection_status}."
    )
    status = receipt.get("status")
    if status == "succeeded" and db_status == "succeeded" and projection_status == "succeeded":
        artifact_refs = (
            (f"workspace://{agent_id}/{workspace_path}",)
            if isinstance(workspace_path, str) and workspace_path
            else ()
        )
        return _typed_success(
            summary,
            result_ref=result_ref,
            artifact_refs=artifact_refs,
            metadata=metadata,
        )
    if status == "unknown" or db_status == "unknown":
        return _typed_unknown(
            "OKR report commit outcome is unknown; reconcile before retrying.",
            str(
                receipt.get("error_code")
                or "okr_report_commit_outcome_unknown"
            ),
            result_ref=result_ref,
            metadata=metadata,
        )
    if db_status == "succeeded" and projection_status == "failed":
        return _typed_failure(
            summary,
            "okr_report_projection_failed",
            result_ref=result_ref,
            metadata=metadata,
        )
    return _typed_failure(
        summary,
        str(receipt.get("error_code") or "okr_report_failed"),
        result_ref=result_ref,
        metadata=metadata,
    )


async def _okr_job_outcome(
    tool_name: str,
    agent_id: uuid.UUID | None,
    arguments: dict,
) -> ToolExecutionOutcome:
    if agent_id is None:
        return _typed_failure(
            "OKR jobs require Agent context.",
            "invalid_tool_arguments",
        )
    report_type: str | None = None
    if tool_name == "generate_okr_report":
        report_type = arguments.get("report_type")
        if report_type not in {"daily", "weekly"}:
            return _typed_failure(
                "report_type must be daily or weekly.",
                "invalid_tool_arguments",
            )
    if not await _agent_is_designated_okr_agent(agent_id):
        return _typed_failure(
            "Only the tenant's designated OKR Agent may run this job.",
            "okr_agent_required",
        )

    agent = await _load_okr_job_agent(agent_id)
    if agent is None:
        return _typed_failure(
            "Agent not found.",
            "source_agent_not_found",
        )

    from app.services import okr_scheduler

    try:
        if tool_name == "collect_okr_progress":
            receipt = await okr_scheduler.collect_all_focus_updates(
                tenant_id=agent.tenant_id,
                okr_agent_id=agent_id,
            )
            if not isinstance(receipt, Mapping):
                return _typed_failure(
                    "OKR collection did not return a structured receipt.",
                    "okr_collection_invalid_receipt",
                )
            return _okr_collection_outcome_from_receipt(receipt)

        if tool_name == "generate_monthly_okr_report":
            receipt = await okr_scheduler.generate_monthly_report(
                agent.tenant_id,
                agent_id,
            )
        elif report_type == "daily":
            receipt = await okr_scheduler.generate_daily_report(
                agent.tenant_id,
                agent_id,
            )
        else:
            receipt = await okr_scheduler.generate_weekly_report(
                agent.tenant_id,
                agent_id,
            )
        if not isinstance(receipt, Mapping):
            return _typed_failure(
                "OKR report job did not return a structured receipt.",
                "okr_report_invalid_receipt",
            )
        return _okr_report_outcome_from_receipt(agent_id, receipt)
    except Exception as exc:
        commit_started = bool(getattr(exc, "commit_started", False))
        if tool_name == "collect_okr_progress":
            operation_id = getattr(exc, "operation_id", None)
            result_ref = (
                f"okr-collection://{operation_id}"
                if operation_id
                else None
            )
            if commit_started:
                return _typed_unknown(
                    "OKR focus collection commit outcome is unknown; reconcile before retrying.",
                    "okr_collection_commit_outcome_unknown",
                    result_ref=result_ref,
                    metadata={"operation_id": operation_id},
                )
            return _typed_failure(
                f"OKR focus collection failed: {type(exc).__name__}.",
                "okr_collection_failed",
                result_ref=result_ref,
            )

        report_id = getattr(exc, "report_id", None)
        workspace_path = getattr(exc, "workspace_path", None)
        result_ref = f"okr-report://{report_id}" if report_id else None
        metadata = {
            "operation_id": getattr(exc, "operation_id", None),
            "report_id": report_id,
            "report_type": getattr(exc, "report_type", report_type),
            "workspace_path": workspace_path,
            "db_status": "unknown" if commit_started else "failed",
            "projection_status": "not_started",
        }
        if commit_started:
            return _typed_unknown(
                "OKR report commit outcome is unknown; reconcile before retrying.",
                "okr_report_commit_outcome_unknown",
                result_ref=result_ref,
                metadata=metadata,
            )
        return _typed_failure(
            f"OKR report generation failed: {type(exc).__name__}.",
            "okr_report_failed",
            result_ref=result_ref,
            metadata=metadata,
        )


async def _collect_okr_progress(agent_id: uuid.UUID | None) -> str:
    """Batch-collect KR progress from legacy team member focus files.

    Delegates to okr_scheduler.collect_all_focus_updates(). The calling agent
    must be the OKR Agent — we look up its tenant from the DB.
    """
    outcome = await _okr_job_outcome(
        "collect_okr_progress",
        agent_id,
        {},
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="OKR collection returned no summary.",
    )


async def _generate_okr_report(agent_id: uuid.UUID | None, arguments: dict) -> str:
    """Generate a daily or weekly OKR report.

    Writes to WorkReport table and returns the markdown content for posting.
    """
    outcome = await _okr_job_outcome(
        "generate_okr_report",
        agent_id,
        arguments,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="OKR report returned no summary.",
    )


async def _generate_monthly_okr_report(agent_id: uuid.UUID | None) -> str:
    """Generate the monthly OKR summary report for the agent's tenant.

    Writes a WorkReport (report_type='monthly') and returns the Markdown
    content. The OKR Agent should forward this to admins via send_platform_message.
    Also triggered automatically by the monthly_okr_report system cron trigger.
    """
    outcome = await _okr_job_outcome(
        "generate_monthly_okr_report",
        agent_id,
        {},
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Monthly OKR report returned no summary.",
    )


async def _get_okr_settings_tool(agent_id: uuid.UUID | None) -> str:
    """Return OKR settings for the agent's tenant as a formatted string.

    The OKR Agent uses this to determine report schedule and period config
    without needing to make HTTP calls to its own API.
    """
    if not agent_id:
        return "OKR tools require agent context."

    try:
        from app.models.agent import Agent as AgentModel
        from app.services.okr_scheduler import get_okr_settings_for_agent
        import json as _json

        async with async_session() as db:
            agent_result = await db.execute(
                select(AgentModel).where(AgentModel.id == agent_id)
            )
            agent = agent_result.scalar_one_or_none()
            if not agent:
                return "Agent not found."

        settings = await get_okr_settings_for_agent(agent.tenant_id)
        return _json.dumps(settings, indent=2, ensure_ascii=False)

    except Exception as e:
        logger.exception(f"[OKR] get_okr_settings failed for agent {agent_id}")
        return f"Failed to get OKR settings: {str(e)[:200]}"


async def _create_objective_outcome(
    agent_id: uuid.UUID | None,
    user_id: uuid.UUID | None,
    arguments: dict,
) -> ToolExecutionOutcome:
    if agent_id is None:
        return _typed_failure(
            "create_objective requires Agent context.",
            "invalid_tool_arguments",
        )
    title = arguments.get("title")
    owner_type = arguments.get("owner_type")
    if not isinstance(title, str) or not title.strip():
        return _typed_failure(
            "title must be a non-empty string.",
            "invalid_tool_arguments",
        )
    if owner_type not in {"company", "user", "agent"}:
        return _typed_failure(
            "owner_type must be company, user, or agent.",
            "invalid_tool_arguments",
        )
    description = arguments.get("description")
    if description is not None and not isinstance(description, str):
        return _typed_failure(
            "description must be a string.",
            "invalid_tool_arguments",
        )
    period_start_raw = arguments.get("period_start")
    period_end_raw = arguments.get("period_end")
    if not isinstance(period_start_raw, str) or not isinstance(period_end_raw, str):
        return _typed_failure(
            "period_start and period_end must use YYYY-MM-DD.",
            "invalid_tool_arguments",
        )
    try:
        period_start = date.fromisoformat(period_start_raw.strip())
        period_end = date.fromisoformat(period_end_raw.strip())
    except ValueError:
        return _typed_failure(
            "period_start and period_end must use YYYY-MM-DD.",
            "invalid_tool_arguments",
        )
    if period_start > period_end:
        return _typed_failure(
            "period_start must be on or before period_end.",
            "invalid_tool_arguments",
        )

    owner_id_raw = arguments.get("owner_id")
    owner_name = arguments.get("owner_name")
    if owner_name is not None and not isinstance(owner_name, str):
        return _typed_failure(
            "owner_name must be a string.",
            "invalid_tool_arguments",
        )
    owner_name = owner_name.strip() if isinstance(owner_name, str) else ""
    owner_id: uuid.UUID | None = None
    if owner_id_raw is not None:
        owner_id, argument_error = _okr_uuid(owner_id_raw, "owner_id")
        if argument_error is not None:
            return argument_error
    if owner_type != "company" and owner_id is None and not owner_name:
        return _typed_failure(
            f"owner_id or owner_name is required for {owner_type} objectives.",
            "invalid_tool_arguments",
        )
    if owner_type == "company":
        owner_id = None

    designated_error = await _require_designated_okr_agent(agent_id)
    if designated_error is not None:
        return designated_error

    commit_started = False
    result_ref: str | None = None
    metadata: dict[str, object] = {
        "owner_type": owner_type,
        "period_start": period_start.isoformat(),
        "period_end": period_end.isoformat(),
    }
    try:
        from app.models.agent import Agent as AgentModel
        from app.models.okr import OKRObjective
        from app.models.org import OrgMember
        from app.models.user import User as UserModel

        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if ctx.get("agent") is None:
                return _typed_failure(
                    "Agent not found.",
                    "source_agent_not_found",
                )
            ctx = dict(ctx)
            ctx["agent_is_designated_okr_agent"] = True
            tenant_id = ctx["tenant_id"]

            resolved_owner_id = owner_id
            if owner_type == "agent":
                if resolved_owner_id is not None:
                    owner_result = await db.execute(
                        select(AgentModel.id).where(
                            AgentModel.id == resolved_owner_id,
                            AgentModel.tenant_id == tenant_id,
                        )
                    )
                    resolved_owner_id = owner_result.scalar_one_or_none()
                else:
                    owner_result = await db.execute(
                        select(AgentModel.id).where(
                            AgentModel.name == owner_name,
                            AgentModel.tenant_id == tenant_id,
                        )
                    )
                    resolved_owner_id = owner_result.scalar_one_or_none()
            elif owner_type == "user":
                if resolved_owner_id is not None:
                    owner_result = await db.execute(
                        select(UserModel.id).where(
                            UserModel.id == resolved_owner_id,
                            UserModel.tenant_id == tenant_id,
                        )
                    )
                    resolved_owner_id = owner_result.scalar_one_or_none()
                    if resolved_owner_id is None:
                        member_result = await db.execute(
                            select(OrgMember.id).where(
                                OrgMember.id == owner_id,
                                OrgMember.tenant_id == tenant_id,
                            )
                        )
                        resolved_owner_id = member_result.scalar_one_or_none()
                else:
                    owner_result = await db.execute(
                        select(UserModel.id).where(
                            UserModel.display_name == owner_name,
                            UserModel.tenant_id == tenant_id,
                        )
                    )
                    resolved_owner_id = owner_result.scalar_one_or_none()
                    if resolved_owner_id is None:
                        member_result = await db.execute(
                            select(OrgMember.id).where(
                                OrgMember.name == owner_name,
                                OrgMember.tenant_id == tenant_id,
                            )
                        )
                        resolved_owner_id = member_result.scalar_one_or_none()

            if owner_type != "company" and resolved_owner_id is None:
                return _typed_failure(
                    f"The requested {owner_type} owner was not found in this tenant.",
                    "okr_owner_not_found",
                )
            permission_error = _can_create_okr_target(
                ctx,
                owner_type,
                resolved_owner_id,
            )
            if permission_error:
                return _typed_failure(
                    permission_error,
                    "okr_permission_denied",
                )

            objective = OKRObjective(
                tenant_id=tenant_id,
                title=title.strip(),
                description=description,
                owner_type=owner_type,
                owner_id=resolved_owner_id,
                period_start=period_start,
                period_end=period_end,
                status="active",
            )
            db.add(objective)
            await db.flush()
            result_ref = str(objective.id)
            metadata.update(
                {
                    "objective_id": result_ref,
                    "owner_id": (
                        str(resolved_owner_id)
                        if resolved_owner_id is not None
                        else None
                    ),
                    "status": objective.status,
                }
            )
            commit_started = True
            await db.commit()
            return _typed_success(
                f"Created Objective {objective.id}: {objective.title}.",
                result_ref=result_ref,
                metadata=metadata,
            )
    except Exception as exc:
        logger.exception("[OKR] typed Objective creation failed")
        if commit_started:
            return _typed_unknown(
                "Objective creation commit acknowledgement was lost; reconcile before retrying.",
                "objective_create_outcome_unknown",
                result_ref=result_ref,
                metadata=metadata,
            )
        return _typed_failure(
            f"Objective creation failed: {type(exc).__name__}.",
            "objective_create_failed",
            result_ref=result_ref,
            metadata=metadata,
        )


async def _create_key_result_outcome(
    agent_id: uuid.UUID | None,
    user_id: uuid.UUID | None,
    arguments: dict,
) -> ToolExecutionOutcome:
    if agent_id is None:
        return _typed_failure(
            "create_key_result requires Agent context.",
            "invalid_tool_arguments",
        )
    objective_id, argument_error = _okr_uuid(
        arguments.get("objective_id"),
        "objective_id",
    )
    if argument_error is not None:
        return argument_error
    title = arguments.get("title")
    if not isinstance(title, str) or not title.strip():
        return _typed_failure(
            "title must be a non-empty string.",
            "invalid_tool_arguments",
        )
    target_value, argument_error = _okr_finite_number(
        arguments.get("target_value"),
        "target_value",
    )
    if argument_error is not None:
        return argument_error
    for field in ("unit", "focus_ref"):
        if field in arguments and not isinstance(arguments[field], str):
            return _typed_failure(
                f"{field} must be a string.",
                "invalid_tool_arguments",
            )
    designated_error = await _require_designated_okr_agent(agent_id)
    if designated_error is not None:
        return designated_error

    assert objective_id is not None and target_value is not None
    commit_started = False
    result_ref: str | None = None
    metadata: dict[str, object] = {
        "objective_id": str(objective_id),
        "target_value": target_value,
    }
    try:
        from app.models.okr import OKRKeyResult, OKRObjective

        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if ctx.get("agent") is None:
                return _typed_failure(
                    "Agent not found.",
                    "source_agent_not_found",
                )
            ctx = dict(ctx)
            ctx["agent_is_designated_okr_agent"] = True
            objective_result = await db.execute(
                select(OKRObjective).where(
                    OKRObjective.id == objective_id,
                    OKRObjective.tenant_id == ctx["tenant_id"],
                )
            )
            objective = objective_result.scalar_one_or_none()
            if objective is None:
                return _typed_failure(
                    f"Objective {objective_id} was not found.",
                    "objective_not_found",
                    result_ref=str(objective_id),
                )
            permission_error = _can_access_existing_okr_target(
                ctx,
                objective.owner_type,
                objective.owner_id,
            )
            if permission_error:
                return _typed_failure(
                    permission_error,
                    "okr_permission_denied",
                    result_ref=str(objective_id),
                )

            key_result = OKRKeyResult(
                objective_id=objective_id,
                title=title.strip(),
                target_value=target_value,
                current_value=0.0,
                unit=(arguments.get("unit") or None),
                focus_ref=(arguments.get("focus_ref") or None),
                status=_okr_progress_status(0.0, target_value),
            )
            db.add(key_result)
            await db.flush()
            result_ref = str(key_result.id)
            metadata.update(
                {
                    "kr_id": result_ref,
                    "current_value": 0.0,
                    "status": key_result.status,
                }
            )
            commit_started = True
            await db.commit()
            return _typed_success(
                f"Created Key Result {key_result.id}: {key_result.title}.",
                result_ref=result_ref,
                metadata=metadata,
            )
    except Exception as exc:
        logger.exception("[OKR] typed Key Result creation failed")
        if commit_started:
            return _typed_unknown(
                "Key Result creation commit acknowledgement was lost; reconcile before retrying.",
                "key_result_create_outcome_unknown",
                result_ref=result_ref,
                metadata=metadata,
            )
        return _typed_failure(
            f"Key Result creation failed: {type(exc).__name__}.",
            "key_result_create_failed",
            result_ref=result_ref,
            metadata=metadata,
        )


async def _create_objective(agent_id: uuid.UUID | None, user_id: uuid.UUID | None, arguments: dict) -> str:
    if not agent_id:
        return "OKR tools require agent context."
    try:
        from app.models.agent import Agent as AgentModel
        from app.models.okr import OKRObjective
        from app.models.user import User as UserModel
        from app.models.org import OrgMember
        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            ag = ctx["agent"]
            if not ag:
                return "Agent not found."

            title = arguments.get("title")
            owner_type = arguments.get("owner_type")
            period_start = arguments.get("period_start")
            period_end = arguments.get("period_end")
            if not all([title, owner_type, period_start, period_end]):
                return "Missing required fields: title, owner_type, period_start, period_end"

            from datetime import date
            p_start = date.fromisoformat(period_start)
            p_end = date.fromisoformat(period_end)

            owner_id_str = arguments.get("owner_id")
            owner_name_hint = arguments.get("owner_name")  # optional name-based fallback
            owner_id: uuid.UUID | None = None

            if owner_id_str:
                try:
                    owner_id = uuid.UUID(owner_id_str)
                except ValueError:
                    owner_id = None

                if owner_id:
                    owner_exists = False
                    if owner_type == "agent":
                        res = await db.execute(select(AgentModel.id).where(AgentModel.id == owner_id))
                        owner_exists = res.scalar_one_or_none() is not None
                    elif owner_type == "user":
                        from app.models.user import User as UserModel
                        from app.models.org import OrgMember
                        res = await db.execute(select(UserModel.id).where(UserModel.id == owner_id))
                        owner_exists = res.scalar_one_or_none() is not None
                        if not owner_exists:
                            # Maybe agent passed OrgMember.id — resolve to linked User.id when available
                            res = await db.execute(
                                select(OrgMember.id, OrgMember.user_id).where(OrgMember.id == owner_id)
                            )
                            member_row = res.first()
                            if member_row:
                                owner_exists = True
                                if member_row.user_id:
                                    # Resolve OrgMember.id → User.id so name lookup in list_objectives works
                                    owner_id = member_row.user_id
                                    logger.info(
                                        f"[OKR] _create_objective: resolved OrgMember.id {owner_id_str} "
                                        f"→ user_id {owner_id}"
                                    )
                                # else: channel-only member, keep OrgMember.id as owner_id

                    if not owner_exists:
                        owner_id = None
                        if not owner_name_hint:
                            return f"owner_id '{owner_id_str}' was not found. Provide a valid UUID, or pass owner_name instead."

            if owner_type != "company" and not owner_id and owner_name_hint:
                # If we don't have a valid UUID but we have a name, look it up
                if owner_type == "agent":
                    res = await db.execute(select(AgentModel.id).where(AgentModel.tenant_id == ag.tenant_id, AgentModel.name == owner_name_hint))
                    owner_id = res.scalar_one_or_none()
                elif owner_type == "user":
                    from app.models.org import OrgMember
                    from app.models.user import User as UserModel
                    # Try platform User.display_name first
                    res = await db.execute(select(UserModel.id).where(UserModel.display_name == owner_name_hint, UserModel.tenant_id == ag.tenant_id))
                    owner_id = res.scalar_one_or_none()
                    if not owner_id:
                        # Fall back to OrgMember.name (Feishu/channel-only users)
                        res = await db.execute(select(OrgMember.id).where(OrgMember.name == owner_name_hint, OrgMember.tenant_id == ag.tenant_id))
                        owner_id = res.scalar_one_or_none()

                if not owner_id:
                    return f"Failed: Could not resolve a valid system UUID for the {owner_type} named '{owner_name_hint}'."

            if owner_type != "company" and not owner_id:
               return f"Failed: owner_id or owner_name is required for {owner_type} OKRs."

            if not ctx["agent_is_system"] and owner_type == "agent" and owner_id is None:
                owner_id = agent_id

            permission_error = _can_create_okr_target(ctx, owner_type, owner_id)
            if permission_error:
                return permission_error

            obj = OKRObjective(
                tenant_id=ag.tenant_id,
                title=title,
                description=arguments.get("description"),
                owner_type=owner_type,
                owner_id=owner_id,
                period_start=p_start,
                period_end=p_end,
                status="active"
            )
            db.add(obj)
            await db.commit()
            owner_info = f"owner={owner_name_hint or owner_id_str or 'unattributed'}"
            return f"Successfully created Objective '{obj.title}' (ID: {obj.id}, {owner_info})"
    except Exception as e:
        logger.exception("[OKR] create_objective failed")
        return f"Failed to create objective: {str(e)[:200]}"


async def _create_key_result(agent_id: uuid.UUID | None, user_id: uuid.UUID | None, arguments: dict) -> str:
    if not agent_id:
        return "OKR tools require agent context."
    import math

    obj_id_str = arguments.get("objective_id")
    title = arguments.get("title")
    raw_target_value = arguments.get("target_value")
    if not isinstance(obj_id_str, str) or not obj_id_str.strip():
        return "Missing objective_id"
    if not isinstance(title, str) or not title.strip():
        return "Missing title"
    if isinstance(raw_target_value, bool):
        return "Invalid target_value: a finite number is required."
    try:
        target_value = float(raw_target_value)
    except (TypeError, ValueError):
        return "Invalid target_value: a finite number is required."
    if not math.isfinite(target_value):
        return "Invalid target_value: a finite number is required."
    try:
        obj_id = uuid.UUID(obj_id_str.strip())
    except ValueError:
        return "Invalid formatted objective_id (must be UUID)"

    try:
        from app.models.okr import OKRObjective, OKRKeyResult
        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if not ctx["agent"]:
                return "Agent not found."

            # Verify objective exists
            obj_res = await db.execute(
                select(OKRObjective).where(
                    OKRObjective.id == obj_id,
                    OKRObjective.tenant_id == ctx["tenant_id"],
                )
            )
            obj = obj_res.scalar_one_or_none()
            if not obj:
                return f"Objective {obj_id} not found."

            permission_error = _can_access_existing_okr_target(ctx, obj.owner_type, obj.owner_id)
            if permission_error:
                return permission_error

            kr = OKRKeyResult(
                objective_id=obj_id,
                title=title.strip(),
                target_value=target_value,
                current_value=0.0,
                unit=arguments.get("unit"),
                focus_ref=arguments.get("focus_ref")
            )
            db.add(kr)
            await db.commit()
            return f"Successfully created Key Result '{kr.title}' (ID: {kr.id})"
    except Exception as e:
        logger.exception("[OKR] create_key_result failed")
        return f"Failed to create key result: {str(e)[:200]}"


async def _update_objective_outcome(
    agent_id: uuid.UUID | None,
    user_id: uuid.UUID | None,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Update Objective metadata.

    Permission rules:
    - Regular agents: can only modify Objectives they own (owner_type='agent', owner_id=agent_id).
    - System agents are constrained by the requesting user's role: admins can modify any OKR,
      non-admins may only modify their own personal OKRs.
    """
    if not agent_id:
        return _typed_failure(
            "update_objective requires Agent context.",
            "invalid_tool_arguments",
        )
    obj_id_str = arguments.get("objective_id")
    if not isinstance(obj_id_str, str) or not obj_id_str.strip():
        return _typed_failure(
            "update_objective requires objective_id.",
            "invalid_tool_arguments",
        )
    try:
        obj_id = uuid.UUID(obj_id_str.strip())
    except ValueError:
        return _typed_failure(
            "update_objective objective_id must be a UUID.",
            "invalid_tool_arguments",
        )

    supported_fields = {
        "title",
        "description",
        "status",
        "period_start",
        "period_end",
    }
    update_fields = [field for field in supported_fields if field in arguments]
    if not update_fields:
        return _typed_failure(
            "update_objective requires at least one supported field to update.",
            "invalid_tool_arguments",
        )
    for field in ("title", "description"):
        if field in arguments and not isinstance(arguments[field], str):
            return _typed_failure(
                f"update_objective {field} must be a string.",
                "invalid_tool_arguments",
            )
    if "title" in arguments and not arguments["title"].strip():
        return _typed_failure(
            "update_objective title must be non-empty.",
            "invalid_tool_arguments",
        )
    if "status" in arguments and arguments["status"] not in {
        "draft",
        "active",
        "completed",
        "archived",
    }:
        return _typed_failure(
            "update_objective status is invalid.",
            "invalid_tool_arguments",
        )
    parsed_dates: dict[str, object] = {}
    try:
        from datetime import date

        for field in ("period_start", "period_end"):
            if field in arguments:
                if not isinstance(arguments[field], str):
                    raise ValueError(field)
                parsed_dates[field] = date.fromisoformat(arguments[field])
    except ValueError:
        return _typed_failure(
            "update_objective period dates must use YYYY-MM-DD.",
            "invalid_tool_arguments",
        )
    if (
        "period_start" in parsed_dates
        and "period_end" in parsed_dates
        and parsed_dates["period_start"] > parsed_dates["period_end"]
    ):
        return _typed_failure(
            "period_start must be on or before period_end.",
            "invalid_tool_arguments",
        )

    commit_started = False
    result_ref = str(obj_id)
    metadata: dict[str, object] = {"objective_id": result_ref}
    try:
        from app.models.okr import OKRObjective
        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if not ctx["agent"]:
                return _typed_failure(
                    "Agent not found.",
                    "source_agent_not_found",
                )

            obj_res = await db.execute(
                select(OKRObjective).where(
                    OKRObjective.id == obj_id,
                    OKRObjective.tenant_id == ctx["tenant_id"],
                )
            )
            obj = obj_res.scalar_one_or_none()
            if not obj:
                return _typed_failure(
                    f"Objective {obj_id} not found.",
                    "objective_not_found",
                    result_ref=result_ref,
                )

            permission_error = _can_access_existing_okr_target(ctx, obj.owner_type, obj.owner_id)
            if permission_error:
                return _typed_failure(
                    permission_error,
                    "objective_permission_denied",
                    result_ref=result_ref,
                )

            next_period_start = parsed_dates.get(
                "period_start", obj.period_start
            )
            next_period_end = parsed_dates.get("period_end", obj.period_end)
            if next_period_start > next_period_end:
                return _typed_failure(
                    "period_start must be on or before period_end.",
                    "invalid_tool_arguments",
                    result_ref=result_ref,
                )

            updates = []
            if "title" in arguments:
                obj.title = arguments["title"]
                updates.append("title")
            if "description" in arguments:
                obj.description = arguments["description"]
                updates.append("description")
            if "status" in arguments:
                obj.status = arguments["status"]
                updates.append("status")
            if "period_start" in arguments:
                obj.period_start = parsed_dates["period_start"]
                updates.append("period_start")
            if "period_end" in arguments:
                obj.period_end = parsed_dates["period_end"]
                updates.append("period_end")

            commit_started = True
            await db.commit()
            metadata.update(
                {
                    "changed_fields": sorted(updates),
                    "status": obj.status,
                    "period_start": obj.period_start.isoformat(),
                    "period_end": obj.period_end.isoformat(),
                }
            )
            return _typed_success(
                f"Updated Objective {obj.id}. Changed fields: {', '.join(updates)}.",
                result_ref=result_ref,
                metadata=metadata,
            )
    except Exception as exc:
        logger.exception("[OKR] update_objective failed")
        if commit_started:
            return _typed_unknown(
                "Objective update outcome is unknown; reconcile before retrying.",
                "objective_update_outcome_unknown",
                result_ref=result_ref,
                metadata=metadata,
            )
        return _typed_failure(
            f"Objective update failed: {type(exc).__name__}.",
            "objective_update_failed",
            result_ref=result_ref,
            metadata=metadata,
        )


async def _update_objective(
    agent_id: uuid.UUID | None,
    user_id: uuid.UUID | None,
    arguments: dict,
) -> str:
    """Legacy display adapter for the typed Objective update."""
    outcome = await _update_objective_outcome(agent_id, user_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Objective update returned no summary.",
    )


async def _update_any_kr_progress(agent_id: uuid.UUID | None, user_id: uuid.UUID | None, arguments: dict) -> str:
    """OKR Agent exclusive version of update_kr_progress."""
    if not agent_id:
        return "OKR tools require agent context."
    try:
        from app.models.okr import OKRKeyResult, OKRObjective, OKRProgressLog
        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if not ctx["agent"]:
                return "Agent not found."

            kr_id_str = arguments.get("kr_id")
            val = arguments.get("value")
            if not kr_id_str or val is None:
                return "Missing kr_id or value"
            try:
                kr_id = uuid.UUID(kr_id_str)
            except ValueError:
                return "Invalid formatted kr_id (must be UUID)"

            kr_res = await db.execute(
                select(OKRKeyResult, OKRObjective)
                .join(OKRObjective, OKRKeyResult.objective_id == OKRObjective.id)
                .where(
                    OKRKeyResult.id == kr_id,
                    OKRObjective.tenant_id == ctx["tenant_id"],
                )
            )
            row = kr_res.first()
            if not row:
                return f"Key Result {kr_id} not found in your organization."

            kr, obj = row
            permission_error = _can_access_existing_okr_target(ctx, obj.owner_type, obj.owner_id)
            if permission_error:
                return permission_error

            old_val = kr.current_value
            kr.current_value = float(val)

            # Auto-compute status if not explicitly given
            explicit_status = arguments.get("status")
            if explicit_status:
                kr.status = explicit_status
            else:
                progress = kr.current_value / kr.target_value if kr.target_value != 0 else 0
                if progress >= 1.0:
                    kr.status = "completed"
                elif progress >= 0.7:
                    kr.status = "on_track"
                elif progress >= 0.4:
                    kr.status = "at_risk"
                else:
                    kr.status = "behind"

            from datetime import datetime
            kr.last_updated_at = datetime.utcnow()

            note = arguments.get("note", "Updated by OKR Agent after check-in")
            log_entry = OKRProgressLog(
                kr_id=kr.id,
                previous_value=old_val,
                new_value=kr.current_value,
                source="okr_agent" if ctx["agent_is_system"] else "agent",
                note=note
            )
            db.add(log_entry)
            await db.commit()

            return f"Successfully updated KR '{kr.title}'. Progress: {old_val} -> {kr.current_value} {kr.unit or ''}. Status: {kr.status}"
    except Exception as e:
        logger.exception(f"[OKR] update_any_kr_progress failed")
        return f"Failed to update kr progress: {str(e)[:200]}"


async def _upsert_member_daily_report_outcome(
    agent_id: uuid.UUID | None,
    user_id: uuid.UUID | None,
    arguments: dict,
) -> ToolExecutionOutcome:
    if agent_id is None:
        return _typed_failure(
            "upsert_member_daily_report requires Agent context.",
            "invalid_tool_arguments",
        )
    report_date_raw = arguments.get("report_date")
    if not isinstance(report_date_raw, str) or not report_date_raw.strip():
        return _typed_failure(
            "report_date must use YYYY-MM-DD.",
            "invalid_tool_arguments",
        )
    try:
        report_date = date.fromisoformat(report_date_raw.strip())
    except ValueError:
        return _typed_failure(
            "report_date must use YYYY-MM-DD.",
            "invalid_tool_arguments",
        )
    content_raw = arguments.get("content")
    if not isinstance(content_raw, str) or not content_raw.strip():
        return _typed_failure(
            "content must be a non-empty string.",
            "invalid_tool_arguments",
        )
    member_type = arguments.get("member_type", "user")
    if member_type not in {"user", "agent"}:
        return _typed_failure(
            "member_type must be user or agent.",
            "invalid_tool_arguments",
        )
    source = arguments.get("source", "okr_agent_assisted")
    if not isinstance(source, str) or not source.strip():
        return _typed_failure(
            "source must be a non-empty string.",
            "invalid_tool_arguments",
        )
    source = source.strip()
    if len(source) > 30:
        return _typed_failure(
            "source must not exceed 30 characters.",
            "invalid_tool_arguments",
        )
    member_id_raw = arguments.get("member_id")
    member_name = arguments.get("member_name")
    if member_name is not None and not isinstance(member_name, str):
        return _typed_failure(
            "member_name must be a string.",
            "invalid_tool_arguments",
        )
    member_name = member_name.strip() if isinstance(member_name, str) else ""
    target_member_id: uuid.UUID | None = None
    if member_id_raw is not None:
        target_member_id, argument_error = _okr_uuid(
            member_id_raw,
            "member_id",
        )
        if argument_error is not None:
            return argument_error
    if target_member_id is None and not member_name:
        return _typed_failure(
            "member_id or member_name is required.",
            "invalid_tool_arguments",
        )

    designated_error = await _require_designated_okr_agent(agent_id)
    if designated_error is not None:
        return designated_error

    stored_content = content_raw[:2000]
    content_truncated = len(content_raw) > len(stored_content)
    content_hash = hashlib.sha256(
        stored_content.encode("utf-8")
    ).hexdigest()
    commit_started = False
    result_ref: str | None = None
    metadata: dict[str, object] = {
        "member_type": member_type,
        "report_date": report_date.isoformat(),
        "content_truncated": content_truncated,
        "okr_content_hash": content_hash,
        "stored_character_count": len(stored_content),
    }
    try:
        from app.models.agent import Agent as AgentModel
        from app.models.okr import MemberDailyReport
        from app.models.user import User as UserModel

        async with async_session() as db:
            ctx = await _load_okr_request_context(db, agent_id, user_id)
            if ctx.get("agent") is None:
                return _typed_failure(
                    "Agent not found.",
                    "source_agent_not_found",
                )
            tenant_id = ctx["tenant_id"]

            member_model = UserModel if member_type == "user" else AgentModel
            if target_member_id is not None:
                member_result = await db.execute(
                    select(member_model).where(
                        member_model.id == target_member_id,
                        member_model.tenant_id == tenant_id,
                    )
                )
            else:
                name_column = (
                    UserModel.display_name
                    if member_type == "user"
                    else AgentModel.name
                )
                member_result = await db.execute(
                    select(member_model).where(
                        name_column == member_name,
                        member_model.tenant_id == tenant_id,
                    )
                )
            member = member_result.scalar_one_or_none()
            if member is None:
                return _typed_failure(
                    f"The requested {member_type} member was not found in this tenant.",
                    "okr_member_not_found",
                )
            target_member_id = member.id
            metadata["member_id"] = str(target_member_id)

            existing_result = await db.execute(
                select(MemberDailyReport).where(
                    MemberDailyReport.tenant_id == tenant_id,
                    MemberDailyReport.member_type == member_type,
                    MemberDailyReport.member_id == target_member_id,
                    MemberDailyReport.report_date == report_date,
                )
            )
            report = existing_result.scalar_one_or_none()
            action = "Updated"
            if report is None:
                action = "Created"
                report = MemberDailyReport(
                    tenant_id=tenant_id,
                    member_type=member_type,
                    member_id=target_member_id,
                    report_date=report_date,
                    content=stored_content,
                    status="submitted",
                    source=source,
                )
                db.add(report)
                await db.flush()
            else:
                report.content = stored_content
                report.source = source
                report.status = "revised"
                report.updated_at = datetime.now(timezone.utc)

            result_ref = str(report.id)
            metadata.update(
                {
                    "report_id": result_ref,
                    "status": report.status,
                    "source": report.source,
                }
            )
            commit_started = True
            await db.commit()
            return _typed_success(
                f"{action} daily report {report.id} for {member_type} {target_member_id} on {report_date.isoformat()}; status={report.status}; stored={len(stored_content)} characters; truncated={str(content_truncated).lower()}.",
                result_ref=result_ref,
                metadata=metadata,
            )
    except Exception as exc:
        logger.exception("[OKR] typed daily report upsert failed")
        if commit_started:
            return _typed_unknown(
                "Daily report commit acknowledgement was lost; reconcile before retrying.",
                "daily_report_outcome_unknown",
                result_ref=result_ref,
                metadata=metadata,
            )
        return _typed_failure(
            f"Daily report upsert failed: {type(exc).__name__}.",
            "daily_report_upsert_failed",
            result_ref=result_ref,
            metadata=metadata,
        )


async def _upsert_member_daily_report(agent_id: uuid.UUID | None, arguments: dict) -> str:
    """OKR Agent exclusive tool for creating or revising a member daily report."""
    if not agent_id:
        return "OKR tools require agent context."

    try:
        from datetime import date as date_cls
        from app.models.agent import Agent as AgentModel
        from app.models.okr import MemberDailyReport
        from app.services.okr_reporting import (
            list_tracked_okr_members,
            upsert_member_daily_report as _upsert,
        )

        report_date_raw = arguments.get("report_date")
        content = (arguments.get("content") or "").strip()
        member_type = arguments.get("member_type") or "user"
        member_id_raw = arguments.get("member_id")
        member_name = (arguments.get("member_name") or "").strip()
        source = (arguments.get("source") or "okr_agent_assisted").strip() or "okr_agent_assisted"

        if not report_date_raw or not content:
            return "Missing report_date or content"

        try:
            report_date = date_cls.fromisoformat(report_date_raw)
        except ValueError:
            return "Invalid report_date format. Use YYYY-MM-DD."

        async with async_session() as db:
            ag_res = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            ag = ag_res.scalar_one_or_none()
            if not ag:
                return "Agent not found."
            if not ag.is_system:
                return "Permission denied: only the OKR Agent can upsert member daily reports."

            target_member_id: uuid.UUID | None = None
            if member_id_raw:
                try:
                    target_member_id = uuid.UUID(member_id_raw)
                except ValueError:
                    return "Invalid member_id format. Use a UUID."

            if not target_member_id:
                if not member_name:
                    return "Provide either member_id or member_name."
                members = await list_tracked_okr_members(ag.tenant_id)
                lowered = member_name.casefold()
                exact_matches = [
                    member for member in members
                    if member.member_type == member_type and member.display_name.casefold() == lowered
                ]
                if len(exact_matches) == 1:
                    target_member_id = exact_matches[0].member_id
                    member_name = exact_matches[0].display_name
                elif len(exact_matches) > 1:
                    return f"Multiple {member_type} members matched '{member_name}'. Please provide member_id."
                else:
                    fuzzy_matches = [
                        member for member in members
                        if member.member_type == member_type and lowered in member.display_name.casefold()
                    ]
                    if len(fuzzy_matches) == 1:
                        target_member_id = fuzzy_matches[0].member_id
                        member_name = fuzzy_matches[0].display_name
                    elif len(fuzzy_matches) > 1:
                        options = ", ".join(member.display_name for member in fuzzy_matches[:5])
                        return f"Multiple {member_type} members matched '{member_name}': {options}. Please provide member_id."
                    else:
                        return f"No {member_type} member matched '{member_name}'."

            existing_res = await db.execute(
                select(MemberDailyReport).where(
                    MemberDailyReport.tenant_id == ag.tenant_id,
                    MemberDailyReport.member_type == member_type,
                    MemberDailyReport.member_id == target_member_id,
                    MemberDailyReport.report_date == report_date,
                )
            )
            existing = existing_res.scalar_one_or_none()
            previous_content = existing.content if existing else ""

        report = await _upsert(
            tenant_id=ag.tenant_id,
            member_type=member_type,
            member_id=target_member_id,
            report_date=report_date,
            content=content,
            source=source,
        )

        resolved_name = member_name or str(target_member_id)
        action = "Updated" if previous_content else "Created"
        details = [
            f"{action} daily report for {resolved_name} on {report.report_date.isoformat()}.",
            f"Stored length: {len(report.content)} characters.",
            f"Status: {report.status}.",
        ]
        if previous_content:
            details.append(f"Previous content: {previous_content}")
        details.append(f"Current content: {report.content}")
        return " ".join(details)
    except Exception as e:
        logger.exception("[OKR] upsert_member_daily_report failed")
        return f"Failed to upsert member daily report: {str(e)[:200]}"


# ── Vercel & Neon Deploy Helper Functions ──

async def _get_vercel_token(agent_id: uuid.UUID, tool_name: str) -> str | None:
    if not tool_name.startswith("vercel_"):
        return None
    # All Vercel operations share one credential source.  Reading a sibling's
    # stale legacy config here would disagree with Runtime readiness and could
    # execute with a different token than the one that made the tool visible.
    config = await _get_tool_config(agent_id, "vercel_deploy")
    return (config or {}).get("vercel_token")


async def _get_vercel_quota_summary(vercel_token: str) -> str:
    import httpx
    headers = {"Authorization": f"Bearer {vercel_token}"}
    async with httpx.AsyncClient() as client:
        try:
            proj_res = await client.get("https://api.vercel.com/v9/projects", headers=headers)
            if proj_res.status_code == 200:
                projects = proj_res.json().get("projects", [])
                project_count = len(projects)
                user_res = await client.get("https://api.vercel.com/v2/user", headers=headers)
                username = "User"
                plan = "Hobby"
                if user_res.status_code == 200:
                    user_data = user_res.json().get("user", {})
                    username = user_data.get("username", username)
                    plan = user_data.get("billing", {}).get("plan", plan)

                quota_str = f"📊 **Vercel Account status ({username} - {plan} Plan)**:\n- Active Projects: {project_count}"
                return quota_str
        except Exception as e:
            logger.warning(f"Error fetching Vercel quota info: {e}")

    return "📊 **Vercel Account status**: Active (Quota details unavailable)"


async def _check_neon_quota_limit(api_key: str) -> tuple[bool, str]:
    import httpx
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Accept": "application/json"
    }
    async with httpx.AsyncClient() as client:
        try:
            res = await client.get("https://console.neon.tech/api/v2/projects", headers=headers)
            if res.status_code == 200:
                projects = res.json().get("projects", [])
                project_count = len(projects)
                if project_count >= 1:
                    return True, f"⚠️ **Neon 免费额度已达上限** (当前项目数: {project_count}/1)。请升级您的 Neon 账户，或者删除已有的旧项目。"
                return False, f"📊 **Neon 账户额度**: {project_count}/1 个项目已使用。"
        except Exception as e:
            logger.warning(f"Error checking Neon quota: {e}")
    return False, "📊 **Neon 账户额度**: 正常 (无法获取详细额度)"


def _prepare_vercel_upload_manifest(
    workspace_root: Path,
    source_dir: object,
) -> list[tuple[str, bytes, str, int]]:
    """Read a complete, workspace-confined upload manifest before provider I/O."""
    source_text = str(source_dir or "").strip()
    if not source_text:
        raise ValueError("source_dir is required for upload deployments")

    relative_source = Path(source_text)
    if relative_source.is_absolute() or ".." in relative_source.parts:
        raise ValueError("source_dir must be a workspace-relative path without '..'")

    resolved_workspace = workspace_root.resolve(strict=True)
    resolved_source = (resolved_workspace / relative_source).resolve(strict=True)
    try:
        resolved_source.relative_to(resolved_workspace)
    except ValueError as exc:
        raise ValueError("source_dir escapes the agent workspace") from exc
    if not resolved_source.is_dir():
        raise ValueError(f"source_dir is not a directory: {source_text}")

    ignored_dirs = {
        ".git",
        ".next",
        ".vercel",
        "build",
        "dist",
        "node_modules",
        "out",
    }
    manifest: list[tuple[str, bytes, str, int]] = []

    def raise_walk_error(error: OSError) -> None:
        raise error

    for root, dirs, files in os.walk(
        resolved_source,
        followlinks=False,
        onerror=raise_walk_error,
    ):
        dirs[:] = sorted(directory for directory in dirs if directory not in ignored_dirs)
        root_path = Path(root)
        for directory in dirs:
            directory_path = root_path / directory
            resolved_directory = directory_path.resolve(strict=True)
            try:
                resolved_directory.relative_to(resolved_source)
            except ValueError as exc:
                raise ValueError(
                    f"directory symlink escapes source_dir: {directory_path}"
                ) from exc
            if directory_path.is_symlink():
                raise ValueError(
                    f"directory symlinks are not supported: {directory_path}"
                )

        for file_name in sorted(files):
            file_path = root_path / file_name
            resolved_file = file_path.resolve(strict=True)
            try:
                resolved_file.relative_to(resolved_source)
            except ValueError as exc:
                raise ValueError(
                    f"file symlink escapes source_dir: {file_path}"
                ) from exc
            if not resolved_file.is_file():
                raise ValueError(f"upload entry is not a regular file: {file_path}")

            relative_path = file_path.relative_to(resolved_source).as_posix()
            file_bytes = file_path.read_bytes()
            digest = hashlib.sha1(
                file_bytes,
                usedforsecurity=False,
            ).hexdigest()
            manifest.append(
                (relative_path, file_bytes, digest, len(file_bytes))
            )

    return manifest


async def _vercel_deploy(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Legacy display adapter for the typed Vercel deployment lifecycle."""
    outcome = await _vercel_deploy_outcome(agent_id, ws, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Vercel deployment returned no summary.",
    )


async def _vercel_read_outcome(
    tool_name: str,
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Execute one Vercel read from explicit HTTP and payload facts."""
    import httpx
    from urllib.parse import quote, urlparse

    if tool_name == "vercel_list_deployments":
        project_name_value = arguments.get("project_name")
        if (
            not isinstance(project_name_value, str)
            or not project_name_value.strip()
        ):
            return _typed_failure(
                "vercel_list_deployments requires project_name.",
                "invalid_tool_arguments",
            )
        project_name = project_name_value.strip()
        request_url = (
            "https://api.vercel.com/v6/deployments?projectId="
            f"{quote(project_name, safe='')}"
        )
        provider_reference = project_name
    elif tool_name == "vercel_get_deploy_logs":
        deployment_value = arguments.get("deployment_id")
        if not isinstance(deployment_value, str) or not deployment_value.strip():
            return _typed_failure(
                "vercel_get_deploy_logs requires deployment_id.",
                "invalid_tool_arguments",
            )
        deployment_reference = deployment_value.strip()
        if deployment_reference.startswith("https://"):
            try:
                parsed = urlparse(deployment_reference)
                parsed_hostname = parsed.hostname
                parsed_port = parsed.port
            except ValueError:
                parsed = None
                parsed_hostname = None
                parsed_port = None
            if (
                parsed is None
                or parsed.scheme != "https"
                or not parsed_hostname
                or parsed.username is not None
                or parsed.password is not None
                or parsed_port is not None
            ):
                return _typed_failure(
                    "deployment_id must be an explicit ID or valid HTTPS URL.",
                    "invalid_tool_arguments",
                )
            deployment_id = parsed_hostname
        elif (
            "://" in deployment_reference
            or "/" in deployment_reference
            or any(character.isspace() for character in deployment_reference)
        ):
            return _typed_failure(
                "deployment_id must be an explicit ID or valid HTTPS URL.",
                "invalid_tool_arguments",
            )
        else:
            deployment_id = deployment_reference
        request_url = (
            "https://api.vercel.com/v2/deployments/"
            f"{quote(deployment_id, safe='')}/events"
        )
        provider_reference = deployment_id
    else:
        return _typed_failure(
            "Unsupported Vercel read tool.",
            "invalid_tool_arguments",
        )

    try:
        token = await _get_vercel_token(agent_id, tool_name)
    except Exception as exc:
        return _typed_failure(
            f"Vercel credential lookup failed: {type(exc).__name__}.",
            "vercel_credentials_lookup_failed",
        )
    if not isinstance(token, str) or not token.strip():
        return _typed_failure(
            "Vercel Access Token is not configured.",
            "vercel_credentials_missing",
        )

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(
                request_url,
                headers={"Authorization": f"Bearer {token.strip()}"},
            )
    except httpx.TimeoutException:
        return _typed_failure(
            "Vercel read timed out.",
            "vercel_read_timeout",
            retryable=True,
        )
    except httpx.TransportError as exc:
        return _typed_failure(
            f"Vercel read transport failed: {type(exc).__name__}.",
            "vercel_read_transport_failed",
            retryable=True,
        )
    except Exception as exc:
        return _typed_failure(
            f"Vercel read failed: {type(exc).__name__}.",
            "vercel_read_failed",
        )

    if not 200 <= response.status_code < 300:
        return _typed_failure(
            f"Vercel read returned HTTP {response.status_code}.",
            f"{tool_name}_http_error",
            retryable=_read_http_status_retryable(response.status_code),
        )
    try:
        data = response.json()
    except Exception:
        return _typed_failure(
            "Vercel read returned invalid JSON.",
            f"{tool_name}_response_invalid",
            retryable=True,
        )

    if tool_name == "vercel_list_deployments":
        if (
            not isinstance(data, Mapping)
            or "error" in data
            or not isinstance(data.get("deployments"), list)
        ):
            return _typed_failure(
                "Vercel returned an invalid deployment collection.",
                "vercel_list_deployments_response_invalid",
                retryable=True,
            )
        deployments = data["deployments"]
        if any(not isinstance(item, Mapping) for item in deployments):
            return _typed_failure(
                "Vercel returned an invalid deployment entry.",
                "vercel_list_deployments_response_invalid",
                retryable=True,
            )
        if not deployments:
            return _typed_success(
                f"No deployments found for project '{project_name}'."
            )

        lines = [
            f"Deployments for {project_name} "
            f"({min(len(deployments), 10)} shown):"
        ]
        evidence_refs: list[str] = []
        for deployment in deployments[:10]:
            deployment_id_value = deployment.get("uid") or deployment.get("id")
            if (
                not isinstance(deployment_id_value, str)
                or not deployment_id_value.strip()
            ):
                return _typed_failure(
                    "Vercel returned a deployment without a stable ID.",
                    "vercel_list_deployments_response_invalid",
                    retryable=True,
                )
            stable_id = deployment_id_value.strip()
            evidence_refs.append(
                f"vercel-deployment://{quote(stable_id, safe='')}"
            )
            deployment_url = str(deployment.get("url") or "").strip()[:500]
            if deployment_url and not deployment_url.startswith(
                ("http://", "https://")
            ):
                deployment_url = f"https://{deployment_url}"
            deployment_state = str(
                deployment.get("state") or deployment.get("readyState") or "unknown"
            ).strip()[:100]
            created_value = deployment.get("created")
            if (
                isinstance(created_value, (int, float))
                and not isinstance(created_value, bool)
            ):
                try:
                    created_text = datetime.fromtimestamp(
                        created_value / 1000,
                        timezone.utc,
                    ).strftime("%Y-%m-%d %H:%M:%S UTC")
                except (OverflowError, OSError, ValueError):
                    created_text = str(created_value)[:100]
            else:
                created_text = str(created_value or "unknown")[:100]
            lines.append(
                f"- ID: {stable_id}; URL: {deployment_url or 'unavailable'}; "
                f"Status: {deployment_state or 'unknown'}; "
                f"Created: {created_text}"
            )
        return _typed_success(
            "\n".join(lines),
            evidence_refs=tuple(evidence_refs),
        )

    if isinstance(data, Mapping):
        if "error" in data or not isinstance(data.get("events"), list):
            return _typed_failure(
                "Vercel returned an invalid deployment log collection.",
                "vercel_get_deploy_logs_response_invalid",
                retryable=True,
            )
        events = data["events"]
    elif isinstance(data, list):
        events = data
    else:
        return _typed_failure(
            "Vercel returned an invalid deployment log collection.",
            "vercel_get_deploy_logs_response_invalid",
            retryable=True,
        )
    if any(not isinstance(event, Mapping) for event in events):
        return _typed_failure(
            "Vercel returned an invalid deployment log entry.",
            "vercel_get_deploy_logs_response_invalid",
            retryable=True,
        )
    evidence_refs = (
        f"vercel-deployment://{quote(provider_reference, safe='')}",
    )
    if not events:
        return _typed_success(
            f"No logs found for deployment '{provider_reference}'.",
            evidence_refs=evidence_refs,
        )

    log_lines: list[str] = []
    for event in events:
        payload = event.get("payload", {})
        if payload is None:
            payload = {}
        if not isinstance(payload, Mapping):
            return _typed_failure(
                "Vercel returned an invalid deployment log payload.",
                "vercel_get_deploy_logs_response_invalid",
                retryable=True,
            )
        text_value = payload.get("text") or event.get("text")
        if text_value is None:
            continue
        if not isinstance(text_value, str):
            return _typed_failure(
                "Vercel returned a non-text deployment log entry.",
                "vercel_get_deploy_logs_response_invalid",
                retryable=True,
            )
        if text_value.strip():
            log_lines.append(text_value.strip()[:4000])
    if not log_lines:
        return _typed_success(
            f"No textual logs found for deployment '{provider_reference}'.",
            evidence_refs=evidence_refs,
        )
    content = "\n".join(log_lines[-100:])
    return _typed_success(
        f"Logs for deployment {provider_reference} (last 100 lines):\n{content}",
        evidence_refs=evidence_refs,
    )


async def _vercel_list_deployments(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter for the typed Vercel deployment list read."""
    outcome = await _vercel_read_outcome(
        "vercel_list_deployments",
        agent_id,
        arguments,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Vercel deployment listing returned no summary.",
    )


async def _vercel_get_deploy_logs(agent_id: uuid.UUID, arguments: dict) -> str:
    """Legacy display adapter for the typed Vercel deployment logs read."""
    outcome = await _vercel_read_outcome(
        "vercel_get_deploy_logs",
        agent_id,
        arguments,
    )
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Vercel deployment logs returned no summary.",
    )


def _deploy_response_object(response) -> Mapping | None:
    """Return one provider JSON object without exposing response text."""
    try:
        payload = response.json()
    except Exception:
        return None
    return payload if isinstance(payload, Mapping) else None


def _deploy_provider_error_code(payload: Mapping | None) -> str | None:
    if payload is None:
        return None
    error = payload.get("error")
    if not isinstance(error, Mapping):
        return None
    code = error.get("code")
    return code.strip() if isinstance(code, str) and code.strip() else None


def _vercel_deployment_https_url(value: object) -> str | None:
    """Normalize only Vercel-style host receipts or explicit HTTPS URLs."""
    from urllib.parse import urlsplit

    if not isinstance(value, str):
        return None
    receipt = value.strip()
    if not receipt or len(receipt.encode("utf-8")) > 2048:
        return None
    candidate = receipt if "://" in receipt else f"https://{receipt}"
    try:
        parsed = urlsplit(candidate)
        port = parsed.port
    except ValueError:
        return None
    if (
        parsed.scheme != "https"
        or not parsed.hostname
        or parsed.username is not None
        or parsed.password is not None
        or port is not None
        or parsed.fragment
    ):
        return None
    return candidate


def _vercel_receipt_metadata_fits_preflight(
    *,
    project_name: str,
    deploy_method: str,
    github_repo: str,
    git_ref: str,
    upload_manifest: list[tuple[str, bytes, str, int]],
) -> bool:
    """Reserve bounded room for receipts before any provider write can occur."""
    metadata: dict[str, object] = {
        "provider": "vercel",
        "operation": "deployment_accepted",
        "project_name": project_name,
        "deploy_method": deploy_method,
        "confirmed_blob_digests": sorted(
            {digest for _path, _content, digest, _size in upload_manifest}
        ),
    }
    if deploy_method == "github":
        metadata["git_ref"] = git_ref
        metadata["linked_repo"] = github_repo
    future_project_id = "p" * 512
    future_deployment_id = "d" * 512
    future_deployment_url = "https://" + "u" * 2040
    metadata.update(
        {
            "project_id": future_project_id,
            "deployment_id": future_deployment_id,
            "deployment_url": future_deployment_url,
            "deployment_state": "S" * 100,
            "error_code": "e" * 200,
            "retryable": False,
            "artifact_refs": [future_deployment_url],
            "evidence_refs": [
                f"vercel-deployment://{future_deployment_id}"
            ],
            "nul_replacements": 0,
            "control_replacements": 0,
            "redaction_count": 0,
            "summary_truncated": False,
            "content_hash": "h" * 64,
            "archive_status": "inline",
        }
    )
    encoded = json.dumps(
        metadata,
        ensure_ascii=False,
        allow_nan=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    # This is the same 16 KiB durable metadata ceiling enforced at settlement;
    # placeholders reserve the largest provider receipts accepted below.
    return len(encoded) <= 16 * 1024


async def _vercel_deploy_outcome(
    agent_id: uuid.UUID,
    workspace_root: Path,
    arguments: dict,
) -> ToolExecutionOutcome:
    """Settle the existing Vercel deployment lifecycle from stage receipts."""
    import httpx
    from urllib.parse import quote

    project_value = arguments.get("project_name")
    method_value = arguments.get("deploy_method", "upload")
    repo_value = arguments.get("github_repo")
    ref_value = arguments.get("git_ref", "main")
    framework_value = arguments.get("framework")
    project_name = (
        project_value.strip() if isinstance(project_value, str) else ""
    )
    deploy_method = (
        method_value.strip() if isinstance(method_value, str) else ""
    )
    github_repo = repo_value.strip() if isinstance(repo_value, str) else ""
    git_ref = ref_value.strip() if isinstance(ref_value, str) else ""
    framework = (
        framework_value.strip()
        if isinstance(framework_value, str)
        else ""
    )
    production = arguments.get("production") is True
    if (
        not project_name
        or deploy_method not in {"upload", "github"}
        or deploy_method == "github"
        and (not github_repo or not git_ref)
    ):
        return _typed_failure(
            "vercel_deploy requires project_name and valid method-specific arguments.",
            "invalid_tool_arguments",
        )

    upload_manifest: list[tuple[str, bytes, str, int]] = []
    if deploy_method == "upload":
        try:
            upload_manifest = _prepare_vercel_upload_manifest(
                workspace_root,
                arguments.get("source_dir"),
            )
        except (OSError, ValueError) as exc:
            return _typed_failure(
                f"Vercel upload preflight failed: {type(exc).__name__}.",
                "vercel_upload_preflight_failed",
            )

    if not _vercel_receipt_metadata_fits_preflight(
        project_name=project_name,
        deploy_method=deploy_method,
        github_repo=github_repo,
        git_ref=git_ref,
        upload_manifest=upload_manifest,
    ):
        return _typed_failure(
            "Vercel deployment receipts would exceed the durable metadata limit.",
            "vercel_deploy_receipt_limit_exceeded",
        )

    try:
        token = await _get_vercel_token(agent_id, "vercel_deploy")
    except Exception as exc:
        return _typed_failure(
            f"Vercel credential lookup failed: {type(exc).__name__}.",
            "vercel_credentials_lookup_failed",
        )
    if not isinstance(token, str) or not token.strip():
        return _typed_failure(
            "Vercel Access Token is not configured.",
            "vercel_credentials_missing",
        )

    project_id: str | None = None
    confirmed_blob_digests: list[str] = []
    linked_repo: str | None = None
    deployment_id: str | None = None
    deployment_url: str | None = None
    deployment_state: str | None = None
    write_stage: str | None = None

    def receipt_metadata(*, operation: str) -> dict:
        metadata: dict[str, object] = {
            "provider": "vercel",
            "operation": operation,
            "project_name": project_name,
            "deploy_method": deploy_method,
            "confirmed_blob_digests": sorted(
                set(confirmed_blob_digests)
            ),
        }
        if project_id:
            metadata["project_id"] = project_id
        if deploy_method == "github":
            metadata["git_ref"] = git_ref
        if linked_repo:
            metadata["linked_repo"] = linked_repo
        if deployment_id:
            metadata["deployment_id"] = deployment_id
        if deployment_url:
            metadata["deployment_url"] = deployment_url
        if deployment_state:
            metadata["deployment_state"] = deployment_state
        return metadata

    def project_stage_failure(
        summary: str,
        error_code: str,
        *,
        unknown: bool,
    ) -> ToolExecutionOutcome:
        if unknown:
            return _typed_unknown(
                summary,
                error_code,
                result_ref=project_id,
                metadata=receipt_metadata(operation=write_stage or "deploy"),
            )
        return _typed_failure(
            summary,
            error_code,
            result_ref=project_id,
            metadata=receipt_metadata(operation=write_stage or "deploy"),
        )

    encoded_project = quote(project_name, safe="")
    headers = {"Authorization": f"Bearer {token.strip()}"}
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            try:
                project_response = await client.get(
                    "https://api.vercel.com/v9/projects/"
                    f"{encoded_project}",
                    headers=headers,
                )
            except Exception as exc:
                return _typed_failure(
                    f"Vercel project lookup failed: {type(exc).__name__}.",
                    "vercel_project_lookup_failed",
                )

            if 200 <= project_response.status_code < 300:
                project_data = _deploy_response_object(project_response)
                project_id_value = (
                    project_data.get("id")
                    if project_data is not None
                    else None
                )
                project_name_receipt = (
                    project_data.get("name")
                    if project_data is not None
                    else None
                )
                if (
                    not isinstance(project_id_value, str)
                    or not project_id_value.strip()
                    or len(project_id_value.strip().encode("utf-8")) > 512
                    or project_name_receipt != project_name
                ):
                    return _typed_failure(
                        "Vercel project lookup returned no matching stable receipt.",
                        "vercel_project_lookup_invalid",
                    )
                project_id = project_id_value.strip()
            elif project_response.status_code == 404:
                create_payload: dict[str, object] = {"name": project_name}
                if framework:
                    create_payload["framework"] = framework
                write_stage = "project_create"
                try:
                    create_response = await client.post(
                        "https://api.vercel.com/v9/projects",
                        headers=headers,
                        json=create_payload,
                    )
                except Exception as exc:
                    return project_stage_failure(
                        f"Vercel project create outcome is unknown: {type(exc).__name__}.",
                        "vercel_project_create_outcome_unknown",
                        unknown=True,
                    )
                if create_response.status_code >= 500:
                    return project_stage_failure(
                        "Vercel project create returned an indeterminate server response.",
                        "vercel_project_create_outcome_unknown",
                        unknown=True,
                    )
                if not 200 <= create_response.status_code < 300:
                    return project_stage_failure(
                        "Vercel rejected the project create request.",
                        "vercel_project_create_rejected",
                        unknown=False,
                    )
                create_data = _deploy_response_object(create_response)
                project_id_value = (
                    create_data.get("id")
                    if create_data is not None
                    else None
                )
                project_name_receipt = (
                    create_data.get("name")
                    if create_data is not None
                    else None
                )
                if (
                    not isinstance(project_id_value, str)
                    or not project_id_value.strip()
                    or len(project_id_value.strip().encode("utf-8")) > 512
                    or project_name_receipt != project_name
                ):
                    return project_stage_failure(
                        "Vercel project create returned no matching stable receipt.",
                        "vercel_project_create_outcome_unknown",
                        unknown=True,
                    )
                project_id = project_id_value.strip()
                write_stage = None
            else:
                return _typed_failure(
                    "Vercel project lookup was rejected; project creation was not attempted.",
                    "vercel_project_lookup_rejected",
                )

            if deploy_method == "upload":
                uploaded: set[str] = set()
                for _path, content, digest, size in upload_manifest:
                    if digest in uploaded:
                        continue
                    write_stage = "blob_upload"
                    try:
                        blob_response = await client.post(
                            "https://api.vercel.com/v2/files",
                            headers={
                                **headers,
                                "Content-Type": "application/octet-stream",
                                "x-vercel-digest": digest,
                                "x-vercel-size": str(size),
                            },
                            content=content,
                        )
                    except Exception as exc:
                        return project_stage_failure(
                            f"Vercel blob upload outcome is unknown: {type(exc).__name__}.",
                            "vercel_blob_upload_outcome_unknown",
                            unknown=True,
                        )
                    if blob_response.status_code >= 500:
                        return project_stage_failure(
                            "Vercel blob upload returned an indeterminate server response.",
                            "vercel_blob_upload_outcome_unknown",
                            unknown=True,
                        )
                    if not 200 <= blob_response.status_code < 300:
                        return project_stage_failure(
                            "Vercel rejected a blob upload.",
                            "vercel_blob_upload_rejected",
                            unknown=False,
                        )
                    uploaded.add(digest)
                    confirmed_blob_digests.append(digest)
                    write_stage = None
            else:
                write_stage = "github_link"
                link_url = (
                    "https://api.vercel.com/v9/projects/"
                    f"{encoded_project}/link"
                )
                try:
                    link_response = await client.post(
                        link_url,
                        headers=headers,
                        json={"type": "github", "repo": github_repo},
                    )
                except Exception as exc:
                    return project_stage_failure(
                        f"Vercel GitHub link outcome is unknown: {type(exc).__name__}.",
                        "vercel_github_link_outcome_unknown",
                        unknown=True,
                    )
                link_data = _deploy_response_object(link_response)
                if link_response.status_code >= 500:
                    return project_stage_failure(
                        "Vercel GitHub link returned an indeterminate server response.",
                        "vercel_github_link_outcome_unknown",
                        unknown=True,
                    )
                if 200 <= link_response.status_code < 300:
                    if (
                        link_data is None
                        or link_data.get("type") != "github"
                        or link_data.get("repo") != github_repo
                    ):
                        return project_stage_failure(
                            "Vercel GitHub link returned no matching receipt.",
                            "vercel_github_link_outcome_unknown",
                            unknown=True,
                        )
                    linked_repo = github_repo
                    write_stage = None
                elif (
                    link_response.status_code == 409
                    and _deploy_provider_error_code(link_data)
                    == "PROJECT_ALREADY_LINKED"
                ):
                    write_stage = None
                    try:
                        reconcile_response = await client.get(
                            "https://api.vercel.com/v9/projects/"
                            f"{encoded_project}",
                            headers=headers,
                        )
                    except Exception as exc:
                        return project_stage_failure(
                            f"Vercel GitHub link reconciliation failed: {type(exc).__name__}.",
                            "vercel_github_link_reconciliation_failed",
                            unknown=False,
                        )
                    reconcile_data = _deploy_response_object(
                        reconcile_response
                    )
                    link_receipt = (
                        reconcile_data.get("link")
                        if reconcile_data is not None
                        else None
                    )
                    if (
                        not 200 <= reconcile_response.status_code < 300
                        or not isinstance(link_receipt, Mapping)
                        or link_receipt.get("type") != "github"
                        or link_receipt.get("repo") != github_repo
                    ):
                        return project_stage_failure(
                            "Vercel project is linked to a different or unverified repository.",
                            "vercel_github_link_mismatch",
                            unknown=False,
                        )
                    linked_repo = github_repo
                else:
                    return project_stage_failure(
                        "Vercel rejected the GitHub link request.",
                        "vercel_github_link_rejected",
                        unknown=False,
                    )

            if deploy_method == "upload":
                deployment_payload: dict[str, object] = {
                    "name": project_name,
                    "files": [
                        {"file": path, "sha": digest, "size": size}
                        for path, _content, digest, size in upload_manifest
                    ],
                }
                if framework:
                    deployment_payload["projectSettings"] = {
                        "framework": framework
                    }
            else:
                deployment_payload = {
                    "name": project_name,
                    "gitSource": {
                        "type": "github",
                        "repo": github_repo,
                        "ref": git_ref,
                    },
                }
            if production:
                deployment_payload["target"] = "production"

            write_stage = "deployment_create"
            try:
                deployment_response = await client.post(
                    "https://api.vercel.com/v13/deployments",
                    headers=headers,
                    json=deployment_payload,
                )
            except Exception as exc:
                return project_stage_failure(
                    f"Vercel deployment create outcome is unknown: {type(exc).__name__}.",
                    "vercel_deployment_create_outcome_unknown",
                    unknown=True,
                )
            if deployment_response.status_code >= 500:
                return project_stage_failure(
                    "Vercel deployment create returned an indeterminate server response.",
                    "vercel_deployment_create_outcome_unknown",
                    unknown=True,
                )
            if not 200 <= deployment_response.status_code < 300:
                return project_stage_failure(
                    "Vercel rejected the deployment create request.",
                    "vercel_deployment_create_rejected",
                    unknown=False,
                )
            deployment_data = _deploy_response_object(deployment_response)
            deployment_id_value = (
                deployment_data.get("id")
                if deployment_data is not None
                else None
            )
            deployment_url_value = (
                deployment_data.get("url")
                if deployment_data is not None
                else None
            )
            if (
                not isinstance(deployment_id_value, str)
                or not deployment_id_value.strip()
                or len(deployment_id_value.strip().encode("utf-8")) > 512
            ):
                return project_stage_failure(
                    "Vercel deployment create returned no stable id/url receipt.",
                    "vercel_deployment_create_outcome_unknown",
                    unknown=True,
                )
            normalized_deployment_url = _vercel_deployment_https_url(
                deployment_url_value
            )
            if normalized_deployment_url is None:
                return project_stage_failure(
                    "Vercel deployment create returned no stable HTTPS URL receipt.",
                    "vercel_deployment_create_outcome_unknown",
                    unknown=True,
                )
            deployment_id = deployment_id_value.strip()
            deployment_url = normalized_deployment_url
            state_value = deployment_data.get("readyState")
            deployment_state = (
                state_value.strip().upper()
                if isinstance(state_value, str)
                and state_value.strip()
                and len(state_value.strip().encode("utf-8")) <= 100
                else "QUEUED"
            )
            write_stage = None

            if deployment_state not in {"READY", "ERROR", "CANCELED"}:
                try:
                    poll_response = await client.get(
                        "https://api.vercel.com/v13/deployments/"
                        f"{quote(deployment_id, safe='')}",
                        headers=headers,
                    )
                except Exception:
                    poll_response = None
                poll_data = (
                    _deploy_response_object(poll_response)
                    if poll_response is not None
                    and 200 <= poll_response.status_code < 300
                    else None
                )
                if (
                    poll_data is not None
                    and poll_data.get("id") == deployment_id
                    and isinstance(poll_data.get("readyState"), str)
                    and str(poll_data.get("readyState")).strip()
                    and len(
                        str(poll_data.get("readyState")).strip().encode("utf-8")
                    )
                    <= 100
                ):
                    deployment_state = str(
                        poll_data["readyState"]
                    ).strip().upper()
                    polled_url = poll_data.get("url")
                    normalized_polled_url = _vercel_deployment_https_url(
                        polled_url
                    )
                    if normalized_polled_url is not None:
                        deployment_url = normalized_polled_url

            metadata = receipt_metadata(operation="deployment_accepted")
            evidence_refs = (
                f"vercel-deployment://{quote(deployment_id, safe='')}",
            )
            if deployment_state in {"ERROR", "CANCELED"}:
                return ToolExecutionOutcome(
                    status="failed",
                    result_summary=f"Vercel deployment reached terminal state {deployment_state}.",
                    result_ref=deployment_id,
                    artifact_refs=(deployment_url,),
                    evidence_refs=evidence_refs,
                    error_code=f"vercel_deployment_{deployment_state.lower()}",
                    metadata=metadata,
                )
            if deployment_state == "READY":
                return _typed_success(
                    f"Vercel deployment {deployment_id} is READY at {deployment_url}.",
                    result_ref=deployment_id,
                    artifact_refs=(deployment_url,),
                    evidence_refs=evidence_refs,
                    metadata=metadata,
                )
            return _typed_success(
                f"Vercel accepted deployment {deployment_id} at {deployment_url}; current state is {deployment_state}.",
                result_ref=deployment_id,
                artifact_refs=(deployment_url,),
                evidence_refs=evidence_refs,
                metadata=metadata,
            )
    except Exception as exc:
        if deployment_id and deployment_url:
            deployment_state = deployment_state or "PENDING"
            return _typed_success(
                f"Vercel accepted deployment {deployment_id}; status polling is pending.",
                result_ref=deployment_id,
                artifact_refs=(deployment_url,),
                evidence_refs=(
                    f"vercel-deployment://{quote(deployment_id, safe='')}",
                ),
                metadata=receipt_metadata(operation="deployment_accepted"),
            )
        if write_stage:
            return project_stage_failure(
                f"Vercel {write_stage} outcome is unknown: {type(exc).__name__}.",
                f"vercel_{write_stage}_outcome_unknown",
                unknown=True,
            )
        return _typed_failure(
            f"Vercel deployment failed before a write was dispatched: {type(exc).__name__}.",
            "vercel_deploy_failed",
            result_ref=project_id,
            metadata=receipt_metadata(operation="deploy_preflight"),
        )


def _deploy_value_storage_key(
    tenant_id: str,
    agent_id: uuid.UUID,
    value_id: str,
) -> str:
    return normalize_storage_key(
        f"runtime/deploy-values/{tenant_id}/{agent_id}/{value_id}.enc"
    )


async def _store_deploy_value_ref(
    agent_id: uuid.UUID,
    value: str,
    **_receipt: object,
) -> str:
    """Encrypt one deploy secret outside Agent-visible workspace paths."""
    from app.core.security import encrypt_data

    if not isinstance(value, str) or not value:
        raise ValueError("deploy value must be a non-empty string")
    tenant_id = await _get_agent_tenant_id(agent_id)
    if not tenant_id:
        raise PermissionError("deploy value requires tenant scope")
    value_id = uuid.uuid4().hex
    encrypted = encrypt_data(value, get_settings().SECRET_KEY)
    await get_storage_backend().write_bytes(
        _deploy_value_storage_key(tenant_id, agent_id, value_id),
        encrypted.encode("ascii"),
        content_type="application/octet-stream",
    )
    return f"deploy-value://{tenant_id}/{agent_id}/{value_id}"


async def _resolve_deploy_value_ref(
    agent_id: uuid.UUID,
    value_ref: str,
) -> str:
    """Resolve only a deploy value owned by the current tenant and Agent."""
    from urllib.parse import urlparse

    from app.core.security import decrypt_data

    if not isinstance(value_ref, str) or not value_ref.strip():
        raise LookupError("deploy value reference is missing")
    tenant_id = await _get_agent_tenant_id(agent_id)
    if not tenant_id:
        raise PermissionError("deploy value requires tenant scope")
    try:
        parsed = urlparse(value_ref.strip())
        path_parts = [part for part in parsed.path.split("/") if part]
    except ValueError as exc:
        raise LookupError("deploy value reference is invalid") from exc
    if (
        parsed.scheme != "deploy-value"
        or parsed.netloc != tenant_id
        or parsed.params
        or parsed.query
        or parsed.fragment
        or len(path_parts) != 2
        or path_parts[0] != str(agent_id)
    ):
        raise PermissionError("deploy value reference scope mismatch")
    try:
        value_id = uuid.UUID(path_parts[1]).hex
    except ValueError as exc:
        raise LookupError("deploy value reference is invalid") from exc
    try:
        encrypted = await get_storage_backend().read_bytes(
            _deploy_value_storage_key(tenant_id, agent_id, value_id)
        )
    except FileNotFoundError as exc:
        raise LookupError("deploy value reference was not found") from exc
    plaintext = decrypt_data(
        encrypted.decode("ascii"),
        get_settings().SECRET_KEY,
    )
    if not plaintext:
        raise LookupError("deploy value reference is empty")
    return plaintext


async def _vercel_set_env_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    import httpx
    from urllib.parse import quote

    project_value = arguments.get("project_name")
    key_value = arguments.get("key")
    project_name = (
        project_value.strip() if isinstance(project_value, str) else ""
    )
    env_key = key_value.strip() if isinstance(key_value, str) else ""
    has_inline = "value" in arguments
    has_ref = "value_ref" in arguments
    target_value = arguments.get("target")
    targets = (
        ["production", "preview", "development"]
        if target_value is None
        else target_value
    )
    valid_targets = {"production", "preview", "development"}
    if (
        not project_name
        or not env_key
        or has_inline == has_ref
        or not isinstance(targets, list)
        or not targets
        or any(
            not isinstance(target, str) or target not in valid_targets
            for target in targets
        )
    ):
        return _typed_failure(
            "vercel_set_env requires project_name, key, exactly one value source, and non-empty valid targets.",
            "invalid_tool_arguments",
        )

    if has_ref:
        value_ref = arguments.get("value_ref")
        if not isinstance(value_ref, str) or not value_ref.strip():
            return _typed_failure(
                "vercel_set_env requires a non-empty value_ref.",
                "invalid_tool_arguments",
            )
        try:
            secret_value = await _resolve_deploy_value_ref(
                agent_id,
                value_ref.strip(),
            )
        except Exception as exc:
            return _typed_failure(
                f"Deploy value reference could not be resolved: {type(exc).__name__}.",
                "deploy_value_ref_unavailable",
            )
    else:
        inline_value = arguments.get("value")
        if not isinstance(inline_value, str) or not inline_value:
            return _typed_failure(
                "vercel_set_env requires a non-empty value.",
                "invalid_tool_arguments",
            )
        secret_value = inline_value

    try:
        token = await _get_vercel_token(agent_id, "vercel_set_env")
    except Exception as exc:
        return _typed_failure(
            f"Vercel credential lookup failed: {type(exc).__name__}.",
            "vercel_credentials_lookup_failed",
        )
    if not isinstance(token, str) or not token.strip():
        return _typed_failure(
            "Vercel Access Token is not configured.",
            "vercel_credentials_missing",
        )

    base_url = (
        "https://api.vercel.com/v9/projects/"
        f"{quote(project_name, safe='')}/env"
    )
    headers = {
        "Authorization": f"Bearer {token.strip()}",
        "Content-Type": "application/json",
    }
    payload = {
        "key": env_key,
        "value": secret_value,
        "type": "encrypted",
        "target": targets,
    }
    create_dispatched = False
    known_env_id: str | None = None
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            create_dispatched = True
            response = await client.post(
                base_url,
                headers=headers,
                json=payload,
            )
            response_data = _deploy_response_object(response)
            if 200 <= response.status_code < 300:
                env_id = (
                    response_data.get("id")
                    if response_data is not None
                    else None
                )
                receipt_key = (
                    response_data.get("key")
                    if response_data is not None
                    else None
                )
                if (
                    not isinstance(env_id, str)
                    or not env_id.strip()
                    or receipt_key != env_key
                ):
                    return _typed_unknown(
                        "Vercel env create returned no matching stable receipt; reconcile before retrying.",
                        "vercel_env_create_outcome_unknown",
                    )
                stable_env_id = env_id.strip()
                known_env_id = stable_env_id
                return _typed_success(
                    f"Environment variable '{env_key}' was created for project '{project_name}'.",
                    result_ref=stable_env_id,
                    evidence_refs=(f"vercel-env://{quote(stable_env_id, safe='')}",),
                    metadata={
                        "provider": "vercel",
                        "operation": "env_create",
                        "env_id": stable_env_id,
                        "env_key": env_key,
                        "project_name": project_name,
                        "targets": list(targets),
                    },
                )

            if not (
                response.status_code == 409
                and _deploy_provider_error_code(response_data)
                == "ENV_ALREADY_EXISTS"
            ):
                if response.status_code >= 500:
                    return _typed_unknown(
                        "Vercel env create returned an indeterminate server response; reconcile before retrying.",
                        "vercel_env_create_outcome_unknown",
                    )
                return _typed_failure(
                    "Vercel rejected the environment variable create request.",
                    "vercel_env_create_rejected",
                )

            try:
                list_response = await client.get(base_url, headers=headers)
            except Exception as exc:
                return _typed_failure(
                    f"Existing Vercel env could not be reconciled: {type(exc).__name__}.",
                    "vercel_env_reconciliation_failed",
                )
            list_data = _deploy_response_object(list_response)
            envs = list_data.get("envs") if list_data is not None else None
            if not 200 <= list_response.status_code < 300 or not isinstance(
                envs,
                list,
            ):
                return _typed_failure(
                    "Existing Vercel env could not be reconciled.",
                    "vercel_env_reconciliation_failed",
                )
            matches = [
                item
                for item in envs
                if isinstance(item, Mapping)
                and item.get("key") == env_key
                and isinstance(item.get("id"), str)
                and str(item.get("id")).strip()
            ]
            if len(matches) != 1:
                return _typed_failure(
                    "Existing Vercel env did not have one stable matching receipt.",
                    "vercel_env_reconciliation_failed",
                )
            env_id = str(matches[0]["id"]).strip()
            known_env_id = env_id
            patch_payload = {
                "value": secret_value,
                "type": "encrypted",
                "target": targets,
            }
            try:
                patch_response = await client.patch(
                    f"{base_url}/{quote(env_id, safe='')}",
                    headers=headers,
                    json=patch_payload,
                )
            except Exception as exc:
                return _typed_unknown(
                    f"Vercel env update outcome is unknown: {type(exc).__name__}; reconcile before retrying.",
                    "vercel_env_update_outcome_unknown",
                    result_ref=env_id,
                    metadata={"env_id": env_id, "env_key": env_key},
                )
            patch_data = _deploy_response_object(patch_response)
            if not 200 <= patch_response.status_code < 300:
                if patch_response.status_code >= 500:
                    return _typed_unknown(
                        "Vercel env update returned an indeterminate server response; reconcile before retrying.",
                        "vercel_env_update_outcome_unknown",
                        result_ref=env_id,
                        metadata={"env_id": env_id, "env_key": env_key},
                    )
                return _typed_failure(
                    "Vercel rejected the existing environment variable update.",
                    "vercel_env_update_rejected",
                    result_ref=env_id,
                    metadata={"env_id": env_id, "env_key": env_key},
                )
            if (
                patch_data is None
                or patch_data.get("id") != env_id
                or patch_data.get("key") != env_key
            ):
                return _typed_unknown(
                    "Vercel env update returned no matching stable receipt; reconcile before retrying.",
                    "vercel_env_update_outcome_unknown",
                    result_ref=env_id,
                    metadata={"env_id": env_id, "env_key": env_key},
                )
            return _typed_success(
                f"Environment variable '{env_key}' was updated for project '{project_name}'.",
                result_ref=env_id,
                evidence_refs=(f"vercel-env://{quote(env_id, safe='')}",),
                metadata={
                    "provider": "vercel",
                    "operation": "env_update",
                    "env_id": env_id,
                    "env_key": env_key,
                    "project_name": project_name,
                    "targets": list(targets),
                },
            )
    except Exception as exc:
        if create_dispatched:
            return _typed_unknown(
                f"Vercel env create outcome is unknown: {type(exc).__name__}; reconcile before retrying.",
                "vercel_env_create_outcome_unknown",
                result_ref=known_env_id,
                metadata=(
                    {"env_id": known_env_id, "env_key": env_key}
                    if known_env_id
                    else None
                ),
            )
        return _typed_failure(
            f"Vercel env request failed before dispatch: {type(exc).__name__}.",
            "vercel_env_create_failed",
        )


async def _vercel_manage_domain_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    import httpx
    from urllib.parse import quote

    action_value = arguments.get("action")
    domain_value = arguments.get("domain")
    project_value = arguments.get("project_name")
    action = action_value.strip() if isinstance(action_value, str) else ""
    domain = domain_value.strip() if isinstance(domain_value, str) else ""
    project_name = (
        project_value.strip() if isinstance(project_value, str) else ""
    )
    if (
        action not in {"check", "bind"}
        or not domain
        or action == "bind"
        and not project_name
    ):
        return _typed_failure(
            "vercel_manage_domain requires a valid action/domain and project_name for bind.",
            "invalid_tool_arguments",
        )
    try:
        token = await _get_vercel_token(agent_id, "vercel_manage_domain")
    except Exception as exc:
        return _typed_failure(
            f"Vercel credential lookup failed: {type(exc).__name__}.",
            "vercel_credentials_lookup_failed",
        )
    if not isinstance(token, str) or not token.strip():
        return _typed_failure(
            "Vercel Access Token is not configured.",
            "vercel_credentials_missing",
        )
    headers = {"Authorization": f"Bearer {token.strip()}"}
    encoded_domain = quote(domain, safe="")

    if action == "check":
        availability_url = (
            "https://api.vercel.com/v1/registrar/domains/"
            f"{encoded_domain}/availability"
        )
        price_url = (
            "https://api.vercel.com/v1/registrar/domains/"
            f"{encoded_domain}/price"
        )
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                availability_response = await client.get(
                    availability_url,
                    headers=headers,
                )
                availability_data = _deploy_response_object(
                    availability_response
                )
                available = (
                    availability_data.get("available")
                    if availability_data is not None
                    else None
                )
                if (
                    not 200 <= availability_response.status_code < 300
                    or not isinstance(available, bool)
                ):
                    return _typed_failure(
                        "Vercel did not return a valid domain availability receipt.",
                        "vercel_domain_availability_failed",
                    )
                price_response = await client.get(price_url, headers=headers)
                price_data = _deploy_response_object(price_response)
                price = (
                    price_data.get("price")
                    if price_data is not None
                    else None
                )
                period = (
                    price_data.get("period")
                    if price_data is not None
                    else None
                )
                if (
                    not 200 <= price_response.status_code < 300
                    or not isinstance(price, (int, float))
                    or isinstance(price, bool)
                    or not isinstance(period, (int, float))
                    or isinstance(period, bool)
                ):
                    return _typed_failure(
                        "Vercel availability was known, but no valid price receipt was returned.",
                        "vercel_domain_price_failed",
                        result_ref=domain,
                        metadata={"domain": domain, "available": available},
                    )
        except Exception as exc:
            return _typed_failure(
                f"Vercel domain check failed: {type(exc).__name__}.",
                "vercel_domain_check_failed",
            )
        availability_text = "available" if available else "unavailable"
        return _typed_success(
            f"Domain '{domain}' is {availability_text}; price is ${price} for period {period}.",
            result_ref=domain,
            evidence_refs=(f"vercel-domain://{encoded_domain}",),
            metadata={
                "provider": "vercel",
                "operation": "domain_check",
                "domain": domain,
                "available": available,
                "price": price,
                "period": period,
            },
        )

    bind_url = (
        "https://api.vercel.com/v9/projects/"
        f"{quote(project_name, safe='')}/domains"
    )
    dispatched = False
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            dispatched = True
            response = await client.post(
                bind_url,
                headers=headers,
                json={"name": domain},
            )
    except Exception as exc:
        if dispatched:
            return _typed_unknown(
                f"Vercel domain bind outcome is unknown: {type(exc).__name__}; reconcile before retrying.",
                "vercel_domain_bind_outcome_unknown",
            )
        return _typed_failure(
            f"Vercel domain bind failed before dispatch: {type(exc).__name__}.",
            "vercel_domain_bind_failed",
        )
    if not 200 <= response.status_code < 300:
        if response.status_code >= 500:
            return _typed_unknown(
                "Vercel domain bind returned an indeterminate server response; reconcile before retrying.",
                "vercel_domain_bind_outcome_unknown",
            )
        return _typed_failure(
            "Vercel rejected the domain bind request.",
            "vercel_domain_bind_rejected",
        )
    data = _deploy_response_object(response)
    if data is None or data.get("name") != domain:
        return _typed_unknown(
            "Vercel domain bind returned no matching receipt; reconcile before retrying.",
            "vercel_domain_bind_outcome_unknown",
        )
    return _typed_success(
        f"Domain '{domain}' was bound to project '{project_name}'.",
        result_ref=domain,
        evidence_refs=(f"vercel-domain://{encoded_domain}",),
        metadata={
            "provider": "vercel",
            "operation": "domain_bind",
            "domain": domain,
            "project_name": project_name,
            "project_id": data.get("projectId"),
            "verified": data.get("verified"),
        },
    )


def _neon_partial_outcome(
    project_id: str,
    *,
    database_name: str,
    region: str,
    error_code: str,
) -> ToolExecutionOutcome:
    return _typed_failure(
        "Neon project was created, but its private connection receipt was not settled; do not recreate the project.",
        error_code,
        result_ref=project_id,
        metadata={
            "provider": "neon",
            "operation": "project_create_partial",
            "project_id": project_id,
            "database_name": database_name,
            "region": region,
        },
    )


async def _neon_create_database_outcome(
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    import httpx

    project_value = arguments.get("project_name")
    database_value = arguments.get("database_name")
    region_value = arguments.get("region", "aws-us-east-1")
    org_value = arguments.get("org_id")
    project_name = (
        project_value.strip() if isinstance(project_value, str) else ""
    )
    database_name = (
        database_value.strip() if isinstance(database_value, str) else ""
    )
    region = region_value.strip() if isinstance(region_value, str) else ""
    org_id = org_value.strip() if isinstance(org_value, str) else ""
    if not project_name or not database_name or not region:
        return _typed_failure(
            "neon_create_database requires project_name, database_name, and region.",
            "invalid_tool_arguments",
        )
    try:
        config = await _get_tool_config(agent_id, "neon_create_database") or {}
    except Exception as exc:
        return _typed_failure(
            f"Neon credential lookup failed: {type(exc).__name__}.",
            "neon_credentials_lookup_failed",
        )
    api_key = config.get("neon_api_key")
    if not isinstance(api_key, str) or not api_key.strip():
        return _typed_failure(
            "Neon API Key is not configured.",
            "neon_credentials_missing",
        )
    try:
        is_blocked, quota_summary = await _check_neon_quota_limit(
            api_key.strip()
        )
    except Exception as exc:
        return _typed_failure(
            f"Neon quota preflight failed: {type(exc).__name__}.",
            "neon_quota_preflight_failed",
        )
    if is_blocked:
        return _typed_failure(
            str(quota_summary or "Neon quota preflight rejected creation."),
            "neon_quota_reached",
        )

    headers = {
        "Authorization": f"Bearer {api_key.strip()}",
        "Content-Type": "application/json",
        "Accept": "application/json",
    }
    create_dispatched = False
    confirmed_project_id: str | None = None
    try:
        async with httpx.AsyncClient(timeout=45.0) as client:
            if not org_id:
                try:
                    org_response = await client.get(
                        "https://console.neon.tech/api/v2/users/me/organizations",
                        headers=headers,
                    )
                except Exception as exc:
                    return _typed_failure(
                        f"Neon organization preflight failed: {type(exc).__name__}.",
                        "neon_org_preflight_failed",
                    )
                org_data = _deploy_response_object(org_response)
                organizations = (
                    org_data.get("organizations")
                    if org_data is not None
                    else None
                )
                if not 200 <= org_response.status_code < 300 or not isinstance(
                    organizations,
                    list,
                ):
                    return _typed_failure(
                        "Neon organization preflight returned no valid collection.",
                        "neon_org_preflight_failed",
                    )
                normalized_orgs = []
                for organization in organizations:
                    if not isinstance(organization, Mapping):
                        return _typed_failure(
                            "Neon organization preflight returned an invalid entry.",
                            "neon_org_preflight_failed",
                        )
                    candidate_id = organization.get("id")
                    if not isinstance(candidate_id, str) or not candidate_id.strip():
                        return _typed_failure(
                            "Neon organization preflight returned an entry without an ID.",
                            "neon_org_preflight_failed",
                        )
                    normalized_orgs.append(
                        (
                            candidate_id.strip(),
                            str(organization.get("name") or "Unnamed")[:100],
                        )
                    )
                if len(normalized_orgs) == 1:
                    org_id = normalized_orgs[0][0]
                elif len(normalized_orgs) > 1:
                    choices = ", ".join(
                        f"{name} ({candidate_id})"
                        for candidate_id, name in normalized_orgs[:20]
                    )
                    return _typed_failure(
                        f"Multiple Neon organizations are available; choose org_id: {choices}.",
                        "neon_org_selection_required",
                    )

            project_payload: dict[str, object] = {
                "project": {
                    "name": project_name,
                    "region_id": region,
                    "pg_version": 15,
                },
                "branch": {"database_name": database_name},
            }
            if org_id:
                project_payload["project"]["org_id"] = org_id  # type: ignore[index]
            create_dispatched = True
            response = await client.post(
                "https://console.neon.tech/api/v2/projects",
                headers=headers,
                json=project_payload,
            )
            if not 200 <= response.status_code < 300:
                if response.status_code >= 500:
                    return _typed_unknown(
                        "Neon project create returned an indeterminate server response; reconcile before retrying.",
                        "neon_project_create_outcome_unknown",
                    )
                return _typed_failure(
                    "Neon rejected the project create request.",
                    "neon_project_create_rejected",
                )
            data = _deploy_response_object(response)
            project = data.get("project") if data is not None else None
            project_id_value = (
                project.get("id") if isinstance(project, Mapping) else None
            )
            if (
                not isinstance(project_id_value, str)
                or not project_id_value.strip()
            ):
                return _typed_unknown(
                    "Neon project create returned no stable project receipt; reconcile before retrying.",
                    "neon_project_create_outcome_unknown",
                )
            project_id = project_id_value.strip()
            confirmed_project_id = project_id
            connection_value = data.get("connection_uri")
            connection_uri = (
                connection_value.strip()
                if isinstance(connection_value, str)
                else ""
            )
            if not connection_uri:
                try:
                    connection_response = await client.get(
                        "https://console.neon.tech/api/v2/projects/"
                        f"{project_id}/connection_string",
                        headers=headers,
                        params={"database_name": database_name},
                    )
                except Exception:
                    return _neon_partial_outcome(
                        project_id,
                        database_name=database_name,
                        region=region,
                        error_code="neon_connection_partial_failure",
                    )
                connection_data = _deploy_response_object(
                    connection_response
                )
                connection_value = (
                    connection_data.get("connection_uri")
                    if connection_data is not None
                    else None
                )
                if (
                    not 200 <= connection_response.status_code < 300
                    or not isinstance(connection_value, str)
                    or not connection_value.strip()
                ):
                    return _neon_partial_outcome(
                        project_id,
                        database_name=database_name,
                        region=region,
                        error_code="neon_connection_partial_failure",
                    )
                connection_uri = connection_value.strip()
    except Exception as exc:
        if confirmed_project_id:
            return _neon_partial_outcome(
                confirmed_project_id,
                database_name=database_name,
                region=region,
                error_code="neon_connection_partial_failure",
            )
        if create_dispatched:
            return _typed_unknown(
                f"Neon project create outcome is unknown: {type(exc).__name__}; reconcile before retrying.",
                "neon_project_create_outcome_unknown",
            )
        return _typed_failure(
            f"Neon project create failed before dispatch: {type(exc).__name__}.",
            "neon_project_create_failed",
        )

    try:
        value_ref = await _store_deploy_value_ref(
            agent_id,
            connection_uri,
            provider="neon",
            resource_id=project_id,
        )
    except Exception:
        return _neon_partial_outcome(
            project_id,
            database_name=database_name,
            region=region,
            error_code="neon_secret_store_partial_failure",
        )
    if not isinstance(value_ref, str) or not value_ref.startswith(
        "deploy-value://"
    ):
        return _neon_partial_outcome(
            project_id,
            database_name=database_name,
            region=region,
            error_code="neon_secret_store_partial_failure",
        )
    return _typed_success(
        f"Neon project '{project_id}' and database '{database_name}' were created; use private value_ref={value_ref} with vercel_set_env.",
        result_ref=project_id,
        evidence_refs=(f"neon-project://{project_id}",),
        metadata={
            "provider": "neon",
            "operation": "project_create",
            "project_id": project_id,
            "database_name": database_name,
            "region": region,
            "value_ref": value_ref,
        },
    )


async def _deploy_simple_write_outcome(
    tool_name: str,
    agent_id: uuid.UUID,
    arguments: dict,
) -> ToolExecutionOutcome:
    if tool_name == "vercel_set_env":
        return await _vercel_set_env_outcome(agent_id, arguments)
    if tool_name == "vercel_manage_domain":
        return await _vercel_manage_domain_outcome(agent_id, arguments)
    if tool_name == "neon_create_database":
        return await _neon_create_database_outcome(agent_id, arguments)
    return _typed_failure(
        "Unsupported simple deploy write.",
        "invalid_tool_arguments",
    )


async def _vercel_set_env(agent_id: uuid.UUID, arguments: dict) -> str:
    outcome = await _vercel_set_env_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Vercel environment write returned no summary.",
    )


async def _vercel_manage_domain(agent_id: uuid.UUID, arguments: dict) -> str:
    outcome = await _vercel_manage_domain_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Vercel domain operation returned no summary.",
    )


async def _neon_create_database(agent_id: uuid.UUID, arguments: dict) -> str:
    outcome = await _neon_create_database_outcome(agent_id, arguments)
    return _legacy_tool_outcome_text(
        outcome,
        fallback="Neon database creation returned no summary.",
    )
